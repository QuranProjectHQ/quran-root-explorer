# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
WK="/sessions/kind-compassionate-feynman/mnt/RootCourse/week02"
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B); GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); SH=prs.slide_height
blank=prs.slide_layouts[6]
def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def slide():
    s=prs.slides.add_slide(blank); b=s.shapes.add_shape(1,0,0,Inches(0.16),SH); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background(); return s
def title(s,t,sz=30):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.28),Inches(12.5),Inches(0.95)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia"); return s
def body(s,lines,top=1.3,size=19,width=12.3,left=0.65,space=8):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(7.5-top-0.2)); tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,bold,col,sz) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space)
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz or size); r.font.bold=bold; r.font.color.rgb=col or INK; sf(r,"Calibri")
def pic(s,relpath,left,top,width=None,height=None):
    p=os.path.join(WK,relpath)
    if os.path.exists(p):
        kw={}
        if width:kw['width']=Inches(width)
        if height:kw['height']=Inches(height)
        s.shapes.add_picture(p,Inches(left),Inches(top),**kw)
def B(t,bold=False,col=None,sz=None): return (t,bold,col,sz)

s=slide()
tb=s.shapes.add_textbox(Inches(0.6),Inches(2.0),Inches(12),Inches(3.4)); tf=tb.text_frame; tf.word_wrap=True
for txt,sz,bold,col,fn in [("WEEK 2",22,True,TEAL,"Georgia"),("Distribution & Concentration",42,True,NAVY,"Georgia"),
  ("Where a root lives, and how unevenly — measured size-true (per 1,000 root-tokens).",21,False,GREY,"Calibri"),
  ("Frequency told us HOW MUCH; today we ask WHERE.  Worked root: ظلم.",19,False,GREY,"Calibri"),
  ("Every value computed from Book6 (6,236 ayahs · 114 surahs).",15,False,GREY,"Calibri")]:
    p=tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]; p.space_after=Pt(8); r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,fn)

s=slide(); title(s,"Recap, and today's question")
body(s,[B("Week 1 — how MUCH a root is named  (frequency).",False,GREY),
 B("A total, though, is silent about geography.",False),
 B("Two roots can have nearly the same count and live completely differently:",True,NAVY),
 B("•  one woven through the whole Qur'an,",False),
 B("•  the other pooled in a few surahs.",False),
 B("Today — WHERE does a root sit, and how evenly is it spread?",True,NAVY),
 B("Two new lenses:  BREADTH (reach)  and  CONCENTRATION (inequality).",True,TEAL),
 B("How much is not where — presence has a geography.",False,GREY),
 B("Worked root: ظلم (injustice).",False,TEAL)],top=1.4,size=19,space=10)

s=slide(); title(s,"What distribution & concentration are")
body(s,[B("DISTRIBUTION — how a root's occurrences spread across the 114 surahs.",True,NAVY,20),
 B("CONCENTRATION — how unevenly they pile up.",True,NAVY,20),
 B("Like a population map: is the root spread evenly across the country,",False),
 B("or crowded into two or three cities?",False,GREY),
 B("Take the per-surah counts — a vector across 114 surahs",False),
 B("(the same vector idea from Week 1) — and summarize its spread and inequality.",False),
 B("•  BREADTH = reach (how many surahs).",False,TEAL),
 B("•  CONCENTRATION = inequality (how pooled).",False,TEAL),
 B("ظلم's 290 ayahs fall across 59 surahs — and not evenly:",False),
 B("its three busiest surahs alone hold 21.7% of them.",False)],top=1.4,size=18,space=8)

s=slide(); title(s,"Breadth — in how many of the 114 surahs")
body(s,[B("The crudest, most intuitive measure of reach:",True,NAVY,20),
 B("how many surahs contain the root at all?",False),
 B("•  Narrow:  عسر (hardship) — 9 surahs.  رشد (right guidance) — 9.",False,RED),
 B("•  Broad:  كفر (disbelief) — 77.  علم (knowledge) — 85.",False,TEAL),
 B("Reach is NOT predicted by frequency:",True,NAVY,20),
 B("a rare root can be wide; a frequent root can be narrow.",False),
 B("Breadth is reach, not amount.",True,TEAL),
 B("But breadth alone can't tell an even spread from a lumpy one —",False),
 B("for that we need concentration.",False,GREY)],top=1.45,size=19,space=10)

s=slide(); title(s,"Concentration — Lorenz, Gini, top-3 share")
pic(s,"fig_concentration.png",7.4,1.65,width=5.7)
body(s,[B("Two roots of equal breadth can be evenly spread OR pooled in one surah.",True,NAVY,18),
 B("Three ways to capture it:",True),
 B("•  top-3 share — the intuitive one.",False),
 B("•  the Lorenz curve — the picture; the harder it bows, the more unequal.",False),
 B("•  the Gini — one number, 0 (even) → 1 (all in one surah).",False),
 B("The app computes the Gini; we read it, we don't hand-calculate it.",False,GREY),
 B("رشد: Gini 0.95, only 9 surahs — concentrated.",False,RED),
 B("كفر: Gini 0.69, 77 surahs — spread.",False,TEAL),
 B("Concentration is NOT importance.",True,NAVY)],top=1.45,size=17.5,width=6.7,space=6)

s=slide(); title(s,"The home surah — two traps")
body(s,[B("Which surah is a root's “home”? Two traps stack.",True,NAVY,20),
 B("TRAP 1 — the surah with the most raw hits is usually just the LONGEST.",True,RED),
 B("al-Baqara has 286 ayahs — the biggest store gets the most foot traffic,",False),
 B("regardless of how good it is. It tops nearly every root's raw count.",False,GREY),
 B("TRAP 2 — dividing by the surah's AYAH count is STILL wrong.",True,RED),
 B("Ayahs vary in length too: a surah of long verses has more room in each.",False),
 B("THE FIX — normalize by ROOT-TOKENS:",True,TEAL,20),
 B("the root's tokens in the surah ÷ the surah's total root-tokens × 1,000.",False),
 B("Measure crowd density per square metre, not per building.",False,GREY)],top=1.4,size=18,space=8)

s=slide(); title(s,"Normalize by ROOT-TOKENS, not ayahs")
pic(s,"fig_normalization_levels.png",0.9,4.7,width=11.5)
body(s,[B("ṣabr (صبر): the “home” changes at EVERY level of normalization —",True,NAVY,20),
 B("•  RAW (ayah-hits)  →  al-Baqara",False,RED),
 B("•  per AYAH  →  al-Kahf",False),
 B("•  per ROOT-TOKENS (size-true)  →  at-Tur",False,TEAL),
 B("Three different surahs. Only per-root-tokens is size-true —",True,NAVY),
 B("ayahs are no more equal than surahs.",False,GREY)],top=1.35,size=18,space=7)

s=slide(); title(s,"ẓulm: the home moves once you normalize")
pic(s,"fig_home_flip.png",0.9,4.5,width=11.5)
body(s,[B("Raw busiest surah for ظلم: al-Baqara (27 hits) —",True,NAVY,20),
 B("but only because al-Baqara is the longest surah in the Qur'an.",False,GREY),
 B("Size-true home: Ibrahim, at 15.8 per 1,000 root-tokens.",True,TEAL,20),
 B("al-Baqara drops to mid-pack once you correct for size.",False)],top=1.4,size=19,space=8)

s=slide(); title(s,"The headline — the length illusion")
body(s,[B("Take the 50 most frequent roots. Their raw busiest surah is al-Baqara for…",False,GREY),
 B("RAW:  30 of 50.",True,RED,30),
 B("Size-true (per root-tokens):  0 of 50.",True,TEAL,30),
 B("The single longest surah masquerades as everyone's home —",True,NAVY,20),
 B("and normalizing dethrones it completely.",True,NAVY,20),
 B("And per-ayah ≠ per-roots — they don't always agree (recall ṣabr).",False),
 B("Trust per-root-tokens.",True,TEAL)],top=1.4,size=19,space=11)

s=slide(); title(s,"The support floor — small-sample reliability")
pic(s,"fig_support_floor.png",7.4,1.7,width=5.7)
body(s,[B("A tiny surah can post a sky-high density off just one or two tokens.",True,NAVY,19),
 B("عسر's highest prevalence is in ash-Sharh —",False),
 B("2 tokens in 16 = 125 per 1k — but that's noise, not a home.",False,RED),
 B("THE FLOOR:",True,TEAL,20),
 B("count ≥ 3 in the surah AND surah ≥ 30 root-tokens.",False),
 B("If no surah qualifies → report “insufficient support.”",False),
 B("عسر has NO reliable home surah.",True,NAVY),
 B("A rate from too little data is a guess.",False,GREY)],top=1.5,size=18,width=6.7,space=8)

s=slide(); title(s,"The size-true rule (locked)")
body(s,[B("Normalize density to per 1,000 ROOT-TOKENS — at every level.",True,TEAL,26),
 B("",False),
 B("Never by ayah-count.  Never by surah-count.  Never raw.",True,RED,24),
 B("",False),
 B("Because containers — ayahs and surahs alike — vary in size.",False,NAVY,22),
 B("It is the same discipline as Week 1's per-1,000-roots rate,",False,GREY),
 B("now applied inside a surah.",False,GREY)],top=1.7,size=20,space=12)

s=slide(); title(s,"What distribution does — and doesn't — tell us")
body(s,[B("GIVES:",True,NAVY,22),
 B("•  a root's geography — where it lives.",False),
 B("•  size-true density (per root-tokens).",False),
 B("•  broad vs pooled — the shape of its presence.",False),
 B("Does NOT give:",True,RED,22),
 B("•  importance — concentration is not centrality.",False),
 B("•  which roots travel together — that is Weeks 3 onward.",False),
 B("كفر across 77 surahs is a pervasive THEME — not the most “central” root.",False,GREY),
 B("So pair every distribution number with a labeled reading.",True,TEAL)],top=1.45,size=19,space=9)

s=slide(); title(s,"The two-sentence reading")
body(s,[B("FACT",True,NAVY,22),
 B("“ظلم spans 59 surahs (Gini 0.74); its size-true home is Ibrahim at",False),
 B("15.8 per 1,000 root-tokens — not al-Baqara, which leads on raw count",False),
 B("only because it is the longest surah.”",False),
 B("INTERPRETATION",True,TEAL,22),
 B("“I read ظلم as a whole-Qur'an concern, not the topic of one surah.”",False),
 B("A classmate should tell instantly which sentence is fact and which is yours.",False,GREY)],top=1.6,size=19,space=10)

s=slide(); title(s,"This week's exercise")
body(s,[B("Part 1 — by hand",True,TEAL,22),
 B("•  Take your root's raw-busiest surah and its size-true home candidate.",False),
 B("•  Compute each one's per-1,000-root-tokens prevalence.",False),
 B("•  Apply the support floor (count ≥ 3, surah ≥ 30 tokens); name the home.",False),
 B("Part 2 — in the app",True,TEAL,22),
 B("•  Record breadth, top-3 share, Gini; confirm your size-true home.",False),
 B("•  One screenshot; one fact + one labeled interpretation.",False),
 B("Submit the night before class — it gates the debrief.",True,NAVY)],top=1.6,size=19,space=10)

s=slide(); title(s,"Where a root lives is a size-true, support-checked claim")
body(s,[B("Never its raw busiest surah.",True,NAVY,24),
 B("",False),
 B("Frequency told us how much; distribution now tells us where —",False,GREY,20),
 B("once we refuse to let the longest surah win by length alone.",False,GREY,20),
 B("",False),
 B("Next week — partners & forms: which roots a root travels with,",True,NAVY,20),
 B("and the morphological shapes it wears.",False,GREY,20)],top=1.9,size=20,space=12)
prs.save(os.path.join(WK,"Week2_Slides.pptx")); print("Week2 deck rebuilt (dense):",len(prs.slides._sldIdLst),"slides")
