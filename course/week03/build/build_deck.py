# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
WK="/sessions/kind-compassionate-feynman/mnt/RootCourse/week03"; FD=os.path.join(WK,"figs")
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B); GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B); GREEN=RGBColor(0x2C,0x7A,0x3F)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); SH=prs.slide_height
blank=prs.slide_layouts[6]
def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def slide():
    s=prs.slides.add_slide(blank); b=s.shapes.add_shape(1,0,0,Inches(0.16),SH); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background(); return s
def title(s,t,sz=30):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.28),Inches(12.5),Inches(0.95)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia"); return s
def body(s,lines,top=1.3,size=19,width=12.3,left=0.65,space=8):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(7.5-top-0.2)); tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,bold,col,sz) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space)
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz or size); r.font.bold=bold; r.font.color.rgb=col or INK; sf(r,"Calibri")
def pic(s,name,left,top,width=None,height=None):
    p=os.path.join(FD,name)
    if os.path.exists(p):
        kw={}
        if width:kw['width']=Inches(width)
        if height:kw['height']=Inches(height)
        s.shapes.add_picture(p,Inches(left),Inches(top),**kw)
def B(t,bold=False,col=None,sz=None): return (t,bold,col,sz)

s=slide()
tb=s.shapes.add_textbox(Inches(0.6),Inches(2.0),Inches(12),Inches(3.4)); tf=tb.text_frame; tf.word_wrap=True
for txt,sz,bold,col,fn in [("WEEK 3",22,True,TEAL,"Georgia"),("Partners & Forms",46,True,NAVY,"Georgia"),
  ("A root's company — its FORMS on the inside, its PARTNERS on the outside.",21,False,GREY,"Calibri"),
  ("Arabic root-and-pattern morphology + collocation.  Worked root: ءمن (believe).",19,False,GREY,"Calibri"),
  ("Every value computed from Book6.",15,False,GREY,"Calibri")]:
    p=tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]; p.space_after=Pt(8); r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,fn)

s=slide(); title(s,"Recap, and today's two questions")
body(s,[B("Week 1 — how MUCH a root is named  (frequency).",False,GREY),
 B("Week 2 — WHERE it sits and how it spreads  (distribution).",False,GREY),
 B("Today — a root's COMPANY, in two directions:",True,NAVY),
 B("•  INSIDE: its surface FORMS — one root, many words (agent, act, attribute).",False),
 B("•  OUTSIDE: its PARTNERS — the roots it travels with in the same ayah.",False),
 B("Both are computed; both need a labeled reading.",True,TEAL),
 B("In Week 1 we melted every word down to its 3 consonants to COUNT it.",False),
 B("This week we pick the shapes back up — to read what we threw away.",False),
 B("Worked root: ءمن (believe).",False,TEAL)],top=1.4,size=19,space=10)

s=slide(); title(s,"The root-and-pattern system (الجذر والوزن)")
pic(s,"fig_root_pattern.png",7.2,1.55,width=6.0)
body(s,[B("Arabic = a 3-consonant ROOT poured into a PATTERN (wazn).",True,NAVY,20),
 B("Root = the raw metal; pattern = the mould you press it into.",False,TEAL),
 B("From ك-ت-ب, “write”:",True),
 B("•  كاتب — writer (the doer)",False),
 B("•  مكتوب — written (the done-to)",False),
 B("•  كِتاب — a book",False),
 B("•  يكتب — he writes",False),
 B("•  كُتِب — was decreed",False),
 B("Same metal, six jobs — the PATTERN assigns the role.",True,NAVY),
 B("So you can't read a word's force from the bare root alone.",False,GREY)],top=1.45,size=18,width=6.5,space=7)

s=slide(); title(s,"Reading a form distribution")
pic(s,"fig_amn_forms.png",7.4,1.7,width=5.7)
body(s,[B("Tally each surface form's share of the root's tokens.",True,NAVY,20),
 B("ءمن: 27 forms across 879 tokens —",False),
 B("•  آمن (verb) — 41%",False),
 B("•  مؤمنين (participle) — 16.5%",False),
 B("•  يؤمن (verb) — 16%",False),
 B("•  إيمان (the noun “faith”) — only 4%",False),
 B("A root is not used evenly — like a coin that lands",False),
 B("far more often on one face than the others.",False,GREY),
 B("Read the dominant form BEFORE you say what a root “means.”",True,TEAL)],top=1.5,size=18,width=6.8,space=7)

s=slide(); title(s,"Pattern families — and the unlearn")
pic(s,"fig_amn_patterns.png",7.0,1.85,width=6.2)
body(s,[B("Group forms into families — where the grammar lives:",True,NAVY,20),
 B("•  Verb (the act): 61%",False,TEAL),
 B("•  Participle (the believers): 26%",False),
 B("•  Masdar (abstract faith): 5%",False),
 B("•  Security branch: 5%",False),
 B("THE UNLEARN:",True,RED,20),
 B("Faith is overwhelmingly a VERB — something DONE —",True),
 B("not the abstract noun إيمان.",True),
 B("The grammar itself encodes the theology.",False,GREY)],top=1.5,size=18,width=6.4,space=7)

s=slide(); title(s,"The intensive forms & the Divine Names")
pic(s,"fig_divine_names.png",7.1,1.85,width=6.1)
body(s,[B("The intensive mould (فعيل / فعّال) turns a quality up to its maximum.",True,NAVY,19),
 B("•  رحم → رحيم   •  سمع → سميع",False),
 B("•  بصر → بصير   •  غفر → غفور",False),
 B("•  حكم → حكيم",False),
 B("These are Names of God.",True,TEAL,20),
 B("Many Divine Names are the SAME mould pressed over different roots —",False),
 B("the attributes are grammatically marked as superlatives.",False),
 B("Learn one pattern and a whole class of the Names lines up.",True,NAVY)],top=1.5,size=18,width=6.4,space=8)

s=slide(); title(s,"Polysemy — one root, divergent meanings")
pic(s,"fig_polysemy.png",6.7,1.7,width=6.5)
body(s,[B("Forms split a root by SENSE — or even by moral VALENCE.",True,NAVY,19),
 B("By sense:",True),
 B("•  ءمن → إيمان (faith) / أمن (security)",False),
 B("•  كتب → كِتاب (a book) / كُتِب (was decreed)",False),
 B("By valence — the deepest case:",True,GREEN),
 B("•  ك-ث-ر → كوثر (blessed abundance, 108:1)",False,GREEN),
 B("           / تكاثر (blameworthy rivalry, 102:1)",False,RED),
 B("One neutral root — “muchness” — pulled to opposite poles",False),
 B("by pattern and context.",False,GREY)],top=1.5,size=18,width=6.0,space=7)

s=slide(); title(s,"Morphological richness varies by root")
pic(s,"fig_form_richness.png",7.3,1.7,width=5.8)
body(s,[B("How many grammatical roles does a root take on?",True,NAVY,20),
 B("•  Rich: ءمن 27 forms, شهد 23, ذكر 23",False),
 B("•  Lean: خلق 12  (78% a single form)",False),
 B("A rich root is one the Qur'an turns over in the hand",False),
 B("many ways — agent, act, attribute, mood.",False),
 B("A lean root it uses for one tight job.",False),
 B("Richness is a clue to how a concept is WORKED.",True,TEAL)],top=1.5,size=18,width=6.8,space=9)

s=slide(); title(s,"Partners — a root's external company")
pic(s,"fig_amn_partners.png",7.3,1.75,width=5.8)
body(s,[B("Partners = roots that co-occur in the same ayah more than chance.",True,NAVY,19),
 B("“You shall know a word by the company it keeps.”",False,GREY),
 B("ءمن's specific partners:",True),
 B("•  صلح (righteous deeds) — z 14",False,TEAL),
 B("•  عمل (works) — z 11",False),
 B("•  رسل (messengers), قلب (heart)",False),
 B("Faith is almost never named alone — it walks arm in arm",False),
 B("with action: آمنوا وعملوا الصالحات.",False),
 B("(The list is already length-controlled — the how is Wk 4–5.)",False,GREY,15)],top=1.5,size=18,width=6.8,space=7)

s=slide(); title(s,"Reading partners honestly — antonyms are partners")
pic(s,"fig_antonym_partners.png",7.4,1.9,width=5.7)
body(s,[B("Three cautions when reading a partner list:",True,NAVY,20),
 B("1.  Significant ≠ meaningful, aligned, or causal.",False),
 B("2.  The list is length-controlled (mechanism = Wk 4–5).",False),
 B("3.  A partner can be an ANTONYM.",True,RED),
 B("Opposites are defined together — like rivals always in the same room:",False),
 B("•  ءمن ↔ كفر  (126 ayahs)",False),
 B("•  هدي ↔ ضلل (52)  ·  ذكر ↔ نسي (12)",False),
 B("A partner is a lead, not a verdict.",True,TEAL)],top=1.5,size=18,width=6.8,space=8)

s=slide(); title(s,"When two roots name each other")
body(s,[B("Some roots are each other's STRONGEST partner — a mutual bond:",True,NAVY,20),
 B("•  رحم ↔ غفر  (mercy & forgiveness, z 25) — the divine pair of pardon.",False,TEAL),
 B("•  سمع ↔ بصر  (hearing & seeing, z 18) — “the All-Hearing, the All-Seeing.”",False,TEAL),
 B("•  حكم ↔ عزز  (wisdom & might, z 24) — “the Mighty, the Wise.”",False,TEAL),
 B("•  خلق ↔ طين / نطف  (creation from clay / a drop).",False),
 B("When two concepts name each other, the text is telling you",False),
 B("they are habitually spoken in one breath.",False),
 B("A strong signal — but still descriptive:",True,NAVY),
 B("the why and the direction wait for Weeks 4–6.",False,GREY)],top=1.45,size=18,space=9)

s=slide(); title(s,"What forms & partners give — and don't")
body(s,[B("FORMS give:  grammar, agency, mood, polysemy — what the root layer hid.",True,NAVY,20),
 B("PARTNERS give:  a concept's lived company — its collocations.",True,NAVY,20),
 B("They do NOT give:",True,RED,20),
 B("•  the full meaning of a word in its ayah  (forms ≠ context).",False),
 B("•  the direction or cause of a partnership — which leads which  (Week 6).",False),
 B("•  which partner shares the MOST ayahs, or why raw counts mislead  (Week 4).",False),
 B("So pair every forms-or-partners claim with a labeled reading.",True,TEAL)],top=1.5,size=18,space=11)

s=slide(); title(s,"The two-sentence reading")
body(s,[B("FACT",True,NAVY,20),
 B("“61% of ء-م-ن's tokens are verb forms, and its strongest partner is صلح (z = 14).”",False),
 B("INTERPRETATION",True,TEAL,20),
 B("“I read faith here as an enacted commitment bound to righteous works,",False),
 B("not an abstract belief.”",False),
 B("The peer-test:",True,NAVY,20),
 B("a classmate should tell instantly which sentence the computer could check,",False),
 B("and which is yours. If not — the two have blurred; rewrite.",False,GREY)],top=1.7,size=19,space=11)

s=slide(); title(s,"This week's exercise")
body(s,[B("Part 1 — by hand",True,TEAL,20),
 B("•  Tally your root's surface forms; record the dominant form.",False),
 B("•  Name its pattern family: verb (act) / participle (agent) / masdar / intensive.",False),
 B("•  Note any polysemy — does the root split by sense or by valence?",False),
 B("Part 2 — in the app",True,TEAL,20),
 B("•  Read your root's significant partners; record the top three.",False),
 B("•  Flag any ANTONYM. Take one screenshot.",False),
 B("Submit one fact + one labeled interpretation per part, the night before class.",True,NAVY)],top=1.6,size=18,space=10)

s=slide(); title(s,"The root says WHAT; the pattern says HOW")
body(s,[B("Forms reveal a concept's grammar; partners reveal its company —",True,NAVY),
 B("both computed, both needing a labeled reading.",True,NAVY),
 B("",False),
 B("The root tells you WHAT, the pattern HOW, the partners WHO WITH.",False,TEAL,20),
 B("",False),
 B("Next week — co-occurrence as a measure:",True,NAVY),
 B("which root shares the MOST ayahs with a target, and why raw counts mislead.",False,GREY)],top=1.8,size=20,space=10)
prs.save(os.path.join(WK,"Week3_Slides.pptx")); print("Week3 deck rebuilt (dense):",len(prs.slides._sldIdLst),"slides")
