# -*- coding: utf-8 -*-
"""Shared Closer-Look slide builder: editable-text panels, canvas-fill, English titles."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B)
GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B); WHITE=RGBColor(0xFF,0xFF,0xFF)
ICE=RGBColor(0xCA,0xDC,0xFC); TINT=RGBColor(0xEA,0xF5,0xF3); TINT2=RGBColor(0xF2,0xF5,0xFB)
FIG="/sessions/kind-compassionate-feynman/mnt/RootCourse/CloserLooks/figs"
def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def newdeck():
    p=Presentation(); p.slide_width=Inches(13.333); p.slide_height=Inches(7.5); return p
def slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); b=s.shapes.add_shape(1,0,0,Inches(0.16),prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background(); return s
def title(s,t,sz=27):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.22),Inches(12.6),Inches(0.92)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia")
def panel(s,left,top,width,height,fill,lines,space=7,align=PP_ALIGN.LEFT,pad=0.32):
    sp=s.shapes.add_shape(5,Inches(left),Inches(top),Inches(width),Inches(height))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background()
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(pad); tf.margin_right=Inches(pad); tf.margin_top=Inches(0.1); tf.margin_bottom=Inches(0.1)
    for i,(txt,sz,bold,col) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space); p.alignment=align
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,"Calibri")
    return sp
def pic(s,name,left,top,width):
    p=os.path.join(FIG,name)
    if os.path.exists(p): s.shapes.add_picture(p,Inches(left),Inches(top),width=Inches(width))
def L(t,sz,b=False,c=INK): return (t,sz,b,c)
