# -*- coding: utf-8 -*-
import sys, os
sys.path.insert(0,"/sessions/kind-compassionate-feynman/mnt/RootCourse/CloserLooks/build")
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from cl_slides import sf,NAVY,TEAL,RED,INK,GREY,TINT,TINT2,ICE
from pptx.dml.color import RGBColor
REDT=RGBColor(0xF7,0xEC,0xEC); AMBERT=RGBColor(0xF7,0xF1,0xDF); FIG="/sessions/kind-compassionate-feynman/mnt/RootCourse/CloserLooks/figs"
DOC="/sessions/kind-compassionate-feynman/mnt/RootCourse/CloserLooks/10_W10_synthesis.pptx"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
def newslide():
    s=prs.slides.add_slide(prs.slide_layouts[6])
    b=s.shapes.add_shape(1,0,0,Inches(0.16),prs.slide_height); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background()
    return s
def title(s,t,sz=23):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.20),Inches(12.6),Inches(0.9)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia")
def panel(s,left,top,width,height,fill,lines,space=9,pad=0.32):
    sp=s.shapes.add_shape(5,Inches(left),Inches(top),Inches(width),Inches(height))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background()
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(pad); tf.margin_right=Inches(pad); tf.margin_top=Inches(0.1); tf.margin_bottom=Inches(0.1)
    for i,(txt,sz,bold,col) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space); p.alignment=PP_ALIGN.LEFT
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,"Calibri")
def pic(s,name,left,top,width):
    p=os.path.join(FIG,name)
    if os.path.exists(p): s.shapes.add_picture(p,Inches(left),Inches(top),width=Inches(width))
def L(t,sz,b=False,c=INK): return (t,sz,b,c)

# ===== SLIDE 1 — conditional logic =====
s=newslide()
title(s,"A Closer Look #10 — the synthesis: the Qur’an’s logic is conditional",sz=23)
panel(s,0.42,1.22,12.5,2.05,TINT2,[
 L("NOT A STATIC FATE — AN IF / THEN / ELSE-IF / THEN, SUSTAINED TO THE END",18,True,NAVY),
 L("The soul is given both vice and virtue, designed optimally. Success is not a single act but a SUSTAINED direction — continual good deeds and thoughts, held to life’s last leg:",17),
 L("IF the soul keeps purifying itself → it SUCCEEDS;   ELSE IF it lets itself be corrupted → it FAILS.   (ash-Shams 91:9–10)",17,True,TEAL),
])
panel(s,0.42,3.42,12.5,2.18,TINT,[
 L("TILTED TOWARD MERCY — AND THE CONDITION IS DURATIVE",18,True,TEAL),
 L("The scales favor the soul: virtue emphasized 6.6×, God named by mercy ~18×, a good deed repaid TENFOLD vs evil once (6:160).",17),
 L("The demand is ongoing: the deed-verb is present-tense (“keeps doing”), the deeds are plural, and perseverance fills 118 verses (“then remained steadfast” — Fussilat 41:30).",17),
])
panel(s,0.42,5.78,12.5,1.42,AMBERT,[
 L("Conditional, so reward and punishment are JUST (a fixed fate would void them) — yet weighted to ENCOURAGE: good repaid tenfold, evil only once, “and they are not wronged.”   (interpretation)",16,True,RED),
],space=4)

# ===== SLIDE 2 — diagnosed + prescribed =====
s=newslide()
title(s,"A Closer Look #10 — the human, diagnosed and then prescribed (doer + deed)",sz=23)
LX=0.42; LW=4.78
panel(s,LX,1.22,LW,2.55,REDT,[
 L("THE DIAGNOSIS  (computed)",18,True,RED),
 L("The human is named by its faults — hasty, forgetful, despairing, grasping, unjust, ungrateful. The cause is placed INSIDE the actor:",17),
 L("heart ↔ disease 21× · heart ↔ faith 2.7×",17,True,NAVY),
])
panel(s,LX,3.95,LW,3.30,TINT,[
 L("THE PRESCRIPTION — two columns",18,True,TEAL),
 L("DOER (inner): faith · taqwā · sincerity · a purified heart · remembrance · knowledge · reflection.",17),
 L("DEED (outer): righteous works · prayer · patience · spending.",17),
 L("Coupled as one: faith → deeds 2.8× (101 verses); sincerity → worship 28.7×.",17,True,NAVY),
])
RX=5.42; RW=7.55
pic(s,"cl10b_doer_deed.png",RX,1.25,RW)
panel(s,RX,4.95,RW,2.30,ICE,[
 L("Both columns are weighted — and the deed rests on the doer",18,True,NAVY),
 L("The inner register rivals the deeds in size: knowledge 728, self 270, remembrance 264, taqwā 237 — vs deeds 313, patience 93, prayer 90.",17),
 L("حسن فاعلي و فعلي — goodness of the doer AND goodness of the deed.",17,True,RED),
])

# ===== SLIDE 3 — thought is the root =====
s=newslide()
title(s,"A Closer Look #10 — thought is the root, action the fruit",sz=23)
panel(s,0.42,1.30,6.05,5.55,TINT2,[
 L("THE VERDICT  (interpretation)",19,True,NAVY),
 L("The object of ash-Shams 91:9–10 is the self (nafs): purify IT, or corrupt IT.",18),
 L("Deeds are the EVIDENCE of which way the self was turned — not the cause itself.",18),
 L("The heart is named as the locus: disease 21×, tranquility 23.5×, faith 2.7×.",18),
 L("So thought is the root, and action the fruit.",18,True,RED),
])
panel(s,6.78,1.30,6.12,4.05,ICE,[
 L("THE SOURCE — the prophetic mission",19,True,NAVY),
 L("Knowledge is not self-made. The messenger’s task is to:",18),
 L("recite the signs → PURIFY (تزكية — the Shams verb) → TEACH the Book and wisdom.",18,True,TEAL),
 L("This formula is complete in 4 verses: 2:129, 2:151, 3:164, 62:2.",18),
])
panel(s,6.78,5.52,6.12,1.33,TINT2,[
 L("Honest limit: نية (explicit “intention”) is ~absent (1 ayah) — it rides on heart, taqwā, and sincerity as a FIELD; the prophetic triad is a 4-verse formula (robust as a phrase, not high-frequency).",14),
],space=4)

# ===== SLIDE 4 — ten lessons =====
s=newslide()
title(s,"Ten Lessons — what careful measurement taught",sz=24)
panel(s,0.42,1.25,12.5,4.55,TINT2,[
 L("#1  A difference is not a finding until it is tested.",17,True,NAVY),
 L("#2  Widespread ≠ evenly spread — control the confound.",17,True,NAVY),
 L("#3  Distinguish a Name from a theme.",17,True,NAVY),
 L("#4  A famous verse is not a frequent pattern.",17,True,NAVY),
 L("#5  Lift finds a real bond — but read its support.",17,True,NAVY),
 L("#6  A shared count hides a direction.",17,True,NAVY),
 L("#7  Motifs surface a text’s recurring scenes.",17,True,NAVY),
 L("#8  A sky-high lift on one verse is not a finding.",17,True,NAVY),
 L("#9  A true number can still be a false claim — audit the reading.",17,True,NAVY),
 L("#10  No single number is a portrait — only the pipeline is.",17,True,RED),
],space=5)
panel(s,0.42,5.95,12.5,1.25,TINT,[
 L("The Qur’an’s direction — the soul toward success — survives measurement; the slogans about it often do not.",17,True,TEAL),
],space=4)

# ===== SLIDES 5 & 6 — Ten Lessons, elaborated =====
def elab_slide(head, blocks):
    s=newslide(); title(s,head,sz=23)
    sp=s.shapes.add_shape(5,Inches(0.42),Inches(1.18),Inches(12.5),Inches(6.05))
    sp.fill.solid(); sp.fill.fore_color.rgb=TINT2; sp.line.fill.background()
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.38); tf.margin_right=Inches(0.38); tf.margin_top=Inches(0.12); tf.margin_bottom=Inches(0.12)
    first=True
    for lesson,lcol,body in blocks:
        p1=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p1.space_before=Pt(0); p1.space_after=Pt(2); p1.alignment=PP_ALIGN.LEFT
        r=p1.add_run(); r.text=lesson; r.font.size=Pt(17); r.font.bold=True; r.font.color.rgb=lcol; sf(r,"Calibri")
        p2=tf.add_paragraph(); p2.space_after=Pt(11); p2.alignment=PP_ALIGN.LEFT
        r=p2.add_run(); r.text=body; r.font.size=Pt(14.5); r.font.color.rgb=INK; sf(r,"Calibri")
    return s

elab_slide("Ten Lessons, elaborated  (1–5)",[
 ("#1  A difference is not a finding until it is tested.",NAVY,
  "Virtue outweighs vice 6.6× — but a raw gap can be an artifact of corpus size or word choice. It only stands after a size-true test (per-1,000 root-tokens) and surface-form filtering.   (CL #1 · frequency)"),
 ("#2  Widespread ≠ evenly spread — control the confound.",NAVY,
  "A root can touch many surahs yet pile up in a few. al-Baqara looked like the “home” of 30 of 50 top roots by raw count — and 0 of 50 once normalized by surah size.   (CL #2 · distribution)"),
 ("#3  Distinguish a Name from a theme.",NAVY,
  "God’s NAMES skew ~18× to mercy (al-Raḥmān, al-Ghafūr); yet punishment is still a frequent THEME (ʿadhāb, 336 verses). A title is not a tally — count the right unit.   (CL #3 · forms)"),
 ("#4  A famous verse is not a frequent pattern.",NAVY,
  "al-ʿAsr’s fourfold (faith, deeds, patience, truth) is iconic, but only faith + deeds actually recurs (101 verses, 2.8×); the other pairs don’t co-occur above chance. Fame ≠ frequency.   (CL #4 · co-occurrence)"),
 ("#5  Lift finds a real bond — but read its support.",NAVY,
  "Lift flags genuine pairs (patience + prayer 5.2×), but a multiplier resting on 2 verses is noise. Always read the × together with the verse-count behind it.   (CL #5 · lift & tiers)"),
])

elab_slide("Ten Lessons, elaborated  (6–10)",[
 ("#6  A shared count hides a direction.",NAVY,
  "Guidance and “the straight path” share verses, but P(guidance | path) = 53% while P(path | guidance) = 9%: the path is a path OF guidance, not the reverse. Symmetry can conceal asymmetry.   (CL #6 · direction)"),
 ("#7  Motifs surface a text’s recurring scenes.",NAVY,
  "Three-root motifs (faith · deeds · reward) recover the Qur’an’s signature moral scene — but only after a length-aware triple-null separates a real motif from coincidence.   (CL #7 · motifs)"),
 ("#8  A sky-high lift on one verse is not a finding.",NAVY,
  "A 576× trio resting on a single verse (21:90) collapses on leave-one-out; a 55× trio spread over 18 verses survives. Support and stability beat raw magnitude.   (CL #8 · significance)"),
 ("#9  A true number can still be a false claim — audit the reading.",NAVY,
  "“Man and woman, 24 each” is arithmetically true for one cherry-picked form but false at the fair root level (66 vs 53). The number is real; the claim built on it is not.   (CL #9 · audit)"),
 ("#10  No single number is a portrait — only the pipeline is.",RED,
  "The human (insān) comes into focus only when frequency, distribution, co-occurrence, direction, motif, significance and audit run together — with computed fact and interpretation kept separate.   (CL #10 · capstone)"),
])

prs.save(DOC)
print("FINAL slides:", len(prs.slides))
