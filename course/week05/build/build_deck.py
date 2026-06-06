# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
WK="/sessions/kind-compassionate-feynman/mnt/RootCourse/week05"; FD=os.path.join(WK,"figs")
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B); GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B); AMBER=RGBColor(0xB8,0x86,0x0B)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); SH=prs.slide_height
blank=prs.slide_layouts[6]
def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def slide():
    s=prs.slides.add_slide(blank); b=s.shapes.add_shape(1,0,0,Inches(0.16),SH); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background(); return s
def title(s,t,sz=30):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.28),Inches(12.5),Inches(0.95)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia"); return s
def body(s,lines,top=1.3,size=19,width=12.3,left=0.65,space=8):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(7.5-top-0.2)); tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,bold,col,sz) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space)
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz or size); r.font.bold=bold; r.font.color.rgb=col or INK; sf(r,"Calibri")
def pic(s,name,left,top,width=None,height=None):
    p=os.path.join(FD,name)
    if os.path.exists(p):
        kw={}
        if width:kw['width']=Inches(width)
        if height:kw['height']=Inches(height)
        s.shapes.add_picture(p,Inches(left),Inches(top),**kw)
def B(t,bold=False,col=None,sz=None): return (t,bold,col,sz)

s=slide()
tb=s.shapes.add_textbox(Inches(0.6),Inches(2.0),Inches(12),Inches(3.4)); tf=tb.text_frame; tf.word_wrap=True
for txt,sz,bold,col,fn in [("WEEK 5",22,True,TEAL,"Georgia"),("Lift & Tiers",46,True,NAVY,"Georgia"),
  ("Putting a bond on trial: is it real, or could it be chance?  Worked pair: صلو ↔ زكو.",21,False,GREY,"Calibri"),
  ("Lift · a length-aware null · Monte-Carlo significance · tiers.",19,False,GREY,"Calibri"),
  ("Every value computed from Book6.",15,False,GREY,"Calibri")]:
    p=tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]; p.space_after=Pt(8); r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,fn)

s=slide(); title(s,"Recap, and today's shift")
body(s,[B("Week 4 RANKED a target's companions and trusted the controlled ranking.",False,GREY),
 B("But a ranking tells you who is closest — not whether even the closest is REAL.",True,NAVY),
 B("Two roots can top each other's lists and still be strangers sharing a few crowded verses.",False),
 B("Today we stop ranking and start JUDGING — for one pair, a VERDICT:",True,TEAL),
 B("•  a LIFT — how much stronger than chance.",False),
 B("•  a SIGNIFICANCE — how sure we can be it is not coincidence.",False),
 B("•  a TIER — structural, borderline, or spurious.",False),
 B("It is the logic of a trial, applied to a pair of words.",False,GREY),
 B("Worked pair: صلو ↔ زكو (prayer ↔ zakat).",False,TEAL)],top=1.4,size=18,space=9)

s=slide(); title(s,"Three questions for every bond")
body(s,[B("To judge a bond between two roots, ask three things — like a court:",True,NAVY,20),
 B("1.  HOW STRONG?  →  the LIFT (× over a fair, length-aware baseline).",False,TEAL),
 B("2.  HOW SURE?  →  the SIGNIFICANCE (could chance alone do this? a p-value).",False,TEAL),
 B("3.  HOW MUCH EVIDENCE?  →  the SUPPORT (at least 5 shared ayahs).",False,TEAL),
 B("A bond worth trusting must pass all three — strong, unlikely-by-chance, and well-supported.",True,NAVY),
 B("Strength alone can be a fluke on 2 ayahs; a big count alone can be pure coincidence.",False,GREY),
 B("The three answers combine into one word: the TIER.",True,NAVY)],top=1.6,size=18,space=10)

s=slide(); title(s,"Two problems Week 4 left unsolved")
pic(s,"fig_two_problems.png",7.0,1.9,width=6.1)
body(s,[B("Problem 1 — long-ayah inflation.",True,RED),
 B("A verse packed with 20 roots is a crowded party where everyone",False),
 B("brushes shoulders; a 3-root verse is a quiet room where",False),
 B("being together MEANS something.",False),
 B("Problem 2 — big counts that are chance.",True,RED),
 B("Two ubiquitous roots share dozens of verses just by being",False),
 B("everywhere. قول↔شيء share 113 — and it means nothing.",False),
 B("The fix: a length-aware null + Monte-Carlo + tiers.",True,TEAL)],top=1.5,size=17,width=6.5,space=6)

s=slide(); title(s,"The length-aware null — a fairer baseline")
pic(s,"fig_lift_deflate.png",7.1,1.9,width=6.0)
body(s,[B("Compare the overlap we SEE to what we'd expect if roots were",True,NAVY,19),
 B("scattered in proportion to ayah LENGTH.",True,NAVY,19),
 B("Lift = observed ÷ this fairer expectation.",True),
 B("Like a handicap in golf: adjust each score for the difficulty",False),
 B("of the course before you compare players.",False,GREY),
 B("Every lift falls a little — the unearned, long-ayah part is trimmed:",False),
 B("صلو↔زكو ×34.6 → ×23.6.  كيل↔وزن ×137 → ×93.",False,TEAL),
 B("Real bonds survive the handicap with room to spare.",True,NAVY)],top=1.55,size=17,width=6.6,space=6)

s=slide(); title(s,"Monte-Carlo significance — testing against chance")
pic(s,"fig_montecarlo.png",7.1,1.95,width=6.0)
body(s,[B("If there were NO bond, how often would chance produce this overlap?",True,NAVY,18),
 B("Keep ayah lengths and root frequencies fixed; scatter the roots",False),
 B("at random; count the overlap. Repeat 3,000 times.",False),
 B("Reshuffle the banquet seating 3,000 times — if two guests still",False,GREY),
 B("end up together far more than any reshuffle manages, they choose each other.",False,GREY),
 B("صلو↔زكو: chance gives 0–3; observed is 28 → p < 0.001.",True,TEAL),
 B("قول↔شيء: chance gives 113+ all the time → p ≈ 0.99.",True,RED)],top=1.6,size=16.5,width=6.6,space=6)

s=slide(); title(s,"The three tiers — a verdict, not a number")
pic(s,"fig_tier_ladder.png",7.0,2.1,width=6.1)
body(s,[B("Combine lift + significance + support, like a court combines evidence:",True,NAVY,18),
 B("Tier 1 — STRUCTURAL",True,TEAL,20),
 B("lift ≥ 3 AND p < 0.001 AND joint ≥ 5  (beyond reasonable doubt).",False),
 B("Tier 2 — BORDERLINE",True,AMBER,20),
 B("significant but modest lift or thin support.",False),
 B("Tier 3 — SPURIOUS",True,RED,20),
 B("p ≥ 0.05 OR lift < 2  (not proven).",False),
 B("Three gates: neither a length artefact nor luck can sneak in.",False,GREY)],top=1.6,size=16.5,width=6.6,space=5)

s=slide(); title(s,"The headline — a big count is NOT a bond")
pic(s,"fig_unlearn.png",6.9,1.85,width=6.2)
body(s,[B("The reversal that should fix the intuition for good:",True,NAVY,19),
 B("قول ↔ شيء — 113 shared ayahs —",True,RED),
 B("lift ×0.8, p ≈ 0.99 → Tier 3, SPURIOUS.",False,RED),
 B("صلو ↔ زكو — 28 shared ayahs —",True,TEAL),
 B("lift ×23.6, p < 0.001 → Tier 1, STRUCTURAL.",False,TEAL),
 B("The pair with the MOST shared verses is the emptiest.",True,NAVY),
 B("113 acquaintances in a crowd vs 28 friends who come to your home.",False,GREY)],top=1.55,size=17,width=6.4,space=6)

s=slide(); title(s,"The verdict — صلو ↔ زكو")
pic(s,"fig_verdict_card.png",7.2,2.2,width=5.7)
body(s,[B("Shared ayahs:  28.",False),
 B("Length-aware lift:  ×23.6 over chance.",False),
 B("Monte-Carlo p:  < 0.001.",False),
 B("Support:  28 of zakat's 56 ayahs — half.",False),
 B("TIER 1 — STRUCTURAL.",True,TEAL,22),
 B("أقيموا الصلاة وآتوا الزكاة —",False),
 B("prayer and almsgiving, inseparable.",False,GREY)],top=1.8,size=18,width=6.6,space=9)

s=slide(); title(s,"Calibration — checking the thresholds")
pic(s,"fig_calibration_table.png",6.7,1.6,width=6.5)
body(s,[B("A classifier you never check is one you can't trust.",True,NAVY,19),
 B("Run all 12 pairs; confirm the known cases land where they should:",False),
 B("•  known-real (صلو↔زكو, كيل↔وزن) → Tier 1.",False,TEAL),
 B("•  known-generic (قول↔شيء, علم↔رحم) → Tier 3.",False,RED),
 B("The thresholds hold — neither too loose nor too strict.",True,NAVY),
 B("Like checking a thermometer against boiling and freezing water",False,GREY),
 B("before you trust it on a fever.",False,GREY)],top=1.6,size=17,width=5.9,space=7)

s=slide(); title(s,"The map — lift vs significance")
pic(s,"fig_scatter.png",6.8,1.5,width=6.4)
body(s,[B("Place every pair by its lift and its significance:",True,NAVY,19),
 B("•  REAL bonds (teal) sit top-right —",False,TEAL),
 B("   high lift, high significance.",False),
 B("•  SPURIOUS (red) sit bottom-left —",False,RED),
 B("   chance-level on both.",False),
 B("•  BORDERLINE (amber) hovers between.",False,AMBER),
 B("The structural cluster is visibly separate from the noise.",True,NAVY)],top=1.6,size=17,width=5.9,space=8)

s=slide(); title(s,"What the verdict can't say (yet)")
body(s,[B("A Tier-1 verdict tells you a bond is REAL. It does NOT tell you:",True,NAVY,20),
 B("•  WHICH root leads the other — does prayer pull in zakat, or the reverse?",False),
 B("   That asymmetry, P(A|B) vs P(B|A), is Week 6.",False,GREY),
 B("•  what the bond MEANS — significance is not meaning.",False),
 B("•  a structural tier is a strong invitation to interpret — never the interpretation.",False),
 B("And remember: most pairs are NOT guilty of any real bond —",True,TEAL),
 B("saying so plainly is as much a finding as convicting the few that are.",False,GREY)],top=1.5,size=18,space=9)

s=slide(); title(s,"The two-sentence reading")
body(s,[B("FACT",True,NAVY,22),
 B("“صلو ↔ زكو share 28 ayahs at lift ×23.6 over a length-aware null,",False),
 B("p < 0.001 — Tier 1, structural.”",False),
 B("INTERPRETATION",True,TEAL,22),
 B("“I read prayer and almsgiving as inseparable halves of one devotion —",False),
 B("the vertical and the horizontal.”",False),
 B("When the statistics say a bond is proven, the temptation to preach is strongest.",True,NAVY),
 B("The statistics earn trust — not meaning.",False,GREY)],top=1.6,size=19,space=10)

s=slide(); title(s,"This week's exercise")
body(s,[B("You are given a PAIR of roots. Deliver a verdict.",True,NAVY,20),
 B("•  Record the joint count and the length-aware lift.",False),
 B("•  Record the Monte-Carlo p-value; check support (joint ≥ 5).",False),
 B("•  Apply the tier rule; name the tier.",False),
 B("•  If your pair has a big count but a high p, explain why it is still spurious.",False),
 B("Submit one fact (lift + p + tier) + one labeled interpretation,",True,TEAL),
 B("the night before class.",False)],top=1.7,size=19,space=10)

s=slide(); title(s,"A verdict comes from lift AND significance")
body(s,[B("Never from the raw number of shared verses.",True,NAVY,24),
 B("",False),
 B("Frequency, distribution, forms, partners, co-occurrence —",False,GREY,18),
 B("and now a tiered VERDICT on a bond: structural, borderline, or spurious.",False,GREY,18),
 B("",False),
 B("Next week — direction & networks:",True,NAVY,20),
 B("which root LEADS the other (P(A|B) vs P(B|A)), and how bonds form hubs.",False,GREY,18)],top=1.8,size=20,space=10)
prs.save(os.path.join(WK,"Week5_Slides.pptx")); print("deck saved:",len(prs.slides._sldIdLst),"slides")
