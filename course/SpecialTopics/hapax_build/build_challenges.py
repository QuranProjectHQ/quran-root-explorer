# -*- coding: utf-8 -*-
"""Special Topic - the Challenges of the Qur'an (al-Tahaddi). Full sec17 spec:
20 slides, >=50% charts (10 chart/figure slides), no filler. Every verse from Book6 by address;
every figure recomputed live from Book6 by gen_figs_chal.py."""
import os,sys
HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
from st_slides import *
from diagrams import fbox,band,harrow
FIG="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics/figs_chal"
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics/SpecialTopic_Quran_Challenges.pptx"
prs=deck()
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def embed(s,ttl,png,cap,cf=TINT2):
    from pptx.util import Inches
    title(s,ttl); w=11.2
    s.shapes.add_picture(os.path.join(FIG,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cf,[L(cap,15,True,TEAL)],space=4)
def gallery(s,ttl,items,fill=TINT,hc=TEAL,intro="the text's own words (address - snippet - gloss):"):
    title(s,ttl); lines=[L(intro,15,True,hc)]
    for ad,ar,gl in items: lines.append(L(ad+"   "+ar+"   -  "+gl,15,False,INK))
    panel(s,0.42,1.15,12.5,6.05,fill,lines,space=7)

A_1788="لَا يَأْتُونَ بِمِثْلِهِ"
A_1113="فَأْتُوا بِعَشْرِ سُوَرٍ مِّثْلِهِ مُفْتَرَيَاتٍ"
A_223="فَأْتُوا بِسُورَةٍ مِّن مِّثْلِهِ"
A_1038="فَأْتُوا بِسُورَةٍ مِّثْلِهِ"
A_5234="فَلْيَأْتُوا بِحَدِيثٍ مِّثْلِهِ"
A_224="فَإن لَّمْ تَفْعَلُوا وَلَن تَفْعَلُوا"
A_482="أَفَلَا يَتَدَبَّرُونَ … لَوَجَدُوا فِيهِ اخْتِلَافًا كَثِيرًا"
A_673="مَّا تَرَى فِي خَلْقِ الرَّحْمَنِ مِن تَفَاوُتٍ … هَلْ تَرَى مِن فُطُورٍ"
A_674="فَارْجِعِ الْبَصَرَ … يَنقَلِبْ إِلَيْكَ الْبَصَرُ خَاسِئًا وَهُوَ حَسِيرٌ"

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.7,TINT2,[L("SPECIAL TOPIC  -  the text that dares its reader",16,True,TEAL),
  L("The Challenges of the Qur'an  (al-Tahaddi)",24,True,NAVY)],space=7)
panel(s,0.42,3.1,12.5,4.0,TINT,[L("A scripture that issues falsifiable dares",18,True,NAVY),
  L("Uniquely, the Qur'an repeatedly CHALLENGES its audience - to imitate it, to find a contradiction in "
    "it, to find a rupture in the creation it describes. These are open, addressable invitations, not "
    "vague claims. This topic gathers them straight from the text, with addresses and live counts.",16.5),
  L("Every verse is quoted from Book6 by address; every chart recomputes from Book6. We PRESENT the "
    "challenges; whether they are 'met' is a theological reading we do not adjudicate.",15.5,True,TEAL)],space=9)

# 2 DIAGRAM - three kinds
s=slide(prs); title(s,"What a 'challenge' is - three kinds, one structure")
band(s,0.42,1.16,12.5,0.42,TINT,"al-tahaddi: an open, addressable, falsifiable dare",NAVY)
three(s,[L("Literary",16,True,TEAL),L("imitate the text - produce its like. Whole -> ten suras -> one sura -> a single statement (17:88; 11:13; 2:23/10:38; 52:34), with the standing prediction 'you never will' (2:24).",14.5)],
 [L("Textual consistency",16,True,AMBER),L("find a contradiction across a corpus revealed over ~23 years - read it through and look (4:82).",14.5)],
 [L("Cosmic observation",16,True,NAVY),L("find a flaw or rupture in creation - inspect the universe, then inspect again (67:3-4).",14.5)])

# 3 FIG escalation
s=slide(prs)
embed(s,"The literary dare escalates DOWN","chal_escalation.png",
  "In the data - the bar is lowered from the whole corpus (~6,236 ayat) to ten suras, to one sura (~55 ayat), to a single statement. Lowering the bar sharpens the dare rather than softening it.",cf=TINT)

# 4 gallery literary verses
s=slide(prs)
gallery(s,"The literary-challenge verses, verbatim",
 [("17:88",A_1788,"jinn & humankind together could not bring its like"),
  ("11:13",A_1113,"then bring ten fabricated suras like it"),
  ("2:23",A_223,"then bring a single sura from its like"),
  ("10:38",A_1038,"then bring one sura like it"),
  ("52:34",A_5234,"then let them bring a statement like it - if truthful")],
 fill=TINT,hc=TEAL)

# 5 FIG mathl haystack
s=slide(prs)
embed(s,"'Bring the like (mithl)' - a needle in the haystack","chal_mathl.png",
  "In the data - the root for 'like' occurs in 148 ayat across the corpus; exactly 5 of them carry the literary challenge. The dare is a precise, locatable subset of a common word - not a vague motif.")

# 6 prediction 2:24
s=slide(prs); title(s,"A standing, falsifiable prediction - 2:24")
panel(s,0.42,1.2,12.5,2.7,AMBERT,[L("2:24   "+A_224,19,True,NAVY),
  L("'And if you do not - and you will NEVER do it -': the text does not merely invite, it PREDICTS the "
    "challenge will go unmet for all time. A prediction is falsifiable: one counter-example would end it.",16.5,True,AMBER)],space=8)
panel(s,0.42,4.1,12.5,3.0,TINT2,[L("Why this is unusual",18,True,NAVY),
  L("Most scriptures assert; this one stakes a claim that could in principle be DISPROVEN by anyone who "
    "produces the like. That structure - an open, repeatable, falsifiable dare - is itself the subject "
    "of this topic, presented here without our verdict.",16.5,True,TEAL)],space=8)

# 7 FIG revorder
s=slide(prs)
embed(s,"When the challenges fall - across the revelation timeline","chal_revorder.png",
  "In the data - the dares appear in both the Meccan period (the literary & cosmic challenges) and the Medinan period (2:23-24, 4:82). The challenge is sustained across the whole revelation, not a one-off.",cf=TINT)

# 8 FIG vocab
s=slide(prs)
embed(s,"The literary-challenge vocabulary","chal_vocab.png",
  "In the data - the dare's words: 'bring!' (486 ayat - the imperative verb of the challenge), 'the like' (148), 'fabricate' (54), 'sura' (16). The challenge is built from the corpus's own core vocabulary.")

# 9 EBAR inventory
s=slide(prs); title(s,"The challenge inventory - verses by kind")
ebar(s,0.42,1.16,12.5,3.7,"Challenge ayat collected from Book6, by kind",
     ["Literary\n(imitate)","Prediction\n(2:24)","Consistency\n(4:82)","Cosmic\n(67:3-4)"],
     [("",[TEAL,AMBER,RED,NAVY],[5,1,1,2])],legend=False,fmt="{:.0f}")
panel(s,0.42,5.0,12.5,2.2,TINT2,[L("Nine challenge ayat, four kinds",16,True,NAVY),
  L("Five literary dares (17:88; 11:13; 2:23; 10:38; 52:34), one standing prediction (2:24), one "
    "consistency dare (4:82), and the two-verse cosmic dare (67:3-4). Few in number, sharp in wording, "
    "and each addressable.",15.5,True,TEAL)],space=7)

# 10 consistency 4:82
s=slide(prs); title(s,"The internal-consistency challenge - 4:82")
panel(s,0.42,1.2,12.5,2.5,TINT,[L("4:82   "+A_482,18,True,NAVY),
  L("'Do they not reflect ...? Had it been from other than God, they would have found in it much "
    "discrepancy.' An invitation to search for internal contradiction across a ~23-year corpus.",16.5,True,TEAL)],space=8)
two(s,[L("The dare",17,True,TEAL),
   L("Read it through (tadabbur) and look for the inconsistencies a long, human-authored, serially-"
     "produced text would accumulate.",16)],
  [L("Where our method touches it",17,True,AMBER),
   L("The Two Books FDR work measured STRUCTURE - contiguity, composition, signal - and found coherent, "
     "non-random organization. That speaks to structure, NOT to this verse's theological claim, which we "
     "do not adjudicate.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 11 FIG fdr
s=slide(prs)
embed(s,"The consistency dare, measured - the Two Books FDR battery","chal_fdr.png",
  "In the data - 6 of 8 representative structural tests survive a single Benjamini-Hochberg 5% FDR (q <= 0.0067); the two borderline tests correctly drop out. Structural coherence is robust - which is NOT proof of 4:82's claim, only that the text is non-random.",cf=TINT)

# 12 cosmic 67:3-4
s=slide(prs); title(s,"The cosmic challenge - look for a flaw  (67:3-4)")
panel(s,0.42,1.2,12.5,2.7,TINT,[L("67:3   "+A_673,17,True,NAVY),
  L("'You see no disparity (tafawut) in the Most Merciful's creation. Return the gaze - do you see any "
    "rupture (futur)?' An empirical dare: inspect the universe for a flaw.",16.5,True,TEAL)],space=8)
panel(s,0.42,4.1,12.5,3.0,AMBERT,[L("67:4   "+A_674,16,True,AMBER),
  L("'Then look again and again - your gaze returns to you humbled and exhausted.' The challenge predicts "
    "the outcome of repeated, careful observation: no rupture found. An invitation to look, framed as a "
    "search for order - which we present, not prove.",15.5,True,NAVY)],space=8)

# 13 FIG cosmic rarity
s=slide(prs)
embed(s,"The cosmic & consistency dares use rare, precise terms","chal_cosmic_rarity.png",
  "In the data - 'disparity' (tafawut) occurs in only 5 ayat, 'rupture' (futur) in 19, 'discrepancy' (ikhtilaf) in 116, 'truthful' (sidq) in 144. The cosmic dare's signature word tafawut is near-unique - a deliberate, exact term.")

# 14 FIG sura context
s=slide(prs)
embed(s,"The challenge suras span the whole corpus","chal_sura_context.png",
  "In the data - the seven challenge suras run from al-Mulk (30 ayat) to al-Baqara (286). The dare is not confined to one register or length - it recurs across short Meccan and long Medinan suras alike.",cf=TINT)

# 15 synthesis diagram
s=slide(prs); title(s,"The family of challenges - one falsifiable structure")
fbox(s,0.55,1.95,3.7,1.55,TINT,"LITERARY","imitate the text - whole->ten->one->a statement  (17:88...52:34)",line=TEAL,tsz=15,ssz=11)
harrow(s,4.35,2.62,0.45,"",color=GREY)
fbox(s,4.95,1.95,3.7,1.55,AMBERT,"CONSISTENCY","find a contradiction across the whole corpus  (4:82)",line=AMBER,tsz=15,ssz=11)
harrow(s,8.75,2.62,0.45,"",color=GREY)
fbox(s,9.35,1.95,3.3,1.55,TINT2,"COSMIC","find a rupture in creation  (67:3-4)",line=NAVY,tsz=15,ssz=11)
panel(s,0.42,3.85,12.5,3.35,TINT,[L("Three targets, one shape",18,True,NAVY),
  L("Whether aimed at the text, its coherence, or the cosmos, every challenge has the same form: an open "
    "invitation, addressable by anyone, that could in principle be defeated by a single counter-example "
    "- a sura, a contradiction, a rupture. That common structure is what makes the tahaddi distinctive.",16,True,TEAL)],space=9)

# 16 honest framing
s=slide(prs); title(s,"How we present these - honestly")
two(s,[L("WHAT WE DO",17,True,TEAL),
   L("- Quote each challenge from Book6 with its address.  - Count the dare-vocabulary from the corpus.  "
     "- Note that each is, by its wording, open and falsifiable.",16,True,NAVY)],
  [L("WHAT WE DO NOT DO",17,True,RED),
   L("- We do NOT declare the challenges met or unmet - that is a theological judgement.  - We do NOT "
     "treat our structural findings as proof of the claims.  - The data is: the text issues these dares, "
     "here.",16,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# 17 audit grid
s=slide(prs); title(s,"Reading the challenges honestly - the audit")
g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
def ar(x,y,mk,mc,hd,bd): panel(s,x,y,w,h,(TINT if mk=="✓" else REDT if mk=="✗" else AMBERT),[L(mk+"  "+hd,15,True,mc),L(bd,13.5)],space=4)
ar(0.42,CY0,"✓",TEAL,"Verses are addressable","every challenge carries a sura:ayah - look it up in Book6.")
ar(0.42+w+g,CY0,"✓",TEAL,"Vocabulary is counted","dare-words tallied live from the corpus.")
ar(0.42,CY0+h+g,"✓",TEAL,"Structure is robust","6/8 FDR tests survive - the corpus is demonstrably non-random.")
ar(0.42+w+g,CY0+h+g,"~",AMBER,"'Met?' is theology","whether a dare stands is a reading we present, not settle.")
ar(0.42,CY0+2*(h+g),"✗",RED,"Not a proof","non-random structure != proof of a divine-origin claim.")
ar(0.42+w+g,CY0+2*(h+g),"✗",RED,"No over-reading","we report verse + count + structure, nothing more.")

# 18 method
s=slide(prs); title(s,"How this was computed (reproducible)")
fbox(s,0.7,1.95,3.7,1.5,TINT,"Book6","6,236 ayat - roots + surface forms - revelation order",line=TEAL,tsz=15,ssz=11)
harrow(s,4.5,2.6,0.5,"locate",color=GREY,lcol=TEAL)
fbox(s,5.4,1.95,3.7,1.5,AMBERT,"challenge ayat","9 verses by address; dare-roots normalized & counted",line=AMBER,tsz=15,ssz=11)
harrow(s,9.25,2.6,0.5,"chart",color=GREY,lcol=RED)
fbox(s,10.2,1.95,2.5,1.5,TINT2,"7 figures","escalation, timeline, vocab, FDR...",line=NAVY,tsz=15,ssz=11)
panel(s,0.42,3.8,12.5,3.4,TINT,[L("Fully reproducible",18,True,NAVY),
  L("Verses are quoted by address; every figure recomputes from Book6.xlsx via gen_figs_chal.py (fixed "
    "seed). Roots are normalized (Persian/Arabic letter variants folded) so counts are stable; the FDR "
    "battery is the live result from the Two Books FDR-summary lecture. Nothing here is hand-tuned.",16,True,TEAL)],space=9)

# 19 takeaway
s=slide(prs); title(s,"Takeaway")
panel(s,0.42,1.2,12.5,5.9,TINT2,[L("A text that invites inspection",18,True,NAVY),
  L("The Qur'an is distinctive in repeatedly daring its reader: imitate me (17:88 -> 52:34), find a "
    "contradiction in me (4:82), find a flaw in the creation I describe (67:3-4) - and it predicts each "
    "dare will stand (2:24).",16.5),
  L("The data shows the dares are few, precise, and addressable; the structure they sit in is "
    "demonstrably non-random (6/8 FDR). Whatever one concludes, the form is remarkable: an open, "
    "falsifiable, repeatable invitation, laid out here from the text itself - the reader inspects, as "
    "the verses ask.",16,True,TEAL),
  L("Presented, not adjudicated: a count of dares and their addresses, faithfully from Book6.",15.5,True,GREY)],space=10)

# 20 quick reference
s=slide(prs); title(s,"Quick reference")
two(s,[L("The challenges (by address)",17,True,TEAL),
   L("Literary: 17:88 - 11:13 - 2:23 - 10:38 - 52:34.  Prediction: 2:24.  Consistency: 4:82.  Cosmic: "
     "67:3-4.  Dare-verb 'bring' in 486 ayat; 'like' in 148.",16)],
  [L("The numbers",17,True,AMBER),
   L("9 challenge ayat - 4 kinds - tafawut only 5 ayat - futur 19 - 6/8 structural tests survive 5% FDR - "
     "7 figures recompute from Book6 (gen_figs_chal.py).",16)],sp=0.5,fa=TINT,fb=AMBERT)

prs.save(OUT)
n=len(list(prs.slides)); print("challenges deck:",OUT,"| slides:",n)
