# -*- coding: utf-8 -*-
import os,sys
HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
from st_slides import *
from diagrams import fbox,band,harrow
FIG="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics/figs_hapax"
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics/SpecialTopic_Hapax_OnceOnly.pptx"
prs=deck()
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def embed(s,ttl,png,cap,cf=TINT2):
    from pptx.util import Inches
    title(s,ttl); w=11.2
    s.shapes.add_picture(os.path.join(FIG,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cf,[L(cap,15,True,TEAL)],space=4)
def gallery(s,ttl,items,fill,hc,intro):
    title(s,ttl); lines=[L(intro,15,True,hc)]+[L(f"{ar}  —  {gl}   ({ad})",15,False,INK) for ar,gl,ad in items]
    panel(s,0.42,1.15,12.5,6.05,fill,lines,space=6)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.7,TINT2,[L("SPECIAL TOPIC  ·  the rarest words",16,True,TEAL),
  L("Hapax Legomena — concepts the Qur'an says exactly ONCE",23,True,NAVY)],space=7)
panel(s,0.42,3.1,12.5,4.0,TINT,[L("Two kinds of 'once', computed from Book6",18,True,NAVY),
  L("A hapax legomenon appears a single time in a corpus. The Qur'an has two levels: a once-only ROOT "
    "(the whole consonantal family occurs once) and a once-only FORM (a specific surface word occurs "
    "once, though its root may be common). 408 roots and 3,027 forms are hapax.",16.5),
  L("Every example carries its address; we report the count and the address, and decline to over-read.",16,True,TEAL)],space=9)

# 2 DIAGRAM — the two levels
s=slide(prs); title(s,"Two levels of 'once-only'")
band(s,0.42,1.2,12.5,0.42,TINT,"the corpus, by how rare its units are",NAVY)
fbox(s,0.7,1.95,5.7,1.6,TINT,"ROOT hapax — 408","the whole 3-letter family occurs once   e.g. صمد (112:2)",line=TEAL,tsz=16,ssz=12)
fbox(s,6.9,1.95,5.7,1.6,AMBERT,"FORM hapax — 3,027","one surface word occurs once; root may be common   e.g. كوثر (108:1, root كثر ×167)",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.8,12.5,3.4,TINT2,[L("Why distinguish them",18,True,NAVY),
  L("A once-only ROOT is a concept the Qur'an names a single time. A once-only FORM is a unique wording of "
    "a possibly-common idea. The two answer different questions — unique concepts vs unique expressions — "
    "so we count and present them separately, both from Book6.",16.5,True,TEAL)],space=9)

# 3 FIG — rarity bands
s=slide(prs)
embed(s,"How rare is 'rare'? — roots by frequency","hapax_rarity.png",
  "In the data — of 1,701 roots: 408 occur once, 212 twice, 127 thrice; only 588 occur 11+ times. Rarity is the NORM, not the exception — the hapax sit at the tip of a steep long tail.",cf=TINT)

# 4 the two flagship examples
s=slide(prs); title(s,"Two flagship examples")
two(s,[L("صمد — a once-only ROOT  (112:2)",17,True,TEAL),
   L("ص-م-د occurs once, in al-Ikhlāṣ: 'Allāh, aṣ-Ṣamad' — the Eternal, Absolute, self-sufficient. A "
     "unique root for a unique description of God. (root freq = 1.)",16)],
  [L("كوثر — a once-only FORM  (108:1)",17,True,AMBER),
   L("كوثر occurs once, opening al-Kawthar — abundance. Its ROOT ك-ث-ر is common (167 āyahs); only this "
     "intensive FORM is unique. (form = 1, root = 167.)",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 5 FIG per-sura roots
s=slide(prs)
embed(s,"Once-only ROOTS by sūra — where they fall","hapax_per_sura.png",
  "In the data — raw counts are highest in the long sūras (al-Baqara 23, Ṭā-Hā 15, al-Ḥajj 14): length alone puts more unique roots there. Length, not mystery, drives the raw tally.")

# 6 FIG density
s=slide(prs)
embed(s,"Density flips it — short sūras are RICHEST","hapax_density.png",
  "In the data — normalize by sūra size and the juz-30 short sūras dominate. Sūra 108 (al-Kawthar) tops it at ~29% hapax. Brevity concentrates the rare word.",cf=TINT)

# 7 FIG cumulative
s=slide(prs)
embed(s,"Where the 408 first appear (cumulative)","hapax_cumulative.png",
  "In the data — new once-only roots accrue steadily through the reading order, with surges in the long early sūras and again in the dense juz-30 tail — the corpus keeps introducing unique vocabulary throughout.")

# 8 FIG forms per sura
s=slide(prs)
embed(s,"Once-only FORMS by sūra","hapax_forms_per_sura.png",
  "In the data — 3,027 forms occur once; like roots, raw counts track length, but the short late sūras are dense with unique wordings (عاديات, كنود, هاويه …).",cf=TINT)

# 9 FIG root length
s=slide(prs)
embed(s,"Are once-only roots different in shape?","hapax_length.png",
  "In the data — hapax roots track the overall root-length profile (mostly 3-letter), with a mild lean toward 4–5-letter roots: rarer roots are slightly likelier to be longer/less-templatic.")

# 10 gallery juz30 roots
s=slide(prs)
gallery(s,"Gallery — once-only ROOTS in the short sūras (juz 30)",
 [("صمد","the Eternal / Absolute","112:2"),("وءد","the infant girl buried alive","81:8"),
  ("عسعس","the night as it recedes","81:17"),("كنس","the stars that sweep / recede","81:16"),
  ("دهق","a cup brim-full","78:34"),("ثجج","water pouring down","78:14"),("وهج","a blazing lamp","78:13"),
  ("طفف","those who give short measure","83:1"),("رين","the rust / stain over hearts","83:14"),
  ("ءرم","Iram of the pillars","89:7")],TINT,TEAL,"once-only in the whole Qur'an:")

# 11 gallery forms
s=slide(prs)
gallery(s,"Gallery — once-only FORMS (common roots, unique words)",
 [("كوثر","abundance (root كثر ×167)","108:1"),("عاديات","the charging war-steeds","100:1"),
  ("كنود","ungrateful to his Lord","100:6"),("هاويه","the abyss / Hell","101:9"),
  ("مسد","a rope of palm-fibre","111:5"),("ضبحا","panting (of the steeds)","100:1"),
  ("مبثوث","scattered (moths)","101:4"),("منفوش","carded (wool)","101:5")],AMBERT,AMBER,"a unique surface form:")

# 12 the foods of 2:61
s=slide(prs); title(s,"Once-only ≠ mystery — often it is SPECIFICITY")
panel(s,0.42,1.15,12.5,2.6,TINT2,[L("One verse, five once-only roots (al-Baqara 2:61)",17,True,NAVY),
  L("بقل herbs · قثّاء cucumber · فوم garlic · عدس lentils · بصل onions — the foods of Egypt, asked of "
    "Moses. Five hapax ROOTS, each naming a SPECIFIC thing mentioned only here. Rarity = lexical "
    "specificity, not hidden code.",16.5,True,TEAL)],space=8)
three(s,[L("Specific items",16,True,NAVY),L("foods, plants, places, materials named once because the topic arises once.",15)],
 [L("Loanwords / names",16,True,AMBER),L("proper nouns (Iram) and borrowings naturally appear a single time.",15)],
 [L("Unique descriptors",16,True,TEAL),L("a few are unique theological terms — صمد, كوثر — singular by design.",15)])

# 13 EBAR — two-level counts
s=slide(prs); title(s,"The corpus by rarity — the headline counts")
ebar(s,0.42,1.16,12.5,3.7,"Book6 inventory — distinct vs once-only",
     ["roots\n(distinct)","ROOT hapax","forms\n(distinct)","FORM hapax"],
     [("",[NAVY,RED,NAVY,RED],[1701,408,6821,3027])],legend=False,fmt="{:.0f}")
panel(s,0.42,5.0,12.5,2.2,TINT2,[L("Roughly 1 in 4 roots, and nearly half of forms, are once-only",16,True,NAVY),
  L("408 of 1,701 roots (24%) and 3,027 of 6,821 forms (44%) occur exactly once — the corpus is built on "
    "a long tail of rare, often single-use vocabulary atop a small core of high-frequency roots.",15.5,True,TEAL)],space=7)

# 14 EBAR — top suras by hapax roots
s=slide(prs); title(s,"Which sūras introduce the most once-only roots")
ebar(s,0.42,1.16,12.5,3.7,"Top sūras by # once-only ROOTS first/only seen there",
     ["al-Baqara\n(2)","Ṭā-Hā\n(20)","al-Ḥajj\n(22)","al-Aʿrāf\n(7)","an-Nisāʾ\n(4)","Yūsuf\n(12)"],
     [("",[TEAL]*6,[23,15,14,13,12,12])],legend=False,fmt="{:.0f}")
panel(s,0.42,5.0,12.5,2.2,AMBERT,[L("Length leads — but read it against size",16,True,AMBER),
  L("The leaders are all long sūras; their raw hapax counts reflect length. The density view (slide 6) "
    "is the size-true reading: short sūras are proportionally richer in once-only words.",15.5,True,NAVY)],space=7)

# 15 DIAGRAM — length vs density resolution
s=slide(prs); title(s,"Resolving the paradox — raw count vs density")
fbox(s,0.7,1.95,5.6,1.6,TINT,"RAW COUNT","long sūras hold more hapax — because they hold more of everything",line=TEAL,tsz=16,ssz=12)
fbox(s,6.9,1.95,5.6,1.6,AMBERT,"DENSITY","short sūras are richer per word — al-Kawthar ~29%",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.8,12.5,3.4,TINT2,[L("Two true statements, one size-true",18,True,NAVY),
  L("Both charts are real. The raw tally tracks length (a confound); the density (hapax ÷ sūra roots) is "
    "the size-true measure and reverses the ranking. The honest claim about 'where rare words live' is "
    "the density one — the same normalization discipline used across the whole course.",16.5,True,TEAL)],space=9)

# 16 DIAGRAM — categories of hapax
s=slide(prs); title(s,"What the 408 once-only roots actually are")
three(s,[L("Concrete specifics",16,True,TEAL),L("foods, animals, materials, body-parts, place-names — most hapax. Named once because mentioned once.",15)],
 [L("Narrative / proper",16,True,AMBER),L("one-off names and story-specific terms (Iram, specific peoples, unique events).",15)],
 [L("Unique descriptors",16,True,RED),L("a small set of singular theological/eschatological terms — aṣ-Ṣamad, al-Kawthar, fuṭūr — unique by design.",15)])

# 17 AUDIT
s=slide(prs); title(s,"Reading hapax honestly — the audit")
g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
def ar(x,y,mk,mc,hd,bd): panel(s,x,y,w,h,(TINT if mk=="✓" else REDT if mk=="✗" else AMBERT),[L(mk+"  "+hd,15,True,mc),L(bd,13.5)],space=4)
ar(0.42,CY0,"✓",TEAL,"Counts are exact","408 roots / 3,027 forms, reproducible from Book6.")
ar(0.42+w+g,CY0,"✓",TEAL,"Each is addressable","every hapax carries a sūra:āyah — look it up.")
ar(0.42,CY0+h+g,"✓",TEAL,"Density is size-true","short-sūra richness holds after normalization.")
ar(0.42+w+g,CY0+h+g,"~",AMBER,"Meaning varies","most hapax are specifics, a few are unique terms.")
ar(0.42,CY0+2*(h+g),"✗",RED,"Not a cipher","a count of one is not a hidden code or miracle.")
ar(0.42+w+g,CY0+2*(h+g),"✗",RED,"No over-reading","we report count + address, nothing more.")

# 18 METHOD
s=slide(prs); title(s,"How this was computed (reproducible)")
fbox(s,0.7,1.95,3.7,1.5,TINT,"Book6","6,236 āyahs · roots + surface forms",line=TEAL,tsz=15)
harrow(s,4.5,2.6,0.5,"count",color=GREY,lcol=TEAL)
fbox(s,5.4,1.95,3.7,1.5,AMBERT,"frequency","normalize forms; tally per root & per form",line=AMBER,tsz=15)
harrow(s,9.25,2.6,0.5,"==1",color=GREY,lcol=RED)
fbox(s,10.2,1.95,2.5,1.5,REDT,"hapax","freq exactly 1",line=RED,tsz=15)
panel(s,0.42,3.8,12.5,3.4,TINT2,[L("Full tables provided",18,True,NAVY),
  L("Every hapax root and form, with its address, is exported to Hapax_roots_full.csv (408) and "
    "Hapax_forms_full.csv (3,027). Roots are normalized (Persian/Arabic letter variants folded) so the "
    "count is stable; the whole pipeline re-runs from Book6.",16.5,True,TEAL)],space=9)

# 19 TAKEAWAY
s=slide(prs); title(s,"Takeaway")
panel(s,0.42,1.2,12.5,5.9,TINT2,[L("The once-only word",18,True,NAVY),
  L("The Qur'an names 408 roots and over 3,000 forms a single time. Some are unique words for unique "
    "things — aṣ-Ṣamad for God, al-Kawthar for the gift; far more are simply the precise word a specific "
    "moment required. Rarity is the corpus's norm, concentrated densest in the short sūras.",16.5),
  L("Read honestly: a count of one, with its address — comprehensive tables attached, nothing over-read.",16,True,TEAL)],space=10)

# 20 QUICK REFERENCE
s=slide(prs); title(s,"Quick reference")
two(s,[L("The numbers",17,True,TEAL),
   L("408 once-only ROOTS · 3,027 once-only FORMS · ~24% of roots, ~44% of forms · densest sūra 108 "
     "(~29%) · flagship صمد (112:2 root), كوثر (108:1 form).",16)],
  [L("Where to find it",17,True,AMBER),
   L("Full lists: Hapax_roots_full.csv, Hapax_forms_full.csv. Galleries above prioritize the salient "
     "short-sūra and theological cases; every figure recomputes from Book6.",16)],sp=0.5,fa=TINT,fb=AMBERT)

prs.save(OUT); print("hapax FULL deck:",OUT,"| slides:",len(list(prs.slides)))
