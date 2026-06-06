# -*- coding: utf-8 -*-
"""One App & Plot Guide per Special Topic (matches the per-module biology lecture guides).
Each guide is generated from the deck itself: its figure slides, their numbers/titles, and
the verse addresses it cites - so every guide is tailored and accurate."""
import os,glob,re
from pptx import Presentation
from docx import Document
from docx.shared import Pt, RGBColor
from docx.oxml.ns import qn
OUTDIR="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics"
ACCENT=RGBColor(0x1F,0x4E,0x79); TEAL=RGBColor(0x0E,0x6D,0x63); GREY=RGBColor(0x55,0x55,0x55)
AR=re.compile(r'[؀-ۿ]'); ADDR=re.compile(r'\b(\d{1,3}:\d{1,3}(?:-\d{1,3})?)\b')
def setcs(r):
    rPr=r._r.get_or_add_rPr(); rf=rPr.get_or_add_rFonts()
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a),"Arial")
def newdoc():
    d=Document(); st=d.styles["Normal"]; st.font.name="Arial"; st.font.size=Pt(11)
    rf=st.element.get_or_add_rPr().get_or_add_rFonts()
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a),"Arial")
    return d
def P(d,text,size=11,bold=False,color=None,after=4,before=0,style=None,italic=False):
    p=d.add_paragraph(style=style); p.paragraph_format.space_after=Pt(after); p.paragraph_format.space_before=Pt(before)
    r=p.add_run(text); r.font.size=Pt(size); r.font.bold=bold; r.italic=italic; setcs(r)
    if color: r.font.color.rgb=color
    return p
def app_for(title):
    t=title.lower()
    if any(k in t for k in ['lift','pair','co-occur','co-occ','bond','attract','needle','haystack','company','partner','keeps']):
        return "Compare & Heatmaps - shared-ayah count AND lift between two roots"
    if any(k in t for k in ['form','voice','number','split','vocabulary','four']):
        return "Morphology / Surface Divergence - separate a root's surface forms by sense"
    if any(k in t for k in ['timeline','revelation','meccan','medinan']):
        return "Statistics - revelation-order (Meccan/Medinan) split"
    if any(k in t for k in ['sura','across the corpus','where','falls','distribution']):
        return "Per-Root Profile / Statistics - per-sura distribution of a root"
    if any(k in t for k in ['fdr','battery','survive']):
        return "Two Books -> FDR Summary - the live Benjamini-Hochberg dashboard"
    if any(k in t for k in ['escalat','inventory','rarity','context','length','field','frequency','principal','schedule']):
        return "Per-Root Profile / Statistics - frequency, fields and counts"
    return "Per-Root Profile / Statistics"

def make_guide(pptx):
    prs=Presentation(pptx); slides=list(prs.slides)
    base=os.path.basename(pptx)[len("SpecialTopic_"):-5]
    # topic main title + subtitle from slide 1
    main=""; sub=""
    for sh in slides[0].shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t=(para.text or "").strip()
                if "SPECIAL TOPIC" in t: sub=t.split("-",1)[-1].strip() if "-" in t else t
                elif t and len(t)>10 and not main: main=t
    figs=[]; verses=[]
    for i,s in enumerate(slides,1):
        has_pic=any(sh.shape_type==13 for sh in s.shapes)
        ttl=""
        for sh in s.shapes:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    t=(para.text or "").strip()
                    if t and not AR.search(t): ttl=t; break
                if ttl: break
        if has_pic and ttl: figs.append((i,ttl))
        for sh in s.shapes:
            if sh.has_text_frame:
                for m in ADDR.findall(sh.text_frame.text or ""): verses.append(m)
    # dedupe verses preserving order
    seen=set(); vlist=[v for v in verses if not (v in seen or seen.add(v))]
    d=newdoc()
    P(d,"Special Topic: "+(main or base)+"  -  App & Plot Guide",15,True,ACCENT,after=2)
    P(d,"Using the Quran Root Explorer app to drive this lecture, and how to read its figures. Every count recomputes from Book6.",10.5,False,GREY,after=8,italic=True)
    P(d,"The app is central to this lecture. Below: the live tasks to run, the screenshots to capture, and the map from each figure to the app feature that reproduces it - so the audience can re-derive every number live.",11,after=8)

    P(d,"Live app tasks (run during the lecture)",12.5,True,TEAL,after=3,before=6)
    if vlist:
        P(d,"Open Ayah Browser and pull this topic's cited verses, vocalized: "+", ".join(vlist[:10])+(" ..." if len(vlist)>10 else "")+".",10.5,style="List Bullet",after=2)
    P(d,"Open Per-Root Profile and search this topic's key root(s) -> read frequency, surface forms, and co-occurring roots.",10.5,style="List Bullet",after=2)
    if figs:
        P(d,"Reproduce the headline figure live: '"+figs[0][1]+"' -> "+app_for(figs[0][1])+".",10.5,style="List Bullet",after=2)
    P(d,"Open Morphology / Surface Divergence to separate the senses behind the key word (the concept-verification guard - §14a).",10.5,style="List Bullet",after=2)
    if any('lift' in t.lower() or 'pair' in t.lower() or 'co-occ' in t.lower() or 'company' in t.lower() for _,t in figs):
        P(d,"Open Compare & Heatmaps for the two roots in play -> read shared count AND lift together.",10.5,style="List Bullet",after=2)
    if any('fdr' in t.lower() or 'battery' in t.lower() for _,t in figs):
        P(d,"Open Two Books -> FDR Summary -> show the live permutation p-values behind the figure.",10.5,style="List Bullet",after=2)

    P(d,"Screenshot capture list",12.5,True,TEAL,after=3,before=8)
    P(d,"sshot-1: a key root's Per-Root Profile (frequency + forms + neighbours).",10.5,style="List Number",after=2)
    P(d,"sshot-2: the app view that reproduces the headline figure"+(": '"+figs[0][1]+"'." if figs else "."),10.5,style="List Number",after=2)
    P(d,"sshot-3: Morphology / Surface Divergence for the key word.",10.5,style="List Number",after=2)
    if vlist:
        P(d,"sshot-4: Ayah Browser - one cited verse ("+vlist[0]+"), vocalized.",10.5,style="List Number",after=2)
    P(d,"sshot-5: Statistics - the relevant distribution / totals.",10.5,style="List Number",after=2)

    P(d,"Plot -> slide map (this deck's "+str(len(figs))+" figures)",12.5,True,TEAL,after=3,before=8)
    for num,ttl in figs:
        P(d,"Slide "+str(num)+" - "+ttl+"   <->   "+app_for(ttl)+".",10.5,style="List Bullet",after=2)

    P(d,"Tip: capture screenshots at 150 dpi or higher; crop to the chart so it reads from the back of the room. Every figure recomputes from Book6 (the app live; the deck via the wk.py kernel) - roots normalized, surface forms used where a concept must be kept apart.",10.5,False,GREY,after=4,before=8,italic=True)
    P(d,"Honest spine: the app and the figures PRESENT the structure; the theological reading is labelled and never adjudicated.",10.5,True,None,after=4)
    out=os.path.join(OUTDIR,"SpecialTopic_"+base+"_App_and_Plot_Guide.docx")
    d.save(out); return out,len(figs),len(vlist)

decks=sorted(glob.glob(OUTDIR+"/SpecialTopic_*.pptx"))
n=0
for p in decks:
    out,nf,nv=make_guide(p); n+=1
print("per-lecture App & Plot Guides written:",n)
