# -*- coding: utf-8 -*-
"""Standard §17 deck-builder helpers for W-series Special Topics. 20-slide template,
>=50% chart/visual, quiz+key. Imports the verified Book6 figures from wk."""
import os, string
from st_slides import *
from diagrams import fbox,band,harrow,chain
from docx import Document
from docx.shared import Pt, RGBColor
from docx.oxml.ns import qn
FIGDIR="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics/figs_w"
OUTDIR="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics"
ACCENT=RGBColor(0x1F,0x4E,0x79); DGREY=RGBColor(0x55,0x55,0x55)

def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def embed(s,ttl,png,cap,cf=TINT2):
    from pptx.util import Inches
    title(s,ttl); w=11.2
    s.shapes.add_picture(os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cf,[L(cap,15,True,TEAL)],space=4)
def gallery(s,ttl,items,fill=TINT,hc=TEAL,intro="the text's own words (address - snippet - gloss):"):
    title(s,ttl); lines=[L(intro,15,True,hc)]
    for ad,ar,gl in items: lines.append(L(ad+"   "+ar+"   -  "+gl,15,False,INK))
    panel(s,0.42,1.15,12.5,6.05,fill,lines,space=7)
def auditgrid(s,ttl,cells):
    """cells = list of (mark,head,body) up to 6; mark in check/tilde/cross"""
    title(s,ttl); g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
    pos=[(0.42,CY0),(0.42+w+g,CY0),(0.42,CY0+h+g),(0.42+w+g,CY0+h+g),(0.42,CY0+2*(h+g)),(0.42+w+g,CY0+2*(h+g))]
    for (mk,hd,bd),(x,y) in zip(cells,pos):
        fill=TINT if mk=="check" else (REDT if mk=="cross" else AMBERT)
        sym="✓" if mk=="check" else ("✗" if mk=="cross" else "~")
        col=TEAL if mk=="check" else (RED if mk=="cross" else AMBER)
        panel(s,x,y,w,h,fill,[L(sym+"  "+hd,15,True,col),L(bd,13.5)],space=4)
def methodslide(s,root_desc,mid_desc,out_desc):
    title(s,"How this was computed (reproducible)")
    fbox(s,0.7,1.95,3.7,1.5,TINT,"Book6",root_desc,line=TEAL,tsz=15,ssz=11)
    harrow(s,4.5,2.6,0.5,"count",color=GREY,lcol=TEAL)
    fbox(s,5.4,1.95,3.7,1.5,AMBERT,"measure",mid_desc,line=AMBER,tsz=15,ssz=11)
    harrow(s,9.25,2.6,0.5,"chart",color=GREY,lcol=RED)
    fbox(s,10.2,1.95,2.5,1.5,TINT2,"figures",out_desc,line=NAVY,tsz=15,ssz=11)
    panel(s,0.42,3.8,12.5,3.4,TINT,[L("Fully reproducible",18,True,NAVY),
      L("Every figure recomputes from Book6.xlsx via the shared wk.py kernel (fixed seed). Roots are "
        "normalized (Persian/Arabic letter variants folded) so counts are stable; surface forms are used "
        "wherever a concept must be kept apart from its root family (per the course's sense-verification rule).",16,True,TEAL)],space=9)

# ---- quiz builder ----
def _setcs(r):
    rPr=r._r.get_or_add_rPr(); rf=rPr.get_or_add_rFonts()
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a),"Arial")
def _newdoc():
    d=Document(); st=d.styles["Normal"]; st.font.name="Arial"; st.font.size=Pt(11)
    rf=st.element.get_or_add_rPr().get_or_add_rFonts()
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a),"Arial")
    return d
def _P(d,segs,size=11,bold=False,color=None,after=4,before=0):
    p=d.add_paragraph(); p.paragraph_format.space_after=Pt(after); p.paragraph_format.space_before=Pt(before)
    if isinstance(segs,str): segs=[(segs,bold)]
    for t,b in segs:
        r=p.add_run(t); r.font.size=Pt(size); r.font.bold=b; _setcs(r)
        if color: r.font.color.rgb=color
    return d
def build_quiz(slug,qtitle,QQ):
    """QQ = [(stem,correct,[distractors],explanation)]; answers rotate A-D."""
    d=_newdoc(); _P(d,[(qtitle+" - Quiz",True)],size=16,color=ACCENT,after=2)
    _P(d,str(len(QQ))+" questions - choose the single best answer. Every value reproducible from Book6.",size=9.5,color=DGREY,after=8)
    KEY=[]
    for qi,(stem,correct,distr,expl) in enumerate(QQ):
        pos=qi%4; opts=list(distr); opts.insert(pos,correct)
        _P(d,[(stem,True)],size=10.5,after=2,before=6)
        for i,o in enumerate(opts): _P(d,string.ascii_uppercase[i]+")  "+o,size=10,after=1)
        KEY.append((qi+1,string.ascii_uppercase[pos],expl))
    d.save(os.path.join(OUTDIR,"SpecialTopic_"+slug+"_Quiz.docx"))
    d=_newdoc(); _P(d,[(qtitle+" - Quiz Answer Key",True)],size=16,color=ACCENT,after=2)
    _P(d,"One point each, "+str(len(QQ))+" total. Every value reproducible from Book6.",size=9.5,color=DGREY,after=8)
    for n,a,ex in KEY: _P(d,[(str(n)+".  "+a+"  ",True),("- "+ex,False)],size=10,after=2)
    d.save(os.path.join(OUTDIR,"SpecialTopic_"+slug+"_Quiz_Answer_Key.docx"))
    return [a for _,a,_ in KEY]
def verify(path):
    from pptx import Presentation
    import re
    prs=Presentation(path); slides=list(prs.slides)
    AR=re.compile(r'[؀-ۿ]'); pics=0; bad=0; arr=0
    for s in slides:
        if any(sh.shape_type==13 for sh in s.shapes): pics+=1
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if AR.search(r.text or ''):
                        arr+=1
                        rpr=r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                        cs=rpr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}cs') if rpr is not None else None
                        if not (cs is not None and cs.get('typeface')): bad+=1
    return len(slides),pics,arr,bad

def standard_deck(spec):
    """Build a 20-slide §17 Special-Topic deck from a spec dict and save it.
    Required keys:
      slug, main, sub, headline, intro1, intro2,
      qhead, qbody, mhead, mpts(list of str),
      figs: list of >=7 dicts {png,cap,cf?}
      gal1,gal2: {title,items:[(ad,ar,gl)],fill?,hc?}
      v1:(head,body), v2:(head,body), v3:(head,body)   # three verdict blocks
      deep:(head,body),
      crit1:(head,body), crit2:(head,body),
      audit: list of (mark,head,body) len 6,
      method:(root_desc,mid_desc,out_desc),
      take:(head,[lines]),
      qr1:(head,body), qr2:(head,body),
      quiz:(qtitle,QQ)
    """
    from pptx import Presentation
    prs=deck()
    F=spec["figs"]
    def cf(i): return F[i].get("cf",TINT2)
    # 1 title
    s=slide(prs)
    panel(s,0.42,1.2,12.5,1.7,TINT2,[L("SPECIAL TOPIC  -  "+spec["sub"],16,True,TEAL),
      L(spec["main"],22,True,NAVY)],space=7)
    panel(s,0.42,3.1,12.5,4.0,TINT,[L(spec["headline"],18,True,NAVY),
      L(spec["intro1"],16.5),L(spec["intro2"],15.5,True,TEAL)],space=9)
    # 2 question + method
    s=slide(prs); title(s,spec.get("s2title","The question, and the method"))
    two(s,[L(spec["qhead"],17,True,NAVY),L(spec["qbody"],16)],
         [L(spec["mhead"],17,True,TEAL)]+[L("- "+p,15) for p in spec["mpts"]],sp=0.42,fa=TINT2,fb=TINT)
    # 3 fig0
    s=slide(prs); embed(s,F[0]["t"],F[0]["png"],F[0]["cap"],cf(0))
    # 4 gallery1
    g=spec["gal1"]; s=slide(prs); gallery(s,g["title"],g["items"],g.get("fill",TINT),g.get("hc",TEAL))
    # 5 fig1
    s=slide(prs); embed(s,F[1]["t"],F[1]["png"],F[1]["cap"],cf(1))
    # 6 verdict three_stack
    v1,v2,v3=spec["v1"],spec["v2"],spec["v3"]
    s=slide(prs); title(s,spec.get("s6title","What the data shows"))
    three(s,[L(v1[0],16,True,TEAL),L(v1[1],15)],[L(v2[0],16,True,AMBER),L(v2[1],15)],[L(v3[0],16,True,NAVY),L(v3[1],15)])
    # 7 fig2
    s=slide(prs); embed(s,F[2]["t"],F[2]["png"],F[2]["cap"],cf(2))
    # 8 gallery2
    g=spec["gal2"]; s=slide(prs); gallery(s,g["title"],g["items"],g.get("fill",AMBERT),g.get("hc",AMBER))
    # 9 fig3
    s=slide(prs); embed(s,F[3]["t"],F[3]["png"],F[3]["cap"],cf(3))
    # 10 fig4
    s=slide(prs); embed(s,F[4]["t"],F[4]["png"],F[4]["cap"],cf(4))
    # 11 deep panel
    s=slide(prs); title(s,spec["deep"][0])
    panel(s,0.42,1.2,12.5,5.9,TINT2,[L(spec["deep"][1],17,False,INK)]+
          [L(x,16,True,TEAL) for x in spec.get("deep_extra",[])],space=10)
    # 12 fig5
    s=slide(prs); embed(s,F[5]["t"],F[5]["png"],F[5]["cap"],cf(5))
    # 13 fig6
    s=slide(prs); embed(s,F[6]["t"],F[6]["png"],F[6]["cap"],cf(6))
    # 14 synthesis chain
    s=slide(prs); title(s,spec["syn"][0])
    ch=spec["syn"][1]  # list of up to 3 (head,sub)
    xs=[0.55,4.95,9.35]; ws=[3.7,3.7,3.3]; fills=[TINT,AMBERT,TINT2]; lines=[TEAL,AMBER,NAVY]
    for i,(hh,ss) in enumerate(ch):
        fbox(s,xs[i],1.95,ws[i],1.55,fills[i],hh,ss,line=lines[i],tsz=15,ssz=11)
        if i<len(ch)-1: harrow(s,xs[i]+ws[i]+0.05,2.62,0.35,"",color=GREY)
    panel(s,0.42,3.85,12.5,3.35,TINT,[L(spec["syn"][2],18,True,NAVY),L(spec["syn"][3],16,True,TEAL)],space=9)
    # 15 critique
    s=slide(prs); title(s,spec.get("s15title","Honest limits and critique"))
    two(s,[L(spec["crit1"][0],17,True,RED),L(spec["crit1"][1],16,True,NAVY)],
         [L(spec["crit2"][0],17,True,NAVY),L(spec["crit2"][1],16,True,NAVY)],sp=0.5,fa=REDT,fb=TINT2)
    # 16 audit
    s=slide(prs); auditgrid(s,spec.get("s16title","Reading this honestly - the audit"),spec["audit"])
    # 17 method
    s=slide(prs); methodslide(s,*spec["method"])
    # 18 appendix snippets
    g=spec.get("gal3")
    s=slide(prs)
    if g: gallery(s,g["title"],g["items"],g.get("fill",TINT2),g.get("hc",NAVY))
    else:
        title(s,"Appendix - input snippets (Book6, vocalized)")
        panel(s,0.42,1.15,12.5,6.05,TINT2,[L("verbatim snippets with addresses:",15,True,NAVY)]+
              [L(ad+"   "+ar+"   -  "+gl,15) for ad,ar,gl in spec["gal1"]["items"][:6]],space=7)
    # 19 takeaway
    s=slide(prs); title(s,"Takeaway")
    panel(s,0.42,1.2,12.5,5.9,TINT2,[L(spec["take"][0],18,True,NAVY)]+
          [L(x,16,(i==len(spec["take"][1])-1),(TEAL if i==len(spec["take"][1])-1 else INK)) for i,x in enumerate(spec["take"][1])],space=10)
    # 20 quickref
    s=slide(prs); title(s,"Quick reference")
    two(s,[L(spec["qr1"][0],17,True,TEAL),L(spec["qr1"][1],16)],
         [L(spec["qr2"][0],17,True,AMBER),L(spec["qr2"][1],16)],sp=0.5,fa=TINT,fb=AMBERT)
    path=os.path.join(OUTDIR,"SpecialTopic_"+spec["slug"]+".pptx")
    prs.save(path)
    qt,QQ=spec["quiz"]; letters=build_quiz(spec["slug"],qt,QQ)
    n,pics,arr,bad=verify(path)
    print(f'{spec["slug"]}: slides={n} pics={pics} arabicRuns={arr} cs_missing={bad} quiz={"".join(letters)}')
    return path
