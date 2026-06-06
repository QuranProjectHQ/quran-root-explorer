# -*- coding: utf-8 -*-
"""Concise instructor scripts (~11 min) for each W-series Special Topic deck.
Talking points are drawn from each built deck's own slide titles and lead lines,
so the script always matches the slides. Style mirrors the Two Books lecture kits."""
import os,glob,re
from pptx import Presentation
from docx import Document
from docx.shared import Pt, RGBColor
from docx.oxml.ns import qn
OUTDIR="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/SpecialTopics"
ACCENT=RGBColor(0x1F,0x4E,0x79); GREY=RGBColor(0x55,0x55,0x55); TEAL=RGBColor(0x0E,0x6D,0x63)
AR=re.compile(r'[؀-ۿ]')
def setcs(r):
    rPr=r._r.get_or_add_rPr(); rf=rPr.get_or_add_rFonts()
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a),"Arial")
def newdoc():
    d=Document(); st=d.styles["Normal"]; st.font.name="Arial"; st.font.size=Pt(11)
    rf=st.element.get_or_add_rPr().get_or_add_rFonts()
    for a in ("w:ascii","w:hAnsi","w:cs","w:eastAsia"): rf.set(qn(a),"Arial")
    return d
def P(d,text,size=11,bold=False,color=None,after=4,before=0,italic=False):
    p=d.add_paragraph(); p.paragraph_format.space_after=Pt(after); p.paragraph_format.space_before=Pt(before)
    r=p.add_run(text); r.font.size=Pt(size); r.font.bold=bold; r.italic=italic; setcs(r)
    if color: r.font.color.rgb=color
    return p

def slide_info(s):
    """Return (title, lead) for a slide: title = first short line, lead = first long non-arabic line."""
    paras=[]
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        for para in sh.text_frame.paragraphs:
            t=(para.text or "").strip()
            if t: paras.append(t)
    if not paras: return ("","")
    title=paras[0]
    lead=""
    for t in paras[1:]:
        if AR.search(t): continue            # skip Arabic verse lines
        if len(t)>40: lead=t; break
    if not lead:
        for t in paras[1:]:
            if not AR.search(t): lead=t; break
    return (title,lead)

# beat plan: (label, slide range inclusive 1-based, minute)
BEATS=[("Frame - the question and method",(1,2),"0:00"),
       ("The opening finding",(3,5),"1:30"),
       ("What the data shows",(6,8),"3:30"),
       ("Going deeper",(9,11),"5:30"),
       ("Distribution and timeline",(12,13),"7:30"),
       ("Synthesis, limits, and the audit",(14,16),"8:30"),
       ("Method, takeaway, quick reference",(17,20),"10:00")]

def make_script(pptx):
    prs=Presentation(pptx); slides=list(prs.slides)
    info=[slide_info(s) for s in slides]
    base=os.path.basename(pptx)[len("SpecialTopic_"):-5]
    # topic title = slide 1's second line (the main title) if present
    main=""
    for sh in slides[0].shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t=(para.text or "").strip()
                if t and "SPECIAL TOPIC" not in t and len(t)>10 and not main: main=t
    d=newdoc()
    P(d,"Special Topic - "+(main or base)+"  -  Instructor Script",16,True,ACCENT,after=2)
    P(d,"Spoken script, ~11 minutes, mapped to the 20-slide deck. Honest spine: present the data, do not adjudicate. Every figure is live from Book6.",9.5,False,GREY,after=8,italic=True)
    for label,(a,b),mins in BEATS:
        P(d,mins+"   "+label,12,True,TEAL,after=2,before=8)
        P(d,"> Slides "+str(a)+("-"+str(b) if b>a else "")+".",10,False,GREY,after=3,italic=True)
        # assemble a talking paragraph from the slides in range
        cues=[]
        for i in range(a-1,min(b,len(info))):
            ttl,lead=info[i]
            ttl=re.sub(r'\s+',' ',ttl).strip()
            lead=re.sub(r'\s+',' ',lead).strip()
            if lead: cues.append(f"{ttl} - {lead}")
            elif ttl: cues.append(ttl)
        para="  ".join(cues)
        P(d,para,11,False,None,after=4)
    P(d,"Close: restate the honest spine - the corpus shows the structure; the reading is labelled, never adjudicated. Point to the quick-reference slide and the companion quiz.",10.5,True,None,after=4,before=8)
    out=os.path.join(OUTDIR,"SpecialTopic_"+base+"_Instructor_Script.docx")
    d.save(out); return out

decks=sorted(glob.glob(OUTDIR+"/SpecialTopic_W*.pptx"))
n=0
for p in decks:
    make_script(p); n+=1
print("instructor scripts written:",n)
