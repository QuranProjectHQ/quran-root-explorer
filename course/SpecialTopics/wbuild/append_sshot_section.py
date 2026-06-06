# -*- coding: utf-8 -*-
import glob, os, re
from pptx import Presentation
ROOT = "/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse"
SP = os.path.join(ROOT, "SpecialTopics")
AR = re.compile(r'[؀-ۿ]')
ADDR = re.compile(r'\b(\d{1,3}:\d{1,3}(?:-\d{1,3})?)\b')
decks = sorted(glob.glob(os.path.join(SP, "SpecialTopic_*.pptx")))
rows = []
for p in decks:
    slug = os.path.basename(p)[len("SpecialTopic_"):-5]
    prs = Presentation(p); slides = list(prs.slides)
    main = ""
    for sh in slides[0].shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = (para.text or "").strip()
                if t and "SPECIAL TOPIC" not in t and len(t) > 8 and not main:
                    main = t
    hfig = ""
    for s in slides:
        if any(sh.shape_type == 13 for sh in s.shapes):
            for sh in s.shapes:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        t = (para.text or "").strip()
                        if t and not AR.search(t):
                            hfig = t; break
                if hfig: break
        if hfig: break
    v = ""
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                m = ADDR.findall(sh.text_frame.text or "")
                if m: v = m[0]; break
        if v: break
    rows.append((slug, main, hfig, v))

L = []
L.append("")
L.append("---")
L.append("")
L.append("## SPECIAL TOPICS (27) - capture set")
L.append("")
L.append("Each Special Topic has its own App & Plot Guide (`SpecialTopic_<slug>_App_and_Plot_Guide.docx`) "
         "listing its sshot-1 to sshot-5. Save captures to `SpecialTopics/shots/<slug>/`. Standard 3-shot "
         "minimum per topic: (a) the key root's Per-Root Profile, (b) the app view reproducing the headline "
         "figure, (c) one cited verse in Ayah Browser. The table gives the headline figure to reproduce and "
         "a suggested first-verse capture per topic.")
L.append("")
L.append("| Topic | Headline figure to reproduce | First cited verse | Suggested folder |")
L.append("|---|---|---|---|")
for slug, main, hfig, v in rows:
    folder = "`SpecialTopics/shots/%s/`" % slug
    L.append("| %s | %s | %s | %s |" % (main[:46], hfig[:50], v or "-", folder))
L.append("")
L.append("Tip: every figure also recomputes from Book6 (each deck via `wbuild/wk.py`, fixed seed); these "
         "screenshots are the live-app counterpart for teaching, not a substitute for the verified charts "
         "already embedded in the decks.")
L.append("")
with open(os.path.join(ROOT, "SCREENSHOT_CAPTURE_GUIDE.md"), "a", encoding="utf-8") as f:
    f.write("\n".join(L))
print("appended", len(rows), "topic rows")
