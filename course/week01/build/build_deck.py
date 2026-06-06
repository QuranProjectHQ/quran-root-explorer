# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
WK="/sessions/kind-compassionate-feynman/mnt/RootCourse/week01"
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B); GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); SH=prs.slide_height
blank=prs.slide_layouts[6]
def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def slide():
    s=prs.slides.add_slide(blank); b=s.shapes.add_shape(1,0,0,Inches(0.16),SH); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background(); return s
def title(s,t,sz=30):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.28),Inches(12.5),Inches(0.95)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia"); return s
def body(s,lines,top=1.3,size=19,width=12.3,left=0.65,space=8):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(7.5-top-0.2)); tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,bold,col,sz) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space)
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz or size); r.font.bold=bold; r.font.color.rgb=col or INK; sf(r,"Calibri")
def pic(s,name,left,top,width=None,height=None):
    p=os.path.join(WK,name)
    if os.path.exists(p):
        kw={}
        if width:kw['width']=Inches(width)
        if height:kw['height']=Inches(height)
        s.shapes.add_picture(p,Inches(left),Inches(top),**kw)
def B(t,bold=False,col=None,sz=None): return (t,bold,col,sz)

s=slide()
tb=s.shapes.add_textbox(Inches(0.6),Inches(2.0),Inches(12),Inches(3.4)); tf=tb.text_frame; tf.word_wrap=True
for txt,sz,bold,col,fn in [("WEEK 1",22,True,TEAL,"Georgia"),("Frequency",46,True,NAVY,"Georgia"),
  ("How often does a root appear — and how honestly can we read that?",21,False,GREY,"Calibri"),
  ("The simplest, most objective layer.  Worked roots: ظلم vs عدل.",19,False,GREY,"Calibri"),
  ("Every value computed from Book6 (6,236 ayahs · ~1,700 roots).",15,False,GREY,"Calibri")]:
    p=tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]; p.space_after=Pt(8); r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,fn)

s=slide(); title(s,"The frame & today's question")
body(s,[B("If I asked which idea matters most in the Qur'an, you'd answer from memory —",False),
 B("and we'd have no way to check each other.",False,GREY),
 B("So we take a smaller, humbler question we CAN settle:",True,NAVY),
 B("not what matters most, but how often a root appears.",True,NAVY),
 B("The course's one rule, kept all term:",True,TEAL,20),
 B("state a computed FACT, then a separately labeled INTERPRETATION — never blur them.",False),
 B("We measure not to replace reading, but to make claims about the text AUDITABLE:",False),
 B("a count can be wrong, but it can be checked.",False,GREY),
 B("Today we claim presence only — not importance, not meaning, not relationships.",False)],top=1.4,size=18,space=9)

s=slide(); title(s,"What frequency analysis is")
body(s,[B("Counting how often a meaning-bearing unit occurs in a text —",True,NAVY,20),
 B("the word-frequency, or “bag-of-words,” tradition in corpus linguistics.",False),
 B("Temporarily set aside order and grammar; ask only: what is present, and how often?",False),
 B("Like emptying a book into a bag and counting the words —",False,GREY),
 B("you lose the sentences, but you can finally tally what's there.",False,GREY),
 B("The working assumption is modest:",True,TEAL,20),
 B("how often a text names something is a SIGNAL of emphasis — a trace, not proof.",False),
 B("We start here because it's the most objective, reproducible layer;",False),
 B("everything later — distribution, partners, co-occurrence — builds on these counts.",False)],top=1.4,size=18,space=9)

s=slide(); title(s,"From text to countable units")
body(s,[B("The same idea wears many clothes; we want to count the idea, not the spelling.",True,NAVY,20),
 B("The pipeline:",True),
 B("•  Tokenize each ayah into words.",False),
 B("•  Reduce each word to its 3-consonant ROOT: كَتَبَ / يَكْتُبُ / كِتاب → كتب.",False),
 B("•  Normalize letters & diacritics so the same root matches everywhere.",False),
 B("•  Drop function words (في، من، الذي) — they carry grammar, not theme.",False),
 B("Count once per ayah (document frequency) — so one long verse can't dominate",False),
 B("by repeating a root.",False,GREY),
 B("The reduction is lossy ON PURPOSE — we trade nuance for comparability.",True,TEAL),
 B("After it, the corpus holds 51,024 root-tokens.",False)],top=1.4,size=18,space=8)

s=slide(); title(s,"Normalization to a rate")
body(s,[B("A bigger text yields bigger counts regardless of emphasis —",True,NAVY,20),
 B("so raw counts can't be compared. Convert to a RATE.",False),
 B("•  per 1,000 AYAHS = ayah-freq ÷ 6,236 × 1,000  (share of verses).",False),
 B("•  per 1,000 ROOTS = term-freq ÷ 51,024 × 1,000  (size-true).",False,TEAL),
 B("Why two? The denominator can flip the ranking — that's not cosmetic:",True,NAVY),
 B("ظلم > هدي per ayah (46.5 vs 43.0), but هدي > ظلم per root (6.19 vs 6.17).",False),
 B("Per-1,000-roots is like deaths-per-100,000 instead of raw deaths —",False,GREY),
 B("the only fair way to compare populations of different sizes.",False,GREY),
 B("Always report both rates: presence is a RATE, not a count.",True,TEAL)],top=1.4,size=18,space=8)

s=slide(); title(s,"The themed-root frequency ladder")
pic(s,"fig_freq_ladder.png",7.1,1.5,width=6.1)
body(s,[B("Rank a themed set of roots by rate — the shape of the vocabulary appears.",True,NAVY,19),
 B("Three descriptive bands:",True),
 B("•  pervasive (≥25/1k): ظلم، نفس، هدي، ضلل",False),
 B("•  mid (5–25): رزق، صبر، شكر",False),
 B("•  rare (<5): عدل، قسط، عسر",False),
 B("The bands orient attention — they are not a law of the text.",False,GREY),
 B("Live in the app: كفر → 465 ayahs (the most-named);  عسر → 12 (the rarest).",False),
 B("Always pair a rate with its raw count and its size-true rate.",True,TEAL)],top=1.5,size=17.5,width=6.6,space=6)

s=slide(); title(s,"Predict, then check — ease vs hardship")
body(s,[B("Before we look — which is named more: ease (يسر) or hardship (عسر)?",True,NAVY,22),
 B("Most rooms guess hardship.",False,GREY),
 B("Check it:",True),
 B("يسر — 40 ayahs.",True,TEAL,28),
 B("عسر — 12 ayahs.",True,RED,28),
 B("About 3.3× more “ease.”",True,NAVY,22),
 B("The habit you just watched: pair the rate with the raw count BEFORE you speak.",False,GREY),
 B("And one number is about to unsettle a far more common assumption…",False)],top=1.45,size=19,space=11)

s=slide(); title(s,"Live in the app — the Per-Root Profile")
body(s,[B("Type a root, press Analyze, open the Per-Root Profile. For ظلم you get:",True,NAVY,20),
 B("•  290 ayahs matched.",False),
 B("•  98.59th percentile — more frequent than ~98.6% of all roots.",False),
 B("•  tier: “ubiquitous” (top 1%).",False),
 B("•  first occurrence 2:17 (al-Baqara); last 76:31 (al-Insan).",False),
 B("Compute the two rates right there:",True,TEAL,20),
 B("46.5 per 1,000 ayahs (share of verses)  ·  6.17 per 1,000 roots (size-true).",False),
 B("Reading habit: always pair the percentile with the raw count and BOTH rates",False),
 B("before you say a single word about the root.",False,GREY)],top=1.45,size=18,space=8)

s=slide(); title(s,"The headline finding & the “unlearn”")
body(s,[B("Predict with me: which does the Qur'an name more —",True,NAVY,22),
 B("عدل (justice) or ظلم (injustice)?",True,NAVY,22),
 B("Most rooms say justice. The data says otherwise:",False,GREY),
 B("ظلم — 290 ayahs   (98.6th percentile, “ubiquitous”).",True,RED,24),
 B("عدل — 24 ayahs.",True,TEAL,24),
 B("About 12× more.",True,NAVY,24),
 B("A genuine unlearning: the corpus names the VIOLATION far more than the IDEAL.",True),
 B("But stay disciplined — this is presence, not endorsement: ظلم is named to condemn it.",False,GREY)],top=1.45,size=19,space=10)

s=slide(); title(s,"Frequency as a vector  [preview]")
body(s,[B("Here is the idea that makes today bigger than it looks.",True,NAVY,20),
 B("List a root's counts across contexts — that ordered list of numbers is a VECTOR.",False),
 B("Do it for every root, and each root becomes a POINT in a space;",False),
 B("roots with similar profiles sit close together.",False),
 B("That is the VECTOR SPACE MODEL — the foundation of how machines “read” text.",True,TEAL,20),
 B("Everything later — partners, co-occurrence, networks — is geometry on these vectors.",False),
 B("Nothing new is computed here: the very counts we tabulated are the coordinates.",False,GREY),
 B("Labeled preview — embeddings & skip-gram are further-study, never asserted against Book6.",False,GREY,15)],top=1.45,size=18,space=9)

s=slide(); title(s,"What frequency gives — and what it loses")
pic(s,"fig1_frequency.png",7.2,1.7,width=6.0)
body(s,[B("GIVES:",True,NAVY,20),
 B("•  objective & reproducible — anyone re-running gets the same numbers.",False),
 B("•  scalable to all 6,236 ayahs.",False),
 B("•  comparable, via rates.",False),
 B("•  scoping — shows WHERE to look deeper.",False),
 B("LOSES — context, above all:",True,RED,20),
 B("•  stance/polarity (كفر counted whether affirmed or refuted).",False),
 B("•  polysemy (ظلم = injustice AND ظلمات = darkness, as one).",False),
 B("•  speaker, syntax, and all relationships between roots.",False),
 B("Frequency SCOPES inquiry; it does not SETTLE it.",True,TEAL)],top=1.5,size=17,width=6.7,space=6)

s=slide(); title(s,"Fact vs. interpretation — the discipline")
body(s,[B("Separate the computed fact from any reading of it.",True,NAVY,20),
 B("One sentence of fact (with its normalization), one labeled sentence of interpretation.",False),
 B("FACT",True,NAVY,22),
 B("“ظلم is named in 290 ayahs (46.5/1k ayahs, 6.17/1k roots) — about 12× عدل.”",False),
 B("INTERPRETATION",True,TEAL,22),
 B("“I read this as the text foregrounding the diagnosis of wrong over the abstract ideal.”",False),
 B("Peer test: can a classmate tell which sentence is fact and which is interpretation?",True,NAVY),
 B("If not — rewrite until they separate cleanly.",False,GREY)],top=1.5,size=18,space=10)

s=slide(); title(s,"Worked example — al-Fatiha (by hand)")
pic(s,"fig_fatiha.png",7.2,1.7,width=6.0)
body(s,[B("Compute frequency by hand on 7 ayahs:",True,NAVY,20),
 B("tokenize → roots → count by ayah, then by token.",False),
 B("The teaching point — repetition:",True,TEAL,20),
 B("رحم has ayah-frequency 2 but term-frequency 4 —",False),
 B("the only root uttered four times.",False),
 B("Ayah-counting hides repetition inside a verse;",False,GREY),
 B("term-counting catches it.",False,GREY),
 B("Then: one fact (“رحم, ayah-freq 2, term-freq 4”)",False),
 B("+ one labeled interpretation.",False)],top=1.5,size=18,width=6.7,space=7)

s=slide(); title(s,"This week's exercise")
body(s,[B("Part 1 — by hand",True,TEAL,22),
 B("•  Count the roots of your assigned short surah, by ayah and by token.",False),
 B("•  Flag any root whose term-frequency exceeds its ayah-frequency.",False),
 B("Part 2 — in the app",True,TEAL,22),
 B("•  Profile your assigned root; record BOTH rates (per 1k ayahs AND per 1k roots).",False),
 B("•  Note percentile and tier; take one screenshot.",False),
 B("Submit one screenshot + two readings (one fact + one interpretation each),",True,NAVY),
 B("the night before class — it gates the debrief.",False)],top=1.6,size=19,space=10)

s=slide(); title(s,"Presence is a rate, not a count — and it is context-blind")
body(s,[B("Frequency is the most objective layer — but the most context-blind.",True,NAVY,22),
 B("It can defy intuition (ظلم 12× عدل) — yet it never tells you why, who, or how.",False,GREY,20),
 B("",False),
 B("Next week — distribution & concentration:",True,NAVY,20),
 B("frequency says HOW MUCH; next we ask WHERE a root sits and how it spreads.",False,GREY,20)],top=1.9,size=20,space=12)
prs.save(os.path.join(WK,"Week1_Slides.pptx")); print("Week1 deck rebuilt (dense):",len(prs.slides._sldIdLst),"slides")
