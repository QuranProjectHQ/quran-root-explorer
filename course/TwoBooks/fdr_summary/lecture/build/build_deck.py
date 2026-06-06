# -*- coding: utf-8 -*-
"""Two Books · FDR Summary capstone deck — >=20 slides, >=half visual."""
import os, sys, math
HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
from st_slides import *
from diagrams import fbox,harrow,vdash,band,matgrid,chain
LEC=os.path.dirname(HERE)
OUT=os.path.join(LEC,"FDR_Summary_Lecture_Slides.pptx")

# the live battery (Book6) — name, domain, raw p, BH q, survives
BAT=[
 ("Contiguity · muṣḥaf","Position",0.0005,0.0010,True),
 ("Contiguity · nuzūl","Position",0.0005,0.0010,True),
 ("Length autocorr.","Signal",0.0005,0.0010,True),
 ("Root entropy","Semantic",0.0005,0.0010,True),
 ("Letter entropy","Sequence",0.002,0.0032,True),
 ("Di-codon adj.","Biology",0.005,0.0067,True),
 ("Shared theme","Semantic",0.049,0.056,False),
 ("Shared length","Position",0.289,0.289,False),
]
def nlp(p): return min(4.0,-math.log10(max(p,1e-4)))
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def audit_row(s,x,y,w,h,mark,mc,head,body):
    panel(s,x,y,w,h,(TINT if mark=="✓" else REDT if mark=="✗" else AMBERT),
          [L(mark+"  "+head,15,True,mc),L(body,14)],space=5)

prs=deck()

FIGDIR=os.path.join(LEC,"figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.20,12.5,1.65,TINT2,[L("TWO BOOKS  ·  FDR SUMMARY — capstone lecture deck",15,True,TEAL),
      L("One Benjamini–Hochberg correction across every domain",22,True,NAVY)],space=7)
panel(s,0.42,3.05,12.5,2.05,TINT,[L("The honest spine",17,True,NAVY),
      L("Collect one representative permutation test from each Two Books domain and correct the whole "
        "battery together. After correction the cross-domain structure is ROBUST — 6 of 8 tests survive "
        "5% FDR — while the borderline shared-theme claim correctly drops out. Crucial caveat: FDR "
        "controls for MULTIPLICITY, not for the sūra-length confound.",15.5,True,TEAL)],space=8)
panel(s,0.42,5.25,12.5,1.95,TINT2,[L("Why this is the capstone",16,True,NAVY),
      L("Disjoint Letters, Signal and Biology each run many tests. Read alone, the best p from a large "
        "battery is how chance masquerades as discovery. This deck reads them ALL together — every q "
        "computed live from Book6.",15.5)],space=7)

# 1b pipeline (diagram)
s=slide(prs); title(s,"How the dashboard is built")
x=0.55; bw=2.65; aw=0.18
for i,(lab,sub,col) in enumerate([("Each page","DL · Signal · Biology",TEAL),
        ("One test\nper domain","8 representative p's",AMBER),
        ("Benjamini–\nHochberg","rank & correct",NAVY),
        ("Dashboard","q-values · survivors",TEAL)]):
    fbox(s,x,2.1,bw,1.5,TINT,lab,sub,line=col,tsz=15,ssz=12)
    if i<3: harrow(s,x+bw,2.6,aw+0.02,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.1,12.5,3.1,TINT2,[L("One pipeline, one correction",18,True,NAVY),
      L("Each Two Books page contributes one representative permutation test; the shared stats kernel "
        "assembles the eight p-values; a single Benjamini–Hochberg pass corrects them together; the "
        "dashboard reports q-values and survivors. No page logic is duplicated.",16.5,True,TEAL)],space=9)

# 2 the multiplicity problem (diagram)
s=slide(prs); title(s,"The problem — run many tests, some pass by luck")
fbox(s,0.7,1.95,5.6,1.5,AMBERT,"At α = 0.05","1 in 20 NULL tests 'passes' by chance alone.",line=AMBER,tsz=15)
fbox(s,6.9,1.95,5.6,1.5,REDT,"Run a battery","With 8 tests, the chance of ≥1 false hit is far above 5%.",line=RED,tsz=15)
panel(s,0.42,3.75,12.5,3.4,TINT2,[L("Why an uncorrected 0.05 is unsafe across a battery",18,True,NAVY),
      L("Every Two Books page runs several permutation tests. Cherry-pick the smallest p and you will "
        "almost always find one under 0.05 — even if nothing real is there. That is the engine of "
        "pseudo-miracles.",16.5),
      L("The shared-theme test lands at p = 0.049 — the textbook borderline result multiplicity should "
        "make us distrust until it is corrected.",16,True,TEAL)],space=9)

# 3 Benjamini–Hochberg (diagram)
s=slide(prs); title(s,"The fix — Benjamini–Hochberg controls the false-discovery rate")
chain(s,1.0,2.0,["p₁","p₂","p₃","…","pₘ"],col=TEAL)
panel(s,0.42,3.2,12.5,3.9,TINT2,[L("How it works",18,True,NAVY),
      L("Rank the m p-values ascending. The i-th passes if it is below (i ÷ m) × α. The q-value is the "
        "smallest α at which a test still passes — a number directly comparable to the 5% threshold.",16.5),
      L("FDR is the right control for an exploratory battery: less brutal than Bonferroni, yet it still "
        "guards against multiplicity. With 8 tests at α = 0.05, a raw p of 0.0005 becomes q = 0.0010, "
        "while p = 0.049 becomes q = 0.056 — now just over the line.",16,True,TEAL)],space=9)

# 3b BH ranking procedure (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 3 — The ranking, made concrete","fdr_staircase.png",
  "In the data","sorted p-values vs the (i/m)·α line. Ranks 1–6 sit below their thresholds and pass; rank 7 (theme, p = 0.049) crosses above and fails. Six of eight survive — exactly the q-values.")

# 4 scope
s=slide(prs); title(s,"What this deck will — and will not — do")
two(s,[L("WILL",18,True,TEAL),
   L("• Assemble one representative test per domain.  • Apply ONE BH-FDR across all of them.  • Read "
     "each q against 5%.  • Report survivors AND casualties.  • State the method's limit.",16.5,True,NAVY)],
  [L("WILL NOT",18,True,RED),
   L("• Read any single p in isolation.  • Promote a borderline uncorrected result.  • Claim FDR fixes "
     "the length confound.  • Treat a surviving test as proof of intent.",16.5,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# 4b Module 1 frame (chart)
s=slide(prs); title(s,"Module 1 — Frame: 8 tests across 5 domains")
from collections import Counter as _C
dc=_C(d for _,d,_,_,_ in BAT)
doms=["Position","Sequence","Semantic","Signal","Biology"]
ebar(s,0.42,1.16,12.5,3.55,"Representative tests per domain (8 total)",
     doms,[("",[TEAL,AMBER,NAVY,TEAL,AMBER],[dc.get(x,0) for x in doms])],legend=False,fmt="{:.0f}")
two(s,[L("What it is / Why",17.5,True,TEAL),
   L("The capstone gathers one representative test from each Two Books domain — eight in all — so a "
     "single correction can speak for the whole project, not one page.",16)],
  [L("In the data · Bridge",17.5,True,AMBER),
   L("Position contributes three (contiguity ×2, shared length), Semantic two (theme, root-entropy), and "
     "Sequence, Signal, Biology one each. Next: the raw p's themselves.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 5 the battery table (chart: raw p by domain)
s=slide(prs); title(s,"Module 4 — The battery: one test per domain")
labels=[n.replace("· ","·\n") for n,_,_,_,_ in BAT]
ebar(s,0.42,1.16,12.5,3.7,"Raw permutation p (−log₁₀ p) — every domain represented; n.s. line at 1.3",
     [n for n,_,_,_,_ in BAT]+["n.s."],
     [("",[TEAL]*len(BAT)+[GREY],[nlp(p) for _,_,p,_,_ in BAT]+[1.30])],legend=False,fmt="{:.1f}")
panel(s,0.42,5.0,12.5,2.2,TINT2,[L("Eight representative tests",16,True,NAVY),
      L("Contiguity in two orderings (Position), shared length & shared theme per tag, letter- and "
        "root-entropy specials (Sequence/Semantic), length autocorrelation (Signal), and di-codon "
        "adjacency (Biology) — one well-chosen test per domain, so the correction never double-counts.",
        15.5,True,TEAL)],space=7)

# 6 the headline: p vs q (chart) [the money slide]
s=slide(prs); title(s,"Modules 5–6 — Raw p vs BH-FDR q across all domains")
ebar(s,0.42,1.16,12.5,3.9,"−log₁₀: light bar = raw p, coloured = BH-FDR q (green survives, red fails); 5% line ≈ 1.3",
     [n for n,_,_,_,_ in BAT],
     [("raw p",[ICE]*len(BAT),[nlp(p) for _,_,p,_,_ in BAT]),
      ("BH-FDR q",[TEAL if sv else RED for _,_,_,_,sv in BAT],[nlp(q) for _,_,_,q,_ in BAT])],
     legend=False,fmt="{:.1f}")
panel(s,0.42,5.2,12.5,2.0,TINT2,[L("The whole verdict in one picture",16,True,NAVY),
      L("Green q-bars clear the 5% line; red q-bars (shared theme, shared length) fall below it. Six of "
        "eight survive after correction — the structure is robust; the borderline theme claim is not.",
        15.5,True,TEAL)],space=7)

# 7 the survivors (chart)
s=slide(prs); title(s,"Module 5 — The survivors (6 of 8)")
surv=[b for b in BAT if b[4]]
ebar(s,0.42,1.16,12.5,3.55,"Survivors — BH-FDR q-value (all ≤ 0.0067)",
     [n for n,_,_,_,_ in surv],[("",[TEAL]*len(surv),[q for _,_,_,q,_ in surv])],legend=False,fmt="{:.4f}")
panel(s,0.42,4.86,12.5,2.34,TINT,[L("In the data — one survivor from every domain",15,True,TEAL),
      L("Contiguity in BOTH orderings (q 0.0010 each), length autocorrelation (q 0.0010, Signal), "
        "root-entropy (q 0.0010, Semantic), letter-entropy (q 0.0032, Sequence), and di-codon adjacency "
        "(q 0.0067, Biology). Every Two Books domain is represented among the survivors.",15.5,True,NAVY)],space=7)

# 8 the casualties (chart)
s=slide(prs); title(s,"Module 6 — The casualties: what drops out")
fail=[b for b in BAT if not b[4]]
finding2(s,
 {"title":"Shared theme — raw p vs BH q","cats":["raw p","BH q","5% line"],
  "series":[("",[ICE,RED,GREY],[0.049,0.056,0.05])],"legend":False,"fmt":"{:.3f}"},
 {"title":"Shared length — raw p","cats":["raw p","5% line"],
  "series":[("",[RED,GREY],[0.289,0.05])],"legend":False,"fmt":"{:.3f}"},
 [L("The correction works",17.5,True,RED),
  L("Shared theme is the textbook case: p = 0.049 looks significant alone, but q = 0.056 — just over the "
    "line. Multiplicity correctly demotes it.",16)],
 [L("Never close",17.5,True,AMBER),
  L("Shared length per tag is p = 0.289 — well inside the null. Both casualties drop out, exactly as a "
    "working guard should.",16)],
 fillA=REDT,fillB=AMBERT)

# 9 anatomy of the borderline (diagram)
s=slide(prs); title(s,"Why p = 0.049 is not a discovery")
band(s,0.42,1.18,12.5,0.4,TINT,"the same number, read alone vs read in a battery",NAVY)
fbox(s,0.9,1.9,5.4,1.6,AMBERT,"Read ALONE","p = 0.049 < 0.05 → looks like a finding.",line=AMBER,tsz=16)
fbox(s,7.0,1.9,5.4,1.6,REDT,"Read in the BATTERY","q = 0.056 > 0.05 → NOT a discovery.",line=RED,tsz=16)
panel(s,0.42,3.8,12.5,3.4,TINT2,[L("Multiplicity changes the honest reading",18,True,NAVY),
      L("A lone borderline p ignores the other seven tests. Benjamini–Hochberg asks how surprising this "
        "p is GIVEN the whole battery — and at rank 7 of 8 it is no longer surprising enough. The shared-"
        "theme claim is reported as not-supported.",16.5,True,TEAL)],space=9)

# 10 cross-domain reading (diagram)
s=slide(prs); title(s,"Module 7 — Cross-domain reading")
x=0.6
for lab,sub,col in [("Position","contiguity ×2\n(two orderings)",TEAL),
                    ("Sequence","letter\nentropy",AMBER),
                    ("Semantic","root\nentropy",NAVY),
                    ("Signal+Bio","autocorr +\ndi-codon",TEAL)]:
    fbox(s,x,1.95,2.95,1.7,TINT,lab,sub,line=col,tsz=15,ssz=12); x+=3.07
panel(s,0.42,4.0,12.5,3.2,TINT2,[L("Structure that survives independent readings",18,True,NAVY),
      L("The six survivors span two INDEPENDENT orderings (muṣḥaf and nuzūl contiguity) and all three "
        "pages plus the Signal and Biology structure tests. Structure that holds across independent "
        "readings and domains is far harder to explain away than any single result.",16.5,True,TEAL)],space=9)

# 11 the dashboard (diagram heatmap)
s=slide(prs); title(s,"The corrected dashboard at a glance")
band(s,0.42,1.18,12.5,0.38,TINT,"green = survives 5% FDR · red = drops out",NAVY)
cells=[TEAL,TEAL,TEAL,TEAL,TEAL,TEAL,RED,RED]
matgrid(s,1.4,1.8,10.5,1.4,1,8,cells)
panel(s,0.42,3.5,12.5,3.7,TINT2,[L("Six green, two red",18,True,NAVY),
      L("Contiguity ×2 · length autocorrelation · root-entropy · letter-entropy · di-codon — survive. "
        "Shared theme · shared length — drop out. One row, the whole Two Books verdict.",16.5,True,TEAL),
      L("Takeaway — robust, multi-domain structure under one joint correction; the borderline claim "
        "correctly excluded. Bridge: the synthesis names the method's limit.",15.5,True,GREY)],space=8)

# 12 synthesis three-stack
s=slide(prs); title(s,"Synthesis — one corrected dashboard")
three(s,[L("Survives 5% FDR (6/8)",17,True,TEAL),
   L("Contiguity ×2, length autocorrelation, root-entropy, letter-entropy, di-codon — q ≤ 0.0067, one "
     "from every domain.",15.5)],
  [L("Drops out",17,True,RED),
   L("Shared theme (p 0.049 → q 0.056) and shared length (p 0.289). The correction does its job.",15.5)],
  [L("The limit",17,True,NAVY),
   L("FDR controls MULTIPLICITY — not the sūra-length confound. A surviving test is reproducible, never "
     "a proof of intent.",15.5)],f=(TINT,REDT,TINT2))

# 12b two orderings agree (chart)
s=slide(prs); title(s,"Robustness — the contiguity signal in two orderings")
finding2(s,
 {"title":"Contiguity q-value — by ordering","cats":["muṣḥaf","nuzūl","5% line"],
  "series":[("",[TEAL,TEAL,GREY],[0.0010,0.0010,0.05])],"legend":False,"fmt":"{:.4f}"},
 {"title":"Significance (−log₁₀ q)","cats":["muṣḥaf","nuzūl","5% line"],
  "series":[("",[TEAL,TEAL,RED],[nlp(0.0010),nlp(0.0010),nlp(0.05)])],"legend":False,"fmt":"{:.1f}"},
 [L("Why two orderings matter",17.5,True,TEAL),
  L("The same contiguity test is run in the muṣḥaf order and the revelation (nuzūl) order — two "
    "INDEPENDENT arrangements of the same sūras.",16)],
 [L("In the data",17.5,True,AMBER),
  L("Both survive at q = 0.0010. A signal that holds in two independent orderings is the hardest kind to "
    "dismiss as an artefact — the backbone of the survivor set.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 AUDIT
s=slide(prs); title(s,"Audit — the method, claim by claim")
g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
audit_row(s,0.42,CY0,w,h,"✓",TEAL,"Multiplicity controlled","One BH-FDR across all 8 tests.")
audit_row(s,0.42+w+g,CY0,w,h,"✓",TEAL,"Cross-domain robustness","6/8 survive, one per domain.")
audit_row(s,0.42,CY0+h+g,w,h,"✓",TEAL,"Borderline demoted","Theme p 0.049 → q 0.056, excluded.")
audit_row(s,0.42+w+g,CY0+h+g,w,h,"✓",TEAL,"Reproducible","Every q re-runs live on the page.")
audit_row(s,0.42,CY0+2*(h+g),w,h,"✗",RED,"FDR ↔ confound control","FDR does NOT fix the length confound.")
audit_row(s,0.42+w+g,CY0+2*(h+g),w,h,"~",AMBER,"Survives ↔ proof","Reproducible ≠ proof of intent.")

# 14 disclaimer
s=slide(prs); title(s,"Not a scientific miracle — and not evidence")
panel(s,0.42,1.20,12.5,2.9,REDT,[L("What we are NOT claiming",18,True,RED),
      L("Surviving 5% FDR does not prove design, intent, or a miracle. FDR controls only for "
        "MULTIPLICITY — it does not correct for the sūra-length confound that several tests share, and it "
        "cannot turn a reproducible pattern into evidence of purpose.",16.5,True,NAVY)],space=9)
panel(s,0.42,4.30,12.5,2.9,TINT,[L("What the dashboard IS for",18,True,NAVY),
      L("A single, honest place to read the whole Two Books battery together, with multiplicity "
        "controlled and every q reproducible live from Book6. It tells us which structure is robust "
        "across domains — and which borderline claim to set aside. Judged by reproducibility, never "
        "offered as proof.",16.5,True,TEAL)],space=9)

# 15 quick reference
s=slide(prs); title(s,"Quick reference — terms & live numbers")
two(s,[L("Terms",17,True,TEAL),
   L("Multiplicity — many tests inflate false positives.  False-discovery rate — expected fraction of "
     "'discoveries' that are null.  Benjamini–Hochberg — rank p's, compare to (i/m)·α.  q-value — "
     "smallest α at which a test still passes.",16)],
  [L("Live Book6 numbers",17,True,AMBER),
   L("8 tests · 6 survive 5% FDR · survivors q ≤ 0.0067 · contiguity p 0.0005 → q 0.0010 · theme p 0.049 "
     "→ q 0.056 (fails) · shared length p 0.289 (fails).",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 16 close
s=slide(prs); title(s,"Close — the Two Books series, corrected")
panel(s,0.42,1.20,12.5,6.0,TINT2,[L("The through-line",18,True,NAVY),
      L("Across Disjoint Letters, Signal and Biology, the corpus shows real geometric, sequential and "
        "compositional structure. Read together under one Benjamini–Hochberg correction, six of eight "
        "representative tests survive 5% FDR — one from every domain — while the borderline shared-theme "
        "claim correctly drops out.",16.5),
      L("That is the disciplined verdict of the whole series: robust, multi-domain, reproducible "
        "structure — with multiplicity controlled, the length confound flagged, and nothing offered as "
        "proof of intent.",16,True,TEAL),
      L("Every q-value on these slides re-runs live on the 📋 Two Books · Global FDR page — anyone can "
        "reproduce the dashboard from Book6.",15.5)],space=10)

prs.save(OUT)
print("FDR deck built:",OUT,"| slides:",len(list(prs.slides)))
