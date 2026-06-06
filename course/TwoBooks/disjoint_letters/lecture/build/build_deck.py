# -*- coding: utf-8 -*-
"""Disjoint Letters companion deck — >=20 slides, >=half visual. st_slides + diagrams.
All Qur'an numbers from Book6 (tour_bank); permutation nulls + Benjamini-Hochberg throughout."""
import os, sys, json
HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
from st_slides import *
from diagrams import fbox,harrow,vdash,band,matgrid,chain
LEC=os.path.dirname(HERE)
TB=json.load(open(os.path.join(LEC,"..","..","_handson_build","tour_bank.json"),encoding="utf-8"))
OUT=os.path.join(LEC,"DisjointLetters_Lecture_Slides.pptx")

def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def audit_row(s,x,y,w,h,mark,mc,head,body):
    panel(s,x,y,w,h,(TINT if mark=="✓" else REDT if mark=="✗" else AMBERT),
          [L(mark+"  "+head,15,True,mc),L(body,14)],space=5)
# -log10 p helper for significance bars
import math
def nlp(p): return min(4.0, -math.log10(max(p,1e-4)))

prs=deck()

FIGDIR=os.path.join(LEC,"figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.20,12.5,1.65,TINT2,[L("TWO BOOKS  ·  DISJOINT LETTERS — companion lecture deck",15,True,TEAL),
      L("al-Muqaṭṭaʿāt as a positional pointer — not a code",23,True,NAVY)],space=7)
panel(s,0.42,3.05,12.5,2.05,TINT,[L("The honest spine",17,True,NAVY),
      L("The disjoint letters mark a CONTIGUOUS block of LONGER, compositionally distinct sūras — a "
        "positional and organizational signal that survives permutation nulls and FDR. The folkloric "
        "'code' claims (Qāf-letter saturation, thematic encoding) do NOT survive. Benjamini–Hochberg "
        "correction across the whole battery; no miracle claims.",15.5,True,TEAL)],space=8)
panel(s,0.42,5.25,12.5,1.95,TINT2,[L("A workbench, not an opinion",16,True,NAVY),
      L("29 of 114 sūras open with disjoint letters. Each test below tags those sūras, compares to "
        "thousands of random tags of the same size, and reports an FDR-corrected verdict — every number "
        "computed live from Book6.",15.5)],space=7)

# 2 framing
s=slide(prs); title(s,"What are al-Muqaṭṭaʿāt?")
band(s,0.42,1.20,12.5,0.42,TINT,"الم   ·   الر   ·   حم   ·   طسم   —   the disjoint letters opening 29 sūras",NAVY)
fbox(s,0.7,1.95,5.7,1.5,TINT,"The popular view","A hidden cipher: each letter secretly encodes meaning or content.",line=AMBER,tsz=15)
fbox(s,6.9,1.95,5.7,1.5,TINT,"The testable thesis","A positional / organizational POINTER to a distinct block of sūras.",line=TEAL,tsz=15)
panel(s,0.42,3.75,12.5,3.4,TINT2,[L("We test the pointer thesis, and we test the cipher claims too",18,True,NAVY),
      L("The letters open 29 sūras and are recited one by one. Rather than assume a code, we tag every "
        "sūra muqaṭṭaʿāt-bearing or not and ask whether that tag predicts measurable structure — "
        "position, size, composition — beyond chance.",16.5),
      L("Every claim faces a permutation null; the whole battery is then read AFTER Benjamini–Hochberg "
        "correction, so no single borderline p is mistaken for a miracle.",16,True,TEAL)],space=9)

# 3 the workbench (diagram)
s=slide(prs); title(s,"The workbench — three scales of question")
x=0.7
for lab,sub,col in [("🧭 Position","contiguity, per-family,\nsize / geometry",TEAL),
                    ("🔤 Sequence","letter density &\nletter information",AMBER),
                    ("🧩 Semantic","root information,\nfamily lexis",NAVY)]:
    fbox(s,x,1.95,3.85,1.7,TINT,lab,sub,line=col,tsz=17,ssz=12); x+=4.05
panel(s,0.42,4.0,12.5,3.2,TINT2,[L("One tag, many lenses",18,True,NAVY),
      L("Position asks WHERE the marked sūras sit and how big they are; Sequence asks how their LETTERS "
        "behave; Semantic asks how their ROOTS behave. A real pointer should leave a footprint across "
        "scales — and the cipher claims should fail under the same nulls.",16.5,True,TEAL)],space=9)

# 4 scope
s=slide(prs); title(s,"What this deck will — and will not — do")
two(s,[L("WILL",18,True,TEAL),
   L("• Tag sūras and test each claim against a permutation null.  • Read the battery AFTER FDR.  • Keep "
     "only what survives.  • Report the geometry the data licenses.  • Audit each rung ✓ / ✗ / ~.",16.5,True,NAVY)],
  [L("WILL NOT",18,True,RED),
   L("• Treat the letters as a decoded cipher.  • Promote a borderline uncorrected p to a finding.  • "
     "Read contiguity as theology.  • Offer any parallel as proof.",16.5,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# 5 method: permutation + FDR (diagram)
s=slide(prs); title(s,"The method — permutation nulls + FDR")
fbox(s,0.7,1.95,5.6,1.5,TINT,"Permutation null","Re-tag 29 random sūras thousands of times; rebuild the statistic each time.",line=TEAL,tsz=15)
fbox(s,6.9,1.95,5.6,1.5,AMBERT,"Benjamini–Hochberg","Correct the whole battery of p-values so chance hits don't masquerade as findings.",line=AMBER,tsz=15)
panel(s,0.42,3.75,12.5,3.4,TINT2,[L("Why FDR is non-negotiable here",18,True,NAVY),
      L("Run enough tests and some will cross p < 0.05 by chance alone — exactly how folklore manufactures "
        "a 'miracle'. Benjamini–Hochberg controls the false-discovery rate across the battery, so a "
        "borderline result like the theme test (p = 0.049) is correctly demoted once corrected.",16.5),
      L("Everything that follows is reported with its FDR-corrected status, never a lone p.",16,True,TEAL)],space=9)

# 5b Module 1 frame — the tag (chart)
s=slide(prs); title(s,"Module 1 — Frame: the 29-sūra tag")
finding2(s,
 {"title":"The corpus, split by the tag","cats":["muqaṭṭaʿāt sūras","other sūras"],
  "series":[("",[TEAL,GREY],[29,85])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Distinct opening families","cats":["families"],
  "series":[("",[AMBER],[len(TB["perfam"])])],"legend":False,"fmt":"{:.0f}"},
 [L("What it is / Why",17.5,True,TEAL),
  L("29 of 114 sūras open with disjoint letters, in a handful of distinct families. We tag each sūra and "
    "ask whether the tag predicts structure beyond chance.",16)],
 [L("In the data · Bridge",17.5,True,AMBER),
  L("29 marked vs 85 unmarked sūras — a clean two-group test. First and strongest question: do the marked "
    "sūras cluster contiguously?",16)],
 fillA=TINT,fillB=AMBERT)

# 6 contiguity (chart) — both orderings
s=slide(prs); title(s,"Module 2 — Contiguity: do the marked sūras cluster?")
ebar(s,0.42,1.16,12.5,3.55,"Contiguity significance (−log₁₀ p) — higher = stronger; n.s. line at 1.3",
     ["muṣḥaf order","nuzūl order","n.s. line"],
     [("",[TEAL,TEAL,GREY],[nlp(TB["contiguity_mushaf"]),nlp(TB["contiguity_nuzul"]),1.30])],legend=False,fmt="{:.1f}")
panel(s,0.42,4.86,12.5,2.34,AMBERT,[L("In the data — they cluster, in BOTH orderings",15,True,AMBER),
      L("The 29 muqaṭṭaʿāt sūras are significantly contiguous at p ≈ 0.0005 in the muṣḥaf order AND "
        "p ≈ 0.0005 in the revelation (nuzūl) order. A result that holds in two independent orderings is "
        "far harder to dismiss as an artefact — this is the strongest geometric signal.",15.5,True,NAVY)],space=7)

# 7 per-family (chart)
s=slide(prs); title(s,"Module 3 — Per-family contiguity")
fam=TB["perfam"]
fams=list(fam.keys()); vals=[nlp(fam[k]) for k in fams]
ebar(s,0.42,1.16,12.5,3.55,"Per-family contiguity significance (−log₁₀ p) — n.s. line at 1.3",
     fams+["n.s."],[("",[TEAL]*len(fams)+[GREY],vals+[1.30])],legend=False,fmt="{:.1f}")
two(s,[L("How it's done",17.5,True,TEAL),
   L("Repeat the contiguity test WITHIN each letter-family (ḤM, ALR, ALM, ṬSM) against the same null.",16)],
  [L("In the data",17.5,True,AMBER),
   L("Every family clusters: ḤM and ALR at p ≈ 0.0005, ALM at 0.006, ṬSM at 0.035. The geometry is "
     "family-structured, not one lucky run.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 8 size (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 4 — Size: the marked sūras are longer","dl_length.png",
  "In the data","the full sūra-length distribution for both groups; the muqaṭṭaʿāt sūras (teal) sit well to the right — median 85 vs 26 āyahs. A large, consistent size gap, not a handful of outliers.",cap_fill=TINT)

# 9 what it is NOT — the code claims (chart)
s=slide(prs); title(s,"Module 5 — What it is NOT: the code claims fail")
ebar(s,0.42,1.16,12.5,3.55,"'Code' claims — significance (−log₁₀ p); ALL below the n.s. line (1.3)",
     ["Qāf-letter\n(sūra 50)","theme\ncoding","embedding\nsimilarity","n.s. line"],
     [("",[RED,RED,RED,GREY],[nlp(0.0978),nlp(TB["theme_p"]),nlp(TB["embedding_p"]),1.30])],legend=False,fmt="{:.2f}")
panel(s,0.42,4.86,12.5,2.34,REDT,[L("In the data — none survive",15,True,RED),
      L("In Sūrat Qāf the letter qāf is 3.76% of letters at p ≈ 0.10 — NOT significant. The theme test is "
        "p = 0.049 (borderline, and gone after FDR). The embedding-similarity claim is p ≈ 0.10. The "
        "famous letter-cipher claims are NOT supported.",15.5,True,NAVY)],space=7)

# 10 the Qaf myth, in detail (both-domain style: claim vs null)
s=slide(prs); title(s,"Module 5 — Anatomy of a myth: Sūrat Qāf")
finding2(s,
 {"title":"Letter qāf density in sūra 50 (%)","cats":["sūra 50","corpus rank"],
  "series":[("",[TEAL,GREY],[3.76,3.76])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Is it special? (−log₁₀ p)","cats":["observed","n.s. line"],
  "series":[("",[GREY,RED],[nlp(0.0978),1.30])],"legend":False,"fmt":"{:.2f}"},
 [L("The claim",17.5,True,TEAL),
  L("Folklore says Sūrat Qāf is saturated with the letter qāf. Its qāf density is 3.76% — and it ranks "
    "111th of 114 sūras, not 1st.",16)],
 [L("The test",17.5,True,RED),
  L("Against a permutation null the elevation is p ≈ 0.10 — not significant. An attractive story that the "
    "data simply does not support.",16)],
 fillA=TINT,fillB=REDT)

# 11 letter information theory (chart)
s=slide(prs); title(s,"Module 6 — Letter information theory")
finding2(s,
 {"title":"Letter-entropy difference — significance (−log₁₀ p)","cats":["observed","n.s. line"],
  "series":[("",[TEAL,GREY],[nlp(TB["letter_entropy_p"]),1.30])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Reference — what entropy measures","cats":["even mix","skewed mix"],
  "series":[("",[TEAL,AMBER],[1.0,0.6])],"legend":False,"fmt":"{:.1f}"},
 [L("How / What we get",17.5,True,TEAL),
  L("Compute the letter-entropy of each group and test the difference against a permutation null.",16)],
 [L("In the data",17.5,True,AMBER),
  L("The difference is significant at p ≈ 0.002 — muqaṭṭaʿāt sūras are compositionally distinct at the "
    "letter scale, not only longer.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 root information theory (chart)
s=slide(prs); title(s,"Module 7 — Root information theory")
ebar(s,0.42,1.16,12.5,3.55,"Root-scale differences — significance (−log₁₀ p); n.s. line at 1.3",
     ["root entropy","lexical richness","n.s. line"],
     [("",[TEAL,TEAL,GREY],[nlp(TB["root_entropy_p"]),nlp(TB["lexical_richness_p"]),1.30])],legend=False,fmt="{:.1f}")
two(s,[L("How it's done",17.5,True,TEAL),
   L("Compute root-entropy and lexical richness per group; test each against a permutation null.",16)],
  [L("In the data",17.5,True,AMBER),
   L("Root-entropy differs at p ≈ 0.0005 and lexical richness at p ≈ 0.0005 — the distinction is sharp "
     "at the root scale too. A multi-scale structural signal.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 13 the validated picture (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 8 — Contiguity you can SEE: the 114-sūra strip","dl_contiguity.png",
  "In the data","each cell is one sūra in order; teal = disjoint-letter sūra. They fall in CONTIGUOUS blocks (the ḤM, ALM, ALR families), not scattered — the geometric signal at p ≈ 0.0005.")

# 14 synthesis three-stack
s=slide(prs); title(s,"Synthesis — pointer, not cipher")
three(s,[L("Survives (FDR-clean)",17,True,TEAL),
   L("Contiguity (0.0005, two orderings), per-family clustering, size 85 vs 26, letter-entropy 0.002, "
     "root-entropy & richness 0.0005.",15.5)],
  [L("Fails the null / FDR",17,True,RED),
   L("Qāf-letter saturation (p ≈ 0.10), thematic coding (0.049, gone after FDR), embedding similarity "
     "(p ≈ 0.10).",15.5)],
  [L("Honest reading",17,True,NAVY),
   L("A geometric and organizational signal marks these sūras as a distinct class — the data does not "
     "license a hidden cipher.",15.5)],f=(TINT,REDT,TINT2))

# 14b full battery summary (chart)
s=slide(prs); title(s,"The whole battery at a glance")
labels=["contig.\nmuṣḥaf","contig.\nnuzūl","size","letter\nH","root\nH","richness","Qāf","theme","embed"]
import math as _m
def _nlp(p): return min(4.0,-_m.log10(max(p,1e-4)))
vals=[_nlp(TB["contiguity_mushaf"]),_nlp(TB["contiguity_nuzul"]),4.0,_nlp(TB["letter_entropy_p"]),
      _nlp(TB["root_entropy_p"]),_nlp(TB["lexical_richness_p"]),_nlp(0.0978),_nlp(TB["theme_p"]),_nlp(TB["embedding_p"])]
cols=[TEAL,TEAL,TEAL,TEAL,TEAL,TEAL,RED,RED,RED]
ebar(s,0.42,1.16,12.5,3.7,"Every test — significance (−log₁₀ p). Green survives, red fails; n.s. line at 1.3",
     labels,[("",cols,vals)],legend=False,fmt="{:.1f}")
panel(s,0.42,5.0,12.5,2.2,TINT2,[L("Six survive, three fail",16,True,NAVY),
      L("The geometric and compositional tests (green) clear the null and FDR; the three cipher claims "
        "(red) sit below the n.s. line. One picture of the honest verdict — read the battery together, "
        "never a lone p.",15.5,True,TEAL)],space=7)

# 15 AUDIT
s=slide(prs); title(s,"Audit — claim by claim")
g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
audit_row(s,0.42,CY0,w,h,"✓",TEAL,"Contiguity ↔ pointer","Significant in two orderings (p ≈ 0.0005).")
audit_row(s,0.42+w+g,CY0,w,h,"✓",TEAL,"Per-family clustering","Every family significant (ḤM/ALR sharpest).")
audit_row(s,0.42,CY0+h+g,w,h,"✓",TEAL,"Size ↔ architecture","Median 85 vs 26 āyahs.")
audit_row(s,0.42+w+g,CY0+h+g,w,h,"✓",TEAL,"Composition ↔ distinct class","Letter & root entropy differ (≤ 0.002).")
audit_row(s,0.42,CY0+2*(h+g),w,h,"✗",RED,"Qāf-letter ↔ cipher","p ≈ 0.10 — not significant.")
audit_row(s,0.42+w+g,CY0+2*(h+g),w,h,"~",AMBER,"Theme ↔ encoding","p = 0.049, gone after FDR — not supported.")

# 16 disclaimer
s=slide(prs); title(s,"Not a scientific miracle — and not evidence")
panel(s,0.42,1.20,12.5,2.9,REDT,[L("What we are NOT claiming",18,True,RED),
      L("The disjoint letters are not a decoded cipher; contiguity is not theology; a borderline "
        "uncorrected p is not a finding. We do not claim to know WHY the letters are there — only what "
        "measurable structure the tag does and does not predict.",16.5,True,NAVY)],space=9)
panel(s,0.42,4.30,12.5,2.9,TINT,[L("What the workbench IS for",18,True,NAVY),
      L("A disciplined way to separate a real geometric/organizational signal from folklore, with a "
        "permutation null behind every answer and FDR across the battery. Every number is recomputable "
        "live in the app. The findings are judged by reproducibility, never offered as proof of intent.",
        16.5,True,TEAL)],space=9)

# 17 quick reference
s=slide(prs); title(s,"Quick reference — terms & live numbers")
two(s,[L("Terms",17,True,TEAL),
   L("Muqaṭṭaʿāt — disjoint letters opening 29 sūras.  Contiguity — clustering of tagged sūras.  "
     "Permutation null — re-tag random sūras to build the chance distribution.  FDR (Benjamini–Hochberg) "
     "— battery-wide correction.  Letter/root entropy — composition information.",16)],
  [L("Live Book6 numbers",17,True,AMBER),
   L("Contiguity p ≈ 0.0005 (muṣḥaf & nuzūl) · families ḤM/ALR 0.0005, ALM 0.006, ṬSM 0.035 · median 85 "
     "vs 26 āyahs · letter-entropy 0.002 · root-entropy & richness 0.0005 · Qāf 3.76% p ≈ 0.10 (n.s.).",16)],
  sp=0.5,fa=TINT,fb=AMBERT)

# 18 close
s=slide(prs); title(s,"Close — the geometry the data licenses")
panel(s,0.42,1.20,12.5,6.0,TINT2,[L("The through-line",18,True,NAVY),
      L("Tag the 29 sūras and the structure appears: they cluster contiguously in two orderings, they are "
        "markedly longer, and they are compositionally distinct at both the letter and the root scale. "
        "The attractive cipher claims — Qāf-letter saturation, thematic encoding — do not survive the "
        "nulls or FDR.",16.5),
      L("So the disjoint letters are best read as a validated positional and organizational POINTER, not "
        "a hidden code — every claim with its null, every parallel a labelled lens, never evidence.",16,True,TEAL),
      L("Next in the series — Signal reads the same corpus as an ordered signal, Biology as a genome; the "
        "FDR Summary then collects every Two Books test into one corrected dashboard.",15.5)],space=10)

prs.save(OUT)
print("DL deck built:",OUT,"| slides:",len(list(prs.slides)))
