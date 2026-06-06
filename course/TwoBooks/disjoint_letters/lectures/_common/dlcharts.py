# -*- coding: utf-8 -*-
"""Native, EDITABLE PowerPoint charts for the DL course, with back-of-room fonts.
Every chart is a real pptx chart object: titles, axis labels, data labels and the
data table are all selectable/editable in PowerPoint."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.oxml.ns import qn
from st_slides import slide, title, panel, L, CY0, CY1, LX, RW, NAVY, TEAL, AMBER, RED, GREY, INK, TINT, TINT2, AMBERT, REDT, ICE, LTEAL, _fill_panel, _tb, sf
PALETTE=[TEAL,NAVY,AMBER,RED,GREY,LTEAL]

# fonts: readable from the back of the room
F_TITLE=20; F_AXIS=15; F_TICK=14; F_LABEL=14; F_LEG=14

def _chfont(chart):
    try:
        chart.font.size=Pt(F_TICK); chart.font.name="Calibri"
    except Exception: pass

def _axis_titles(chart,xlab,ylab):
    for ax,lab in ((getattr(chart,'category_axis',None),xlab),(getattr(chart,'value_axis',None),ylab)):
        if ax is None or not lab: continue
        try:
            ax.has_title=True; ax.axis_title.text_frame.text=lab
            r=ax.axis_title.text_frame.paragraphs[0].runs[0]; r.font.size=Pt(F_AXIS); r.font.bold=True; r.font.color.rgb=NAVY
        except Exception: pass
        try:
            ax.tick_labels.font.size=Pt(F_TICK); ax.tick_labels.font.color.rgb=INK
        except Exception: pass

def _title(chart,t):
    if not t: chart.has_title=False; return
    chart.has_title=True; chart.chart_title.text_frame.text=t
    r=chart.chart_title.text_frame.paragraphs[0].runs[0]; r.font.size=Pt(F_TITLE); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Calibri")

def _legend(chart,on=True):
    chart.has_legend=on
    if on:
        chart.legend.position=XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout=False
        chart.legend.font.size=Pt(F_LEG)

def _color_series(plot):
    for i,s in enumerate(plot.series):
        try: s.format.fill.solid(); s.format.fill.fore_color.rgb=PALETTE[i%len(PALETTE)]
        except Exception: pass

def _dlabels(plot,fmt="0",on=True,color=NAVY):
    if not on: return
    try:
        dl=plot.data_labels; dl.number_format=fmt; dl.number_format_is_linked=False
        dl.font.size=Pt(F_LABEL); dl.font.bold=True; dl.font.color.rgb=color
        dl.position=XL_LABEL_POSITION.OUTSIDE_END
    except Exception: pass

def _frame(prs,t):
    s=slide(prs); title(s,t); return s

def _addchart(s,xl_type,cd,x=1.4,y=1.18,w=10.5,h=3.5):
    gf=s.shapes.add_chart(xl_type,Inches(x),Inches(y),Inches(w),Inches(h),cd); return gf.chart

def chart_slide(prs,ttl,kind,ctitle,cats,series,insight,ylab="",xlab="",
                datalabels=True,legend=None,fmt="0",fill=TINT,ymax=None,ymin=None):
    """kind: col,bar,line,pie. series=[(name,[vals]),...]. Native editable chart + insight panel."""
    s=_frame(prs,ttl)
    cd=CategoryChartData(); cd.categories=cats
    for nm,vals in series: cd.add_series(nm,vals)
    xt={"col":XL_CHART_TYPE.COLUMN_CLUSTERED,"bar":XL_CHART_TYPE.BAR_CLUSTERED,
        "line":XL_CHART_TYPE.LINE_MARKERS,"pie":XL_CHART_TYPE.PIE}[kind]
    ch=_addchart(s,xt,cd)
    _chfont(ch); _title(ch,ctitle)
    multi=len(series)>1
    _legend(ch, legend if legend is not None else multi)
    if kind=="pie":
        ch.plots[0].has_data_labels=True
        dl=ch.plots[0].data_labels; dl.number_format='0"%"' if fmt=="pct" else "0"; dl.number_format_is_linked=False
        dl.font.size=Pt(F_LABEL); dl.font.bold=True; dl.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
        dl.show_percentage=True if fmt=="pct" else False
        for i,pt in enumerate(ch.plots[0].series[0].points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb=PALETTE[i%len(PALETTE)]
    else:
        _color_series(ch.plots[0])
        _axis_titles(ch,xlab,ylab)
        _dlabels(ch.plots[0],fmt=("0.0" if fmt=="f1" else "0"),on=datalabels)
        if ymax is not None:
            try: ch.value_axis.maximum_scale=ymax
            except Exception: pass
        if ymin is not None:
            try: ch.value_axis.minimum_scale=ymin
            except Exception: pass
    pt=CY0+3.62; ph=CY1-pt
    panel(s,LX,pt,RW,ph,fill,insight,space=8)
    return s

def scatter_slide(prs,ttl,ctitle,series,insight,xlab="",ylab="",fill=TINT2,legend=True):
    """series=[(name,[(x,y),...]),...] -> native editable XY scatter."""
    s=_frame(prs,ttl)
    cd=XyChartData()
    for nm,pts in series:
        sr=cd.add_series(nm)
        for (xx,yy) in pts: sr.add_data_point(xx,yy)
    ch=_addchart(s,XL_CHART_TYPE.XY_SCATTER,cd)
    _chfont(ch); _title(ch,ctitle); _legend(ch,legend and len(series)>1)
    for i,sr in enumerate(ch.plots[0].series):
        try: sr.marker.format.fill.solid(); sr.marker.format.fill.fore_color.rgb=PALETTE[i%len(PALETTE)]; sr.marker.size=9
        except Exception: pass
    _axis_titles(ch,xlab,ylab)
    pt=CY0+3.62; ph=CY1-pt
    panel(s,LX,pt,RW,ph,fill,insight,space=8)
    return s
