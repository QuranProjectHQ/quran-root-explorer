# -*- coding: utf-8 -*-
import os,sys
HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
from st_slides import *
from st_slides import _tb,_rect,_fill_panel
from diagrams import fbox,harrow,vdash,band,sigrow
from pptx.util import Inches,Pt
ROOT=os.path.abspath(os.path.join(HERE,".."))
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def Tt(prs,t,sz=21):
    s=slide(prs); title(s,t,sz); return s
def titleslide(prs,tag,headline,body1,body2):
    s=slide(prs)
    panel(s,0.42,1.05,12.5,1.5,TINT2,[L(tag,15,True,TEAL),L(headline,23,True,NAVY)],space=6)
    panel(s,0.42,2.85,12.5,4.35,TINT,[L(body1,16.5),L(body2,16,True,TEAL)],space=9)
    return s
def numline(s,x,y,w,families,xmax=115,bh=2.2):
    ax=s.shapes.add_shape(1,Inches(x),Inches(y+bh),Inches(w),Inches(0.02)); ax.fill.solid(); ax.fill.fore_color.rgb=GREY; ax.line.fill.background()
    for i,(name,nums,col) in enumerate(families):
        yy=y+0.1+i*((bh-0.2)/max(1,len(families)-1))
        for n in nums:
            d=s.shapes.add_shape(9,Inches(x+ (n/xmax)*w -0.06),Inches(yy-0.06),Inches(0.12),Inches(0.12)); d.fill.solid(); d.fill.fore_color.rgb=col; d.line.fill.background()
        _tb(s,x-0.02,yy-0.16,1.6,0.32,[(name,11,True,col)])
def audit(s,ok,bad,sil):
    title(s,"Audit — supported, broken, silent")
    three(s,[L("✓ SUPPORTED",17,True,TEAL),L(ok,16)],[L("✗ BREAKS",17,True,RED),L(bad,16)],[L("~ SILENT",17,True,AMBER),L(sil,16)],f=(TINT,REDT,AMBERT))
def takeaway(s,rel,take):
    title(s,"Real-world relevance & takeaway")
    two(s,[L("RELEVANCE",18,True,NAVY),L(rel,17,True,TEAL)],[L("TAKEAWAY",18,True,AMBER),L(take,17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
def appslide(prs,steps,body):
    s=slide(prs); title(s,"The app — explore it live")
    band(s,0.42,1.2,12.5,0.4,TINT,"disjoint-letter pointer explorer",TEAL)
    x=0.55; bw=2.98; aw=0.1
    for i,(t,sub,fl,ln) in enumerate(steps):
        fbox(s,x,1.9,bw,1.35,fl,t,sub,line=ln,tsz=15,ssz=11)
        if i<len(steps)-1: harrow(s,x+bw-0.02,2.45,aw+0.05,"",color=GREY)
        x+=bw+aw
    panel(s,0.42,4.65,12.5,2.55,TINT,[L("Hands on the data",18,True,NAVY),L(body,16.5,True,TEAL)],space=7)
FIGDL=os.path.join(ROOT,"figs_dl")
def figslide(prs,ttl,png,insight,fill=TINT,split=0.62):
    from PIL import Image
    s=slide(prs); title(s,ttl)
    path=os.path.join(FIGDL,png); band_h=CY1-CY0; img_h=band_h*split-0.10
    if os.path.exists(path):
        iw,ih=Image.open(path).size; ar=iw/ih; h=img_h; w=h*ar
        if w>12.0: w=12.0; h=w/ar
        x=(13.333-w)/2; y=CY0+(img_h-h)/2+0.02
        s.shapes.add_picture(path,Inches(x),Inches(y),width=Inches(w),height=Inches(h))
    pt=CY0+img_h+0.14; ph=CY1-pt
    panel(s,LX,pt,RW,ph,fill,insight,space=8)
    return s
def figtwo(prs,ttl,pngL,pngR,caption,fill=TINT2):
    from PIL import Image
    s=slide(prs); title(s,ttl); cw=6.15; ih=3.45
    for i,png in enumerate((pngL,pngR)):
        path=os.path.join(FIGDL,png)
        if os.path.exists(path):
            iw,ihh=Image.open(path).size; ar=iw/ihh; h=ih; w=h*ar
            if w>cw: w=cw; h=w/ar
            x=(0.42 if i==0 else 6.79)+(cw-w)/2; y=CY0+(ih-h)/2
            s.shapes.add_picture(path,Inches(x),Inches(y),width=Inches(w),height=Inches(h))
    pt=CY0+ih+0.16; ph=CY1-pt
    panel(s,LX,pt,RW,ph,fill,caption,space=8)
    return s
def save(prs,fn):
    prs.save(os.path.join(ROOT,fn)); return len(prs.slides)
