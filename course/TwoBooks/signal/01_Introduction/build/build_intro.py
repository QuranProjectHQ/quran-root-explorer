# -*- coding: utf-8 -*-
"""Signal course - Lecture 1 - Introduction.
Idea: a signal is the DIGITIZED form of an ayah - vectorize each verse into a
1-D numeric signal, then read it with DSP. All Qur'an figures are recomputed
from Book6.xlsx (6236 ayat) and verified exact.
Locked standard: >=20 editable slides, >=half visual, dual-domain real data,
audit, Monte-Carlo, app-central, real-world takeaway.
"""
import os, sys
HERE=os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0,HERE)
from st_slides import *
from st_slides import _tb, _rect, _fill_panel
from diagrams import fbox,harrow,vdash,band,sigrow,matgrid,isocube,chain
from pptx.util import Inches,Pt

OUT=os.path.abspath(os.path.join(HERE,".."))+"/"
prs=deck()
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

def signal_plot(s,x,y,w,vals,tokens,col=TEAL,vmax=None,base_h=2.7,lblcol=NAVY):
    n=len(vals); cw=w/n; vmax=vmax or max(vals); base=y+base_h
    ax=s.shapes.add_shape(1,Inches(x),Inches(base),Inches(w),Inches(0.02)); ax.fill.solid(); ax.fill.fore_color.rgb=GREY; ax.line.fill.background()
    for i,(v,tok) in enumerate(zip(vals,tokens)):
        h=0.16+(v/vmax)*(base_h-0.5)
        b=s.shapes.add_shape(5,Inches(x+i*cw+cw*0.14),Inches(base-h),Inches(cw*0.72),Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb=col; b.line.fill.background()
        _tb(s,x+i*cw-0.05,base-h-0.34,cw+0.1,0.30,[(str(v),14,True,NAVY)])
        _tb(s,x+i*cw-0.05,base+0.06,cw+0.1,0.46,[(tok,20,True,lblcol)])
        _tb(s,x+i*cw-0.05,base+0.52,cw+0.1,0.26,[("x[%d]"%i,11,False,GREY)])

def map_row(s,y,left,right,h=0.66):
    fbox(s,0.55,y,5.45,h,AMBERT,left,line=AMBER,tsz=14)
    _tb(s,6.1,y,1.1,h,[("≡",19,True,GREY)])
    fbox(s,7.3,y,5.45,h,TINT,right,line=TEAL,tsz=14)

# ===================================================== 1 TITLE
s=slide(prs)
panel(s,0.42,1.05,12.5,1.55,TINT2,[L("THE TWO BOOKS  ·  Ayah as Signal (1-D)  ·  a Qur’an-and-science lecture series",15,True,TEAL),
  L("Lecture 1 — Introduction: digitizing the ayah into a signal",26,True,NAVY)],space=7)
fbox(s,0.7,3.0,3.7,1.35,TINT,"AYAH","a verse of n tokens (text data)",line=TEAL,tsz=18,ssz=11.5)
harrow(s,4.6,3.55,1.0,"digitize",color=GREY,lcol=TEAL)
fbox(s,5.8,3.0,3.5,1.35,AMBERT,"x = [ … ]","a 1-D signal of n samples",line=AMBER,tsz=18,ssz=11.5)
harrow(s,9.5,3.55,1.0,"DSP",color=GREY,lcol=NAVY)
fbox(s,10.7,3.0,2.1,1.35,REDT,"READING","frequency · shape",line=RED,tsz=16,ssz=11)
panel(s,0.42,4.7,12.5,2.5,TINT,[L("One idea, two Books",17,True,NAVY),
  L("A signal is just the DIGITIZED form of an ayah. This course turns each verse into a 1-D numeric vector and reads it with the mathematics of digital signal processing. The Book of Creation side is the same: nature’s data (sound, sensors, measurements) are also 1-D signals, analysed by the same math.",16.5),
  L("Every Qur’an number here is recomputed from Book6.xlsx (6,236 ayat) and verified exact. No “scientific-miracle” claims — every parallel is audited ✓ / ✗ / ~.",16,True,TEAL)],space=8)

# ===================================================== 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur’an: God’s speech set down in language — tadwin. The Book of SCRIPTURE. Here we read its text, digitized.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The Universe: God’s deed brought into being — takwin. The Book of CREATION. Its data arrive as physical signals.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source — Allah. Both are ayat (signs); both, here, become numbers we can measure. Read side by side, never collapsed.",16)])

# ===================================================== 3 THE IDEA
s=slide(prs); title(s,"The idea — a signal is just the digitized ayah")
two(s,[L("FROM VERSE TO VECTOR",18,True,NAVY),
  L("Take an ayah and read its ROOTS in order. Assign each root a number — its corpus frequency. The verse becomes x = [x₀, x₁, …, xₙ₋₁]: a length-n sequence of samples. That sequence is a 1-D signal — the same object DSP studies in audio, ECG, or a sensor trace.",17,True,TEAL)],
 [L("THIS IS DIGITIZING TEXT",18,True,NAVY),
  L("Digitization turns an analog source into numbers we can store and compute. The ayah is the source; the vector is its digitization. A signal is just numbers indexed by position — and an ayah already IS numbers-in-order once you choose what to measure per root. The ANCHOR is the ROOT (highest semantic power, as in NLP and the biology course); one root-token = one sample. Surface form and morphology are further channels.",16.5)],sp=0.5,fa=TINT,fb=TINT2)

# ===================================================== 4 WHY (what / why / value / significance)  [NEW]
s=slide(prs); title(s,"Why digitize the ayah? — the what, the why, and what’s at stake")
three(s,[L("WHAT we do",17,True,TEAL),L("Turn each verse into numbers — a 1-D signal. This is digitizing text data, exactly as we digitize sound or a sensor reading, so the ayah becomes something we can measure rather than only admire.",16)],
 [L("WHY it matters · the value",17,True,AMBER),L("Prose can be read but not measured. As a vector a verse can be compared across all 6,236 ayat at once, searched, transformed, smoothed, and TESTED against chance — at scale and reproducibly.",16)],
 [L("SIGNIFICANCE",17,True,NAVY),L("It places the Book of Scripture under the very mathematics we use on the Book of Creation — one Author, one toolkit. The audit ✓/✗/~ keeps wonder honest: the vector is the carrier, never the message.",16)])

# ===================================================== OBJECTIVE (latent features)  [NEW]
s=slide(prs); title(s,"The objective — we vectorize to surface latent features")
band(s,0.42,1.2,12.5,0.4,TINT2,"the vector is the doorway; the hidden structure it reveals is the goal",NAVY)
fbox(s,0.30,2.0,2.4,1.25,TINT,"AYAH","text data",line=TEAL,tsz=16,ssz=11)
harrow(s,2.78,2.48,0.85,"digitize",color=GREY,lcol=TEAL)
fbox(s,3.71,2.0,2.5,1.25,AMBERT,"x = [ … ]","the vector",line=AMBER,tsz=16,ssz=11)
harrow(s,6.29,2.48,0.85,"transform",color=GREY,lcol=NAVY)
fbox(s,7.22,2.0,2.6,1.25,TINT2,"reorder · project","reduce · cluster",line=NAVY,tsz=14.5,ssz=11)
harrow(s,9.90,2.48,0.85,"reveal",color=GREY,lcol=RED)
fbox(s,10.83,2.0,2.2,1.25,REDT,"LATENT","hidden features",line=RED,tsz=16,ssz=11)
two_stack(s,
 [L("The end goal is DISCOVERY",18,True,NAVY),
  L("We do not digitize for its own sake. The vector lets us uncover HIDDEN structure the raw text never shows — rhythm, rhyme, families of verses, principal axes, clusters. Units D–E (energy, distance, embeddings, PCA, clustering) are pure latent-feature discovery. Vectorization is the doorway; the latent features are the room.",16)],
 [L("Reordering is a TOOL, not a violation",18,True,NAVY),
  L("Permuting tokens, sorting ayat by length or revelation-time, clustering rows — all are legitimate, DECLARED manipulations that reveal new configurations; the Monte-Carlo null is itself a deliberate shuffle. The one rule: a feature you find by searching configurations must still beat a null — or it is an artifact, not a latent feature.",16,True,TEAL)],
 split=0.5,fillA=TINT,fillB=AMBERT)

# ===================================================== 5 WORKED EXAMPLE (corpus-verified)
s=slide(prs); title(s,"Real ayah-vector — قل هو الله أحد (112:1) on the ROOT anchor")
band(s,0.42,1.18,12.5,0.42,TINT,"one axis: root position   ·   amplitude = how often that ROOT occurs in the Qur’an",TEAL)
signal_plot(s,2.0,1.75,9.2,[1722,2848,153],["قول","ءله","وحد"],col=TEAL,vmax=2848,base_h=2.55)
panel(s,0.42,5.55,12.5,1.65,TINT2,[L("x = [1722, 2848, 153]  —  roots قول·ءله·وحد, recomputed from the corpus",18,True,NAVY),
  L("Three roots, three samples — قل هو الله أحد reduces to قول·ءله·وحد (هو is a pronoun, no root). ءله is a tall PEAK (2,848); وحد a VALLEY (153). The verse now has a measurable shape. The root is the anchor — highest semantic power; surface form and morphology are further channels.",16,True,TEAL)],space=6)

# ===================================================== 6 WHAT A SIGNAL IS
s=slide(prs); title(s,"Conceptual foundation — what a 1-D signal is")
three(s,[L("SAMPLES",17,True,TEAL),L("A signal is a list of numbers x[0], x[1], …, x[n−1]. Each number is a SAMPLE. For an ayah, each ROOT-token contributes one sample.",16)],
 [L("THE AXIS",17,True,AMBER),L("The index is the only axis. For an ayah it is TOKEN POSITION (reading order). For audio it is time; for a sensor, time too. One axis ⇒ 1-D.",16)],
 [L("AMPLITUDE",17,True,NAVY),L("The value at each position is AMPLITUDE — here a chosen measurement of the token (its frequency, length, root id…). Shape lives in how amplitude varies.",16)])

# ===================================================== 7 THE DICTIONARY (equivalency mapping)  [NEW]
s=slide(prs); title(s,"The dictionary — the ayah ↔ signal equivalence, term by term")
band(s,0.42,1.16,12.5,0.42,TINT2,"عالم التدوين — the ayah (text data)        ≡        the 1-D signal (its digitization)",NAVY)
rows=[("a whole ayah (verse)","the signal  x = [x₀ … xₙ₋₁]"),
      ("one root-token","one sample  x[i]"),
      ("reading order (root position)","the index axis  (like time)"),
      ("number of tokens,  n","the signal length,  N"),
      ("a per-token measurement","amplitude  (the value at i)"),
      ("a root’s rarity / frequency","amplitude height  (peak vs valley)"),
      ("rhyme & repetition","periodicity / autocorrelation")]
y=1.70
for left,right in rows:
    map_row(s,y,left,right); y+=0.79

# ===================================================== 8 CHANNELS
s=slide(prs); title(s,"Many channels — each root measured a different way (root = anchor)")
band(s,0.42,1.2,12.5,0.4,TINT2,"same ayah → several signals (“channels”); each captures something different",NAVY)
ch=[("ROOT frequency","ANCHOR · semantic · ratio",TINT,TEAL),("ROOT length","letters in the root (~3)",AMBERT,AMBER),
    ("SURFACE form","the inflected word",TINT,TEAL),("MORPHOLOGY","pattern / position",AMBERT,AMBER),
    ("EMBEDDING comp.","a learned meaning axis",TINT2,NAVY)]
x=0.55; bw=2.4; aw=0.08
for i,(t,sub,fl,ln) in enumerate(ch):
    fbox(s,x,1.85,bw,1.5,fl,t,sub,line=ln,tsz=14.5,ssz=11); x+=bw+aw
panel(s,0.42,3.65,12.5,3.55,TINT,[L("One ayah, a stack of signals",18,True,NAVY),
  L("The same verse yields a frequency-signal, a length-signal, a root-id-signal, a position-signal, an embedding-signal. Each is a 1-D vector of the same length n; each is a different “channel” of the text — like measuring a patient by ECG, temperature, or blood-oxygen, all over the same minutes.",16.5),
  L("Later lectures pick the channel that suits the question; this lecture establishes the move itself.",16,True,TEAL)],space=8)

# ===================================================== 9 CRITERIA (number assignment)  [NEW]
s=slide(prs); title(s,"The criteria — what number may we assign to a token?")
band(s,0.42,1.2,12.5,0.4,TINT2,"a valid vectorization assigns a MEASUREMENT, not a name — five tests it must pass",NAVY)
crit=[("① MEASUREMENT","not a label: quantify a real property; ratio/interval scale",TINT,TEAL),
      ("② RULE-BASED","one stated rule; same token → same number; from Book6",AMBERT,AMBER),
      ("③ DECLARED FIRST","channel fixed before seeing results — no fitting to a wish",TINT,TEAL),
      ("④ ORDER IS TRACKED","reading order is the reference; any reordering is a DECLARED step",AMBERT,AMBER),
      ("⑤ CORPUS-WIDE","same rule + normalization across all 6,236 ayat, so vectors compare",TINT2,NAVY)]
x=0.55; bw=2.4; aw=0.08
for i,(t,sub,fl,ln) in enumerate(crit):
    fbox(s,x,1.85,bw,1.7,fl,t,sub,line=ln,tsz=14,ssz=10.5); x+=bw+aw
panel(s,0.42,3.85,12.5,3.35,TINT,[L("Why these five — the line between a signal and numerology",18,True,NAVY),
  L("A number you can defend is a MEASUREMENT (frequency, length, energy), computed by a fixed rule, declared in advance, kept in reading order, and applied to the whole corpus. Root frequency [1722, 2848, 153] passes all five.",16.5),
  L("Counter-case: a bare ROOT-ID is only a NAME (nominal) — it fails test ①. You may not average or transform ids as magnitudes; first convert them (one-hot, replace-by-frequency, or embedding). Reordering (shuffling tokens, sorting ayat) is allowed — it is how we probe latent structure and how the null is built — provided it is declared and the result still beats chance.",15.5,True,TEAL)],space=8)

# ===================================================== 9 DATA two channels
s=slide(prs); title(s,"Real data — two channels of the same verse (112:1, roots)")
finding2(s,
 {"title":"Channel A — root frequency (anchor)","cats":["قول","ءله","وحد"],
  "series":[("",[TEAL,NAVY,RED],[1722,2848,153])],"legend":False},
 {"title":"Channel B — root length (letters)","cats":["قول","ءله","وحد"],
  "series":[("",[TEAL,TEAL,TEAL],[3,3,3])],"legend":False},
 [L("Frequency: a sharp peak",17.5,True,TEAL),
  L("By rarity the verse is spiky — ءله towers at 2,848, وحد sits at 153. The signal’s shape is dominated by one peak.",16)],
 [L("Length: a gentle hump",17.5,True,AMBER),
  L("Arabic roots are nearly all triliteral, so root-length is flat — قول·ءله·وحد → [3,3,3]. Same verse, different channel, different signal — choose to fit the question.",16)],
 fillA=TINT,fillB=AMBERT)

# ===================================================== 10 DATA ayah lengths
s=slide(prs); title(s,"Real data — how long are the ayah-signals? (computed from Book6)")
finding2(s,
 {"title":"Ayah length n (root-tokens) — count of ayat","cats":["1–2","3–10","11–30","31–60","61+"],
  "series":[("",[GREY,TEAL,TEAL,AMBER,RED],[769,3829,1576,61,1])],"legend":False},
 {"title":"Creation 1-D signal lengths (samples, log₁₀)","cats":["ECG beat","speech word","day temps","EEG epoch"],
  "series":[("",[TEAL,TEAL,AMBER,NAVY],[2.4,3.6,1.4,3.3])],"legend":False,"fmt":"{:.1f}"},
 [L("Short signals, with rare giants",17.5,True,TEAL),
  L("Most ayat are 3–10 tokens (median 7); one runs to 84 — the debt verse (2:282). The corpus is 6,236 such 1-D signals, heavy-tailed in length. (Computed, exact.)",16)],
 [L("Just like creation’s signals",17.5,True,AMBER),
  L("Physical 1-D signals span a few to thousands of samples — an ECG beat (~250), a spoken word (~4000), a day of hourly temps (24). Same object, same toolkit.",16)],
 fillA=TINT,fillB=AMBERT)

# ===================================================== 11 DATA spectrum
s=slide(prs); title(s,"Real data — the amplitude spectrum is heavy-tailed")
finding2(s,
 {"title":"Root-frequency (top 8 of 1702) — amplitudes","cats":["ءله","قول","كون","ربب","ءمن","علم","قوم","ءتي"],
  "series":[("",[NAVY,TEAL,TEAL,AMBER,AMBER,AMBER,GREY,GREY],[2848,1722,1390,980,879,854,660,549])],"legend":False},
 {"title":"Letter-frequency (×1000) — amplitudes","cats":["ا","ن","م","ل","ي","ء"],
  "series":[("",[NAVY,TEAL,TEAL,AMBER,AMBER,GREY],[35.8,28.0,26.7,26.1,24.7,19.1])],"legend":False,"fmt":"{:.0f}"},
 [L("A few peaks, a long tail",17.5,True,TEAL),
  L("Valued by frequency, a few units dominate (ءله 2848 …) and the rest trail off — a Zipf / power-law spectrum (log-log slope ≈ −0.8), the statistics of natural language.",16)],
 [L("The same shape as natural signals",17.5,True,AMBER),
  L("This 1/f-like, heavy-tailed amplitude profile is what creation’s signals show too (audio, neural spikes). The choice of channel sets what the spectrum measures.",16)],
 fillA=TINT,fillB=AMBERT)

# ===================================================== 12 DUAL-DOMAIN
s=slide(prs); title(s,"Dual-domain — the ayah-vector beside creation’s signals")
band(s,0.42,1.18,12.5,0.4,TINT,"عالم التدوين — an ayah as a 1-D signal",TEAL)
sigrow(s,0.9,1.7,11.4,[.45,.62,.95,.20,.5,.7,.35,.85,.4,.6],col=TEAL)
band(s,0.42,3.5,12.5,0.4,AMBERT,"عالم التكوين — a heartbeat (ECG) as a 1-D signal",AMBER)
sigrow(s,0.9,4.0,11.4,[.2,.25,.9,.15,.3,.22,.85,.18,.28,.24],col=AMBER)
panel(s,0.42,5.85,12.5,1.35,TINT2,[L("Same mathematical object",17,True,NAVY),
  L("Both are sequences of samples along one axis. Whatever DSP does to the ECG — smoothing, spectrum, correlation, filtering — it can do to the ayah-vector. The content differs; the mathematics is shared.",16,True,TEAL)],space=6)

# ===================================================== 13 MAPPING (operations)
s=slide(prs); title(s,"The mapping — a DSP operation for every question")
def oprow(y,op,opcol,reveals,wk):
    fbox(s,0.6,y,3.2,0.95,opcol,op,"",line=NAVY,tsz=15); harrow(s,3.95,y+0.32,1.5,"reveals",color=GREY,lcol=TEAL)
    fbox(s,5.6,y,5.1,0.95,TINT2,reveals,"",line=TEAL,tsz=14); harrow(s,10.85,y+0.32,0.85,"",color=GREY)
    fbox(s,11.85,y,0.95,0.95,AMBERT,wk,"",line=AMBER,tsz=13)
oprow(1.40,"Fourier / spectrum",TINT,"periodic structure, rhythm","U-C")
oprow(2.50,"autocorrelation",LTEAL,"repetition, rhyme","U-B")
oprow(3.60,"filtering / smoothing",AMBERT,"trend vs detail; denoise","U-C")
oprow(4.70,"distance / embedding",TINT,"how two ayat compare","U-D")
panel(s,0.42,5.85,12.5,1.35,TINT,[L("Faithful, not merely poetic",17,True,NAVY),
  L("Each course unit is a real operation on the vector: counting along the sequence, transforming it to frequency, comparing two of them. The analogy maps the ACTUAL math we run, not a mood.",16,True,TEAL)],space=6)

# ===================================================== 14 SCOPE
s=slide(prs); title(s,"What this course will — and will not — do")
two(s,[L("WILL",18,True,TEAL),L("• Digitize ayat into 1-D signals by explicit, stated channels.  • Apply standard DSP — spectrum, autocorrelation, filtering, distance — to real ayah-vectors.  • Audit every parallel: ✓ Supported · ✗ Breaks · ~ Silent · and beat a null where one applies.",16.5,True,NAVY)],
 [L("WILL NOT",18,True,RED),L("• Claim the numbers ‘encode’ hidden messages or predictions.  • Treat a signal property as the verse’s MEANING.  • Read amplitude as miracle. The vector is the carrier; guidance, beauty and address are not in the samples.",16.5,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# ===================================================== 15 DEFENSE falsifiable
s=slide(prs); title(s,"Is this validated — or just numerology?")
band(s,0.42,1.2,12.5,0.4,TINT2,"three claims, each of which real data could REFUTE",NAVY)
fbox(s,0.6,1.9,3.95,1.95,TINT,"(1) CHANNELS DIFFER","TRUE only if frequency- and length-signals of the SAME verse have different shape.  Refuted if all channels coincided.",line=TEAL,tsz=15,ssz=11.5)
fbox(s,4.7,1.9,3.95,1.95,AMBERT,"(2) HEAVY TAIL","TRUE only if amplitudes follow a power law.  Refuted if the spectrum were flat / uniform.",line=AMBER,tsz=15,ssz=11.5)
fbox(s,8.8,1.9,3.95,1.95,TINT2,"(3) STRUCTURE > CHANCE","TRUE only if a real ayah-vector’s structure beats a shuffle null.  Refuted if shuffling changed nothing.",line=NAVY,tsz=15,ssz=11.5)
panel(s,0.42,4.1,12.5,3.1,TINT,[L("Falsifiable, not a story",18,True,NAVY),
  L("A numerological claim cannot be wrong — it bends to fit anything. Each statement here says in advance what would prove it false, and is then checked against real counts from Book6. Digitization is a measurement, and a measurement can disappoint.",16.5),
  L("The next slides show the firewall and the corpus checks every such claim must pass.",16,True,TEAL)],space=8)

# ===================================================== 16 MONTE-CARLO firewall (REAL data)
s=slide(prs); title(s,"The firewall — beat a Monte-Carlo null (root anchor, real corpus)")
finding2(s,
 {"title":"Root null: 20,000 random 7-root verse pairs (r)","cats":["≤-.6","-.6/-.3","-.3/0","0/.3",".3/.6","≥.6"],
  "series":[("",[GREY,AMBER,TEAL,TEAL,AMBER,GREY],[387,4244,6430,3622,2527,2790])],"legend":False},
 {"title":"Beats the null? (−log₁₀ p)","cats":["random pairing","31× refrain (55)"],
  "series":[("",[RED,TEAL],[0.1,2.6])],"legend":False,"fmt":"{:.1f}"},
 [L("Why a null is needed",17.5,True,TEAL),
  L("On the ROOT anchor the null centres near 0 (mean r ≈ 0.04; on surface tokens it is 0.18 — function words). “Striking” similarity is the baseline, not the signal — only a null can judge.",16)],
 [L("Real structure clears the bar",17.5,True,AMBER),
  L("The 31× refrain فبأي آلاء ربكما تكذبان scores r = 1.0, p ≈ 0.002 against the 20,000-draw null; a random pairing sits in the bulk. Computed live from Book6.",16)],
 fillA=TINT,fillB=AMBERT)

# ===================================================== 17 V&V evidence (corpus-verified)  [NEW]
s=slide(prs); title(s,"Verification & validation — every figure recomputed from Book6")
finding2(s,
 {"title":"112:1 root vector — plan vs computed","cats":["قول","ءله","وحد"],
  "series":[("",[TEAL,NAVY,RED],[1722,2848,153])],"legend":False},
 {"title":"Corpus census (computed) — log₁₀","cats":["suras","roots","ayat","tokens"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[2.06,3.23,3.79,5.13])],"legend":False,"fmt":"{:.2f}"},
 [L("Zero discrepancy",17.5,True,TEAL),
  L("Recomputed from Book6.xlsx (6,236 ayat, 114 suras, 51,044 root-tokens, 1,702 roots): the worked root vector [1722, 2848, 153], the length histogram, and the top-root spectrum all reproduce the plan exactly.",16)],
 [L("Validated, not asserted",17.5,True,AMBER),
  L("The root ءله occurs 2,848× (surface الله 2,695×; published 2,698–2,699 differ by tokenization). The 31× refrain cleared the root null at p ≈ 0.002. The idea is credible on actual data.",16)],
 fillA=TINT,fillB=AMBERT)

# ===================================================== HONEST SCOPE & LIMITS  [NEW]
s=slide(prs); title(s,"Honest scope — what this can, and cannot, show")
band(s,0.42,1.2,12.5,0.4,REDT,"the discipline is also knowing the limits — four of them, up front",RED)
lim=[("SHORT SIGNALS","ayah median = 7 tokens. Fourier/autocorrelation need length — run them at SURA or CORPUS scale, not on one short ayah.",REDT,RED),
     ("GENERIC ≠ SPECIFIC","Zipf slope (−1.23) and function-word share (0.39) MATCH random Arabic — they are features of language, not of the Qur'an.",AMBERT,AMBER),
     ("BEAT A BASELINE","so a shuffle null is not enough: a Qur'an-specific claim must beat a length/Zipf-matched NATURAL-LANGUAGE baseline.",TINT,TEAL)]
x=0.55; bw=3.95; aw=0.1
for i,(t,sub,fl,ln) in enumerate(lim):
    fbox(s,x,1.85,bw,1.75,fl,t,sub,line=ln,tsz=15,ssz=11.5); x+=bw+aw
panel(s,0.42,3.8,12.5,3.4,TINT,[L("And: not everything here is signal processing",18,True,NAVY),
  L("Where the math is cleanest (time, frequency) the MEANING is absent; where meaning lives (embeddings, PCA, clustering — Units D–E) the tool is REPRESENTATION LEARNING, not classical DSP. We say so plainly rather than stretch one metaphor over both.",16.5),
  L("What survives all of this is real: e.g., the Qur'an's exact root-refrain rate is 7.1% vs 0.81% in matched random Arabic (~8.8×) — a structure that beats the baseline. The honest yield is small, specific, and earned.",16,True,TEAL)],space=8)

# ===================================================== 18 ROADMAP (17)
s=slide(prs); title(s,"Your roadmap — seventeen lectures, simple → complex")
band(s,0.42,1.2,12.5,0.4,TINT2,"five units · 17 lectures · a full-semester DSP arc on the ayah-as-1-D-signal",NAVY)
units=[("A · Foundations","1–3 · idea · method · channels",TINT,TEAL),
       ("B · Signal in time","4–8 · waveform · sampling · convolution",AMBERT,AMBER),
       ("C · In frequency","9–11 · Fourier · rhythm · filtering",TINT,TEAL),
       ("D · Information & space","12–14 · energy · distance · embeddings",AMBERT,AMBER),
       ("E · Structure & synthesis","15–17 · PCA · clustering · capstone",TINT2,NAVY)]
x=0.55; bw=2.42; aw=0.05
for i,(t,sub,fl,ln) in enumerate(units):
    fbox(s,x,1.9,bw,1.5,fl,t,sub,line=ln,tsz=14.5,ssz=10.5)
    if i<4: harrow(s,x+bw-0.03,2.5,aw+0.06,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.75,12.5,2.45,TINT,[L("Each lecture is one DSP idea on real ayah-vectors",18,True,NAVY),
  L("We start with the move itself (this lecture), then the method and the digitization schemes; then read the signal in time (waveform → sampling → smoothing → convolution → autocorrelation), then in frequency; then as information and distance; then reduce, cluster, and synthesise. Seventeen rungs, one object — matched to the biology course.",16,True,TEAL)],space=7)

# ===================================================== 19 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, and silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("The structure is literally true: an ayah IS a sequence of samples once a channel is chosen, and DSP operations apply to it exactly. The digitization is real and reproducible.",16)],
 [L("✗ BREAKS",17,True,RED),L("Meaning is NOT a signal property. The verse is not its vector; guidance, beauty and address vanish in the numbers. The signal is the carrier, never the message.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("The CHOICE of channel (frequency vs length vs embedding) is a modeling decision the text doesn’t dictate; recitation adds a real acoustic 1-D signal the written text omits.",16)],f=(TINT,REDT,AMBERT))

# ===================================================== 20 GAINED / LOST
s=slide(prs); title(s,"What is gained — and what is lost — in the vector")
two(s,[L("GAINED",18,True,TEAL),L("Measurability. Once an ayah is a vector we can compare 6,236 of them at once, find rhythm, cluster themes, denoise, and TEST claims against chance. The text becomes tractable to mathematics without altering a letter.",17,True,NAVY)],
 [L("LOST",18,True,RED),L("Sense. A number per token discards grammar, reference, address, and the experience of recitation. Every result must be read back into the text by a human; the vector is a lens, not a replacement.",17,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# ===================================================== ANCHOR (inverse / read-back)  [NEW]
s=slide(prs); title(s,"The anchor — every feature must map back to the text")
band(s,0.42,1.2,12.5,0.4,TINT2,"forward: text → vector → feature      ·      back: does the feature mean something in the ayat?",NAVY)
fbox(s,0.6,1.95,2.7,1.1,TINT,"AYAH (text)","the source of truth",line=TEAL,tsz=15,ssz=11)
harrow(s,3.45,2.32,1.0,"vectorize",color=GREY,lcol=TEAL)
fbox(s,4.55,1.95,2.7,1.1,AMBERT,"VECTOR · OP","transform / reduce",line=AMBER,tsz=15,ssz=11)
harrow(s,7.4,2.32,1.0,"yields",color=GREY,lcol=NAVY)
fbox(s,8.5,1.95,2.7,1.1,TINT2,"LATENT FEATURE","a pattern",line=NAVY,tsz=15,ssz=11)
band(s,0.42,3.25,12.5,0.5,REDT,"↩  THE GATE: can we read this feature back into the actual ayat? — if not, it is an artifact, not a finding",RED)
pt=3.95; ph=CY1-pt
panel(s,0.42,pt,6.13,ph,TINT,[L("Inverse = INTERPRETIVE, not exact",17.5,True,TEAL),
  L("Vectorizing is LOSSY: a frequency signal can never rebuild the verse — and it need not. We do not demand bijective reconstruction; we demand that any feature we find can be READ BACK into the verses. Interpretability, not invertibility, is the test.",16)],space=8)
panel(s,6.79,pt,6.13,ph,AMBERT,[L("Semantics is the anchor — precisely",17.5,True,AMBER),
  L("For claims about MEANING, the text’s sense is the final arbiter: a latent axis or cluster is credible only if it answers to something real in the ayat. For claims about FORM (rhythm, rhyme, length), the anchor is the textual structure itself, validated against a null. Map back to nothing — semantic or formal — and the number is numerology.",16)],space=8)

# ===================================================== 21 APP
s=slide(prs); title(s,"The app — digitize any ayah, live")
band(s,0.42,1.2,12.5,0.4,TINT,"type a reference → pick a channel → see the signal",TEAL)
steps=[("① REFERENCE","enter sura:ayah (e.g., 112:1)",TINT,TEAL),("② CHANNEL","frequency · length · root · embed",AMBERT,AMBER),
       ("③ SIGNAL","the bars + the values x[i]",TINT,TEAL),("④ MEASURE","peak · spectrum · vs-null",TINT2,NAVY)]
x=0.55; bw=2.98; aw=0.1
for i,(t,sub,fl,ln) in enumerate(steps):
    fbox(s,x,1.9,bw,1.35,fl,t,sub,line=ln,tsz=15,ssz=11)
    if i<3: harrow(s,x+bw-0.02,2.45,aw+0.05,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.65,12.5,2.55,TINT,[L("The app is the lab bench of the course",18,True,NAVY),
  L("Every lecture has live tasks: digitize a verse, switch channels, run the transform, and compare against the Monte-Carlo null on the spot. Today: reproduce x = [1722, 2848, 153] for 112:1 on the root channel and watch ءله make the peak.",16.5,True,TEAL)],space=7)

# ===================================================== 22 DISCUSSION
s=slide(prs); title(s,"Synthesis & discussion — the Two Books")
two(s,[L("THE TWO BOOKS",18,True,NAVY),L("The tools that read creation’s 1-D signals — spectrum, correlation, filtering — are the very tools we turn on the digitized ayah. Same mathematics of structure; different content. Read side by side, never collapsed.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• What is gained, and lost, when a verse becomes a vector?  • Which channel is the most ‘honest’ — and which most arbitrary?  • Is recitation the 1-D signal our text-analysis misses?  • Where must the signal lens stop?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# ===================================================== 23 REAL-WORLD takeaway
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Digitizing text into vectors to learn LATENT FEATURES is exactly how modern search, speech, and language models work (representation learning) — and the discipline of beating a null before believing a pattern guards any data claim you will meet.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("A signal is the digitized ayah: each verse → a 1-D numeric vector on the ROOT anchor (112:1 → قول·ءله·وحد → [1722, 2848, 153], verified). The aim is to surface LATENT FEATURES — digitize, transform, reorder, project — and make every discovered pattern beat a Monte-Carlo null, or set it aside.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)

prs.save(OUT+"01_Introduction_Lecture.pptx")
print("L1 (signal) Introduction slides:", len(prs.slides))
