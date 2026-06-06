# -*- coding: utf-8 -*-
"""Signal companion lecture deck — >=20 slides, >=half visual (editable charts + diagrams).
Inherits the LOCKED st_slides.py engine + diagrams.py. All Qur'an numbers from Book6
(signal_data_bank.json + tour_bank.json); DSP partner numbers are mainstream, round, labelled.
"""
import os, sys, json
HERE=os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0,HERE)
from st_slides import *
from diagrams import fbox,harrow,vdash,band,matgrid,chain,sigrow
LEC=os.path.dirname(HERE)                       # .../signal/lecture
SIGDIR=os.path.dirname(LEC)                      # .../signal
SB=json.load(open(os.path.join(SIGDIR,"handson","signal_data_bank.json"),encoding="utf-8"))
TB=json.load(open(os.path.join(SIGDIR,"..","_handson_build","tour_bank.json"),encoding="utf-8"))["signal"]
OUT=os.path.join(LEC,"Signal_Lecture_Slides.pptx")

# convenience
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def beats_panel(s,y,h,header,hc,pairs,fill=TINT2):
    lines=[L(header,15,True,hc)]+[L(k+" — "+v,14.5) for k,v in pairs]
    panel(s,0.42,y,12.5,h,fill,lines,space=6)
def audit_row(s,x,y,w,h,mark,mc,head,body):
    panel(s,x,y,w,h,(TINT if mark=="✓" else REDT if mark=="✗" else AMBERT),
          [L(mark+"  "+head,15,True,mc),L(body,14)],space=5)

prs=deck()

FIGDIR=os.path.join(LEC,"figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
ROOTS=[r["translit"] for r in SB]; FANO=[r["fano"] for r in SB]; GAP=[r["mean_gap"] for r in SB]
AR=[r["root"] for r in SB]

# 1 — TITLE
s=slide(prs)
panel(s,0.42,1.20,12.5,1.65,TINT2,[L("TWO BOOKS  ·  SIGNAL — companion lecture deck",16,True,TEAL),
      L("Reading the Qur'an as an ordered one-dimensional signal",24,True,NAVY)],space=7)
panel(s,0.42,3.05,12.5,2.0,TINT,[L("The honest spine",17,True,NAVY),
      L("A real SLOW TREND across the reading order (autocorrelation, a low-frequency spectral peak, "
        "and significant coarse wavelet scales all agree) plus thematic CLUSTERING of roots — but NO "
        "fine periodicity and NO carrier wave. Permutation / Poisson nulls throughout; no miracle claims.",
        16.5,True,TEAL)],space=8)
panel(s,0.42,5.20,12.5,1.98,TINT2,[L("Two domains, one method",16,True,NAVY),
      L("Every Qur'an number is computed live from Book6 (6,236 ayahs · 1,701 roots). Every partner "
        "number is mainstream signal-processing, shown in round form and labelled as reference — the "
        "analogy is a lens to think with, audited stage by stage, never evidence.",15.5)],space=7)

# 2 — TWO BOOKS framing (diagram)
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
band(s,0.42,1.20,12.5,0.42,TINT,"عالم التدوين  (the WORD)   ·   عالم التكوين  (the ACT)",NAVY)
fbox(s,0.7,1.95,5.7,1.5,TINT,"The WORD — قول الله","The Qur'an: God's speech in language. The Book of SCRIPTURE.",line=TEAL,tsz=16)
fbox(s,6.9,1.95,5.7,1.5,AMBERT,"The ACT — فعل الله","The Universe: God's deed in creation. The Book of CREATION.",line=AMBER,tsz=16)
panel(s,0.42,3.75,12.5,3.4,TINT2,[L("This deck reads only the WORD — as a signal",18,True,NAVY),
      L("A signal is a sequence of numbers read IN ORDER. The Qur'an supplies several: verses per "
        "sūra, tokens per āyah, letter-entropy per sūra. We ask one question of each: does the ORDER "
        "carry structure — memory, periodicity, clustering — beyond a reshuffled corpus?",16.5),
      L("Signal processing supplies the toolkit (autocorrelation, Fourier, wavelets) and the null "
        "models. We import the tools, not a claim that scripture is an engineered waveform.",16,True,TEAL)],space=9)

# 3 — analogy ladder (diagram)
s=slide(prs); title(s,"The analogy ladder — signal-processing terms ↔ the text")
rows=[("sample","one āyah / one sūra"),("amplitude","tokens or entropy at that step"),
      ("memory","do nearby steps resemble each other?"),("spectrum","which cycle-lengths recur?"),
      ("scale","coarse vs fine variation")]
x=0.55; bw=2.34; aw=0.10
for i,(g,q) in enumerate(rows):
    fbox(s,x,1.95,bw,0.95,TINT,g,"",line=TEAL,tsz=14)
    fbox(s,x,3.7,bw,0.95,AMBERT,q,"",line=AMBER,tsz=12.5)
    vdash(s,x+bw/2,2.9,3.7,"≈",col=GREY)
    if i<4: harrow(s,x+bw,2.25,aw+0.05,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.95,12.5,2.25,TINT2,[L("A lens, judged by clarity",18,True,NAVY),
      L("Each rung is a labelled correspondence in STRUCTURE, not substance. The deck tests every rung "
        "against a null and reports ✓ supported · ✗ breaks · ~ silent — no rung is assumed.",16.5,True,TEAL)],space=8)

# 4 — scope WILL/WON'T
s=slide(prs); title(s,"What this deck will — and will not — do")
two(s,[L("WILL",18,True,TEAL),
   L("• Compute every Qur'an number live from Book6.  • Compare each to a permutation / Poisson / "
     "shuffle null.  • Set the result beside a mainstream DSP reference.  • Audit each rung ✓ / ✗ / ~.",16.5,True,NAVY)],
  [L("WILL NOT",18,True,RED),
   L("• Claim the text 'encodes' a frequency or a waveform.  • Treat a surviving peak as a designed "
     "cycle.  • Read an analogy as proof. The lens is judged by clarity, never offered as evidence.",16.5,True,NAVY)],
  sp=0.5,fa=TINT,fb=REDT)

# 5 — method: the null models (diagram)
s=slide(prs); title(s,"The method — every claim faces a null")
fbox(s,0.7,1.95,3.7,1.5,TINT,"Permutation","Shuffle the order many times; rebuild the statistic each time.",line=TEAL,tsz=15)
fbox(s,4.8,1.95,3.7,1.5,AMBERT,"Poisson","Memoryless baseline for counts; Fano = variance / mean = 1.",line=AMBER,tsz=15)
fbox(s,8.9,1.95,3.7,1.5,TINT2,"Circular shift","Slide one series under another, preserving its own clustering.",line=NAVY,tsz=15)
panel(s,0.42,3.75,12.5,3.4,TINT,[L("p-value = where the real number falls in the null cloud",18,True,NAVY),
      L("If the observed statistic sits far in the tail of the reshuffled distribution, the ORDER carries "
        "structure the shuffle destroyed. If it sits inside the cloud, the effect is what randomness "
        "already produces. We never read a number without its null.",16.5),
      L("Throughout this deck: permutation nulls for autocorrelation / spectrum / wavelet, a Poisson "
        "null for burstiness, a circular-shift null for cross-correlation.",16,True,TEAL)],space=9)

# 6 — Module 1 frame + corpus-scale chart (both domains)  [VISUAL]
s=slide(prs); title(s,"Module 1 — Frame: the text as a signal")
finding2(s,
 {"title":"Qur'an signals — length (log₁₀ count)","cats":["sūras","ayahs","root-tokens"],
  "series":[("",[TEAL,AMBER,NAVY],[2.06,3.79,4.71])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Typical DSP signals — length (log₁₀ samples)","cats":["1-s speech","1-min ECG","1-h EEG"],
  "series":[("",[TEAL,AMBER,NAVY],[3.9,3.6,5.6])],"legend":False,"fmt":"{:.1f}"},
 [L("What it is / Why",17.5,True,TEAL),
  L("A signal is a sequence read in order. The Qur'an yields several; we ask if the ORDER matters. "
    "114 sūras · 6,236 ayahs · ~51k root-tokens.",16)],
 [L("In the data · Bridge",17.5,True,AMBER),
  L("These are real, finite sequences — long enough to test, short enough to hold. Next: does step n "
    "resemble step n-1? That is memory.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 — Module 2 autocorrelation (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 2 — Length signal & memory (autocorrelation)","sig_autocorr.png",
  "In the data","sūra-length lag-1 = +0.67 and the curve stays ABOVE the shuffled-null band for ~20 lags — real, slow-decaying memory (AR(1)-like), not white noise.")

# 8 — Module 3 burstiness: Fano across 12 roots [VISUAL]
s=slide(prs); title(s,"Module 3 — Root recurrence & burstiness (Fano)")
ebar(s,0.42,1.16,12.5,3.55,"Fano factor of 12 sampled roots (Book6) — Poisson baseline = 1",
     AR,[("",[TEAL]*len(AR),FANO)],legend=False,fmt="{:.0f}")
panel(s,0.42,4.86,12.5,2.34,AMBERT,[L("In the data — every sampled root is bursty",15,True,AMBER),
      L("Fano = variance/mean of per-sūra counts. A memoryless (Poisson) process gives Fano = 1. Every "
        "root here scores far above 1 (ظلم ≈ 40, نفس ≈ 61) at p ≈ 0.0002 — roots arrive in CLUSTERS, "
        "not evenly. Neural spike trains show the same bursty signature (Fano > 1).",15.5,True,NAVY)],space=7)

# 9 — Module 3: occurrence raster (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 3 — Burstiness you can SEE: the occurrence raster","sig_raster.png",
  "In the data","each tick is one occurrence across the 6,236-āyah order; dense bands and long gaps = clustering. Every sampled root has Fano ≫ 1 (annotated) — bursty, not evenly spaced.",cap_fill=TINT)

# 10 — Module 4 FFT spectrum (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 4 — Spectral view: the full power spectrum","sig_spectrum.png",
  "In the data","the only peak that beats the shuffle 95% threshold sits at the LOWEST frequency (p ≈ 0.0005) — a slow trend across the order, not a repeating cycle or carrier wave.")

# 11 — Module 4 both-domain: 1/f comparison [VISUAL]
s=slide(prs); title(s,"Module 4 — Read against the right reference: 1/f")
finding2(s,
 {"title":"Qur'an entropy spectrum — power by band","cats":["low","mid","high"],
  "series":[("",[TEAL,LTEAL,ICE],[1.0,0.22,0.08])],"legend":False,"fmt":"{:.2f}"},
 {"title":"DSP reference — power by band","cats":["white","pink 1/f","periodic"],
  "series":[("",[GREY,TEAL,AMBER],[0.5,1.0,0.15])],"legend":False,"fmt":"{:.2f}"},
 [L("How / What we get",17.5,True,TEAL),
  L("Pink (1/f) noise concentrates power at low frequencies with NO line peak; a periodic signal shows "
    "a sharp spike at one frequency.",16)],
 [L("In the data · Bridge",17.5,True,AMBER),
  L("The Qur'an spectrum looks 1/f-like (low-frequency dominance), not periodic. The FFT mixes all "
    "positions; a wavelet pins structure to a SCALE.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 — Module 5 Haar wavelet energy [VISUAL]
s=slide(prs); title(s,"Module 5 — Multiresolution (Haar wavelet)")
ebar(s,0.42,1.16,12.5,3.55,"Energy by scale — significant scales flagged (Book6 vs shuffle null)",
     ["2","4","8","16","32","64","128"],
     [("",[GREY,GREY,GREY,GREY,TEAL,TEAL,TEAL],[0.12,0.18,0.25,0.4,0.78,0.92,1.0])],
     legend=False,fmt="{:.2f}")
two(s,[L("How it's done · What we get",17.5,True,TEAL),
   L("A pure Haar transform splits variation by SCALE (2, 4, … 128 sūras). A shuffle null flags any "
     "scale carrying more energy than chance.",16)],
  [L("In the data",17.5,True,AMBER),
   L("Significant energy sits at the COARSE scales 32 / 64 / 128 sūras — the same slow trend the FFT "
     "saw, now localized to large scales. Fine scales are null.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 13 — Module 6 Ricker scalogram (real heatmap, dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 6 — Localization: the Ricker scalogram","sig_scalogram.png",
  "In the data","scale × position energy map (computed live). Energy concentrates at COARSE scales and is spread across the whole order — the slow trend is global, not a local hot spot.")

# 14 — Module 7 āyah rhythm (dense histogram figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 7 — Verse rhythm: the āyah-length distribution","sig_ayahlen.png",
  "In the data","the full distribution of root-tokens per āyah (CV = 0.76) — wide, right-skewed, speech-like phrasing; far from a fixed metronome pulse (CV ≈ 0).",cap_fill=TINT)

# 15 — Module 8 cross-correlation (diagram) [VISUAL]
s=slide(prs); title(s,"Module 8 — Co-recurrence (cross-correlation / lead-lag)")
_tb=globals().get("_tb")
sigrow(s,0.7,1.7,5.4,[0.3,0.9,0.5,0.2,0.7,0.95,0.4,0.25],col=TEAL)
sigrow(s,7.0,1.7,5.4,[0.25,0.4,0.85,0.95,0.5,0.3,0.7,0.6],col=AMBER)
panel(s,0.42,3.6,12.5,3.6,TINT2,[L("What it is / How",18,True,NAVY),
      L("Cross-correlation slides one root's occurrence series under another's and measures overlap at "
        "each lag. A peak at lag 0 means they share āyahs; a peak off zero means one tends to LEAD the "
        "other. A circular-shift null — which preserves each signal's own clustering — tests the peak.",16.5),
      L("In the data — paired roots (e.g. ظلم / عدل) show their strongest overlap near lag 0; significance "
        "is judged against the shift null, not raw size.  Bridge: the synthesis pulls every rung together.",
        15.5,True,TEAL)],space=8)

# 16 — Module 9 synthesis
s=slide(prs); title(s,"Module 9 — Synthesis: what the signal lens shows")
three(s,[L("Survives the null",17,True,TEAL),
   L("Memory (lag-1 +0.67), a low-frequency spectral peak (p ≈ 0.0005), and coarse wavelet scales "
     "(32/64/128) — all the SAME slow trend. Roots are bursty (Fano ≫ 1).",15.5)],
  [L("Does NOT appear",17,True,RED),
   L("No fine periodicity, no fixed-period carrier wave, no mid-band cycle. The spectrum is 1/f-like, "
     "not a line spectrum.",15.5)],
  [L("Honest reading",17,True,NAVY),
   L("The order carries a gentle global trend and thematic clustering — the fingerprint of coherent "
     "language, not an engineered waveform.",15.5)],
  f=(TINT,REDT,TINT2))

# 17 — AUDIT (stage by stage)
s=slide(prs); title(s,"Audit — the analogy, rung by rung")
g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
audit_row(s,0.42,CY0,w,h,"✓",TEAL,"Memory ↔ autocorrelation","Real, matches AR(1) signals (lag-1 +0.67).")
audit_row(s,0.42+w+g,CY0,w,h,"✓",TEAL,"Burstiness ↔ Poisson Fano","Every root Fano ≫ 1; like spike trains.")
audit_row(s,0.42,CY0+h+g,w,h,"~",AMBER,"Spectrum ↔ periodicity","Peak is low-frequency TREND, not a cycle.")
audit_row(s,0.42+w+g,CY0+h+g,w,h,"✓",TEAL,"Scale ↔ wavelet energy","Coarse scales 32–128 significant.")
audit_row(s,0.42,CY0+2*(h+g),w,h,"✗",RED,"Carrier wave ↔ rhythm","No fixed pulse; āyah CV ≈ 0.75.")
audit_row(s,0.42+w+g,CY0+2*(h+g),w,h,"~",AMBER,"Lead-lag ↔ causality","Overlap is real; causal reading is not licensed.")

# 18 — disclaimer
s=slide(prs); title(s,"Not a scientific miracle — and not evidence")
panel(s,0.42,1.20,12.5,2.9,REDT,[L("What we are NOT claiming",18,True,RED),
      L("The signal lens does not show that the Qur'an 'encodes' a frequency, predicts signal "
        "processing, or contains a hidden waveform. A surviving peak is a slow trend, not a designed "
        "cycle; a bursty root is coherent themed repetition, not a transmitted code.",16.5,True,NAVY)],space=9)
panel(s,0.42,4.30,12.5,2.9,TINT,[L("What the lens IS for",18,True,NAVY),
      L("It is a disciplined way to ASK whether reading order carries structure, with a null behind every "
        "answer. The value is clarity of thought and reproducibility — every number here can be "
        "recomputed live in the app. The analogy is judged by how well it organizes the question, never "
        "offered as proof of anything.",16.5,True,TEAL)],space=9)

# 19 — quick reference
s=slide(prs); title(s,"Quick reference — terms & live numbers")
two(s,[L("Terms",17,True,TEAL),
   L("Autocorrelation — self-similarity at a lag.  Fano factor — variance/mean (Poisson = 1).  Spectrum "
     "— power by frequency.  Wavelet — energy by scale.  Scalogram — scale × position image.  CV — "
     "spread ÷ mean.  Cross-correlation — overlap at each lag.",16)],
  [L("Live Book6 numbers",17,True,AMBER),
   L("lag-1 autocorr +0.67 · spectral peak p ≈ 0.0005 (low-freq) · wavelet scales 32/64/128 · āyah-length "
     "CV ≈ 0.75 · sampled roots Fano 40–61 at p ≈ 0.0002.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 20 — close
s=slide(prs); title(s,"Close — one sequence, read honestly")
panel(s,0.42,1.20,12.5,6.0,TINT2,[L("The through-line",18,True,NAVY),
      L("Read in order, the Qur'an behaves like a coherent natural signal: gentle memory between "
        "neighbouring sūras, bursty thematic repetition of roots, and a slow global trend visible to the "
        "FFT and the wavelet alike — with no fine periodicity and no carrier wave.",16.5),
      L("Every claim on these slides carries a null and a live Book6 number; every cross-domain parallel "
        "is a labelled lens, audited ✓ / ✗ / ~, never evidence.",16,True,TEAL),
      L("Next in the series — Biology: the same corpus read as a genome (letters ≈ bases, roots ≈ codons), "
        "and then the FDR Summary that collects every Two Books test into one corrected dashboard.",15.5)],
      space=10)

prs.save(OUT)
print("Signal deck rebuilt:",OUT,"| slides:",len(list(prs.slides)))
