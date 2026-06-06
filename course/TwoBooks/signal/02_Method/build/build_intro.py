# -*- coding: utf-8 -*-
"""Signal course - Lecture 2 - The Method. ANCHOR = ROOTS (col ریشه نحوی),
the unit of highest semantic power (as in the biology course and in NLP);
surface forms and morphology are complementary channels. All figures from Book6."""
import os, sys
HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
from st_slides import *
from st_slides import _tb,_rect,_fill_panel
from diagrams import fbox,harrow,vdash,band,sigrow
from pptx.util import Inches,Pt
OUT=os.path.abspath(os.path.join(HERE,".."))+"/"
prs=deck()
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def vecbars(s,x,y,w,vals,toks,col,vmax,bh=1.7):
    n=len(vals); cw=w/n; base=y+bh
    ax=s.shapes.add_shape(1,Inches(x),Inches(base),Inches(w),Inches(0.02)); ax.fill.solid(); ax.fill.fore_color.rgb=GREY; ax.line.fill.background()
    for i,(v,t) in enumerate(zip(vals,toks)):
        h=0.1+(v/vmax)*(bh-0.35)
        b=s.shapes.add_shape(5,Inches(x+i*cw+cw*0.15),Inches(base-h),Inches(cw*0.7),Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb=col; b.line.fill.background()
        _tb(s,x+i*cw-0.05,base-h-0.28,cw+0.1,0.26,[(str(v),11,True,NAVY)])
        _tb(s,x+i*cw-0.05,base+0.04,cw+0.1,0.34,[(t,15,True,NAVY)])

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.05,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Ayah as Signal (1-D)  ·  Lecture 2",15,True,TEAL),
  L("The Method — testing every claim on real Qur’an data",25,True,NAVY)],space=6)
steps=[("VERIFY","is the number right?",TINT,TEAL),("VALIDATE","is the pattern real?",AMBERT,AMBER),
       ("BEAT NULL + BASELINE","better than chance & Arabic?",TINT,TEAL),("AUDIT","✓ / ✗ / ~",REDT,RED)]
x=0.55; bw=2.92; aw=0.12
for i,(t,sub,fl,ln) in enumerate(steps):
    fbox(s,x,2.95,bw,1.25,fl,t,sub,line=ln,tsz=15.5,ssz=10.5)
    if i<3: harrow(s,x+bw-0.02,3.45,aw+0.02,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.6,12.5,2.6,TINT,[L("A pipeline, run on actual verses — anchored on ROOTS",17,True,NAVY),
  L("We vectorize on the ROOT (ریشه) — the unit of highest semantic power, as in the biology course (root↔codon) and in NLP (the stem/lemma). قل هو الله أحد → roots قول·ءله·وحد → [1722, 2848, 153]. Surface forms and morphology stay as complementary channels.",16.5),
  L("This lecture decides which root-signal findings to trust. Every number is computed from Book6.xlsx (6,236 āyāt).",16,True,TEAL)],space=7)

# 2 VERIFY (root vectors, concrete)
s=slide(prs); title(s,"Verify — recompute the root vectors from the text")
band(s,0.42,1.18,12.5,0.4,TINT,"anchor = ROOT (ریشه); each verse re-derived from Book6, root by root",TEAL)
vecbars(s,1.2,1.75,5.0,[1722,2848,153],["قول","ءله","وحد"],TEAL,2848)
vecbars(s,7.0,1.75,5.6,[381,2848,339,339],["سمو","ءله","رحم","رحم"],AMBER,2848)
panel(s,0.42,4.05,12.5,3.15,TINT2,[L("112:1 → [1722, 2848, 153]   ·   1:1 → [381, 2848, 339, 339]",18,True,NAVY),
  L("قل هو الله أحد reduces to roots قول·ءله·وحد — هو is a pronoun with no root, so it carries no sample (one root-token = one sample). The Basmala’s roots are سمو·ءله·رحم·رحم (رحمن and رحيم share the root رحم). ءله is the peak.",16.5),
  L("Surface forms (col 5) and full morphology (col 7) remain available as additional channels; the root is the anchor because it carries the meaning.",16,True,TEAL)],space=7)

# 3 WHY ROOTS (surface artifact -> motivation)
s=slide(prs); title(s,"Why roots — surface tokens fake a perfect match")
band(s,0.42,1.18,12.5,0.4,REDT,"on SURFACE tokens, 70:8 and 37:153 correlate r=1.00 — only because both carry ال (8374) twice",RED)
vecbars(s,0.9,1.7,5.4,[420,81,8374,120,1454,8374,5],["يوم","تكون","ال","سماء","ك","ال","مهل"],GREY,8374)
vecbars(s,6.9,1.7,5.6,[511,5,8374,17,715,8374,12],["أ","اصطفى","ال","بنات","على","ال","بنين"],GREY,8374)
panel(s,0.42,3.95,12.5,3.25,TINT,[L("Function words inflate surface similarity — roots remove it",18,True,NAVY),
  L("These two unrelated verses hit r=1.00 on surface tokens because the two huge ال peaks (8,374) dominate. Drop to the ROOT anchor and ال (no root) disappears — the artifact is gone. Mean pairwise r across length-7 verses falls from 0.18 (surface) to 0.04 (roots).",16.5),
  L("The anchor choice is itself a validation decision: roots carry meaning AND behave better statistically.",16,True,TEAL)],space=7)

# 4 NULL MODEL
s=slide(prs); title(s,"The null model — shuffle the roots, then compare")
three(s,[L("① STATE the null",17,True,TEAL),L("‘Order carries nothing.’ Model it by SHUFFLING the roots of a verse — قول ءله وحد → وحد قول ءله — same samples, scrambled.",16)],
 [L("② MEASURE under it",17,True,AMBER),L("Compute the same statistic on thousands of shuffles and on length-matched random root-verses. This is the null distribution.",16)],
 [L("③ COMPARE",17,True,NAVY),L("Observed in the bulk → chance (drop it). In the tail → a candidate, p = tail fraction. The refrain roots ءلی·ربب·کذب land in the tail (p≈0.002).",16)])

# 5 SAMPLING THE NULL
s=slide(prs); title(s,"Sampling the null — when enumeration is impossible")
band(s,0.42,1.2,12.5,0.4,TINT2,"a 7-root verse has 5,040 orderings; the space of reorderings/channels is ~10^N",NAVY)
fbox(s,0.6,2.0,3.6,1.55,REDT,"ALL configurations","reorderings · channels — uncountable",line=RED,tsz=15.5,ssz=11)
harrow(s,4.4,2.7,1.6,"draw 10⁴",color=GREY,lcol=NAVY)
fbox(s,6.2,2.0,3.6,1.55,AMBERT,"a RANDOM SAMPLE","20,000 draws stand in for the whole",line=AMBER,tsz=15.5,ssz=11)
harrow(s,10.0,2.7,1.4,"estimate",color=GREY,lcol=TEAL)
fbox(s,11.5,2.0,1.65,1.55,TINT,"p","tail frac.",line=TEAL,tsz=16,ssz=11)
panel(s,0.42,3.9,12.5,3.3,TINT,[L("Monte-Carlo: a sampled null model",18,True,NAVY),
  L("We cannot enumerate every reordering of every verse — so we DRAW a large random sample and read the observed value against it. For the refrain similarity, 20,000 draws give p ≈ 0.002. More draws → a sharper p (error ∝ 1/√draws).",16.5),
  L("Sampling tames the impossible; the null sets the bar. This judges every testable claim in the course.",16,True,TEAL)],space=8)

# 6 THE ROOT ANCHOR PAYS OFF (real null, contrast)
s=slide(prs); title(s,"The root anchor pays off — a cleaner null")
finding2(s,
 {"title":"Mean pairwise r (7-token verses)","cats":["surface tokens","ROOT anchor"],
  "series":[("",[RED,TEAL],[0.18,0.04])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Root null (20,000 random 7-root pairs)","cats":["≤-.6","-.6/-.3","-.3/0","0/.3",".3/.6","≥.6"],
  "series":[("",[GREY,AMBER,TEAL,TEAL,AMBER,GREY],[387,4244,6430,3622,2527,2790])],"legend":False},
 [L("Roots centre the null near zero",17.5,True,TEAL),
  L("On surface tokens random verses correlate at mean r ≈ 0.18 (function words). On the ROOT anchor that falls to 0.04 — the null is honest, so a real signal can actually stand out.",16)],
 [L("And a real one still clears it",17.5,True,AMBER),
  L("The exact refrain (identical roots ءلی·ربب·کذب, 31×) sits at r=1.0, p ≈ 0.002 against this root null. A clean null + a real tail = a finding.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 NATURAL-LANGUAGE BASELINE (root numbers)
s=slide(prs); title(s,"Beyond the null — beat random Arabic, not just randomness")
finding2(s,
 {"title":"Generic root features: Qur’an vs random Arabic","cats":["Zipf slope","top-10 share"],
  "series":[("Qur’an",NAVY,[0.757,0.214]),("baseline",AMBER,[0.764,0.214])],"legend":True,"fmt":"{:.2f}"},
 {"title":"Exact root-refrain rate (%)","cats":["random Arabic","Qur’an"],
  "series":[("",[GREY,TEAL],[0.81,7.13])],"legend":False,"fmt":"{:.2f}"},
 [L("Most structure is GENERIC to Arabic",17.5,True,AMBER),
  L("Root-Zipf slope (−0.76) and top-10 root share (0.21) are IDENTICAL in length/Zipf-matched random Arabic roots. Beating a shuffle here proves nothing — these belong to the language.",16)],
 [L("The residue that is SPECIFIC",17.5,True,TEAL),
  L("Exact root-verse repetition is 7.1% in the Qur’an vs 0.81% in matched random Arabic — ~8.8× (p≈0.03). THIS beats the baseline; the next slide shows the verses.",16)],
 fillA=AMBERT,fillB=TINT)

# 8 REFRAIN TABLE (concrete)
s=slide(prs); title(s,"The refrains — the actual repeated āyāt (from Book6)")
band(s,0.42,1.16,12.5,0.4,TINT2,"exact whole-verse repetition (same roots) · the verse · ×count · sūra",NAVY)
refs=[("فَبِأَيِّ آلَاءِ رَبِّكُمَا تُكَذِّبَانِ","×31","55 Ar-Raḥmān"),
      ("وَيْلٌ يَوْمَئِذٍ لِّلْمُكَذِّبِينَ","×11","77 Al-Mursalāt"),
      ("وَإِنَّ رَبَّكَ لَهُوَ الْعَزِيزُ الرَّحِيمُ","×8","26 Ash-Shuʿarāʾ"),
      ("فَاتَّقُوا اللَّهَ وَأَطِيعُونِ","×8","26 Ash-Shuʿarāʾ"),
      ("إِنَّ فِي ذَٰلِكَ لَآيَةً ۖ وَمَا كَانَ أَكْثَرُهُم مُّؤْمِنِينَ","×6","26 Ash-Shuʿarāʾ"),
      ("الم  /  حم  (openings)","×6/7","muqaṭṭaʿāt")]
y=1.72
for ar,cnt,su in refs:
    fbox(s,0.55,y,8.7,0.62,TINT,ar,line=TEAL,tsz=15.5)
    fbox(s,9.4,y,1.55,0.62,AMBERT,cnt,line=AMBER,tsz=14.5)
    fbox(s,11.1,y,2.0,0.62,TINT2,su,line=NAVY,tsz=11.5)
    y+=0.72
panel(s,0.42,6.12,12.5,1.08,TINT,[L("7.1% vs 0.81% — ~8.8× the matched-Arabic rate (the clearest Qur’an-specific structure so far)",16.5,True,NAVY),
  L("Refrains are a real rhetorical device (Ar-Raḥmān’s repeated question, al-Mursalāt’s repeated warning) — and they survive the null AND the baseline.",15.5,True,TEAL)],space=5)

# 9 MULTIPLE COMPARISONS
s=slide(prs); title(s,"The trap of many tests — and the correction")
finding2(s,
 {"title":"400 random channels, true effect = 0","cats":["'sig' p<.05","after BH-FDR"],
  "series":[("",[RED,TEAL],[18,0])],"legend":False},
 {"title":"A hidden confound (length coupling)","cats":["random×length 'sig'","expected"],
  "series":[("",[RED,GREY],[110,8])],"legend":False},
 [L("Search enough and you ‘find’ things",17.5,True,RED),
  L("Test 400 meaningless channels and ~5% (18) cross p<0.05 by chance. Benjamini-Hochberg FDR removes ALL 18 — none were real.",16)],
 [L("And beware confounds",17.5,True,AMBER),
  L("110 of 150 random channels ‘correlated’ with verse length — not chance but a CONFOUND (averaging over more roots shrinks variance). Rule out the boring explanation first.",16)],
 fillA=REDT,fillB=AMBERT)

# 10 EFFECT SIZE
s=slide(prs); title(s,"p is not enough — report the effect size")
two(s,[L("A small p can be a tiny effect",18,True,RED),L("With 6,236 āyāt, trivial gaps turn ‘significant.’ p answers ‘non-zero?’ — not ‘big enough to matter?’ A pattern can beat the null yet be negligible.",16.5,True,NAVY)],
 [L("So pair p with magnitude",18,True,TEAL),L("The refrain result is not just p≈0.03 — it is 7.1% vs 0.81%, a ~8.8× ratio (31 repeats of one verse alone). Big effect, clear story. Always give both.",16.5,True,NAVY)],sp=0.5,fa=REDT,fb=TINT)

# 11 PRE-REGISTRATION
s=slide(prs); title(s,"Declare first — the cure for the search")
two(s,[L("DECIDE BEFORE YOU LOOK",18,True,TEAL),L("Fix the channel (root? surface? morphology?), the statistic, the null, the baseline and the threshold IN ADVANCE. Choosing them after seeing the data is how a search of millions of configurations manufactures a ‘finding.’",16.5,True,NAVY)],
 [L("AND COUNT YOUR TESTS",18,True,AMBER),L("If exploration is the goal, say so — then correct for every test (FDR) and confirm survivors on held-out āyāt. Exploratory and confirmatory work are both fine; conflating them is not.",16.5,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)

# 12 ALIASING & OVERFITTING
s=slide(prs); title(s,"Two classic ways to fool yourself")
two(s,[L("ALIASING — under-sampling lies",18,True,RED),L("Sample too coarsely and a fast pattern masquerades as a slow one. In text: binning rare roots together invents periodicities that are artifacts of the binning, not the verse.",16.5,True,NAVY)],
 [L("OVERFITTING — explaining noise",18,True,RED),L("With enough parameters you can fit any wiggle, even noise. A model that ‘explains’ one āyah but predicts no other has learned the noise. Guard: hold out data; prefer the simplest model that still beats the null.",16.5,True,NAVY)],sp=0.5,fa=REDT,fb=REDT)

# 13 SCALE RULE
s=slide(prs); title(s,"The scale rule — short verses are unstable")
finding2(s,
 {"title":"Lag-1 autocorr estimate spread vs length","cats":["n≈4","n≈7","n≈15","n≈30","n≈60"],
  "series":[("",[RED,AMBER,TEAL,TEAL,NAVY],[0.484,0.307,0.155,0.119,0.087])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Min n where each tool is valid","cats":["amplitude","distance","autocorr","Fourier"],
  "series":[("",[TEAL,TEAL,AMBER,RED],[1,1,15,30])],"legend":False,"fmt":"n≥{:.0f}"},
 [L("From الم (1 root) to 2:282 (84)",17.5,True,TEAL),
  L("The shortest āyah carries 1 root; the longest is 2:282 (84 root-tokens). Median is 7. A lag-1 autocorrelation has std 0.48 at n≈4 but 0.09 at n≈60 — short verses give unstable spectra.",16)],
 [L("So match the tool to the scale",17.5,True,AMBER),
  L("Robust at āyah scale: amplitude, energy, distance. Length-hungry (Fourier, autocorrelation, spectrogram): run at SŪRA or CORPUS scale. State the scale on every claim.",16)],
 fillA=TINT,fillB=AMBERT)

# 14 PIPELINE
s=slide(prs); title(s,"The pipeline — every claim runs this gauntlet")
gates=[("VERIFY","recompute from Book6",TINT,TEAL),("VALIDATE","is it real?",AMBERT,AMBER),
       ("NULL","beat chance (sampled)",TINT,TEAL),("BASELINE","beat random Arabic",AMBERT,AMBER),
       ("CORRECT","FDR for many tests",TINT,TEAL),("READ-BACK","map to the roots/text",REDT,RED)]
x=0.42; bw=1.98; aw=0.05
for i,(t,sub,fl,ln) in enumerate(gates):
    fbox(s,x,2.0,bw,1.5,fl,t,sub,line=ln,tsz=13.5,ssz=9.5)
    if i<5: harrow(s,x+bw-0.02,2.6,aw+0.04,"",color=GREY)
    x+=bw+aw
panel(s,0.42,3.9,12.5,3.3,TINT,[L("Survive all six, then audit",18,True,NAVY),
  L("A finding that is correctly computed, beats a sampled null, beats a natural-language baseline, survives multiple-comparison correction, and reads back into the roots/text earns a ✓. Fail any gate and it is ✗ or ~. Most candidates die here; that is the point.",16.5,True,TEAL)],space=8)

# 15 CASE STUDY
s=slide(prs); title(s,"A claim, end to end — فبأي آلاء ربكما تكذبان")
band(s,0.42,1.2,12.5,0.4,TINT2,"claim: the Qur’an repeats whole āyāt far more than chance — does it survive?",NAVY)
g=[("VERIFY ✓","7.1% computed"),("NULL ✓","p≈0.002"),("BASELINE ✓","7.1 vs 0.81%"),("FDR ✓","1 declared test"),("READ-BACK ✓","real refrain")]
x=0.55; bw=2.4; aw=0.06
for i,(t,sub) in enumerate(g):
    fbox(s,x,1.9,bw,1.45,TINT,t,sub,line=TEAL,tsz=14,ssz=11)
    if i<4: harrow(s,x+bw-0.03,2.5,aw+0.05,"",color=GREY)
    x+=bw+aw
panel(s,0.42,3.7,12.5,3.5,TINT,[L("This one earns a ✓",18,True,NAVY),
  L("فَبِأَيِّ آلَاءِ رَبِّكُمَا تُكَذِّبَانِ (roots ءلی·ربب·کذب) recurs 31× in Sūrat ar-Raḥmān. The claim is correctly computed, beats a sampled root null (p≈0.002), beats a matched-Arabic baseline ~8.8×, is a single declared test, and reads back as a real rhetorical refrain.",16.5),
  L("Modest and specific — exactly the honest yield the method is built to find, and to separate from the generic.",16,True,TEAL)],space=7)

# 16 REPRODUCIBILITY
s=slide(prs); title(s,"Reproducibility — anyone can re-run it")
three(s,[L("DATA",17,True,TEAL),L("One stated source — Book6.xlsx (6,236 āyāt); the ROOT column (ریشه) is the anchor. Every figure traces to a column and a rule.",16)],
 [L("CODE",17,True,AMBER),L("Each deck is generated by a script, each statistic by a function. Re-running reproduces the number — and the slide.",16)],
 [L("SEED",17,True,NAVY),L("Monte-Carlo uses a fixed random seed, so the null, the p-value and the histogram are identical on every run.",16)])

# 17 DUAL-DOMAIN
s=slide(prs); title(s,"Dual-domain — the same method guards creation’s data")
two(s,[L("عالم التكوين — science runs this too",18,True,AMBER),L("A gene-expression spike, a candidate exoplanet dip, a clinical effect — each must beat a null, a baseline and a multiple-testing correction. Genomics calls it FDR across 20,000 genes; astronomy, the look-elsewhere effect.",16.5,True,NAVY)],
 [L("عالم التدوين — scripture-as-data",18,True,TEAL),L("We apply the identical discipline to the root-vector. Same method, two Books: a pattern is a finding only when chance, baseline and search are ruled out — in the cell and in the verse alike.",16.5,True,NAVY)],sp=0.5,fa=AMBERT,fb=TINT)

# 18 APP
s=slide(prs); title(s,"The app — run the method live")
band(s,0.42,1.2,12.5,0.4,TINT,"pick a verse · choose channel (root/surface/morph) · set the null · toggle baseline",TEAL)
st=[("① VERSE","e.g. 55:13 refrain",TINT,TEAL),("② CHANNEL","root anchor (+others)",AMBERT,AMBER),
    ("③ NULL+DRAWS","shuffle, 10³→10⁵",TINT,TEAL),("④ VERDICT","p · baseline · FDR",REDT,RED)]
x=0.55; bw=2.98; aw=0.1
for i,(t,sub,fl,ln) in enumerate(st):
    fbox(s,x,1.9,bw,1.35,fl,t,sub,line=ln,tsz=15,ssz=11)
    if i<3: harrow(s,x+bw-0.02,2.45,aw+0.05,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.65,12.5,2.55,TINT,[L("The method is a button, not a sermon",18,True,NAVY),
  L("Students pick 55:13 on the ROOT channel, set a null, slide the draws and watch the p-value sharpen, then toggle the natural-language baseline to see a ‘significant’ result collapse into ‘generic.’ Today: reproduce the refrain — 7.1% vs 0.81%, p≈0.03.",16.5,True,TEAL)],space=7)

# 19 FALSIFIABILITY
s=slide(prs); title(s,"Falsifiable by construction")
two(s,[L("A METHOD that can say NO",18,True,TEAL),L("The apparatus exists to KILL claims: the root anchor kills the surface artifact (r=1.0 from ال), the null kills the random-looking, the baseline kills the merely-Arabic (Zipf), FDR kills the lucky, the read-back kills the meaningless.",16.5,True,NAVY)],
 [L("The opposite of numerology",18,True,RED),L("Numerology never says no — it bends to fit anything. Here most candidates fail; the refrain survived, countless others will not. That asymmetry separates a finding from a fancy.",16.5,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# 20 AUDIT
s=slide(prs); title(s,"Audit — the verdict the method produces")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Exact root-verse repetition (refrains): 7.1% vs 0.81% baseline, p≈0.002, reads back to فبأي آلاء ربكما تكذبان. Survives every gate.",16)],
 [L("✗ BREAKS",17,True,RED),L("Surface ‘70:8 ≈ 37:153’ (r=1.0) — an artifact of ال, fixed by the root anchor; and root-Zipf / share — generic to Arabic. Both fail.",16)],
 [L("~ SILENT",17,True,AMBER),L("A spectrum of a 1–3 root āyah — untestable at this scale; set aside until run at sūra/corpus scale, not asserted.",16)],f=(TINT,REDT,AMBERT))

# 21 DISCUSSION
s=slide(prs); title(s,"Synthesis & discussion — trusting a result")
two(s,[L("THE HABIT",18,True,NAVY),L("Before believing any root-signal finding, ask in order: computed right? beats a sampled null? beats random Arabic? survives the search? reads back into the roots/text? Only yes to all earns belief.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Why does the root anchor beat surface tokens here?  • When is sampling the null unavoidable?  • Is the refrain result a confound — how would you check?  • At what scale should ‘rhythm’ be tested?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 22 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Root/stem vectorization, sampled nulls, baselines and FDR are the daily tools of NLP, genomics and clinical trials. Learn them here and you can audit any data claim you meet.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Anchor on the ROOT; then verify → beat a sampled null → beat random Arabic → correct for many tests → read back. On real āyāt: the refrain survives (7.1% vs 0.81%); the surface r=1.0 pair and Zipf do not.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)

prs.save(OUT+"02_Method_Lecture.pptx")
print("L2 Method slides:", len(prs.slides))
