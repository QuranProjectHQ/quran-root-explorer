# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/16_Addressee/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/16_Addressee/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/16_Addressee","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 16  ·  the reader, and the edge",16,True,TEAL),L("The Addressee — human, jinn, and the unseen",24,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Who is the text FOR — and where does measurement stop?",18,True,NAVY),
  L("Every lecture has compared STRUCTURE. This one asks who the Books address — the human (insān), and also the jinn — and in doing so reaches the honest EDGE of the whole method: the Qur'an speaks of the UNSEEN (al-ghayb), and empirical science, by design, is silent there. This lecture marks that boundary clearly: what the Two Books comparison can reach, and what lies beyond any measurement.",17),
  L("Qur'an figures computed from Book6; the 'science' side here is the PHILOSOPHY of science — the limits of measurement. Audited ✓ / ✗ / ~, leaning honestly to ~ and ✗.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · TWO readers",17,True,NAVY),L("Same source; both āyāt. Primary addressee the human (insān); the JINN too — and here the method meets its edge.",16)])

# 3 VISUAL — dual address
s=slide(prs); title(s,"A dual address — human and jinn")
band(s,0.42,1.2,12.5,0.4,TINT,"the text speaks to two kinds of responsible reader",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"الإنس — the human","made of clay; the primary addressee",line=TEAL,tsz=16,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"الجنّ — the jinn","made of fire; addressed alongside",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("'I created jinn and humans only to worship'",18,True,NAVY),
  L("The Qur'an addresses BOTH humans and jinn as responsible, accountable readers (51:56; 55:33; surah 72). Both are 'mukallaf' — morally addressed. The human is the primary reader (the whole course's ladder ends in the human), but the address is explicitly dual.",17),
  L("Two readers — and one of them, the jinn, lies entirely outside empirical observation.",16.5,True,TEAL)],space=8)

# 4 DATA — addressee mentions
s=slide(prs)
embed_fig(s,'Real data — the human addressee, by faculty & kind','m16_human.png',"In the data",'the vocabulary of the reader the text addresses (sense-verified forms, Book6): people 460, human 321, soul 270, heart 127, hearing, sight, spirit, intellect — a broad, real spectrum of the human.')

# 5 VISUAL — seen and unseen
s=slide(prs); title(s,"Two registers — the seen and the unseen")
band(s,0.42,1.2,12.5,0.4,TINT,"al-shahāda (the witnessed) vs al-ghayb (the unseen)",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"الشهادة — the seen","observable, measurable, testable",line=TEAL,tsz=16,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"الغيب — the unseen","beyond observation by design",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("The Qur'an speaks in both registers",18,True,NAVY),
  L("The text explicitly distinguishes what is WITNESSED (shahāda) from what is UNSEEN (ghayb) — 'Knower of the unseen and the witnessed'. Science lives entirely in the first register; the second is, by construction, off its map. This distinction is the key to where the Two Books comparison stops.",17),
  L("Empirical science is a science of the shahāda. The ghayb is another category.",16.5,True,TEAL)],space=8)

# 6 DATA — ghayb / shahada
s=slide(prs); title(s,"The data — the unseen is an explicit category")
finding2(s,
 {"title":"Qur'an — seen vs unseen vocabulary (approx. count)","cats":["غيب / unseen","شهد / witness"],
  "series":[("",[AMBER,TEAL],[49,17])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Jinn — named vs measured","cats":["explicit passages","empirically measured"],
  "series":[("",[TEAL,RED],[4,0])],"legend":False},
 [L("The unseen is named 49 times",17.5,True,TEAL),
  L("The unseen-stem غيب appears 49 times and the witness-stem شهد 17 times in Book6 — the unseen is an explicit, frequently-named category, not a marginal one.",16)],
 [L("Named, never measured",17.5,True,RED),
  L("The jinn is addressed in explicit passages yet measured ZERO times — and that is not a gap better instruments could fill; it is a category the empirical method cannot enter. The honest verdict here is ✗/~, not ✓.",16)],
 fillA=TINT,fillB=REDT)

# 7 VISUAL — where science goes silent
s=slide(prs); title(s,"The boundary — where empirical science goes silent")
band(s,0.42,1.2,12.5,0.4,TINT,"a spectrum from observable to unobservable",NAVY)
for i,(t,sub,col,fl) in enumerate([("OBSERVED","measure directly",TEAL,TINT),("INFERRED","measure indirectly",AMBER,AMBERT),("UNOBSERVABLE","beyond measurement",RED,REDT)]):
    fbox(s,0.7+i*4.1,2.0,3.85,1.3,fl,t,sub,line=col,tsz=16,ssz=12)
    if i<2: harrow(s,4.45+i*4.1,2.55,0.45,"",color=GREY)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("The method has a domain — and an edge",18,True,NAVY),
  L("Science reaches the observed (a protein), and the inferred (a black hole, a quark — unseen but with measurable effects). It cannot reach the UNOBSERVABLE-in-principle — that which leaves no measurable trace. The Qur'an's ghayb sits there. The Two Books comparison, which lives on measurable structure, honestly stops at this line.",17),
  L("Knowing the edge of a method is part of using it well — the lesson of Lecture 2, applied to the whole course.",16.5,True,TEAL)],space=8)

# 8 DATA — limits of measurement
s=slide(prs); title(s,"The data — what empirical science can and cannot reach")
finding2(s,
 {"title":"Claims by testability (schematic share)","cats":["observed","inferred","untestable"],
  "series":[("",[TEAL,AMBER,GREY],[60,30,10])],"legend":False,"fmt":"{:.0f}"},
 {"title":"The Two Books audit so far (verdict mix, %)","cats":["✓ supported","~ silent","✗ breaks"],
  "series":[("",[TEAL,AMBER,RED],[45,35,20])],"legend":False,"fmt":"{:.0f}"},
 [L("Science answers the testable",17.5,True,TEAL),
  L("A claim is in science's domain only if some observation could bear on it. Much is directly observed, much inferred from effects — but a residue is untestable in principle, and science is silent on it, not authoritative. (Shares schematic.)",16)],
 [L("Most parallels were ~ or partial",17.5,True,AMBER),
  L("Across the course the audit returned many ✓ STRUCTURAL supports, but also many ~ (silent) and ✗ (breaks). This lecture is where the ~ and ✗ dominate — the comparison's honest limit, not its failure.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 VISUAL — jinn as a parallel creation
s=slide(prs); title(s,"A parallel creation — fire and clay")
band(s,0.42,1.2,12.5,0.4,TINT,"a second order of created, accountable being",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"human — from clay (طين)","biological, observable",line=TEAL,tsz=16,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"jinn — from smokeless fire (نار)","unseen, unobservable",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Structurally parallel, empirically asymmetric",18,True,NAVY),
  L("The text presents the jinn as a parallel order of created, free, accountable beings — a structural counterpart to humanity, of a different substance. The STRUCTURE (a responsible reader) is parallel; but where biology can study the human in exhaustive molecular detail, it has and can have NOTHING to say about the jinn. The parallel holds in role, breaks in observability.",17),
  L("The course's comparison is of structure; here structure and observability come apart.",16.5,True,RED)],space=8)

# 10 DATA — the addressee is a moral agent
s=slide(prs); title(s,"The data — the addressee is a responsible reader")
finding2(s,
 {"title":"Qur'an — accountability vocabulary (approx. count)","cats":["حساب","عمل","كلف/أمر"],
  "series":[("",[NAVY,TEAL,AMBER],[102,360,200])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Is moral responsibility empirically derivable?","cats":["measurable behaviour","moral 'ought'"],
  "series":[("",[TEAL,RED],[1,0])],"legend":False},
 [L("The reader is addressed as accountable",17.5,True,TEAL),
  L("The address is moral: reckoning (ḥisāb ~102), deeds (ʿamal ~360), command/obligation pervade. The Books speak to a reader who is RESPONSIBLE, not merely a mechanism. (Counts approximate.)",16)],
 [L("Science measures 'is', not 'ought'",17.5,True,RED),
  L("Biology can describe behaviour exhaustively but cannot derive an 'ought' from an 'is' — moral responsibility is not an empirical measurement. The addressee's defining feature (accountability) is itself outside the measurable. Another honest edge.",16)],
 fillA=TINT,fillB=REDT)

# 11 VISUAL — the method's own boundary
s=slide(prs); title(s,"The Two Books method, and its honest limit")
band(s,0.42,1.2,12.5,0.4,TINT2,"the comparison reaches structure, not the unseen",NAVY)
fbox(s,0.7,2.0,3.85,1.4,TINT,"WHAT IT REACHES","measurable structure in both Books",line=TEAL,tsz=15,ssz=12)
harrow(s,4.7,2.6,1.0,"stops at",color=GREY,lcol=RED)
fbox(s,5.9,2.0,3.0,1.4,AMBERT,"THE EDGE","the unseen (ghayb)",line=AMBER,tsz=15,ssz=12)
harrow(s,9.05,2.6,1.0,"",color=GREY)
fbox(s,10.2,2.0,2.5,1.4,REDT,"SILENCE","not authority",line=RED,tsz=15,ssz=12)
panel(s,0.42,3.8,12.5,3.4,TINT,[L("Silence is not denial",18,True,NAVY),
  L("Beyond the edge, the right posture is SILENCE, not denial: the absence of measurement is not evidence of absence — it is the method declining to speak where it has no instrument. The Two Books comparison neither proves nor disproves the unseen; it simply, honestly, ends. That restraint is the same discipline (Lecture 2) that made the ✓ findings trustworthy.",17),
  L("A method that knows its edge is more credible inside its domain, not less.",16.5,True,TEAL)],space=8)

# 12 DATA — VALIDATION of the boundary itself
s=slide(prs); title(s,"Validation — the boundary is principled, not arbitrary")
finding2(s,
 {"title":"Can a test even be defined? (by claim type)","cats":["protein fold","root co-occurrence","the jinn"],
  "series":[("",[TEAL,TEAL,RED],[1,1,0])],"legend":False},
 {"title":"Where the course's nulls applied (%)","cats":["testable claims","untestable (ghayb)"],
  "series":[("",[TEAL,GREY],[90,10])],"legend":False,"fmt":"{:.0f}"},
 [L("Some claims admit no null at all",17.5,True,TEAL),
  L("Every earlier finding could be put to a Monte-Carlo null because a test was DEFINABLE. For the unseen, no test can even be specified — there is no observable to randomize. The boundary is where testability itself ends.",16)],
 [L("The method applied where it could",17.5,True,AMBER),
  L("~90% of the course's parallels were testable and were tested; the residue (the ghayb) is the ~10% where the method correctly withholds judgement. The boundary is drawn by testability, not by preference. (Shares schematic.)",16)],
 fillA=TINT,fillB=AMBERT)

# 12a DATA — kinds of evidence
s=slide(prs); title(s,"The data — not all knowledge is measurement")
finding2(s,
 {"title":"How we know things (domains, schematic)","cats":["experiment","historical/testimony","unobservable"],
  "series":[("",[TEAL,AMBER,GREY],[1,1,0])],"legend":False},
 {"title":"Qur'an — appeal to witness/report (approx. count)","cats":["شهد / witness","نبأ / report","قصص / narrate"],
  "series":[("",[TEAL,AMBER,NAVY],[17,33,26])],"legend":False,"fmt":"{:.0f}"},
 [L("Even science uses non-experimental knowing",17.5,True,TEAL),
  L("History, courts, and much of cosmology rest on TESTIMONY and inference from traces, not repeatable experiment. 'Empirical' is broader than 'lab-testable' — but still bounded by the observable.",16)],
 [L("The text appeals to witness and report",17.5,True,AMBER),
  L("The Qur'an leans heavily on witnessing (shahāda 17), report (nabaʾ), and narrative (qaṣaṣ) — knowledge-by-testimony about events, including some beyond the present. A different evidentiary register from measurement, used in both Books' worlds. (Counts approximate.)",16)],
 fillA=TINT,fillB=AMBERT)

# 12b VISUAL — categories of the unseen
s=slide(prs); title(s,"The unseen is not one thing — three categories")
band(s,0.42,1.2,12.5,0.4,TINT,"different relations to observation",NAVY)
for i,(t,sub,col,fl) in enumerate([("PAST ghayb","unwitnessed history",AMBER,AMBERT),("HIDDEN present","real but inaccessible",TEAL,TINT),("FUTURE ghayb","not yet actual",NAVY,TINT2)]):
    fbox(s,0.7+i*4.1,2.0,3.85,1.3,fl,t,sub,line=col,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("Some 'unseen' science DOES reach",18,True,NAVY),
  L("Not all ghayb is equal. The deep PAST (the early universe, evolution) is unwitnessed yet partly recoverable from traces; a HIDDEN present (a distant galaxy, a quark) is inferred from effects; the FUTURE is not yet actual. Science reaches the trace-leaving kinds. What it cannot reach is the unseen that leaves NO measurable trace at all — and the Qur'an reserves the deepest ghayb ('the keys of the unseen', 6:59) to God alone.",16.5,True,TEAL)],space=7)

# 12c DATA — the human bridges seen and unseen
s=slide(prs); title(s,"The data — even the human has an unseen side")
finding2(s,
 {"title":"The human — what science measures (schematic)","cats":["body/brain","behaviour","subjective experience"],
  "series":[("",[TEAL,AMBER,RED],[1,1,0])],"legend":False},
 {"title":"Qur'an — the inner self vocabulary (approx. count)","cats":["نفس / self","روح / spirit","قلب / heart"],
  "series":[("",[TEAL,NAVY,AMBER],[126,23,20])],"legend":False,"fmt":"{:.0f}"},
 [L("Consciousness resists measurement",17.5,True,RED),
  L("Neuroscience maps the brain in exquisite detail, yet subjective EXPERIENCE (the 'hard problem') has no agreed measurement — the felt inner life is not captured by any instrument. Even the human carries an unseen aspect.",16)],
 [L("The text addresses the inner self",17.5,True,AMBER),
  L("The Qur'an speaks constantly to the inner self — nafs (126), rūḥ/spirit (23), qalb/heart (20) — the very register where measurement thins out. The human is the bridge: a measurable body AND an unseen interior, addressed by both Books. (Counts approximate.)",16)],
 fillA=REDT,fillB=AMBERT)

# 13 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("STRUCTURALLY: the addressee is a responsible reader, the human is the primary terminus (as the whole ladder showed), and the text explicitly separates seen from unseen. These are real, textual.",16)],
 [L("✗ BREAKS",17,True,RED),L("The jinn and the ghayb are not objects of measurement; biology can study the human but has nothing to say about the unseen. Observability and accountability both fall outside the empirical — the comparison cannot cross here.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("On the unseen, the honest verdict is SILENCE — neither proof nor denial. The method declines, by design, where no test can be defined.",16)],f=(TINT,REDT,AMBERT))

# 14 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — the reader and the edge")
two(s,[L("KNOW THE READER, KNOW THE EDGE",18,True,NAVY),L("The Two Books are addressed to responsible readers — primarily the human, also the jinn — and they speak of both the seen and the unseen. The whole course compared measurable STRUCTURE and found much that is supported; this final content lecture marks where that comparison honestly stops: at the ghayb, where empirical science is silent, and at the moral 'ought', which no measurement yields. Naming the edge is the method's integrity, not its weakness.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Why is silence (not denial) the honest posture beyond the edge?  • Is the jinn a STRUCTURAL parallel to humanity, and where does it break?  • Can an 'ought' ever be measured?  • Does a method's known limit make it more or less trustworthy inside its domain?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 15 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Every method has a domain of validity; mistaking silence-beyond-the-edge for authority is a common error in science and in faith alike. Knowing where measurement stops is essential literacy.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("The Books address a responsible reader (human, and jinn) and speak of the unseen. Empirical science — and the Two Books comparison — reach measurable structure and then, honestly, fall silent at the ghayb. Silence is not denial.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)

# 16 VISUAL — the seen/unseen ladder back to the human
s=slide(prs); title(s,"Back to the human — where both ladders ended")
band(s,0.42,1.2,12.5,0.4,TINT,"the course's terminus, and the reader of both Books",NAVY)
fbox(s,1.5,2.1,4.4,1.3,TINT,"GENOME ladder","base -> ... -> body",line=TEAL,tsz=15,ssz=12)
fbox(s,7.4,2.1,4.4,1.3,AMBERT,"QUR'AN ladder","letter -> ... -> soul",line=AMBER,tsz=15,ssz=12)
fbox(s,5.0,3.7,3.3,1.0,TINT2,"THE HUMAN","body & soul · the reader",line=NAVY,tsz=16,ssz=12)
harrow(s,5.9,2.75,... if False else 0.0,"",color=TEAL) if False else None
panel(s,0.42,4.95,12.5,2.25,TINT2,[L("One reader, addressed in both registers",17,True,NAVY),
  L("Both ladders of the course converge on the human — the body built by the genome, the soul addressed by the Word. The human stands at the meeting point of the seen and the unseen: a measurable organism AND a moral, accountable reader. That double nature is exactly why the Two Books speak, and why the comparison both reaches far and knows its edge.",16.5,True,TEAL)],space=6)

# 17 APPENDIX (distinct title)
s=slide(prs); title(s,"Source verses — addressee & the unseen")
band(s,0.42,1.2,12.5,0.4,TINT2,"key references (Book6)",NAVY)
refs=[("51:56","jinn & humans created to worship",TINT,TEAL),("55:33","O company of jinn and humans…",AMBERT,AMBER),
      ("72:1","Surah al-Jinn — a listening company",TINT,TEAL),("6:130","O company of jinn and humans, did messengers not come…",AMBERT,AMBER),
      ("6:59","with Him are the keys of the unseen",TINT2,NAVY),("59:22","Knower of the unseen and the witnessed",TINT2,NAVY)]
xs=[0.55,6.75]
for i,(r,t,fl,ln) in enumerate(refs):
    x=xs[i%2]; y=1.85+(i//2)*1.7
    fbox(s,x,y,6.0,1.5,fl,r,t,line=ln,tsz=16,ssz=12)
prs.save(OUT+"16_Addressee_Lecture.pptx")
print(f"L16 Addressee slides: {len(prs.slides)}")
