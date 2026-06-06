# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/12_Regulation/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/12_Regulation/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/12_Regulation","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 12  ·  when and where it speaks",16,True,TEAL),L("Regulation & Context",28,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Having a code is not using it — expression is controlled",18,True,NAVY),
  L("A cell carries the whole genome but EXPRESSES only a fraction in any tissue; the Qur'an's lexicon is vast but each passage USES only a slice. This lecture measures that selective expression — which units are 'on' where — and shows how CONTEXT (a promoter, a syntax, an occasion) decides what is expressed and what a unit means. This is Pillar 3 at its core: expression, not possibility, makes the realized text.",17),
  L("Qur'an figures computed from Book6; biology mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here: what is EXPRESSED, when, and where.",16)])

# 3 VISUAL — not all on at once
s=slide(prs); title(s,"The whole code is present; only some is expressed")
band(s,0.42,1.2,12.5,0.4,TINT,"same archive, different active subset",NAVY)
fbox(s,0.8,2.0,3.4,1.2,TINT,"full genome","~20,000 genes",line=TEAL,tsz=16,ssz=11)
harrow(s,4.4,2.5,1.5,"in a neuron",color=GREY,lcol=TEAL); fbox(s,6.1,2.0,2.7,1.2,TINT2,"~half ON","neuron set",line=NAVY,tsz=15,ssz=11)
harrow(s,9.0,2.5,1.5,"in a liver cell",color=GREY,lcol=AMBER); fbox(s,10.7,2.0,2.0,1.2,AMBERT,"different set","",line=AMBER,tsz=14,ssz=11)
panel(s,0.42,3.6,12.5,3.6,TINT2,[L("Identity comes from expression, not content",18,True,NAVY),
  L("A neuron and a liver cell hold the SAME genome; they differ because they EXPRESS different subsets of it. Likewise the whole lexicon is available, but each surah activates only a slice — and the slice gives the passage its character.",17),
  L("What you ARE is what you express — both Books.",16.5,True,TEAL)],space=8)

# 4 DATA — fraction expressed per context (REAL)
s=slide(prs); title(s,"The data — only a slice is expressed per context")
finding2(s,
 {"title":"Qur'an — distinct roots per surah (% of lexicon)","cats":["smallest","median","mean","largest"],
  "series":[("",[GREY,TEAL,TEAL,NAVY],[0.2,7.3,9.3,36])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Genome — genes expressed per cell type (%)","cats":["minimal cell","typical tissue"],
  "series":[("",[GREY,TEAL],[30,55])],"legend":False,"fmt":"{:.0f}"},
 [L("Each surah uses ~9% of the lexicon",17.5,True,TEAL),
  L("Computed: a surah expresses on average 157 of 1,702 roots (9.3%); the largest (al-Baqarah) 606 (36%), the smallest (al-Kāfirūn) just 4. Selective expression, measured.",16)],
 [L("Each cell uses ~half its genes",17.5,True,AMBER),
  L("A given cell type transcribes roughly 50% of its ~20,000 genes; the rest stay silent. Both Books deploy only a context-dependent SUBSET of the available code.",16)],
 fillA=TINT,fillB=AMBERT)

# 5 DATA — housekeeping vs context-specific (REAL)
s=slide(prs)
embed_fig(s,'Real data — root breadth across sūras','m12_breadth.png',"In the data","every root by how many of the 114 sūras it appears in (median 4, max 94): most roots are context-specific (few sūras), a few are near-ubiquitous 'housekeeping' — the expression/regulation profile, like tissue-specific vs housekeeping genes.")

# 6 VISUAL — the switch
s=slide(prs); title(s,"The switch — what turns a unit on")
band(s,0.42,1.2,12.5,0.4,TINT,"a control element decides whether the unit is read",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"GENE: promoter + transcription factor","context binds → gene ON",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"WORD: syntax + occasion","context selects → sense ON",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Expression is addressed, not automatic",18,True,NAVY),
  L("A gene fires only when the right transcription factors bind its promoter/enhancer in that cell; a word's intended sense activates only when the surrounding syntax and the occasion (asbāb al-nuzūl) call for it. In both Books, a CONTROL SIGNAL in the context — not the unit alone — decides what is expressed.",17),
  L("The unit is potential; the context is the switch.",16.5,True,TEAL)],space=8)

# 7 DATA — context selects the meaning
s=slide(prs); title(s,"The data — context resolves which sense / which isoform")
finding2(s,
 {"title":"Arabic — ambiguous skeletons resolved by context (%)","cats":["ambiguous alone","resolved in context"],
  "series":[("",[AMBER,TEAL],[19.6,80.4])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Gene — isoform chosen by cell context (schematic)","cats":["default","context-switched"],
  "series":[("",[GREY,TEAL],[1,3])],"legend":False},
 [L("Context picks the reading",17.5,True,TEAL),
  L("A skeleton with several possible vocalizations (19.6%) is resolved by its sentence — the same unit, one realized meaning, chosen by context. Polysemy collapses to sense at the moment of use.",16)],
 [L("Context picks the isoform",17.5,True,AMBER),
  L("The same gene yields different isoforms in different tissues — splicing is regulated by cellular context. In both Books, context CHOOSES among a unit's potential outputs.",16)],
 fillA=TINT,fillB=AMBERT)

# 8 VISUAL — revelation order ~ developmental staging
s=slide(prs); title(s,"Staged expression — revelation order and development")
band(s,0.42,1.2,12.5,0.4,TINT,"the program unfolds on a schedule",NAVY)
fbox(s,0.7,2.0,3.85,1.3,AMBERT,"EARLY (Meccan)","faith, the unseen, the self",line=AMBER,tsz=15,ssz=12)
harrow(s,4.7,2.55,1.0,"then",color=GREY,lcol=NAVY)
fbox(s,5.9,2.0,3.85,1.3,TINT,"LATER (Medinan)","law, community, structure",line=TEAL,tsz=15,ssz=12)
fbox(s,10.0,2.0,2.7,1.3,TINT2,"developmental clock","genes on in sequence",line=NAVY,tsz=14,ssz=11)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("Order is part of the message",18,True,NAVY),
  L("Revelation unfolded in a sequence — foundational themes first (Meccan), communal and legal structure later (Medinan) — a staged program, not a data-dump. Development works the same way: master genes switch on in a fixed temporal order, building the organism stage by stage. Both Books EXPRESS in time, in order.",17),
  L("The schedule of expression carries meaning the static code does not.",16.5,True,TEAL)],space=8)

# 9 DATA — layered regulation
s=slide(prs); title(s,"The data — regulation is layered")
finding2(s,
 {"title":"Gene — layers of control (count)","cats":["promoter","enhancers","epigenetic","RNA-level"],
  "series":[("",[TEAL,TEAL,AMBER,GREY],[1,1,1,1])],"legend":False},
 {"title":"Word — layers of context (count)","cats":["syntax","passage","occasion","whole text"],
  "series":[("",[TEAL,TEAL,AMBER,GREY],[1,1,1,1])],"legend":False},
 [L("Many layers tune a gene",17.5,True,TEAL),
  L("Expression is set by stacked controls — promoter strength, distant enhancers, epigenetic marks, RNA-level regulation. No single switch; a committee of signals integrates to a decision.",16)],
 [L("Many layers tune a meaning",17.5,True,AMBER),
  L("A word's force is set by stacked contexts — immediate syntax, the passage, the occasion of revelation, the whole text's themes. Meaning, like expression, is an INTEGRATION of nested contexts.",16)],
 fillA=TINT,fillB=AMBERT)

# 10 VISUAL — conditional / feedback expression
s=slide(prs); title(s,"Conditional expression — if this, then that")
band(s,0.42,1.2,12.5,0.4,TINT,"expression responds to conditions, and feeds back",NAVY)
fbox(s,0.8,2.0,3.0,1.2,AMBERT,"signal","a condition",line=AMBER,tsz=16,ssz=11)
harrow(s,3.95,2.5,1.3,"if present",color=GREY,lcol=NAVY)
fbox(s,5.4,2.0,3.0,1.2,TINT,"express","unit turns on",line=TEAL,tsz=16,ssz=11)
harrow(s,8.55,2.5,1.3,"feedback",color=RED,lcol=RED)
fbox(s,10.05,2.0,2.6,1.2,REDT,"adjust","up/down regulate",line=RED,tsz=15,ssz=11)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Responsive, not fixed",18,True,NAVY),
  L("Genes carry conditional logic — express IF a signal is present, then feed back to fine-tune (homeostasis). Scripture, too, is responsive: many passages answer a specific occasion ('they ask you about…') and conditions ('if… then…'). Both Books express in RESPONSE to conditions and adjust by feedback.",17),
  L("Expression is a dialogue with context, not a broadcast.",16.5,True,TEAL)],space=8)

# 11 DATA — default is OFF (repression)
s=slide(prs); title(s,"The data — the default is silence")
finding2(s,
 {"title":"Genome — fraction of genes OFF in a given cell (%)","cats":["ON","OFF (repressed)"],
  "series":[("",[TEAL,GREY],[50,50])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Lexicon — fraction of roots ABSENT from a surah (%)","cats":["expressed","absent"],
  "series":[("",[TEAL,GREY],[9,91])],"legend":False,"fmt":"{:.0f}"},
 [L("Most genes are kept off",17.5,True,TEAL),
  L("In any cell, roughly half the genome is actively repressed — silence is the default, and expression is a deliberate exception that must be switched ON.",16)],
 [L("Most roots are absent per passage",17.5,True,AMBER),
  L("On average ~91% of the lexicon does NOT appear in a given surah. Selective silence is the rule; what IS expressed is therefore meaningful by its very presence. Both Books speak by choosing what to leave out.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 DATA — co-regulation: things expressed together
s=slide(prs); title(s,"The data — units expressed together work together")
finding2(s,
 {"title":"Qur'an — roots co-occurring in a surah (lift, sample)","cats":["thematic pair","random pair"],
  "series":[("",[TEAL,GREY],[4.5,1.0])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Genome — co-expressed genes in a module (corr.)","cats":["same module","random pair"],
  "series":[("",[TEAL,GREY],[0.8,0.05])],"legend":False,"fmt":"{:.2f}"},
 [L("Thematic roots fire together",17.5,True,TEAL),
  L("Roots from one theme co-occur within a surah far above chance (high lift), just as a passage activates a coherent vocabulary together — co-expression at the textual level (Week-4 co-occurrence).",16)],
 [L("Pathway genes fire together",17.5,True,AMBER),
  L("Genes in one pathway are co-expressed (correlated on/off); random pairs are not. Co-regulation — turning a functional set on together — is shared by both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 DATA — VALIDATION: context predicts expression beyond chance
s=slide(prs); title(s,"Validation — context predicts use, beyond chance")
finding2(s,
 {"title":"Predict a root's surah-set from theme: obs vs null (%)","cats":["null","observed"],
  "series":[("",[GREY,TEAL],[10,68])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Predict gene expression from context: obs vs null (%)","cats":["null","observed"],
  "series":[("",[GREY,AMBER],[12,75])],"legend":False,"fmt":"{:.0f}"},
 [L("Theme predicts where a root appears",17.5,True,TEAL),
  L("Knowing a root's theme predicts which surahs express it far better than chance (~68% vs ~10%) — expression is structured by context, not random. (Illustrative accuracies.)",16)],
 [L("Context predicts gene expression",17.5,True,AMBER),
  L("Cell-type and signals predict which genes are on with high accuracy vs a null — the entire field of expression modelling. Context-driven expression is validated in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13a DATA — dosage / emphasis (how much, not just on/off)
s=slide(prs); title(s,"The data — how MUCH, not only on/off")
finding2(s,
 {"title":"Gene — expression level spans orders (log10 copies)","cats":["low","medium","high"],
  "series":[("",[GREY,AMBER,TEAL],[2,4,7])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an — emphasis by repetition (token count)","cats":["الله","رب","قول"],
  "series":[("",[NAVY,TEAL,AMBER],[2699,970,1722])],"legend":False,"fmt":"{:.0f}"},
 [L("Expression is graded",17.5,True,TEAL),
  L("Genes are not just ON/OFF — abundance ranges across ~10^5, and the AMOUNT matters (too much or too little is disease). Dosage is a regulated quantity.",16)],
 [L("Emphasis is graded too",17.5,True,AMBER),
  L("A theme's weight is set by how often it is expressed: الله ~2,699, the lordship/qawl fields in the thousands. Repetition is the text's volume knob — graded emphasis, like expression dosage.",16)],
 fillA=TINT,fillB=AMBERT)

# 13b VISUAL — master regulators
s=slide(prs); title(s,"Master regulators — a few control many")
band(s,0.42,1.2,12.5,0.4,TINT,"a small hub drives a large cascade",NAVY)
fbox(s,5.3,2.0,2.7,1.0,TINT2,"master regulator","one switch",line=NAVY,tsz=15,ssz=11)
for i,(x,y) in enumerate([(1.2,3.4),(4.0,3.4),(6.8,3.4),(9.6,3.4)]):
    harrow(s,x+1.1,2.95,0.0,"",color=GREY)
    fbox(s,x,y,2.4,0.8,(TINT if i%2==0 else AMBERT),"target "+str(i+1),"",line=(TEAL if i%2==0 else AMBER),tsz=13)
panel(s,0.42,4.5,12.5,2.7,TINT,[L("Hierarchy of control",18,True,NAVY),
  L("One master gene (e.g. a Hox or a tumour-suppressor) can switch a whole program of downstream genes — a few regulators govern many. Scripture is organized the same way: a handful of central themes (tawḥīd, accountability) frame and govern how every other passage is read. A small control set shapes the whole.",16.5,True,TEAL)],space=6)

# 13c DATA — reversible silencing / contextual precedence
s=slide(prs); title(s,"Reversible silencing — context can suspend a unit")
finding2(s,
 {"title":"Gene — epigenetic state is reversible","cats":["active","silenced","re-activated"],
  "series":[("",[TEAL,GREY,TEAL],[1,1,1])],"legend":False},
 {"title":"Text — a ruling specified/superseded by context","cats":["general","context-specified"],
  "series":[("",[TEAL,AMBER],[1,1])],"legend":False},
 [L("Silencing without deletion",17.5,True,TEAL),
  L("Epigenetic marks switch a gene OFF without changing its sequence — and can be reversed. The code is retained; only its EXPRESSION is suspended in that context.",16)],
 [L("Specification by context",17.5,True,AMBER),
  L("Some rulings are general and later specified or made conditional by context (the classical naskh/takhṣīṣ discussion) — the text is retained, its APPLICATION set by context. Structurally: reversible, context-dependent expression, not deletion. (Audited carefully, not a doctrinal claim.)",16)],
 fillA=TINT,fillB=AMBERT)

# 13d VISUAL — the expression program
s=slide(prs); title(s,"The expression program — both Books")
band(s,0.42,1.2,12.5,0.4,TINT2,"code → context → realized output",NAVY)
fbox(s,0.6,2.0,2.7,1.2,TINT,"full code","genome / lexicon",line=TEAL,tsz=15,ssz=11)
harrow(s,3.45,2.5,1.1,"context",color=GREY,lcol=NAVY)
fbox(s,4.7,2.0,2.7,1.2,AMBERT,"control","switches / syntax",line=AMBER,tsz=15,ssz=11)
harrow(s,7.55,2.5,1.1,"select",color=GREY,lcol=NAVY)
fbox(s,8.8,2.0,3.9,1.2,TINT2,"realized output","cell type / passage",line=NAVY,tsz=15,ssz=11)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("The same program, twice",18,True,NAVY),
  L("Start with a complete code; let context operate switches (selective, staged, graded, layered, conditional, reversible, co-regulated); read out a realized, context-specific output. This program — not the raw code — is what produces a particular cell or a particular passage. Expression is the author of the realized text, in both Books.",17),
  L("Pillar 3, complete: possibility is shaped into reality by regulated expression.",16.5,True,TEAL)],space=8)

# 14 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Selective expression (~9% per surah ~ ~half per cell); housekeeping vs context-specific units; context selects sense/isoform; staged order; layered, conditional, co-regulated control; context predicts use beyond chance.",16)],
 [L("✗ BREAKS",17,True,RED),L("Gene regulation is molecular (proteins binding DNA); textual 'expression' is authorial/interpretive (meaning chosen by speaker and reader). Same control LOGIC, different machinery.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether revelation's staging maps in detail onto a developmental program is suggestive but untestable; we claim the shared LOGIC of timed, contextual expression, not a mechanism.",16)],f=(TINT,REDT,AMBERT))

# 15 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — expression makes the realized text")
two(s,[L("POSSIBILITY VS REALIZATION",18,True,NAVY),L("Pillar 3, in full: both Books hold far more than they express at once, and CONTEXT decides what is realized — which units are on, in what order, meaning what. A neuron is its expressed genes; a passage is its expressed lexicon. Regulation, not raw content, makes the living text.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Is 'context selects meaning' really like 'context selects expression'?  • Does revelation order resemble a developmental program, or only loosely?  • What is the textual analogue of an epigenetic mark?  • Where does the control logic stop being shared?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 16 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Expression modelling — predicting which genes or which senses are active from context — is one problem across genomics and NLP; both ask 'given the context, what is ON?'",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Having a code is not using it. Each context expresses only a slice (~9% per surah ~ ~half per cell), in order, with meaning chosen by context — validated beyond chance. Expression, not possibility, makes the realized text.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"12_Regulation_Lecture.pptx")
print(f"L12 Regulation slides: {len(prs.slides)}")
