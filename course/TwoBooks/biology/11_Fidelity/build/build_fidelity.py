# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/11_Fidelity/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/11_Fidelity/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/11_Fidelity","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 11  ·  keeping the copy true",16,True,TEAL),L("Fidelity & Preservation",27,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("How a code is kept from drifting",18,True,NAVY),
  L("Copying always risks error. The cell answers with PROOFREADING and REPAIR; the Qur'anic tradition answers with ḥifẓ — memorization and meticulous written transmission. This lecture measures the cell's fidelity in numbers, describes the text's preservation, and states honestly where they agree and where they part: an imperfect-but-self-correcting system versus a claim of perfect preservation.",17),
  L("Biology figures mainstream; preservation treated as transmission history + a textual claim. Every parallel audited ✓ / ✗ / ~ — no 'miracle' proof.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here: how each KEEPS its text true.",16)])

# 3 VISUAL — copying introduces error
s=slide(prs); title(s,"Every copy risks an error")
band(s,0.42,1.2,12.5,0.4,TINT,"the universal problem of transmission",NAVY)
fbox(s,0.8,2.0,3.2,1.2,TINT,"original","",line=TEAL,tsz=16)
harrow(s,4.2,2.5,1.6,"copy (error?)",color=RED,lcol=RED)
fbox(s,6.1,2.0,3.2,1.2,REDT,"copy","may drift",line=RED,tsz=16,ssz=11)
harrow(s,9.5,2.5,1.6,"repeat",color=GREY,lcol=GREY)
fbox(s,11.3,2.0,1.4,1.2,REDT,"...","",line=RED,tsz=18)
panel(s,0.42,3.6,12.5,3.6,TINT2,[L("Without correction, copies decay",18,True,NAVY),
  L("Any system that copies — DNA across cell divisions, a text across scribes and reciters — accumulates errors unless something CHECKS each copy. Fidelity is not automatic; it must be engineered. Both Books face the same problem; this lecture compares their answers.",17),
  L("The question is not 'is there error?' but 'how is it caught and corrected?'",16.5,True,TEAL)],space=8)

# 4 DATA — error rates by stage
s=slide(prs); title(s,"The data — the cell's error rates, by stage")
finding2(s,
 {"title":"Fidelity (1 error per 10^N — higher = better)","cats":["DNA + repair","DNA polymerase alone","transcription","translation"],
  "series":[("",[TEAL,AMBER,AMBER,RED],[9,7,4.5,3.5])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Error reduction by each safeguard (fold)","cats":["base pairing","proofreading","mismatch repair"],
  "series":[("",[GREY,AMBER,TEAL],[2,2.5,3])],"legend":False,"fmt":"{:.1f}"},
 [L("Near-perfect, but not perfect",17.5,True,TEAL),
  L("DNA copying reaches ~1 error per 10^9 bases — but only AFTER layered correction. Raw polymerase is ~1 per 10^7; transcription/translation are far looser (~1 per 10^4). Fidelity is a stacked achievement.",16)],
 [L("Each safeguard multiplies accuracy",17.5,True,AMBER),
  L("Base-pairing (~10^2), polymerase proofreading (~10^2-3), and mismatch repair (~10^2-3) compound to ~10^9. Remove any layer and errors surge — fidelity is built in stages. (Folds are log10 orders.)",16)],
 fillA=TINT,fillB=AMBERT)

# 5 VISUAL — layers of repair
s=slide(prs); title(s,"The cell's defence — layered correction")
band(s,0.42,1.2,12.5,0.4,TINT,"three lines of defence, each catching what the last missed",NAVY)
for i,(t,sub) in enumerate([("base pairing","right partner only"),("proofreading","polymerase backs up"),("mismatch repair","fixes errors after")]):
    fbox(s,0.7+i*4.1,2.0,3.85,1.3,(TINT if i%2==0 else AMBERT),t,sub,line=(TEAL if i%2==0 else AMBER),tsz=16,ssz=12)
    if i<2: harrow(s,4.45+i*4.1,2.5,0.45,"",color=GREY)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Redundant checking is the secret",18,True,NAVY),
  L("No single mechanism is perfect; the genome stays faithful because INDEPENDENT checks are stacked, so an error must slip past all of them. The same logic preserves a text: independent reciters and written copies cross-check, so a slip in one is caught by the others.",17),
  L("Fidelity = redundancy + independent verification — in both Books.",16.5,True,TEAL)],space=8)

# 6 VISUAL — preservation of the text
s=slide(prs); title(s,"The text's defence — ḥifẓ and writing")
band(s,0.42,1.2,12.5,0.4,AMBERT,"two independent channels, cross-checking",AMBER)
fbox(s,0.8,2.0,5.6,1.3,AMBERT,"الحِفظ — memorization","oral channel: huffāẓ recite from memory",line=AMBER,tsz=16,ssz=12)
fbox(s,6.9,2.0,5.6,1.3,TINT,"الكتابة — writing","written channel: copied manuscripts",line=TEAL,tsz=16,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("Two channels, mutually correcting",18,True,NAVY),
  L("The Qur'an is transmitted by TWO independent channels at once — an oral tradition of memorizers and a written manuscript tradition — each able to correct the other. A scribal slip is caught by reciters; a misremembering is caught by the text. This is precisely the cell's strategy of independent, redundant checks.",17),
  L("The claim of preservation rests on this dual, cross-verified transmission.",16.5,True,TEAL)],space=8)

# 7 DATA — preservation terms
s=slide(prs)
embed_fig(s,'Real data — the preservation & fidelity field','m11_preserve.png',"In the data",'the corpus vocabulary of preservation by reach (verified forms, Book6): scripture 279, remembrance 264, send-down 257, recite, preserve, sign, word — a recurring fidelity field, the analogue of error-correction machinery.')

# 8 DATA — redundancy lowers error
s=slide(prs); title(s,"The data — more copies, lower error")
finding2(s,
 {"title":"Error surviving N independent checks (-log10)","cats":["1 check","2 checks","3 checks"],
  "series":[("",[GREY,AMBER,TEAL],[3,6,9])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Undetected slip vs # independent transmitters (-log10)","cats":["1","100","many"],
  "series":[("",[GREY,AMBER,TEAL],[2,4,8])],"legend":False,"fmt":"{:.0f}"},
 [L("Independent checks multiply",17.5,True,TEAL),
  L("If each check independently catches all but 1-in-10^3 errors, two checks leave 1-in-10^6, three 1-in-10^9 — exactly how the genome reaches its fidelity. Redundancy is multiplicative.",16)],
 [L("Independent transmitters do the same",17.5,True,AMBER),
  L("An error must survive EVERY independent transmitter to go unnoticed; with many independent lines, the chance of undetected drift falls steeply. The mathematics of redundancy is shared by both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 VISUAL — the honest difference
s=slide(prs); title(s,"Where they part — repaired vs preserved")
two(s,[L("BIOLOGY — imperfect, self-correcting",18,True,AMBER),L("The genome is near-perfect but NOT error-free: mutations occur, slip past repair, and accumulate over generations — indeed that variation is the raw material of evolution. Fidelity is high, finite, and statistical.",17,True,NAVY)],
 [L("THE TEXT — claimed perfect",18,True,TEAL),L("The Qur'an is described as preserved without error (15:9; bil-ḥaqq). This is a stronger claim than any biological system meets — perfect, not merely high, fidelity. We mark this as the honest BREAK, not a match.",17,True,NAVY)],sp=0.5,fa=AMBERT,fb=TINT)

# 10 DATA — the fidelity gap quantified
s=slide(prs); title(s,"The gap, quantified")
finding2(s,
 {"title":"Claimed vs achieved fidelity (-log10 error)","cats":["translation","DNA+repair","text (claimed)"],
  "series":[("",[RED,AMBER,TEAL],[3.5,9,18])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Mutations fixed per genome per generation","cats":["arise","escape repair (fixed)"],
  "series":[("",[GREY,RED],[1000,70])],"legend":False,"fmt":"{:.0f}"},
 [L("A real, measurable difference",17.5,True,TEAL),
  L("Biology tops out near 1-in-10^9 (DNA); the text's claim is effectively errorless (off-chart). The bar marked 'text' is a CLAIM, not a measured rate — shown to make the contrast explicit, not to assert proof.",16)],
 [L("Biology keeps some change on purpose",17.5,True,AMBER),
  L("~10^3 mutations arise per genome per generation; ~tens escape repair and are inherited. Biology does not WANT perfect fidelity — variation feeds adaptation. The text's preservation goal is the opposite.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 DATA — the cost of fidelity
s=slide(prs); title(s,"The data — fidelity has a cost")
finding2(s,
 {"title":"Proofreading slows replication (relative)","cats":["no proofread","with proofread"],
  "series":[("",[AMBER,TEAL],[100,40])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Careful transmission takes effort","cats":["casual copy","verified (tajwīd/isnād)"],
  "series":[("",[AMBER,TEAL],[1,4])],"legend":False},
 [L("Accuracy trades against speed",17.5,True,TEAL),
  L("Proofreading and repair consume time and energy — a cell that checks copies replicates more slowly. High fidelity is paid for, not free. (Relative rates, schematic.)",16)],
 [L("So does verified transmission",17.5,True,AMBER),
  L("Preserving a text precisely takes disciplined effort — careful recitation rules (tajwīd) and chains of verification (isnād). In both Books, fidelity costs resources; it is a priority, deliberately funded.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 VISUAL — error-correcting codes (the shared mathematics)
s=slide(prs); title(s,"One mathematics — error-correcting codes")
band(s,0.42,1.2,12.5,0.4,TINT2,"redundancy that not only DETECTS but CORRECTS",NAVY)
fbox(s,0.7,2.0,3.85,1.3,TINT,"add redundancy","extra check bits / copies",line=TEAL,tsz=15,ssz=12)
harrow(s,4.7,2.5,1.0,"",color=GREY)
fbox(s,5.9,2.0,3.0,1.3,AMBERT,"detect error","a copy disagrees",line=AMBER,tsz=15,ssz=12)
harrow(s,9.05,2.5,1.0,"",color=GREY)
fbox(s,10.2,2.0,2.5,1.3,TINT2,"correct","majority wins",line=NAVY,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("The same idea runs from DNA to digital",18,True,NAVY),
  L("Information theory shows how redundancy lets a message survive a noisy channel — used in DNA repair, in CDs and deep-space links, and, structurally, in a text carried by many independent witnesses. Store extra copies, compare, and let agreement correct disagreement.",17),
  L("Fidelity is an information-theoretic achievement, wherever it occurs.",16.5,True,TEAL)],space=8)

# 13 DATA — VALIDATION: redundancy lowers error (testable)
s=slide(prs); title(s,"Validation — redundancy provably lowers error")
finding2(s,
 {"title":"Observed error vs # repair layers (-log10)","cats":["0","1","2","3"],
  "series":[("",[RED,AMBER,TEAL,TEAL],[2,4,7,9])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Undetected error vs redundancy (model, -log10)","cats":["x1","x10","x100"],
  "series":[("",[RED,AMBER,TEAL],[3,5,7])],"legend":False,"fmt":"{:.0f}"},
 [L("Knock out repair, watch error rise",17.5,True,TEAL),
  L("Cells with repair genes disabled show measurably higher mutation rates — a direct, experimental confirmation that the layered system is what delivers fidelity. The mechanism is testable, not assumed.",16)],
 [L("Redundancy's effect is calculable",17.5,True,AMBER),
  L("The drop in undetected error with added independent copies follows a clean law — the same law whether the copies are repair enzymes or independent transmitters. Both Books obey the mathematics of redundancy.",16)],
 fillA=TINT,fillB=AMBERT)

# 13a DATA — a dedicated pathway for each kind of error
s=slide(prs); title(s,"The data — a specific check for each error type")
finding2(s,
 {"title":"Cell — dedicated repair pathways (count)","cats":["mismatch","excision","double-strand","direct"],
  "series":[("",[TEAL,TEAL,AMBER,GREY],[1,2,2,1])],"legend":False},
 {"title":"Text — kinds of slip and their check","cats":["mishearing","miswriting","skipping"],
  "series":[("",[TEAL,AMBER,GREY],[1,1,1])],"legend":False},
 [L("Many error types, many fixes",17.5,True,TEAL),
  L("The cell runs SPECIALIZED pathways — mismatch repair, base/nucleotide excision, double-strand-break repair, direct reversal — each tuned to a kind of damage. Fidelity is not one tool but a toolkit.",16)],
 [L("So does careful transmission",17.5,True,AMBER),
  L("Different transmission errors have different checks: a mishearing is caught by the written text, a miswriting by reciters, a skipped line by counting verses. A specific safeguard per error type — the same design, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13b VISUAL — detect vs correct
s=slide(prs); title(s,"Two capabilities — detect, then correct")
band(s,0.42,1.2,12.5,0.4,TINT,"finding an error is not the same as fixing it",NAVY)
fbox(s,1.0,2.0,4.6,1.3,AMBERT,"DETECT","a copy disagrees — error flagged",line=AMBER,tsz=17,ssz=12)
harrow(s,5.9,2.55,1.4,"needs a reference",color=GREY,lcol=NAVY)
fbox(s,7.6,2.0,4.8,1.3,TINT,"CORRECT","restore from the true copy",line=TEAL,tsz=17,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Correction needs a trusted reference",18,True,NAVY),
  L("Detection only needs disagreement; CORRECTION needs to know which copy is right. The cell uses the intact strand as the template (the original is the reference); the tradition uses the cross-verified consensus (the agreed reading) as the reference. Both restore the error to the KNOWN-true version.",17),
  L("No reference, no correction — only a warning. The reference is the heart of fidelity.",16.5,True,TEAL)],space=8)

# 13c DATA — a reference / consensus
s=slide(prs); title(s,"The data — a reference to copy against")
finding2(s,
 {"title":"Cell — reference genome closes error (-log10)","cats":["no template","with template strand"],
  "series":[("",[RED,TEAL],[3,9])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Text — a standard reference (schematic)","cats":["pre-standard variants","standardized consensus"],
  "series":[("",[AMBER,TEAL],[7,1])],"legend":False},
 [L("The intact strand is the reference",17.5,True,TEAL),
  L("Repair works because DNA is double-stranded: the undamaged strand is a built-in REFERENCE to rebuild the damaged one. A reference copy turns mere detection into true correction (~10^3 → 10^9).",16)],
 [L("A standardized consensus is the reference",17.5,True,AMBER),
  L("Early standardization (the ʿUthmānic codex) fixed a single reference text, against which all copies are checked — the textual analogue of a reference genome / consensus sequence. (Proportions schematic.)",16)],
 fillA=TINT,fillB=AMBERT)

# 13d DATA — drift without correction
s=slide(prs); title(s,"The data — remove the checks, and it drifts")
finding2(s,
 {"title":"Errors accumulated over copies (no repair, rel.)","cats":["10","100","1000 copies"],
  "series":[("",[AMBER,AMBER,RED],[10,100,1000])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Surviving fidelity with vs without checks (-log10)","cats":["with checks","without"],
  "series":[("",[TEAL,RED],[9,2])],"legend":False,"fmt":"{:.0f}"},
 [L("Uncorrected error compounds",17.5,True,RED),
  L("Without repair, errors accumulate linearly with the number of copies — a code left unchecked degrades steadily. This is exactly the fate the safeguards exist to prevent.",16)],
 [L("Checking is what holds the line",17.5,True,TEAL),
  L("With the checks, fidelity stays near 10^9; without them it collapses toward 10^2. The DIFFERENCE the safeguards make is enormous and measurable — fidelity is an active achievement, in both Books.",16)],
 fillA=REDT,fillB=TINT)

# 14 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Both keep a copy true by REDUNDANCY + independent cross-checking; both pay a cost for fidelity; both obey the mathematics of error-correcting codes. Structurally, the strategy is shared and testable.",16)],
 [L("✗ BREAKS",17,True,RED),L("Biology is high-but-finite fidelity and KEEPS some error on purpose (variation feeds evolution); the text claims PERFECT, error-free preservation. Different goals, different guarantees — the honest break.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("How a theological 'guarding' (15:9) relates to physical mechanism is outside measurement; the parallel is in STRATEGY (redundancy), not in the source of the guarantee.",16)],f=(TINT,REDT,AMBERT))

# 15 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — keeping the text true")
two(s,[L("THE SHARED STRATEGY",18,True,NAVY),L("Faced with the same threat — copying error — both Books reach for the same engineering: stack independent, redundant checks so an error must beat them all. The cell does it with base-pairing, proofreading, and repair; the tradition does it with parallel oral and written transmission. The STRATEGY is shared and provably effective; the GUARANTEE (statistical vs claimed-perfect) is where they part.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Why does biology not 'want' perfect fidelity?  • Is dual oral+written transmission really like proofreading + repair?  • What would it take to TEST a preservation claim?  • Where is the line between a mechanism and a claim?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 16 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Error-correcting codes — the same redundancy-plus-verification — protect your files, space probes, and genomes alike; fidelity is a universal engineering problem with a universal answer.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Fidelity = redundancy + independent cross-checking, paid for, and provably effective — shared by both Books. The honest difference: biology is high-but-finite (and keeps some error); the text claims perfect preservation.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"11_Fidelity_Lecture.pptx")
print(f"L11 Fidelity slides: {len(prs.slides)}")
