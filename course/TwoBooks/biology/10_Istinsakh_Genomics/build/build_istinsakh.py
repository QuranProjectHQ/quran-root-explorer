# -*- coding: utf-8 -*-
import sys, json
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/10_Istinsakh_Genomics/build")
from st_slides import *
from diagrams import fbox,harrow,vdash,band,dna,scales
from pptx.util import Inches,Pt
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/10_Istinsakh_Genomics/"
GN=json.load(open("snip_genome.json",encoding="utf-8"))
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/10_Istinsakh_Genomics","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
SUP=("✓ AUDIT — SUPPORTED by the text",16,True,TEAL)
BRK=("✗ AUDIT — the text BREAKS the parallel",16,True,RED)

def pipeline(s,y,boxes,arrows,line):
    # boxes: [(title,sub,fill)] x4 ; arrows: [label]x3 ; returns center xs
    bw=2.45; aw=0.82; x=0.55; cx=[]
    for i,(t,sub,fl) in enumerate(boxes):
        fbox(s,x,y,bw,1.05,fl,t,sub,line=line,tsz=16,ssz=12.5); cx.append(x+bw/2)
        if i<3:
            harrow(s,x+bw+0.0,y+0.38,aw,arrows[i],color=GREY,lcol=line); 
        x+=bw+aw
    return cx

# 1 TITLE  (text + mini two-books motif)
s=slide(prs)
panel(s,0.42,1.15,12.5,1.7,TINT2,[L("THE TWO BOOKS  ·  a Qur’an-and-science lecture",17,True,TEAL),
  L("Istinsākh & the Genome",33,True,NAVY)],space=8)
fbox(s,1.4,3.3,4.2,1.5,TINT,"Book of CREATION","the genome",line=TEAL,tsz=18,ssz=13)
harrow(s,5.9,3.9,1.4,"āyāt (signs)",color=AMBER,lcol=AMBER,h=0.34)
fbox(s,7.6,3.3,4.2,1.5,AMBERT,"Book of SCRIPTURE","the record of deeds",line=AMBER,tsz=18,ssz=13)
panel(s,0.42,5.2,12.5,1.95,REDT,[L("“…We were TRANSCRIBING (نستنسخ) what you used to do.” — 45:29",17,True,NAVY),
  L("A labelled analogy, audited stage by stage: ✓ Supported · ✗ Breaks · ~ Silent-but-surmisable. No “scientific-miracle” claims.",16,False,INK)],space=9)

s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three_stack(s,
 [L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur’an: God’s speech, revealed in language. The composed Book — tadwīn, “what is set down.” The Book of SCRIPTURE.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The Universe: God’s deed, revealed in creation. The Book of CREATION — takwīn, “what is brought into being.”",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same SOURCE — Allah. Same primary ADDRESSEE — the human being (insān); the jinn too (a later lecture). Both are āyāt (signs); this series reads them side by side, never collapsed.",16)],
 fills=(TINT,AMBERT,TINT2))
# 2 QUESTION (text)
s=slide(prs); title(s,"The question — a life recorded, to what end?")
two(s,[L("THE HOOK",18,True,NAVY),L("The Qur’an: deeds are written down, by watchers, omitting nothing. Biology: cells continually TRANSCRIBE coded instructions, expressed later as traits.",17),L("One shape: copy now, express later, lose nothing. We test it honestly.",17,True,TEAL)],
 [L("THE RULE",18,True,RED),L("For every parallel: does the text SUPPORT it, BREAK it, or is it SILENT but surmisable? Never claim the Qur’an “contains” genetics.",17,True,NAVY)],sp=0.45,fa=TINT2,fb=REDT)

# 3 VISUAL — central dogma
s=slide(prs); title(s,"The genome’s pipeline (central dogma)")
band(s,0.42,1.2,12.5,0.46,TINT,"BOOK OF CREATION — molecular biology",TEAL)
dna(s,0.7,2.2,0.9,1.6,TEAL)
pipeline(s,2.6,[("DNA","genome",TINT),("mRNA","transcript",TINT),("protein","",TINT),("phenotype","trait",TINT)],
         ["transcription","translation","express"],TEAL)
panel(s,0.42,4.55,12.5,2.65,TINT2,[L("Read left to right",18,True,NAVY),
  L("An inherited archive (DNA) is COPIED into a working message (transcription), the message is READ into a protein (translation), and proteins make the organism’s visible traits — the phenotype, expressed later, in context.",17),
  L("Three moves: copy → read → express. Information is preserved throughout.",17,True,TEAL)],space=9)

# 4 VISUAL — deeds pipeline
s=slide(prs); title(s,"The deeds pipeline (the Qur’an’s record)")
band(s,0.42,1.2,12.5,0.46,AMBERT,"BOOK OF SCRIPTURE — the record of deeds",AMBER)
pipeline(s,2.4,[("deeds","free acts",AMBERT),("record","الكتاب",AMBERT),("reckoning","اقرأ كتابك",AMBERT),("outcome","hereafter",AMBERT)],
         ["istinsākh 45:29","read & weigh 17:14","requital"],AMBER)
scales(s,10.9,1.9,1.0,col=NAVY)
panel(s,0.42,4.2,12.5,3.0,TINT,[L("The same three moves",18,True,NAVY),
  L("Free deeds are TRANSCRIBED into a record (نستنسخ, 45:29), the record is READ — “Read your record; you suffice as your own reckoner” (17:14) — and weighed on the just scales (21:47), and the consequence is EXPRESSED in the hereafter.",17),
  L("copy → read → express — stated in the text’s own words.",17,True,TEAL)],space=9)

# 4b VISUAL DATA — rich charts
s=slide(prs)
embed_fig(s,'Real data — the recording / preservation field','m10_recording.png',"In the data","a spectrum of istinsākh (transcription) concepts by āyah-reach: write كتب 279, remember ذكر 264, send-down نزل 257 … transcribe نسخ 4 — the Qur'an's own vocabulary of recording and preservation, computed from Book6.")

# 5 VISUAL — THE PARALLEL MAPPING (centerpiece)
s=slide(prs); title(s,"The two pipelines, side by side")
band(s,0.42,1.18,12.5,0.42,TINT,"GENOME",TEAL)
cxt=pipeline(s,1.72,[("DNA","genome",TINT),("mRNA","transcript",TINT),("protein","",TINT),("phenotype","trait",TINT)],
         ["transcription","translation","express"],TEAL)
band(s,0.42,4.95,12.5,0.42,AMBERT,"DEEDS",AMBER)
cxb=pipeline(s,5.5,[("deeds","free acts",AMBERT),("record","الكتاب",AMBERT),("reckoning","اقرأ كتابك",AMBERT),("outcome","hereafter",AMBERT)],
         ["istinsākh","read & weigh","requital"],AMBER)
for xc,tag in zip(cxt,["✓","✓","✓","✓"]):
    vdash(s,xc,2.77,5.5,tag,col=TEAL)
panel(s,0.42,6.95,12.5,0.42,TINT2,[L("Same structure, stage for stage — copy → read → express, lose nothing.  (Where it BREAKS comes next.)",14.5,True,NAVY)],space=2)

# 6 PRIMER — analogy + verdicts (text)
s=slide(prs); title(s,"What an analogy is — and the three verdicts")
two(s,[L("MAPS STRUCTURE, NOT IDENTITY",18,True,TEAL),L("An analogy lines up the SHAPE of two systems to think with. It is not a claim of identity, shared mechanism, or prediction.",17,True,NAVY)],
 [L("THREE VERDICTS  (the Week-9 audit)",18,True,RED),L("✓ SUPPORTED — the text states it.",16.5,True,TEAL),L("✗ BREAKS — the text contradicts it.",16.5,True,RED),L("~ SILENT — but surmisable from the text.",16.5,True,AMBER)],sp=0.42,fa=TINT,fb=TINT2)

# 7 VISUAL — fidelity break
s=slide(prs); title(s,"Where it BREAKS — fidelity")
fbox(s,0.7,1.8,5.4,1.6,AMBERT,"BIOLOGY","copying has ERRORS → repair fixes most",line=AMBER,tsz=17,ssz=13)
fbox(s,7.2,1.8,5.4,1.6,TINT,"THE RECORD","bil-ḥaqq — ERROR-FREE",line=TEAL,tsz=17,ssz=13)
harrow(s,6.15,2.45,1.0,"vs",color=RED,lcol=RED)
panel(s,0.42,3.9,12.5,3.3,REDT,[L("✗ The parallel cracks here",18,True,RED),
  L("Biology is accurate but imperfect — mutations happen, and cells run proofreading and repair. The Qur’anic record is described as flawless: “in truth,” “you are not wronged a single thread.”",17),
  L("Perfect fidelity vs an error-prone, self-correcting system. Honest verdict: BREAKS.",17,True,NAVY)],space=9)

# 7b BIOLOGY DATA — fidelity & code size
s=slide(prs); title(s,"Molecular biology — by the numbers")
finding2(s,
 {"title":"Fidelity (1 error per 10^N — higher = better)","cats":["DNA copy","transcription","translation"],
  "series":[("",[TEAL,AMBER,RED],[9,4.5,3.5])],"legend":False,"fmt":"{:.1f}"},
 {"title":"The genetic code is tiny","cats":["bases","amino acids","codons"],
  "series":[("",[TEAL,AMBER,NAVY],[4,20,64])],"legend":False},
 [L("Finite fidelity — quantified",17.5,True,TEAL),
  L("DNA copying is near-perfect — ~1 error per 10^9 bases (proofreading + repair); transcription/translation are looser (~1 per 10^4). The record’s claim is PERFECT (bil-ḥaqq) — off this chart. That is the fidelity break, in numbers.",16)],
 [L("A small code, vast output",17.5,True,AMBER),
  L("4 bases → 64 triplet codons → 20 amino acids build ALL proteins; the Qur’an: ~28 letters → ~1700 roots. Small alphabets, unbounded output.",16)],
 fillA=TINT,fillB=AMBERT)
# 8 VISUAL — one-way arrow
s=slide(prs); title(s,"A one-way bias toward the good")
fbox(s,0.8,2.0,3.4,1.4,TINT,"الحسنات","good deeds",line=TEAL,tsz=20,ssz=13)
harrow(s,4.4,2.55,2.2,"drive away (11:114)",color=TEAL,lcol=TEAL,h=0.34)
fbox(s,6.9,2.0,3.4,1.4,REDT,"السيئات","bad deeds",line=RED,tsz=20,ssz=13)
fbox(s,10.6,2.0,2.1,1.4,REDT,"✗","no auto-reverse",line=RED,tsz=22,ssz=12)
panel(s,0.42,3.9,12.5,3.3,AMBERT,[L("✓ SUPPORTED — and directional",18,True,TEAL),
  L("“The good deeds drive away the bad” (11:114); repentance converts bad into good (25:70). There is NO symmetric rule that bad erases good — the arrow points one way.",17),
  L("Biology echoes this only as ERROR-CORRECTION (DNA repair; error-correcting codes; smaller chips → more capacity). CAVEAT: do NOT read it as “evolution → perfection” — that overreach is excluded.",16,True,RED)],space=8)

# 9 VISUAL — editable transcript / umm al-kitab
s=slide(prs); title(s,"The transcript is editable — repentance")
fbox(s,0.8,1.9,3.6,1.4,TINT,"the record","سيئات (draft)",line=NAVY,tsz=17,ssz=13)
harrow(s,4.6,2.45,2.0,"repentance 25:70",color=TEAL,lcol=TEAL,h=0.34)
fbox(s,6.9,1.9,3.6,1.4,TINT,"rewritten","حسنات",line=TEAL,tsz=17,ssz=13)
fbox(s,11.0,1.7,1.8,1.8,AMBERT,"أم الكتاب","master archive (13:39)",line=AMBER,tsz=15,ssz=11)
panel(s,0.42,3.9,12.5,3.3,TINT2,[L("The text OVER-answers biology",18,True,NAVY),
  L("Biology’s pipeline is largely write-once. But the text is not silent: for the repentant, “God CONVERTS their evil deeds into good” (25:70); “He erases what He wills and confirms, and with Him is the Mother of the Record” (13:39).",17),
  L("Hope is built into the pipeline — the draft can be rewritten before the final read.",17,True,TEAL)],space=9)

# 10 VISUAL — scorecard
s=slide(prs); title(s,"The scorecard — going back and forth")
def cols3(items):
    n=len(items); gap=0.3; cw=(12.5-(n-1)*gap)/n
    for i,(hd,hc,fl,lines) in enumerate(items):
        x=0.42+i*(cw+gap); L0=[(hd,18,True,hc)]+[(t,15,False,INK) for t in lines]
        panel(s,x,1.18,cw,7.28-1.18,fl,L0,space=8)
cols3([("✓ SUPPORTED",TEAL,TINT,["Transcription (deeds→record) 45:29","Translation (record→verdict) 17:14","Phenotype (deferred outcome)","Nothing lost / fidelity 78:29","Weighed, not tallied 21:47","Editable by repentance 25:70","One-way bias toward good 11:114"]),
       ("✗ BREAKS",RED,REDT,["Error-free vs error-prone biology","Consequence NOT inherited 35:18","Free deeds AUTHOR the code","Expressed in ANOTHER realm","Moral, not blind expression","“Evolution→perfection” — overreach"]),
       ("~ SILENT (surmisable)",AMBER,AMBERT,["Mechanism of recording — unstated","HOW deeds are “weighed” — unstated","Descendants’ share by grace 52:21","umm al-kitāb as “master archive”"])])

# 10a DATA — exhaustive recording vs complete copy
s=slide(prs); title(s,"The data — nothing omitted")
finding2(s,
 {"title":"Qur'an — completeness terms (occurrences)","cats":["كتب","حفظ","احصى"],
  "series":[("",[NAVY,TEAL,AMBER],[279,42,11])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Cell — every division copies the whole genome (log10)","cats":["bases copied","errors left"],
  "series":[("",[TEAL,GREY],[9.51,0.5])],"legend":False,"fmt":"{:.1f}"},
 [L("'It omits nothing'",17.5,True,TEAL),
  L("The record is exhaustive — 'leaving out nothing small or great' (18:49); writing (كتب) 279, preserving (ḥifẓ) 42, enumerating (aḥṣā) 11. Completeness is a stated property.",16)],
 [L("Every cell, the whole archive",17.5,True,AMBER),
  L("Each cell division copies all ~3.2 billion bases — the COMPLETE genome — into the daughter cell, with proofreading leaving ~1 error per 10^9. Exhaustive copying is a real biological fact too.",16)],
 fillA=TINT,fillB=AMBERT)

# 10b DATA — dedicated machinery / agents
s=slide(prs); title(s,"The data — a dedicated apparatus does the copying")
finding2(s,
 {"title":"Qur'an — recording agents/terms (occurrences)","cats":["رقب (watcher)","شهد (witness)","عدد"],
  "series":[("",[NAVY,TEAL,AMBER],[20,160,57])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Cell — machines of the central dogma (proofread?)","cats":["polymerase","ribosome","repair"],
  "series":[("",[TEAL,AMBER,TEAL],[1,1,1])],"legend":False},
 [L("Recording is staffed, not casual",17.5,True,TEAL),
  L("The text names dedicated agents — watchers (raqīb), witnesses (shahida ~160), noble scribes (kirāman kātibīn). Recording is performed by an apparatus, not left to chance.",16)],
 [L("So is transcription",17.5,True,AMBER),
  L("The cell runs dedicated machines — RNA polymerase, the ribosome, and repair enzymes — each with proofreading. Both Books copy with purpose-built machinery, not by accident.",16)],
 fillA=TINT,fillB=AMBERT)

# 10c VISUAL — deferred expression (timing)
s=slide(prs); title(s,"Copied now, expressed later — the timing")
band(s,0.42,1.2,12.5,0.4,TINT,"the record is written before its consequence appears",NAVY)
fbox(s,0.7,2.0,2.7,1.1,TINT,"gene copied","at conception",line=TEAL,tsz=15,ssz=11)
harrow(s,3.55,2.45,1.5,"developmental clock",color=GREY,lcol=TEAL)
fbox(s,5.2,2.0,2.7,1.1,TINT,"trait expressed","years later",line=TEAL,tsz=15,ssz=11)
fbox(s,8.6,2.0,4.1,1.1,AMBERT,"deed → record → reckoning","outcome deferred (hereafter)",line=AMBER,tsz=14,ssz=11)
panel(s,0.42,3.5,12.5,3.7,TINT2,[L("Expression is deferred in both Books",18,True,NAVY),
  L("A gene is copied at conception but many traits switch on much later, on a developmental schedule. A deed is recorded now but its consequence is expressed later, at the reckoning. In both, there is a gap — copy first, express on a timetable — and the gap is where REGULATION lives (Lecture 12).",17),
  L("Same structure: archive early, express on schedule, in context.",16.5,True,TEAL)],space=8)

# 11 DISCIPLINE (text)
s=slide(prs); title(s,"What this lecture does NOT claim")
two(s,[L("NOT A “SCIENTIFIC MIRACLE”",18,True,RED),L("No verse is read as biology; no biology is offered as proof. iʿjāz-by-science is exactly the unsubstantiated move this course rejects.",17,True,NAVY)],
 [L("WHAT IT IS",18,True,TEAL),L("A structurally honest mental model — 7 parallels supported, 6 broken, 4 left silent but surmisable.",17,True,NAVY)],sp=0.5,fa=REDT,fb=TINT)

# 12 SYNTHESIS — Two Books (text)
s=slide(prs); title(s,"Synthesis & discussion — the Two Books")
two(s,[L("THE TWO BOOKS",18,True,NAVY),L("The genome is a page of the Book of Creation; istinsākh (45:29) is from the Book of Scripture — both called āyāt. Read side by side (never collapsed), each lends the other a grammar of consequence: copy now, express later, lose nothing, bias toward the good.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Where do the Two Books agree, and break?  • Is repentance an edit or a new transcript (25:70 / 13:39)?  • Why a one-way arrow toward the good that biology mirrors only as repair?  • What does reading creation as an āyah ask of a scientist?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 13 APPENDIX
s=slide(prs); title(s,"Appendix — the verses (Book6, vocalized col 11)")
def pk(refs): return [(e["ref"],e["snip"],e["tag"]) for e in GN if e["ref"] in refs]
appendix(s,[("TRANSCRIPTION & fidelity",TEAL,TINT,pk({"45:29","78:29","18:49"})),
            ("TRANSLATION / weighing",AMBER,AMBERT,pk({"17:14","21:47","99:7"})),
            ("EDITABLE & directional",NAVY,TINT2,pk({"25:70","13:39","11:114"}))])
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("A vivid, accountable mental model for how a life is recorded and its consequence deferred — and a guard against “scientific-miracle” overreach.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Deeds are transcribed now, expressed later, nothing lost — a labelled analogy, audited, never offered as proof.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"10_Istinsakh_Genomics_Lecture.pptx")
nvis=9  # diagram/chart-driven slides
print(f"slides: {len(prs.slides)} ; visual diagram slides: {nvis}")
