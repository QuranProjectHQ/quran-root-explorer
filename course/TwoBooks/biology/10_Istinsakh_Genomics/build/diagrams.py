# -*- coding: utf-8 -*-
"""Editable pptx diagram primitives (boxes, arrows, dashed links, scales, DNA ladder)."""
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from st_slides import sf, NAVY, TEAL, RED, AMBER, GREY, INK, WHITE, TINT, TINT2, AMBERT, REDT, LTEAL
def fbox(s,x,y,w,h,fill,title,sub="",line=NAVY,tsz=14,ssz=11):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.color.rgb=line; sp.line.width=Pt(1.5)
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_top=Inches(0.04); tf.margin_bottom=Inches(0.04); tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=title; r.font.size=Pt(tsz); r.font.bold=True; r.font.color.rgb=line; sf(r,"Calibri")
    if sub:
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r=p2.add_run(); r.text=sub; r.font.size=Pt(ssz); r.font.color.rgb=INK; sf(r,"Calibri")
    return sp
def harrow(s,x,y,w,label="",color=GREY,lcol=TEAL,h=0.3):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb=color; a.line.fill.background()
    if label:
        tb=s.shapes.add_textbox(Inches(x-0.25),Inches(y-0.42),Inches(w+0.5),Inches(0.4)); tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=label; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=lcol; sf(r,"Calibri")
    return a
def vdash(s,x,y1,y2,tag="",col=AMBER):
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x),Inches(y1),Inches(x),Inches(y2))
    cn.line.color.rgb=col; cn.line.width=Pt(1.5)
    ln=cn.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    if tag:
        tb=s.shapes.add_textbox(Inches(x-0.55),Inches((y1+y2)/2-0.16),Inches(1.1),Inches(0.32)); tf=tb.text_frame
        tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=tag; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=col; sf(r,"Calibri")
        # white halo behind tag
    return cn
def band(s,x,y,w,h,fill,label,lcol):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background()
    tb=s.shapes.add_textbox(Inches(x+0.1),Inches(y+0.03),Inches(w-0.2),Inches(0.34)); tf=tb.text_frame
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=lcol; sf(r,"Calibri")
    return sp
def dna(s,x,y,w,h,col=TEAL):
    # simple double-ladder: two verticals + rungs
    for off in (0,w):
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x+off),Inches(y),Inches(x+off),Inches(y+h)); ln.line.color.rgb=col; ln.line.width=Pt(2.5)
    n=6
    for i in range(n):
        yy=y+0.12+i*(h-0.24)/(n-1)
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x),Inches(yy),Inches(x+w),Inches(yy)); ln.line.color.rgb=(AMBER if i%2 else NAVY); ln.line.width=Pt(2)
def scales(s,x,y,w,col=NAVY):
    # fulcrum triangle + beam + two pans
    beam=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(0.08)); beam.fill.solid(); beam.fill.fore_color.rgb=col; beam.line.fill.background()
    post=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+w/2-0.03),Inches(y),Inches(0.06),Inches(0.7)); post.fill.solid(); post.fill.fore_color.rgb=col; post.line.fill.background()
    for px,c in [(x-0.05,TEAL),(x+w-0.55,AMBER)]:
        pan=s.shapes.add_shape(MSO_SHAPE.TRAPEZOID,Inches(px),Inches(y+0.5),Inches(0.6),Inches(0.3)); pan.rotation=180; pan.fill.solid(); pan.fill.fore_color.rgb=c; pan.line.fill.background()
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(px+0.3),Inches(y+0.06),Inches(px+0.3),Inches(y+0.5)); ln.line.color.rgb=col; ln.line.width=Pt(1.5)
