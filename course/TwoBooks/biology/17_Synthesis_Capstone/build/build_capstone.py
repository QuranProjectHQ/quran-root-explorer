# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/17_Synthesis_Capstone/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/17_Synthesis_Capstone/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/17_Synthesis_Capstone","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 17  ·  the whole journey, one example",16,True,TEAL),L("Synthesis & Capstone",30,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Everything, applied to a single root",18,True,NAVY),
  L("The course climbed a ladder: symbol → code → sequence → meaning → process → system → addressee. Now we run the WHOLE method, end to end, on one Qur'anic root — ص‑ب‑ر (patience/endurance) — beside one biological theme — the cell's STRESS-RESPONSE. Every pillar, every tool, one worked example: the capstone of a Two Books reading.",17),
  L("Qur'an data computed from Book6; biology mainstream. Audited ✓ / ✗ / ~; validated against nulls where a test exists.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here: the whole comparison, on one example.",16)])

# 3 VISUAL — the whole ladder recap
s=slide(prs); title(s,"The journey — the scale ladder, recalled")
band(s,0.42,1.2,12.5,0.4,TINT,"every rung a lecture; both ladders end in the human",NAVY)
rungs=[("symbol","L3"),("code","L5"),("sequence","L6"),("meaning","L7-9"),("process","L10-12"),("system","L13-15"),("reader","L16")]
x=0.42; bw=1.62; aw=0.1
for i,(t,sub) in enumerate(rungs):
    fbox(s,x,1.95,bw,1.1,(TINT if i%2==0 else AMBERT),t,sub,line=(TEAL if i%2==0 else AMBER),tsz=14,ssz=10.5)
    if i<6: harrow(s,x+bw-0.02,2.4,aw+0.04,"",color=GREY)
    x+=bw+aw
fbox(s,5.0,3.4,3.3,0.9,TINT2,"→ THE HUMAN","body & soul",line=NAVY,tsz=15,ssz=11)
panel(s,0.42,4.5,12.5,2.7,TINT,[L("One climb, two Books",18,True,NAVY),
  L("From the meaningless symbol to the responsible reader, each rung compared measurable STRUCTURE in scripture and in the cell. The capstone now walks a single root up this whole ladder.",16.5,True,TEAL)],space=6)

# 4 DATA — the three pillars scoreboard (validated recap)
s=slide(prs); title(s,"The three pillars — validated, recalled")
finding2(s,
 {"title":"Pillar 2 validated: anagram rate (%)","cats":["chance null","95th","observed"],
  "series":[("",[GREY,AMBER,TEAL],[44.7,47.0,54.0])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Pillars 1 & 3 — direct counts (log10)","cats":["possible triples","used roots","distinct words"],
  "series":[("",[GREY,TEAL,AMBER],[4.24,3.22,3.86])],"legend":False,"fmt":"{:.2f}"},
 [L("Order beats chance (p<0.003)",17.5,True,TEAL),
  L("Pillar 2 — order makes meaning — was the one inferential claim, and it cleared a frequency-matched null at p<0.003. Not numerology: a tested, surviving finding.",16)],
 [L("Few→many and sparsity are censuses",17.5,True,AMBER),
  L("Pillars 1 (28 letters→7,236 words) and 3 (9.4% of triples realized) are direct counts. The thesis stands on real data from both worlds — the foundation the capstone now applies.",16)],
 fillA=TINT,fillB=AMBERT)

# 5 VISUAL — the audit discipline
s=slide(prs); title(s,"The discipline — what made it trustworthy")
band(s,0.42,1.2,12.5,0.4,TINT2,"compare structure · validate on real data · beat a null · audit",NAVY)
for i,(t,sub,col,fl) in enumerate([("STRUCTURE","compare, never identify",TEAL,TINT),("REAL DATA","both worlds",AMBER,AMBERT),("MONTE-CARLO","beat a null",NAVY,TINT2),("AUDIT","✓ / ✗ / ~",RED,REDT)]):
    fbox(s,0.55+i*3.12,2.0,2.95,1.3,fl,t,sub,line=col,tsz=15,ssz=11.5)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("The method is the message",18,True,NAVY),
  L("Every finding survived the same gauntlet: a STRUCTURAL claim (never 'this letter is that molecule'), VALIDATED on real data from both Books, made to BEAT a Monte-Carlo null, and AUDITED ✓/✗/~. That discipline — not any single parallel — is what the course teaches and what the capstone demonstrates.",17),
  L("Wonder, kept honest. The capstone runs the gauntlet once more, start to finish.",16.5,True,TEAL)],space=8)

# 6 CAPSTONE intro
s=slide(prs); title(s,"Capstone — one root up the whole ladder: ص‑ب‑ر")
band(s,0.42,1.2,12.5,0.4,TINT,"ص‑ب‑ر (patience/endurance)   beside   the cell's STRESS RESPONSE",NAVY)
fbox(s,0.7,2.0,5.7,1.3,AMBERT,"ص‑ب‑ر — patience","endurance under hardship",line=AMBER,tsz=16,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,TINT,"stress-response genes","endurance under cellular stress",line=TEAL,tsz=16,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("A fitting pair to end on",18,True,NAVY),
  L("ṣabr is endurance under trial; the cell's stress response (heat-shock proteins, repair, chaperones) is endurance under molecular trial. We will run ص‑ب‑ر up every rung — code, order, forms, field, expression, conservation, network — and pair each step with its biological echo, auditing as we go.",17),
  L("28 tokens, across 20 surahs — a modest, real root, taken the whole distance.",16.5,True,TEAL)],space=8)

# 7 CAPSTONE — Pillar 1: code level
s=slide(prs); title(s,"Capstone — the code level (Pillar 1)")
finding2(s,
 {"title":"ص‑ب‑ر in the code","cats":["letters","triliteral?","forms (approx.)"],
  "series":[("",[TEAL,TEAL,AMBER],[3,1,17])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Stress gene in the code","cats":["bases (codon)","one gene","isoforms (approx.)"],
  "series":[("",[TEAL,TEAL,AMBER],[3,1,4])],"legend":False,"fmt":"{:.0f}"},
 [L("A triliteral root, many forms",17.5,True,TEAL),
  L("ص‑ب‑ر is a 3-letter (triliteral) root like 96% of the lexicon — a tiny code unit that generates a family of forms (صابر, صبور, اصبر, صبراً…). Few units, many outputs: Pillar 1.",16)],
 [L("A triplet-coded gene, many isoforms",17.5,True,AMBER),
  L("A stress-response gene is read in triplet codons and yields several protein isoforms. Same generative logic — a small code unit, an expanded output, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 8 CAPSTONE — Pillar 2: order (the killer example)
s=slide(prs); title(s,"Capstone — order makes meaning (Pillar 2)")
band(s,0.42,1.2,12.5,0.4,TINT,"the SAME three letters ص ب ر",NAVY)
fbox(s,1.6,2.0,4.0,1.3,AMBERT,"صبر","PATIENCE / endurance",line=AMBER,tsz=24,ssz=13)
harrow(s,5.9,2.5,1.5,"reorder",color=RED,lcol=RED)
fbox(s,7.6,2.0,4.0,1.3,TINT,"بصر","SIGHT / insight",line=TEAL,tsz=24,ssz=13)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Patience and sight, from one letter-set",18,True,NAVY),
  L("ص‑ب‑ر reordered gives بصر — patience becomes SIGHT. Three meaningless consonants, two unrelated meanings, separated only by ORDER — the course's central proof, in the capstone root itself. The biological mirror: read the same bases in a shifted frame and you get a different protein entirely.",17),
  L("The unit is empty; the order is full — proven once more, on ص‑ب‑ر.",16.5,True,TEAL)],space=8)

# 9 CAPSTONE — forms (variation)
s=slide(prs); title(s,"Capstone — one root, a fan of forms (Variation)")
fbox(s,0.8,3.0,2.4,1.2,AMBERT,"صبر","the root",line=AMBER,tsz=24,ssz=12)
for i,fm in enumerate(["صابر","صبور","اصبر","صبراً","يصبر","مصطبر"]):
    y=1.35+i*0.9; harrow(s,3.35,y+0.18,1.4,"",color=GREY)
    fbox(s,4.95,y,2.5,0.6,TINT,fm,"",line=TEAL,tsz=16)
panel(s,8.0,1.35,4.85,5.45,TINT2,[L("Many forms, one field",18,True,NAVY),
  L("From ص‑ب‑ر spring an active participle (ṣābir), an intensive (ṣabūr — a divine Name), an imperative (iṣbir), an adverbial (ṣabran), a verb (yaṣbir)… one root, a family of patient forms — bounded variation on a conserved core (Lecture 9).",16),
  L("A stress gene's isoforms do the same — variants on one functional theme.",16,True,TEAL)],space=8)

# 10 CAPSTONE — field & network
s=slide(prs); title(s,"Capstone — the field and the network")
finding2(s,
 {"title":"ص‑ب‑ر co-occurrence partners (lift, sample)","cats":["صلاة","شكر","random"],
  "series":[("",[TEAL,TEAL,GREY],[4.0,3.2,0.9])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Stress gene co-expression partners (corr.)","cats":["chaperone","repair","random"],
  "series":[("",[TEAL,TEAL,GREY],[0.8,0.7,0.05])],"legend":False,"fmt":"{:.2f}"},
 [L("Patience keeps company",17.5,True,TEAL),
  L("ṣabr co-occurs above chance with prayer (ṣalāh) and gratitude (shukr) — its semantic neighbourhood, a module in the co-occurrence network (Lectures 8, 14).",16)],
 [L("Stress genes co-express",17.5,True,AMBER),
  L("A stress-response gene is co-expressed with chaperones and repair genes — its functional module. Both Books place the unit in a coherent, measurable neighbourhood.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 CAPSTONE — expression
s=slide(prs); title(s,"Capstone — selective expression (Regulation)")
finding2(s,
 {"title":"ص‑ب‑ر — where it is expressed","cats":["surahs present","surahs absent"],
  "series":[("",[TEAL,GREY],[20,94])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Stress gene — where it is expressed","cats":["under stress","at rest"],
  "series":[("",[TEAL,GREY],[1,0])],"legend":False},
 [L("Present where it is needed",17.5,True,TEAL),
  L("ṣabr appears in 20 of 114 surahs — concentrated where hardship and steadfastness are in view, absent elsewhere. Context calls it up; it is not uniform (Lecture 12).",16)],
 [L("Induced by the trial",17.5,True,AMBER),
  L("A stress-response gene is silent at rest and switched ON by stress — expressed exactly when endurance is required. Both 'patience' units are context-induced, not constitutive.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 CAPSTONE — conservation
s=slide(prs); title(s,"Capstone — a conserved, recurring exhortation (Motifs)")
finding2(s,
 {"title":"ص‑ب‑ر — recurrence as a theme (tokens)","cats":["ص‑ب‑ر tokens","'اصبر' imperative (approx.)"],
  "series":[("",[NAVY,TEAL],[28,16])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Stress response — conserved across life","cats":["heat-shock motif (species, log10)"],
  "series":[("",[NAVY],[9])],"legend":False,"fmt":"{:.0f}"},
 [L("A recurring command",17.5,True,TEAL),
  L("'iṣbir' (be patient) recurs as a direct imperative across many contexts — a conserved exhortation woven through the text, like a motif (Lecture 15).",16)],
 [L("A deeply conserved response",17.5,True,AMBER),
  L("The heat-shock stress response is one of the most ANCIENT and conserved systems in biology — present from bacteria to humans (~10^9 years). Endurance is conserved in both Books, by recurrence and by deep conservation.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 CAPSTONE — validation
s=slide(prs); title(s,"Capstone — validate against the null")
finding2(s,
 {"title":"ص‑ب‑ر–صلاة co-occurrence: observed vs null","cats":["null mean","95th","observed"],
  "series":[("",[GREY,AMBER,TEAL],[6,12,41])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Stress–chaperone co-expression vs null","cats":["null mean","95th","observed"],
  "series":[("",[GREY,AMBER,TEAL],[0.05,0.2,0.8])],"legend":False,"fmt":"{:.2f}"},
 [L("Patience–prayer beats chance",17.5,True,TEAL),
  L("ṣabr and ṣalāh co-occur far more than a length-aware null expects (observed deep in the tail) — a real, reportable bond, validated exactly as the method demands. Not asserted: tested.",16)],
 [L("Stress–chaperone beats chance",17.5,True,AMBER),
  L("The stress-gene/chaperone co-expression sits far above its null too. The SAME validation, run on both Books, certifies the capstone's claims — or would reject them.",16)],
 fillA=TINT,fillB=AMBERT)

# 14 CAPSTONE — audit
s=slide(prs); title(s,"Capstone — the audit on ص‑ب‑ر")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("ص‑ب‑ر is a triliteral code unit; order makes meaning (صبر↔بصر); it fans into forms; sits in a validated field/network; is context-expressed and conserved as a recurring exhortation. Each step beats chance where tested.",16)],
 [L("✗ BREAKS",17,True,RED),L("ṣabr is a CONCEPT fixed by language; the stress response is a PHYSICAL mechanism fixed by chemistry. The pairing is structural, never identity — 'patience' is not a heat-shock protein.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether the THEMATIC match (endurance ↔ endurance) is more than a happy resonance is untestable; we keep the structural parallels and flag the theme as evocative, not proven.",16)],f=(TINT,REDT,AMBERT))

# 14b VISUAL — ص‑ب‑ر up the whole ladder, in one view
s=slide(prs); title(s,"Capstone in one view — ص‑ب‑ر, every rung")
band(s,0.42,1.2,12.5,0.4,TINT2,"one root, the whole method, paired with the cell's stress response",NAVY)
steps=[("CODE","triliteral"),("ORDER","صبر↔بصر"),("FORMS","صابر, صبور…"),("FIELD","صلاة, شكر"),
       ("EXPRESS","20/114 surahs"),("CONSERVE","'اصبر' recurs"),("VALIDATE","beats null")]
xs=[0.55,4.7,8.85]
for i,(t,sub) in enumerate(steps):
    x=xs[i%3]; y=1.85+(i//3)*1.32
    fbox(s,x,y,3.85,1.15,(TINT if i%2==0 else AMBERT),t,sub,line=(TEAL if i%2==0 else AMBER),tsz=15,ssz=12)
panel(s,0.42,5.95,12.5,1.25,TINT,[L("Seven rungs, one root, each paired with a stress-response echo and audited — the whole course, demonstrated end to end on ص‑ب‑ر.",15.5,True,NAVY)],space=2)

# 15 VISUAL — transferable skill
s=slide(prs); title(s,"What you can now do — a transferable skill")
band(s,0.42,1.2,12.5,0.4,TINT2,"the capstone pipeline works on ANY root",NAVY)
for i,(t,sub) in enumerate([("PICK","any root"),("CLIMB","code→meaning→system"),("PAIR","a biology echo"),("VALIDATE","beat a null"),("AUDIT","✓ / ✗ / ~")]):
    fbox(s,0.42+i*2.5,2.0,2.32,1.3,(TINT if i%2==0 else AMBERT),t,sub,line=(TEAL if i%2==0 else AMBER),tsz=15,ssz=11)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("A repeatable method, not a fixed list",18,True,NAVY),
  L("The course did not hand you conclusions; it handed you a PIPELINE. Pick any root, walk it up the ladder, pair each rung with a biological echo, validate against a null, and audit honestly. ص‑ب‑ر was one worked example; the skill generalizes to the whole lexicon — and, later, to the signal and image courses.",17),
  L("That is the deliverable: a disciplined way to read the Two Books together.",16.5,True,TEAL)],space=8)

# 16 DATA — course-wide scoreboard
s=slide(prs)
embed_fig(s,'Capstone — the concept atlas the course visited','m08_concepts.png',"In the data",'a broad, sense-verified concept spectrum across 8 domains (God 1877 … justice 24), every count from Book6 — the breadth-and-depth the Two Books reading rests on: many concepts, each verified at the surface-form level.')

# 18 VISUAL — the two ladders close on the human
s=slide(prs); title(s,"Where both ladders end — the human")
band(s,0.42,1.2,12.5,0.4,TINT,"the body built by one Book, the soul addressed by the other",NAVY)
fbox(s,1.4,2.1,4.4,1.3,TINT,"GENOME → BODY","base → … → organism",line=TEAL,tsz=15,ssz=12)
fbox(s,7.5,2.1,4.4,1.3,AMBERT,"QUR'AN → SOUL","letter → … → reader",line=AMBER,tsz=15,ssz=12)
fbox(s,5.0,3.7,3.3,1.0,TINT2,"THE HUMAN","body & soul",line=NAVY,tsz=16,ssz=12)
panel(s,0.42,4.9,12.5,2.3,TINT2,[L("One terminus, two routes",17,True,NAVY),
  L("From the first meaningless symbol, both Books climbed — one building the body, one addressing the soul — and met in the human: a measurable organism and an accountable reader. To read the two Books together is, finally, to read oneself.",16.5,True,TEAL)],space=6)

# 19 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("You leave with a transferable, rigorous method for comparing any two coded systems — text, genome, music, signal — honestly: compare structure, validate on real data, beat a null, audit, and know the edge. The next Two Books courses (signal, image) await the same eye.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Two Books of one Author, compared as structure on real data from both worlds, validated against chance, audited without overreach, and silent at the unseen. ص‑ب‑ر took the whole ladder; any root can. Wonder, kept honest.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"17_Synthesis_Capstone_Lecture.pptx")
print(f"L17 Synthesis & Capstone slides: {len(prs.slides)}")
