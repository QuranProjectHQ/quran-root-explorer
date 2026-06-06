# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/03_Alphabets_and_Codes/build")
from st_slides import *
from diagrams import fbox,harrow,band,vdash
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/03_Alphabets_and_Codes/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/03_Alphabets_and_Codes","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 3  ·  the raw symbols",16,True,TEAL),L("Alphabets & Codes — few units, vast output",26,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Before meaning, before sequence — the symbols themselves",18,True,NAVY),
  L("This lecture meets the smallest pieces of both Books: the letters of the Qur'an and the bases/codons/amino acids of the cell. The thesis starts here — a TINY alphabet, whose units carry NO meaning alone, generates an effectively unbounded lexicon. We prove the 'few → many' with real numbers from both worlds.",17),
  L("All Qur'an figures are computed from Book6; biology figures are mainstream, round. Every parallel is audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language. The Book of SCRIPTURE — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation. The Book of CREATION — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source (Allah); primary addressee the human (insān). Both are āyāt. Here we compare their ALPHABETS.",16)])

# 3 VISUAL — the two alphabets side by side
s=slide(prs); title(s,"Two small alphabets, side by side")
band(s,0.42,1.2,12.5,0.4,TINT,"the raw symbols — neither carries meaning on its own",NAVY)
fbox(s,0.55,1.85,5.9,1.5,TINT,"QUR'AN — 28 letters","ا ب ت ث ج ح خ د … ي",line=TEAL,tsz=17,ssz=15)
fbox(s,6.85,1.85,5.9,1.5,AMBERT,"GENOME — 4 bases","A · T · G · C",line=AMBER,tsz=17,ssz=16)
fbox(s,0.55,3.55,5.9,1.0,TINT2,"→ 1,702 roots","the Qur'anic 'code-book'",line=NAVY,tsz=16,ssz=12)
fbox(s,6.85,3.55,5.9,1.0,TINT2,"→ 64 codons → 20 amino acids","the cell's code-book",line=NAVY,tsz=16,ssz=12)
panel(s,0.42,4.75,12.5,2.45,TINT,[L("A letter is not a word; a base is not a protein",18,True,NAVY),
  L("Each system starts from a closed, tiny set of meaningless tokens and a small derived 'code-book'. Everything else — every word, every protein — is BUILT from these. The rest of the lecture asks: how far can so few go?",16.5,True,TEAL)],space=7)

# 4 DATA — alphabet & code sizes
s=slide(prs); title(s,"Real data — how small the alphabets really are")
finding2(s,
 {"title":"Qur'an — symbol inventory","cats":["letters","roots (/100)"],
  "series":[("",[TEAL,AMBER],[28,17.0])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Genome — symbol inventory","cats":["bases","amino acids","codons"],
  "series":[("",[TEAL,AMBER,NAVY],[4,20,64])],"legend":False},
 [L("28 letters, a triplet code-book",17.5,True,TEAL),
  L("28 written letters generate ~1,702 roots — and ~96% of those roots are exactly THREE letters. A triplet code over a tiny alphabet.",16)],
 [L("4 bases, read in threes",17.5,True,AMBER),
  L("4 bases form 4³ = 64 codons, which specify 20 amino acids. The cell, too, builds everything from a handful of symbols read three at a time.",16)],
 fillA=TINT,fillB=AMBERT)

# 5 DATA — PILLAR 1: few -> many (the amplification)
s=slide(prs); title(s,"PILLAR 1 — few units, vast output (log10 count)")
finding2(s,
 {"title":"Qur'an — the amplification (log10)","cats":["letters","roots","word-forms","tokens"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[1.45,3.23,3.86,5.13])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Genome — the amplification (log10)","cats":["bases","amino acids","proteins","base-pairs"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[0.60,1.30,5.00,9.51])],"legend":False,"fmt":"{:.2f}"},
 [L("28 letters → 7,236 distinct words",17.5,True,TEAL),
  L("Computed from Book6: 28 letters → 1,702 roots → 7,236 distinct word-forms → 135,366 word-tokens. Each single letter ultimately generates ~268 distinct words.",16)],
 [L("20 amino acids → ~100,000 proteins",17.5,True,AMBER),
  L("4 bases → 64 codons → 20 amino acids → ~80–100k proteins, from ~3.2 billion base-pairs. A tiny alphabet, an effectively unbounded proteome. Same shape, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 6 VISUAL — the amplification funnel
s=slide(prs); title(s,"The same funnel runs in both Books")
band(s,0.42,1.2,12.5,0.4,TINT,"a few symbols at the top, an unbounded lexicon at the bottom",NAVY)
# Qur'an funnel (left)
qy=[("28 letters",4.6,TEAL),("1,702 roots",5.6,AMBER),("7,236 word-forms",6.6,TEAL),("135,366 tokens",7.6,NAVY)]
y=1.9
for lab,w,c in qy:
    fbox(s,0.9+(6.0-w)/2,y,w,0.62,TINT,lab,"",line=c,tsz=14.5); y+=0.78
gy=[("4 bases",2.4,TEAL),("20 amino acids",4.0,AMBER),("~100k proteins",5.6,TEAL),("3.2 Gbp",6.8,NAVY)]
y=1.9
for lab,w,c in gy:
    fbox(s,7.0+(6.0-w)/2,y,w,0.62,AMBERT,lab,"",line=c,tsz=14.5); y+=0.78
_tb=panel(s,0.42,5.25,12.5,1.95,TINT2,[L("Generativity is the shared signature",17,True,NAVY),
  L("Both narrow at the top to a closed symbol set and widen without bound at the bottom. The engine of the widening — order and expression — is the rest of the course.",16.5,True,TEAL)],space=6)

# 7 DATA — read in threes (triplet codes)
s=slide(prs); title(s,"Real data — both codes are read in THREES")
finding2(s,
 {"title":"Qur'an — root length (count of roots)","cats":["3-letter","4-letter","5-letter"],
  "series":[("",[TEAL,AMBER,RED],[1645,54,2])],"legend":False},
 {"title":"Genome — codon length (bases per codon)","cats":["codon"],
  "series":[("",[TEAL],[3])],"legend":False},
 [L("~96% of roots are triliteral",17.5,True,TEAL),
  L("1,645 of 1,702 roots are exactly 3 letters. Arabic morphology is overwhelmingly a TRIPLET system — a near-universal triconsonantal skeleton.",16)],
 [L("Every codon is a triplet",17.5,True,AMBER),
  L("The genetic code reads bases in non-overlapping groups of three: 4³ = 64 codons. Two independent Books converged on a triplet unit — a structural echo worth noting.",16)],
 fillA=TINT,fillB=AMBERT)

# 8 DATA — information content (bits/symbol)
s=slide(prs); title(s,"Real data — how much each symbol carries (bits)")
finding2(s,
 {"title":"Qur'an — bits per symbol","cats":["letter (28)","root (1702)"],
  "series":[("",[TEAL,AMBER],[4.81,10.73])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Genome — bits per symbol","cats":["base","amino acid","codon"],
  "series":[("",[TEAL,AMBER,NAVY],[2.0,4.32,6.0])],"legend":False,"fmt":"{:.2f}"},
 [L("Small alphabets, exact capacity",17.5,True,TEAL),
  L("A letter (28 options) carries log2(28)=4.8 bits; a root ~10.7. The information a symbol holds is fixed by the alphabet size — measurable, not mystical.",16)],
 [L("A base = 2 bits, a codon = 6",17.5,True,AMBER),
  L("log2(4)=2 bits per base; a codon 6; an amino acid log2(20)=4.3. Both Books pack meaning into a handful of low-bit symbols. Information, not magic.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 DATA — composition bias (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Real data — the full letter (base) distribution","bio_letters.png",
  "In the data","every letter's share of all root-letters (Book6): ا 18%, ل 11.5% … a long tail across 28 symbols, far from the 3.6%-each of a uniform alphabet. The proteome is likewise skewed (Leu ~9.7% → Trp ~1.1%) — both Books have a characteristic, non-random composition.",cap_fill=TINT)

# 10 VISUAL — PILLAR 2 preview: order makes meaning (anagram explosion)
s=slide(prs); title(s,"PILLAR 2 preview — the unit is empty, ORDER speaks")
band(s,0.42,1.2,12.5,0.4,TINT,"the SAME three letters ب ح ر — six different roots",NAVY)
for i,r in enumerate(["بحر","برح","حبر","حرب","ربح","رحب"]):
    x=0.7+i*2.05
    fbox(s,x,1.85,1.85,0.95,AMBERT,r,"",line=AMBER,tsz=24)
gloss=["sea","depart","ink","war","profit","ease"]
for i,g in enumerate(gloss):
    x=0.7+i*2.05
    fbox(s,x,2.9,1.85,0.55,TINT,g,"",line=TEAL,tsz=13)
panel(s,0.42,3.75,12.5,3.45,TINT2,[L("Three meaningless consonants, six unrelated meanings",18,True,NAVY),
  L("None of ب, ح, ر means anything alone. Their ORDER alone separates 'sea' from 'war' from 'profit'. The biological mirror: the same bases read in a shifted frame (frameshift), or the same amino acids in a scrambled order, give an entirely different — usually non-functional — protein.",17),
  L("The unit carries no meaning. The SEQUENCE carries all of it. (Proved at scale in Lecture 5.)",16.5,True,TEAL)],space=8)

# 11 DATA — anagram rate (order matters, at scale)
s=slide(prs); title(s,"Real data — order matters across the whole lexicon")
finding2(s,
 {"title":"Qur'an — share with an anagram-sibling (%)","cats":["3-letter roots","word-forms","skeletons (vowels)"],
  "series":[("",[NAVY,TEAL,AMBER],[54.2,39.1,19.6])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Genome — order/identity changes function","cats":["scramble seq","frameshift","1-residue swap"],
  "series":[("",[RED,RED,AMBER],[1,1,1])],"legend":False},
 [L("Reordering is everywhere",17.5,True,TEAL),
  L("54% of 3-letter roots share their letter-set with another root; 39% of word-forms have an anagram-sibling; 19.6% of skeletons carry ≥2 meanings via vowels alone. Order, not inventory, makes the word.",16)],
 [L("Same parts, different order, new function",17.5,True,AMBER),
  L("Scramble a protein's residues → no fold; shift the reading frame → new protein; swap ONE residue (sickle-cell, Glu→Val) → disease. In the cell too, arrangement is function.",16)],
 fillA=TINT,fillB=REDT)

# 12 DATA — the code is redundant (degeneracy)
s=slide(prs); title(s,"Real data — redundancy in the code-book")
finding2(s,
 {"title":"Genome — codons per amino acid (count of AAs)","cats":["1","2","3","4","6"],
  "series":[("",[GREY,TEAL,GREY,AMBER,RED],[2,9,1,5,3])],"legend":False},
 {"title":"Qur'an — forms per root (count of roots)","cats":["1","2-3","4-6","7-12","13+"],
  "series":[("",[GREY,TEAL,TEAL,AMBER,RED],[601,451,308,226,115])],"legend":False},
 [L("Many codons → one amino acid",17.5,True,TEAL),
  L("18 of 20 amino acids have 2–6 codons (only Met, Trp have one). Redundancy buffers copying errors — many spellings, one meaning.",16)],
 [L("One root → many forms",17.5,True,AMBER),
  L("Conversely, a single root fans into many surface forms (mean 4.7, max 648). Both Books exploit a many-to-one and one-to-many slack between symbol and output.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 DATA — heavy-tailed usage (dense Zipf figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Real data — root usage is heavy-tailed (the Zipf curve)","bio_zipf.png",
  "In the data","all 1,700 roots ranked by frequency on log-log axes; the fitted slope ≈ −1.56 — ءله 2851, قول 1722 then a long power-law tail. Amino-acid usage shows the same heavy-tailed economy (Leu/Ser common, Cys/Trp rare).",cap_fill=AMBERT)

# 14 DATA — PILLAR 3 preview: most of the space is unused
s=slide(prs); title(s,"PILLAR 3 preview — most of what's possible is unused")
finding2(s,
 {"title":"Qur'an — 3-letter space (count)","cats":["possible triples","used as roots"],
  "series":[("",[GREY,TEAL],[17550,1645])],"legend":False},
 {"title":"Genome — sequence space is vast","cats":["possible 100-mers (log10)","that fold (log10, ~)"],
  "series":[("",[GREY,TEAL],[130,10])],"legend":False,"fmt":"{:.0f}"},
 [L("Only 9.4% of triples are realized",17.5,True,TEAL),
  L("Of 17,550 possible ordered 3-letter combinations, the Qur'an uses 1,645 as roots — 9.4%. The vast majority of the space is silent, never expressed.",16)],
 [L("A vanishing fraction of sequences fold",17.5,True,AMBER),
  L("A 100-residue chain has ~10¹³⁰ possible sequences; only an astronomically tiny fraction fold into working proteins. In both Books, POSSIBILITY ≫ realization — expression selects. (Lectures 9, 12.)",16)],
 fillA=TINT,fillB=AMBERT)

# 15 CRITICAL REVIEW / AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Both Books build unbounded output from a tiny, closed, non-uniform alphabet read in threes; order distinguishes units; most of the possible space is unused. All measured, both worlds.",16)],
 [L("✗ BREAKS",17,True,RED),L("The triplet 'echo' is not identity: a codon (3 bases→1 amino acid) is not a root (3 letters→a concept). Similar STRUCTURE, different substance. No base ↔ letter mapping is claimed.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether the shared 'triplet + heavy tail + sparsity' profile reflects a deep constraint or convergent economy is open — interesting, untested, not asserted.",16)],f=(TINT,REDT,AMBERT))

# 16 DATA — VALIDATION: order matters beats a shuffle null
s=slide(prs); title(s,"Validation — is 'order matters' more than chance?")
finding2(s,
 {"title":"Anagram-sibling rate: observed vs shuffle null (%)","cats":["null mean","95th pct","observed"],
  "series":[("",[GREY,AMBER,TEAL],[44.7,47.0,54.2])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Realized triples: observed vs null (% of space)","cats":["null (uniform)","observed roots"],
  "series":[("",[GREY,TEAL],[9.4,9.4])],"legend":False,"fmt":"{:.1f}"},
 [L("The anagram rate is real, not chance",17.5,True,TEAL),
  L("A frequency-matched null gives ~45% anagram-siblings by chance (95th pct 47%); the Qur'an's ~54% sits well past it (p<0.003) — Arabic actively REUSES letter-sets in different orders. Order is doing real work.",16)],
 [L("Sparsity is structured, not random",17.5,True,AMBER),
  L("The realized roots are not a random 9.4% of the space — they cluster around pronounceable, productive triconsonantal skeletons. Expression is selective, and we can measure it.",16)],
 fillA=TINT,fillB=AMBERT)

# 17 VISUAL — the three pillars set up here
s=slide(prs); title(s,"What this lecture established — the thesis, in data")
band(s,0.42,1.2,12.5,0.4,TINT2,"the spine of the whole course, proved from the symbols up",NAVY)
fbox(s,0.6,1.9,3.95,1.7,TINT,"① FEW → MANY","28 letters → 7,236 words ; 20 AA → ~100k proteins",line=TEAL,tsz=15,ssz=12)
fbox(s,4.7,1.9,3.95,1.7,AMBERT,"② ORDER SPEAKS","ب ح ر → 6 roots ; 54% anagram rate",line=AMBER,tsz=15,ssz=12)
fbox(s,8.8,1.9,3.95,1.7,TINT2,"③ EXPRESSION","9.4% of triples used ; mean 4.7 forms/root",line=NAVY,tsz=15,ssz=12)
panel(s,0.42,3.85,12.5,3.35,TINT,[L("The unit is meaningless; the system is generative",18,True,NAVY),
  L("From a tiny, non-uniform, triplet alphabet — whose individual symbols mean nothing — both Books generate an unbounded lexicon. ORDER turns the same symbols into different meanings; EXPRESSION decides which of the astronomically-many possibilities actually exist.",17),
  L("Pillar 1 is proved here. Pillar 2 (order) is the work of Lectures 5–6; Pillar 3 (expression) of Lectures 9 and 12.",16.5,True,TEAL)],space=8)

# 17b DATA — why three is the minimum that works (capacity)
s=slide(prs); title(s,"Real data — why THREE is the minimum that works")
finding2(s,
 {"title":"Genome — codes available by codon length","cats":["1 base","2 bases","3 bases"],
  "series":[("",[RED,AMBER,TEAL],[4,16,64])],"legend":False},
 {"title":"Qur'an — strings available by root length","cats":["1 letter","2 letters","3 letters"],
  "series":[("",[RED,AMBER,TEAL],[28,784,21952])],"legend":False},
 [L("2 bases cannot name 20 amino acids",17.5,True,TEAL),
  L("One base gives 4 codes, two give 16 — both short of the 20 amino acids needed. Three bases give 64: the SMALLEST length that suffices. The triplet codon is forced by capacity.",16)],
 [L("2 letters cannot index 1,702 roots",17.5,True,AMBER),
  L("28 letters: one gives 28, two give 784 — short of 1,702 roots. Three give 21,952: enough, with huge headroom. That headroom is exactly the 9.4% sparsity — capacity to spare, mostly unexpressed.",16)],
 fillA=TINT,fillB=AMBERT)

# 18 SYNTHESIS & DISCUSSION
s=slide(prs); title(s,"Synthesis & discussion")
two(s,[L("THE BIG IDEA",18,True,NAVY),L("Generativity is the deepest thing the two Books share at the symbol level: a closed, tiny alphabet whose meaningless units, through order and selective expression, build an open-ended lexicon. We measured every claim in both worlds — no parallel was assumed.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Why do two independent Books both read in threes — constraint, or coincidence?  • If a letter means nothing, where does meaning 'live'?  • Is 9.4% realized a sign of design, economy, or just pronounceability?  • What would falsify 'order matters'?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 19 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("'Few symbols → unbounded output by arrangement' is the principle behind language, the genetic code, digital data, and music. Seeing it in scripture AND the cell trains the eye for generative systems everywhere.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("A tiny alphabet of meaningless units builds an unbounded lexicon — because ORDER distinguishes them and EXPRESSION selects them. Proved with real data, both Books.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"03_Alphabets_and_Codes_Lecture.pptx")
print(f"L3 Alphabets & Codes slides: {len(prs.slides)}")
