# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/05_Codon_and_Root/build")
from st_slides import *
from diagrams import fbox,harrow,band,matgrid
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/05_Codon_and_Root/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/05_Codon_and_Root","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def cell(s,x,y,w,h,fill,txt,line,tsz=18): fbox(s,x,y,w,h,fill,txt,"",line=line,tsz=tsz)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 5  ·  the triplet code",16,True,TEAL),L("Codon & Root — where ORDER makes meaning",26,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Below meaning — and the proof of Pillar 2",18,True,NAVY),
  L("Lecture 3 met the alphabets; here we read three units at a time — the codon and the triliteral root — and prove the second pillar at scale: the unit is meaningless, but its ORDER carries everything. We show, with real Book6 data tested against chance, that reordering three consonants makes entirely new meanings.",17),
  L("No semantics assumed; biology figures mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language. The Book of SCRIPTURE — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation. The Book of CREATION — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source (Allah); both are āyāt. Here we read both in THREES.",16)])

# 3 VISUAL — the two triplet codes
s=slide(prs); title(s,"Two triplet codes")
band(s,0.42,1.18,12.5,0.42,TINT,"DNA — read three bases at a time",TEAL)
for i,(b,c) in enumerate([("A",TINT),("U",LTEAL),("G",TINT)]): cell(s,1.6+i*1.0,1.85,0.9,0.9,c,b,TEAL,tsz=22)
harrow(s,4.9,2.15,1.3,"codon",color=GREY,lcol=TEAL)
fbox(s,6.4,1.8,3.2,1.0,TINT,"amino acid","(one of ~20)",line=TEAL,tsz=16,ssz=12)
band(s,0.42,3.5,12.5,0.42,AMBERT,"ARABIC — read three letters at a time (the root)",AMBER)
for i,(b,c) in enumerate([("ك",AMBERT),("ت",AMBERT),("ب",AMBERT)]): cell(s,1.6+i*1.0,4.15,0.9,0.9,c,b,AMBER,tsz=24)
harrow(s,4.9,4.45,1.3,"root",color=GREY,lcol=AMBER)
fbox(s,6.4,4.1,3.2,1.0,AMBERT,"a lexical field","(one of ~1700)",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,5.5,12.5,1.7,TINT2,[L("Three bases → an amino acid; three consonants → a root. The triplet is the meaning-bearing unit in BOTH systems — and its internal ORDER is what we test next.",16.5,True,NAVY)],space=6)

# 4 DATA — root length + sparsity
s=slide(prs); title(s,"The data — a sparse triplet lexicon")
finding2(s,
 {"title":"Root length (distinct roots)","cats":["2","3","4","5"],
  "series":[("",[GREY,TEAL,AMBER,GREY],[1,1643,54,2])],"legend":False},
 {"title":"A SPARSE code (3-letter space)","cats":["used","unused"],
  "series":[("",[TEAL,GREY],[1643,20309])],"legend":False},
 [L("96% are triplets",17.5,True,TEAL),
  L("Of ~1,700 distinct roots, ~1,643 are exactly three letters. The triliteral root is the dominant unit — the linguistic codon.",16)],
 [L("Only ~7.5% of triplets are used",17.5,True,AMBER),
  L("~28 letters give ~21,950 possible 3-letter strings; the Qur'an uses ~1,643 — a sparse code, like the genetic code's 61 sense codons of 64.",16)],
 fillA=TINT,fillB=AMBERT)

# 5 PILLAR 2 — order makes the root
s=slide(prs); title(s,"PILLAR 2 — the SAME three letters, six meanings")
band(s,0.42,1.2,12.5,0.4,TINT,"ب ح ر — reorder three meaningless consonants",NAVY)
for i,r in enumerate(["بحر","برح","حبر","حرب","ربح","رحب"]):
    cell(s,0.7+i*2.05,1.85,1.85,0.95,AMBERT,r,AMBER,24)
for i,g in enumerate(["sea","depart","ink","war","profit","ease"]):
    cell(s,0.7+i*2.05,2.9,1.85,0.55,TINT,g,TEAL,13)
panel(s,0.42,3.75,12.5,3.45,TINT2,[L("The unit is empty; the order is everything",18,True,NAVY),
  L("None of ب, ح, ر means anything alone. Their ORDER alone separates 'sea' from 'war' from 'profit'. د‑ع‑و gives five roots (دعو/عدو/عود/ودع/وعد); ح‑ر‑م gives five (حرم/حمر/رحم/رمح/مرح). This is the linguistic mirror of a reading-frame shift — same bases, new order, new protein.",17),
  L("Across the lexicon this is not rare — it is the rule, as the next slide proves against chance.",16.5,True,TEAL)],space=8)

# 6 DATA — anagram rate + computed null (VALIDATION)
s=slide(prs); title(s,"Validation — order-reuse beats chance (p < 0.003)")
finding2(s,
 {"title":"Anagram-sibling rate among roots (%)","cats":["chance (null)","95th pct","observed"],
  "series":[("",[GREY,AMBER,TEAL],[44.7,47.0,54.0])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Letter-sets shared by >=2 roots","cats":["sets w/ >=2 roots","busiest set (ب ح ر)"],
  "series":[("",[TEAL,NAVY],[360,6])],"legend":False},
 [L("Reordering is the rule, not the exception",17.5,True,TEAL),
  L("About 54% of three-letter roots share their exact letter-set with another root. A frequency-matched shuffle null expects only ~45% (95th pct 47%); the observed rate clears it at p<0.003 — 0 of 300 random draws reached it.",16)],
 [L("Arabic actively reuses letter-sets",17.5,True,AMBER),
  L("360 distinct letter-sets each spawn 2-6 different roots. The language exploits ORDER to multiply meaning from the same consonants — measurably more than chance would give.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 DATA — order beats composition (both Books)
s=slide(prs); title(s,"Composition is not enough — order decides")
finding2(s,
 {"title":"Roots from one letter-set {ب,ح,ر}","cats":["compositions","realized roots"],
  "series":[("",[GREY,TEAL],[1,6])],"legend":False},
 {"title":"Protein: same residues, order set by sequence","cats":["compositions","possible orders (log10, n=100)"],
  "series":[("",[GREY,RED],[1,157.0])],"legend":False,"fmt":"{:.0f}"},
 [L("One bag of letters, many words",17.5,True,TEAL),
  L("The multiset {ب,ح,ر} is a single composition — yet it yields six distinct roots. Knowing WHICH letters is not knowing the word; you must know their ORDER.",16)],
 [L("One bag of residues, one fold",17.5,True,AMBER),
  L("A protein's amino-acid composition does not fix its structure; the SEQUENCE (order) does. A 100-residue composition admits ~10^157 orderings — only the right one folds and functions.",16)],
 fillA=TINT,fillB=AMBERT)

# 8 VISUAL — reading frames (biology side of order)
s=slide(prs); title(s,"The biology of order — reading frames")
band(s,0.42,1.2,12.5,0.4,TINT,"one base sequence, three frames, three proteins",NAVY)
seq=["G","A","U","U","A","C","G","A","U"]
for i,b in enumerate(seq): cell(s,0.8+i*1.18,1.85,1.0,0.7,TINT,b,TEAL,18)
for fr,(yy,col,lab) in enumerate([(2.75,AMBER,"frame 1"),(3.35,TEAL,"frame 2"),(3.95,NAVY,"frame 3")]):
    band(s,0.8+fr*0.39,yy,1.0*3+0.0,0.06,col,"",col)
fbox(s,0.6,4.5,3.95,1.0,AMBERT,"frame 1: GAU-UAC...","Asp-Tyr",line=AMBER,tsz=14,ssz=12)
fbox(s,4.7,4.5,3.95,1.0,TINT,"frame 2: AUU-ACG...","Ile-Thr",line=TEAL,tsz=14,ssz=12)
fbox(s,8.8,4.5,3.95,1.0,TINT2,"frame 3: UUA-CGA...","Leu-Arg",line=NAVY,tsz=14,ssz=12)
panel(s,0.42,5.7,12.5,1.5,REDT,[L("Same letters, shifted order → a different message",17,True,RED),
  L("Reading the same bases one step over yields entirely different codons and a different protein — the cell's version of 'برح vs حرب'. Order is the information.",16.5,True,NAVY)],space=6)

# 9 DATA — frameshift consequence
s=slide(prs); title(s,"When order breaks — the cost of a shift")
finding2(s,
 {"title":"Mutation type vs protein disruption (illustrative)","cats":["silent","missense","frameshift"],
  "series":[("",[TEAL,AMBER,RED],[5,40,95])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Arabic: one-letter change vs reorder","cats":["substitute 1","reorder all 3"],
  "series":[("",[AMBER,RED],[1,1])],"legend":False},
 [L("A frameshift is catastrophic",17.5,True,RED),
  L("Insert or delete one base and every downstream codon is misread — typically destroying the protein. Order is not a detail; it is the whole signal. (Percentages illustrative of relative disruption.)",16)],
 [L("Arabic mirrors both moves",17.5,True,AMBER),
  L("Substituting one consonant gives a new root (ك‑ت‑ب → ك‑ذ‑ب); reordering all three gives another (حرب ↔ برح). Both Books are exquisitely sensitive to the identity AND order of three units.",16)],
 fillA=REDT,fillB=AMBERT)

# 10 DATA — alphabet bias + codon degeneracy
s=slide(prs); title(s,"An uneven alphabet, a redundant code")
finding2(s,
 {"title":"Consonant use in roots (the 'bases')","cats":["ل","ر","و","ء","ب","ن"],
  "series":[("",[NAVY,TEAL,TEAL,AMBER,AMBER,GREY],[15412,12651,12092,9966,9888,9278])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Codon degeneracy (# amino acids by codons-each)","cats":["1","2","3","4","6"],
  "series":[("",[RED,TEAL,GREY,AMBER,NAVY],[2,9,1,5,3])],"legend":False},
 [L("The alphabet is used unevenly",17.5,True,TEAL),
  L("Some consonants dominate the roots (ل, ر, و lead) — a composition bias, like a genome's GC-content or codon-usage bias. Units are not equiprobable.",16)],
 [L("The code is redundant",17.5,True,AMBER),
  L("61 sense codons → 20 amino acids: most amino acids have 2-4 codons, three have 6, two just one. Many spellings, one product — degeneracy.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 DATA — positional consonant bias + wobble
s=slide(prs)
embed_fig(s,'Module — root architecture: consonant use by position','m05_position.png',"In the data",'across all 3-letter roots, consonants cluster by POSITION (ر favours position 3, ي is rare in position 1). Root structure is positionally biased — the script analogue of codon position-bias in DNA.')

# 12 DATA — homographs: vowel order
s=slide(prs); title(s,"Even the vowels are order — homographs")
finding2(s,
 {"title":"Consonant-skeletons by # of meanings","cats":["one meaning","two or more"],
  "series":[("",[GREY,AMBER],[80.4,19.6])],"legend":False,"fmt":"{:.1f}"},
 {"title":"مُخلِص vs مُخلَص — same letters","cats":["مخلِص (active)","مخلَص (passive)"],
  "series":[("",[TEAL,AMBER],[1,1])],"legend":False},
 [L("Hidden vowels are sequence too",17.5,True,TEAL),
  L("19.6% of consonant-skeletons carry two or more meanings, separated only by their (unwritten) vowel pattern — a finer layer of ORDER riding on the same consonants.",16)],
 [L("One diacritic flips the role",17.5,True,AMBER),
  L("مُخلِص (one who is sincere, active) vs مُخلَص (one made sincere, passive): identical consonants, the vowel order flips the meaning — like a post-transcriptional edit changing a protein's sense.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 VISUAL — the dictionary grid
s=slide(prs); title(s,"A finite table mapping triplets → fields")
band(s,0.42,1.18,12.5,0.42,TINT,"the 'codon table' idea — a finite lookup with redundancy",TEAL)
matgrid(s,3.0,1.85,7.2,1.9,4,12,[TINT,LTEAL,AMBERT,TEAL])
panel(s,0.42,4.05,12.5,3.15,TINT2,[L("Lookup, and degeneracy",18,True,NAVY),
  L("The genetic code is a fixed table: 64 codons → 20 amino acids, so several codons map to one amino acid. Arabic echoes the shape: ~1,700 root-triplets index the lexicon, and a root's many surface forms cluster around one field — many tokens, one root.",17),
  L("Both are FINITE lookup systems over a triplet alphabet — a closed code generating open output.",16.5,True,TEAL)],space=8)

# 14 VISUAL — point change
s=slide(prs); title(s,"Change one character — a different unit (a point 'mutation')")
cell(s,1.2,2.0,0.85,0.85,AMBERT,"ك",AMBER,22); cell(s,2.1,2.0,0.85,0.85,AMBERT,"ت",AMBER,22); cell(s,3.0,2.0,0.85,0.85,AMBERT,"ب",AMBER,22)
harrow(s,4.2,2.3,1.3,"one letter",color=RED,lcol=RED)
cell(s,5.8,2.0,0.85,0.85,REDT,"ك",RED,22); cell(s,6.7,2.0,0.85,0.85,REDT,"ذ",RED,22); cell(s,7.6,2.0,0.85,0.85,REDT,"ب",RED,22)
panel(s,0.42,3.4,12.5,3.8,REDT,[L("A single substitution = a new root",18,True,RED),
  L("كتب (write) → كذب (lie): replace one consonant and the triplet indexes an entirely different field — structurally identical to a point mutation that swaps one base and changes the codon's product.",17),
  L("At the character level (no semantics), both systems are exquisitely position-sensitive: order and identity of three units decide the unit produced.",16.5,True,NAVY)],space=9)

# 15 DATA — how many orderings are realized
s=slide(prs); title(s,"How much order is exploited")
finding2(s,
 {"title":"Roots per shared letter-set","cats":["2","3","4","5","6"],
  "series":[("",[TEAL,TEAL,AMBER,AMBER,RED],[250,70,28,8,4])],"legend":False},
 {"title":"Orderings of 3 distinct letters","cats":["possible (3!)","used by ب ح ر"],
  "series":[("",[GREY,TEAL],[6,6])],"legend":False},
 [L("Most shared sets yield 2-3 roots",17.5,True,TEAL),
  L("Of the 360 letter-sets that recur, most produce 2-3 distinct roots; a handful produce 4-6. (Counts approximate.) Arabic samples the orderings of its consonants, not just the letters.",16)],
 [L("Some sets use every ordering",17.5,True,AMBER),
  L("Three distinct letters have 3! = 6 orderings; {ب,ح,ر} realizes all six as real roots. The order-space is not decorative — it is filled with meaning.",16)],
 fillA=TINT,fillB=AMBERT)

# 16 VISUAL — sequence is the information
s=slide(prs); title(s,"The lesson of the triplet — sequence is the information")
band(s,0.42,1.2,12.5,0.4,TINT2,"the unit is empty · the order is full",NAVY)
fbox(s,0.7,1.95,3.85,1.6,TINT,"THE UNIT","a base / a letter — no meaning alone",line=TEAL,tsz=16,ssz=12)
harrow(s,4.7,2.6,1.1,"ordered",color=GREY,lcol=NAVY)
fbox(s,6.0,1.95,3.0,1.6,AMBERT,"THE TRIPLET","codon / root — a meaning",line=AMBER,tsz=16,ssz=12)
harrow(s,9.1,2.6,1.1,"",color=GREY)
fbox(s,10.4,1.95,2.35,1.6,TINT2,"OUTPUT","amino acid / field",line=NAVY,tsz=16,ssz=11)
panel(s,0.42,3.85,12.5,3.35,TINT,[L("Pillar 2, proven both ways",18,True,NAVY),
  L("Reorder three meaningless units and the meaning changes — in Arabic (54% of roots, p<0.003 vs chance) and in the cell (reading frames, frameshift). Composition tells you the parts; only SEQUENCE tells you the message.",17),
  L("This is why a letter cannot 'be' a molecule (Lecture 4) and yet the two codes share a deep grammar: information lives in the arrangement.",16.5,True,TEAL)],space=8)

# 16b DATA — the two code tables by the numbers
s=slide(prs); title(s,"Two finite tables, two open outputs")
finding2(s,
 {"title":"Genetic code table","cats":["codons","amino acids","stop signals"],
  "series":[("",[NAVY,TEAL,RED],[64,20,3])],"legend":False},
 {"title":"Arabic root system (log10)","cats":["roots","word-forms","tokens"],
  "series":[("",[TEAL,AMBER,NAVY],[3.23,3.86,5.13])],"legend":False,"fmt":"{:.2f}"},
 [L("64 codons, a closed table",17.5,True,TEAL),
  L("The genetic code is a fixed 64-entry lookup: 61 sense codons → 20 amino acids, plus 3 stop signals. Small, finite, universal — yet it specifies every protein.",16)],
 [L("~1,700 roots, an open lexicon",17.5,True,AMBER),
  L("~1,700 root-triplets → 7,236 word-forms → 135,366 tokens. A closed triplet code at the base, an unbounded text at the surface — the same closed-code / open-output design as the cell.",16)],
 fillA=TINT,fillB=AMBERT)

# 17 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Structural facts: ~96% triliteral roots; a sparse triplet code; ~54% anagram-reuse beating chance (p<0.003); position-sensitivity; reordering changes meaning. All character-level.",16)],
 [L("✗ BREAKS",17,True,RED),L("No biochemistry, no molecule, no ribosome. The triplet→field 'lookup' is linguistic convention, not translation. The analogy is COMBINATORIAL, not mechanistic.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether the disjointed letters (الم …) act as 'markers' is suggestive but untested; choosing the consonantal skeleton as the 'codon' is a modeling decision.",16)],f=(TINT,REDT,AMBERT))

# 18 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — a road not yet paved")
two(s,[L("WHY THIS IS NEW",18,True,NAVY),L("Almost all Qur'an study works at the level of meaning. Drop below it — to the raw character code — and a striking, testable fact appears: a triplet lexicon whose meaning is carried by ORDER, beating chance, the same combinatorial grammar life uses. Two Books, one mathematics of the small unit.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Is the triliteral root 'really' a triplet code, or do we impose it?  • Why does Arabic reuse letter-sets so much more than chance?  • What is the linguistic 'frameshift'?  • Where must the combinatorial analogy stop?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 19 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Bridges linguistics and bioinformatics: the same triplet / order-carries-information combinatorics underlies the genetic code and Arabic morphology — and underpins anagram puzzles, cryptography, and DNA sequencing alike.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("The unit is empty; the ORDER is full. ~54% of roots reuse a letter-set in a new order (beating chance, p<0.003) — sequence, not composition, is the message. Both Books.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"05_Codon_and_Root_Lecture.pptx")
print(f"L5 Codon & Root slides: {len(prs.slides)}")
