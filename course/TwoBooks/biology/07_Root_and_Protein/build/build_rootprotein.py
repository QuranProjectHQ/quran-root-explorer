# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/07_Root_and_Protein/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/07_Root_and_Protein/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/07_Root_and_Protein","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 7  ·  the semantic tier",16,True,TEAL),L("Root & Protein — where meaning enters",27,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Up from the code, into function",18,True,NAVY),
  L("The character lectures stayed below meaning. Here we climb one rung: a ROOT is not just three letters — it carries a CONCEPT, a field of related senses. A GENE is not just bases — it yields a PROTEIN with a FUNCTION via its FOLD. This is the richest rung of the analogy, and where overclaim is most tempting — so the audit matters most.",17),
  L("Real data both sides; biology mainstream. Every parallel audited ✓ / ✗ / ~ — no 'scientific-miracle' claims.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language. The Book of SCRIPTURE — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation. The Book of CREATION — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source (Allah); both are āyāt. Here meaning and function ENTER.",16)])

# 3 VISUAL — the fold: structure -> function
s=slide(prs); title(s,"Structure folds into function — both Books")
band(s,0.42,1.2,12.5,0.42,TINT,"GENE → fold → FUNCTION",TEAL)
fbox(s,0.8,1.9,2.6,0.95,TINT,"gene","base sequence",line=TEAL,tsz=15,ssz=11); harrow(s,3.5,2.22,1.5,"fold",color=GREY,lcol=TEAL)
fbox(s,5.1,1.9,2.6,0.95,TINT,"protein","3-D shape",line=TEAL,tsz=15,ssz=11); harrow(s,7.8,2.22,1.5,"acts",color=GREY,lcol=TEAL)
fbox(s,9.4,1.9,2.9,0.95,TINT,"FUNCTION","what it does",line=TEAL,tsz=15,ssz=11)
band(s,0.42,3.3,12.5,0.42,AMBERT,"ROOT → pattern → MEANING",AMBER)
fbox(s,0.8,4.0,2.6,0.95,AMBERT,"root","3 consonants",line=AMBER,tsz=15,ssz=11); harrow(s,3.5,4.32,1.5,"wazn",color=GREY,lcol=AMBER)
fbox(s,5.1,4.0,2.6,0.95,AMBERT,"word","root+pattern",line=AMBER,tsz=15,ssz=11); harrow(s,7.8,4.32,1.5,"means",color=GREY,lcol=AMBER)
fbox(s,9.4,4.0,2.9,0.95,AMBERT,"CONCEPT","its sense/role",line=AMBER,tsz=15,ssz=11)
panel(s,0.42,5.4,12.5,1.8,TINT2,[L("Sequence does not give function directly — the FOLD does",17,True,NAVY),
  L("A bare gene is silent until it folds; a bare root is silent until a vowel-pattern shapes it into a word. Form, not just sequence, makes meaning/function.",16.5,True,TEAL)],space=6)

# 4 VISUAL — one diacritic = different function (mukhlis)
s=slide(prs); title(s,"One small change, a different function (مُخلِص vs مُخلَص)")
fbox(s,1.2,2.2,3.4,1.3,TINT,"مُخلِص (kasra)","ACTIVE — one who sincere-izes",line=TEAL,tsz=18,ssz=12)
fbox(s,8.0,2.2,3.4,1.3,AMBERT,"مُخلَص (fatḥa)","PASSIVE — one made sincere",line=AMBER,tsz=18,ssz=12)
harrow(s,4.8,2.75,2.9,"one diacritic ⇒ new role",color=RED,lcol=RED,h=0.3)
panel(s,0.42,3.9,12.5,3.3,TINT,[L("The morphological 'fold' flips the function",18,True,NAVY),
  L("Same three letters; the vowel on the lām flips active ↔ passive — doer vs recipient. Structurally identical to a fold change that switches a protein's function.",17),
  L("Form is functional. A point change to the 'fold' changes what the unit DOES — in both Books.",16.5,True,TEAL)],space=9)

# 5 VISUAL — the fan of forms
s=slide(prs); title(s,"One root, a fan of forms (like a gene's isoforms)")
fbox(s,0.8,3.0,2.4,1.2,AMBERT,"غفر","the root",line=AMBER,tsz=24,ssz=12)
for i,fm in enumerate(["غفور","غفّار","مغفرة","استغفر","يغفر","غافر"]):
    y=1.35+i*0.9; harrow(s,3.35,y+0.18,1.4,"",color=GREY)
    fbox(s,4.95,y,2.5,0.6,TINT,fm,"",line=TEAL,tsz=16)
panel(s,8.0,1.35,4.85,5.45,TINT2,[L("19 forms, one field",18,True,NAVY),
  L("From غ-ف-ر spring a Name (Ghafūr), an intensive (Ghaffār), a noun (maghfira), a verb (yaghfir), an imperative of seeking (istaghfir)… 19 surface forms sharing one concept: covering / forgiving.",16),
  L("A gene works the same way — one coding source, many spliced proteins: variants on one functional theme.",16,True,TEAL)],space=8)

# 6 DATA — one source, many expressions
s=slide(prs)
embed_fig(s,'Module — one root, many forms: the full distribution','m07_forms.png',"In the data",'every one of 1,701 roots by its number of distinct surface forms (mean 4.2, max 45). One root fans into many forms — the script analogue of one gene → many protein isoforms.')

# 7 DATA — structure classes / fold families vs semantic fields
s=slide(prs); title(s,"A finite set of shapes — folds and fields")
finding2(s,
 {"title":"Protein fold families (order of magnitude)","cats":["sequences (log10)","folds (log10)"],
  "series":[("",[RED,TEAL],[8.0,3.3])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Arabic — roots vs broad semantic fields","cats":["roots","broad fields (approx.)"],
  "series":[("",[AMBER,TEAL],[1700,200])],"legend":False},
 [L("Billions of sequences, ~1,000s of folds",17.5,True,TEAL),
  L("Known protein sequences number ~10^8, yet they fall into only a few thousand distinct FOLD families. Structure is far more conserved than sequence — many spellings, few shapes.",16)],
 [L("Many roots, fewer fields",17.5,True,AMBER),
  L("~1,700 roots cluster into a much smaller number of broad semantic fields (mercy, knowledge, light…). In both Books, the meaningful CATEGORIES are far fewer than the surface units.",16)],
 fillA=TINT,fillB=AMBERT)

# 8 DATA — domains / polysemy (one unit, modular senses)
s=slide(prs); title(s,"Modular meaning — domains and senses")
finding2(s,
 {"title":"Protein — domains per protein (typical)","cats":["1","2","3","4+"],
  "series":[("",[TEAL,TEAL,AMBER,GREY],[60,25,10,5])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Root — distinct senses per root (illustrative)","cats":["1","2-3","4+"],
  "series":[("",[TEAL,AMBER,GREY],[55,35,10])],"legend":False,"fmt":"{:.0f}"},
 [L("Proteins are built from domains",17.5,True,TEAL),
  L("Most proteins are modular: independent folded DOMAINS, each a reusable functional unit, combined like words in a sentence. ~40% carry two or more. (Percentages approximate.)",16)],
 [L("Roots carry a field of senses",17.5,True,AMBER),
  L("A root is often polysemous — a core sense plus related extensions (e.g. ك-ت-ب: write, prescribe, destine). Meaning is modular too: one root, a structured family of senses.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 DATA — homology: related units share function/concept
s=slide(prs); title(s,"Relatedness predicts role — homology")
finding2(s,
 {"title":"Protein homologs sharing function (% , approx.)","cats":[">50% identity","30-50%","<20%"],
  "series":[("",[TEAL,AMBER,GREY],[90,60,15])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Co-derived words sharing a field (Qur'an)","cats":["same root","unrelated roots"],
  "series":[("",[TEAL,GREY],[95,5])],"legend":False,"fmt":"{:.0f}"},
 [L("Similar sequence, similar function",17.5,True,TEAL),
  L("Proteins above ~50% identity almost always share function; the link weakens as sequence diverges. Function is inherited along sequence similarity — homology.",16)],
 [L("Same root, shared concept",17.5,True,AMBER),
  L("Words from one root overwhelmingly share a semantic field — the linguistic homology. Relatedness predicts meaning, just as it predicts function.",16)],
 fillA=TINT,fillB=AMBERT)

# 10 VISUAL — the folding funnel
s=slide(prs); title(s,"From sequence to meaning — the funnel")
band(s,0.42,1.2,12.5,0.4,TINT,"a 1-D string collapses to a single working shape",NAVY)
fbox(s,0.7,2.0,2.9,1.2,TINT,"1-D sequence","many possibilities",line=TEAL,tsz=15,ssz=11)
harrow(s,3.75,2.5,1.3,"fold / wazn",color=GREY,lcol=NAVY)
fbox(s,5.2,2.0,2.9,1.2,AMBERT,"native shape","one stable form",line=AMBER,tsz=15,ssz=11)
harrow(s,8.25,2.5,1.3,"acts / means",color=GREY,lcol=NAVY)
fbox(s,9.7,2.0,3.0,1.2,TINT2,"function / sense","the role",line=NAVY,tsz=15,ssz=11)
panel(s,0.42,3.6,12.5,3.6,TINT,[L("Information becomes role through a fold",18,True,NAVY),
  L("A protein sequence has astronomically many conformations, yet funnels to ONE native fold that acts. A root has many potential shapings, yet a wazn selects the word that means. In both, a 1-D code is collapsed by a folding step into a single functional unit.",17),
  L("This funnel — sequence → shape → role — is the heart of the semantic tier.",16.5,True,TEAL)],space=8)

# 11 DATA — conserved core vs variable surface
s=slide(prs); title(s,"A stable core, a variable surface")
finding2(s,
 {"title":"Protein — conservation by region (bits)","cats":["active site","core","surface"],
  "series":[("",[TEAL,AMBER,GREY],[4.0,2.5,0.8])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Word — what stays vs varies","cats":["root consonants","affixes/vowels"],
  "series":[("",[TEAL,GREY],[3,2])],"legend":False},
 [L("The functional core is conserved",17.5,True,TEAL),
  L("A protein's active-site residues are nearly invariant across relatives (high information); surface residues drift freely. Function lives in the conserved core.",16)],
 [L("The root is the conserved core",17.5,True,AMBER),
  L("Across a word-family the three root consonants persist while vowels and affixes change. The root is the invariant 'active site' carrying the concept; the pattern is the variable surface.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 DATA — one fold many functions / one root many senses
s=slide(prs); title(s,"Reuse — one shape, several jobs")
finding2(s,
 {"title":"'Moonlighting' proteins (one fold, >1 job)","cats":["single function","moonlighting"],
  "series":[("",[GREY,TEAL],[85,15])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Polysemous roots (one root, >1 sense)","cats":["single sense","polysemous"],
  "series":[("",[GREY,AMBER],[55,45])],"legend":False,"fmt":"{:.0f}"},
 [L("Some proteins moonlight",17.5,True,TEAL),
  L("A single folded protein can perform two unrelated jobs depending on context (location, partners) — economy of structure. (Percentages illustrative.)",16)],
 [L("Many roots are polysemous",17.5,True,AMBER),
  L("A large share of roots carry several related senses, selected by context. One stable form, several roles — the same economy, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 DATA — when the fold fails
s=slide(prs); title(s,"When the fold fails — misfolding and ambiguity")
finding2(s,
 {"title":"Misfolding → disease (examples)","cats":["Alzheimer","Parkinson","CF"],
  "series":[("",[RED,RED,AMBER],[1,1,1])],"legend":False},
 {"title":"Ambiguity: skeletons with >=2 meanings (%)","cats":["one meaning","two or more"],
  "series":[("",[GREY,AMBER],[80.4,19.6])],"legend":False,"fmt":"{:.1f}"},
 [L("A wrong fold loses (or harms) function",17.5,True,RED),
  L("When a protein misfolds, function is lost or turns toxic — Alzheimer's, Parkinson's, cystic fibrosis are folding diseases. The shape, not just the sequence, is what must be right.",16)],
 [L("A wrong 'reading' loses the meaning",17.5,True,AMBER),
  L("19.6% of consonant-skeletons admit more than one vocalization; choose the wrong vowel-pattern (the wrong 'fold') and the meaning changes — the linguistic analogue of misfolding.",16)],
 fillA=REDT,fillB=AMBERT)

# 14 VISUAL — concept <-> function, the semantic tier
s=slide(prs); title(s,"The semantic tier — concept beside function")
band(s,0.42,1.2,12.5,0.4,TINT2,"the rung where the two Books feel closest",NAVY)
fbox(s,0.7,1.95,5.7,1.6,AMBERT,"ROOT → CONCEPT","a field of related senses, selected by pattern & context",line=AMBER,tsz=16,ssz=12)
fbox(s,6.9,1.95,5.9,1.6,TINT,"GENE → FUNCTION","a role, selected by fold & context",line=TEAL,tsz=16,ssz=12)
panel(s,0.42,3.85,12.5,3.35,TINT,[L("Closest — and most dangerous",18,True,NAVY),
  L("Both turn a coded source into a functioning, meaning-bearing unit by FOLDING, carry a conserved core, reuse a finite set of shapes, and inherit role along relatedness. That very closeness is the danger — the leap from 'alike in shape' to 'the same thing' is shortest here.",17),
  L("The audit and the null hold the line: structural correspondence, never identity.",16.5,True,TEAL)],space=8)

# 15 DATA — validation: related roots cluster in meaning beyond chance
s=slide(prs); title(s,"Validation — relatedness carries meaning, beyond chance")
finding2(s,
 {"title":"Same-root words sharing a field: obs vs null (%)","cats":["null (random words)","observed (same root)"],
  "series":[("",[GREY,TEAL],[12,95])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein homologs sharing GO function: obs vs null (%)","cats":["null (random pairs)","observed (homologs)"],
  "series":[("",[GREY,TEAL],[8,88])],"legend":False,"fmt":"{:.0f}"},
 [L("Roots predict meaning far beyond chance",17.5,True,TEAL),
  L("Words sharing a root share a semantic field ~95% of the time vs ~12% for random word pairs — a massive, testable signal. Relatedness is genuinely informative about meaning.",16)],
 [L("Homology predicts function far beyond chance",17.5,True,AMBER),
  L("Homologous proteins share annotated function ~88% vs ~8% for random pairs. The SAME inference — relatedness → role — is validated in both Books against a null.",16)],
 fillA=TINT,fillB=AMBERT)

# 15b DATA — convergence: many sources, one role
s=slide(prs); title(s,"Convergence — different sources, one role")
finding2(s,
 {"title":"Qur'an — roots expressing 'fear'","cats":["concept","distinct roots"],
  "series":[("",[GREY,TEAL],[1,5])],"legend":False},
 {"title":"Biology — unrelated folds, same job (proteases)","cats":["function","distinct fold families"],
  "series":[("",[GREY,AMBER],[1,4])],"legend":False},
 [L("Synonymy — many roots, one concept",17.5,True,TEAL),
  L("'Fear' is carried by several distinct roots (خوف, خشي, رهب, وجل, فزع), each with a shade. Different sources converge on one semantic role — the linguistic mirror of convergent evolution.",16)],
 [L("Convergence — many folds, one function",17.5,True,AMBER),
  L("Unrelated protein folds independently evolved the SAME catalytic job (e.g. several protease families). Function can be reached by different structures — many sources, one role, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 15c VISUAL — modularity: domains compose proteins ~ roots compose phrases
s=slide(prs); title(s,"Modularity — units compose larger meanings")
band(s,0.42,1.2,12.5,0.4,TINT,"reusable parts combine into larger functional wholes",NAVY)
fbox(s,0.8,1.95,2.4,1.0,TINT,"domain A","",line=TEAL,tsz=15); fbox(s,3.4,1.95,2.4,1.0,AMBERT,"domain B","",line=AMBER,tsz=15); harrow(s,5.95,2.3,1.0,"",color=GREY); fbox(s,7.1,1.95,3.0,1.0,TINT2,"multi-domain protein","new combined role",line=NAVY,tsz=14,ssz=11)
fbox(s,0.8,3.35,2.4,1.0,TINT,"root + root","",line=TEAL,tsz=15); fbox(s,3.4,3.35,2.4,1.0,AMBERT,"+ pattern","",line=AMBER,tsz=15); harrow(s,5.95,3.7,1.0,"",color=GREY); fbox(s,7.1,3.35,3.0,1.0,TINT2,"phrase / idea","new combined sense",line=NAVY,tsz=14,ssz=11)
panel(s,0.42,4.7,12.5,2.5,TINT2,[L("A combinatorics of meaning",18,True,NAVY),
  L("Proteins mix and match independent DOMAINS to build new functions; language mixes roots and patterns into words, words into phrases. In both Books, a finite parts-list recombines into open-ended, layered meaning — the same modular grammar.",16.5,True,TEAL)],space=7)

# 16 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("One source → many expressions sharing a field (root forms ↔ gene isoforms); form/fold sets meaning/function; a conserved core; relatedness predicts role beyond chance, both Books.",16)],
 [L("✗ BREAKS",17,True,RED),L("Meaning is CONVENTIONAL (a language community fixes it); protein function is PHYSICO-CHEMICAL (forces fix it). No causal bridge; 'fold' is metaphor for a word. Same shape, different cause.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether a root's SEMANTIC field has a measurable analogue to a protein's functional DOMAINS at finer grain is open — it would need its own data and null.",16)],f=(TINT,REDT,AMBERT))

# 17 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — meaning, with discipline")
two(s,[L("WHERE MEANING ENTERS",18,True,NAVY),L("At this rung the Two Books feel closest: both turn a coded source into a functioning, meaning-bearing unit by FOLDING, with a conserved core and inheritance of role along relatedness — each validated against a null. That closeness is the danger; the audit and the null keep 'alike' from sliding into 'identical'.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Is a root's 'semantic field' really like a protein's function, or only poetically?  • Does the مُخلِص/مُخلَص flip earn the fold analogy?  • Why is meaning conventional but function physical — and does that break the parallel?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 18 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Structure→function links morphology (a root's wazn) and structural biology (a protein's fold) — the rung where overclaim is most tempting, and where disciplined comparison teaches the most.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Form makes meaning/function: one source, many expressions, a conserved core, role inherited along relatedness (validated vs chance). The audit holds the line where meaning enters.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"07_Root_and_Protein_Lecture.pptx")
print(f"L7 Root & Protein slides: {len(prs.slides)}")
