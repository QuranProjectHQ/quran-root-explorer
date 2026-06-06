# -*- coding: utf-8 -*-
"""LOCKED Special-Topics slide system: full-canvas fill, back-of-room fonts, no empty bottoms.
Grid: title 0.18–1.05; content 1.18–7.28 (always filled). Body >=17pt, headers 18, titles 20–22.
Appendix columns distribute snippets into EQUAL vertical slots -> no empty space, ever."""
import os, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B)
GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B); AMBER=RGBColor(0xB8,0x86,0x0B)
WHITE=RGBColor(0xFF,0xFF,0xFF); ICE=RGBColor(0xCA,0xDC,0xFC); LTEAL=RGBColor(0x7F,0xCA,0xBD)
TINT=RGBColor(0xEA,0xF5,0xF3); TINT2=RGBColor(0xF2,0xF5,0xFB); AMBERT=RGBColor(0xF7,0xF1,0xDF); REDT=RGBColor(0xF7,0xEC,0xEC)
FIG="/sessions/kind-compassionate-feynman/mnt/RootCourse/SpecialTopics/figs"
CY0, CY1 = 1.18, 7.28          # content band (locked)
LX, RW = 0.42, 12.50           # left x, usable width

def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def deck():
    p=Presentation(); p.slide_width=Inches(13.333); p.slide_height=Inches(7.5); return p
def slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    b=s.shapes.add_shape(1,0,0,Inches(0.16),prs.slide_height); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background()
    return s
def title(s,t,sz=21):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.16),Inches(12.55),Inches(0.95)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia")
def _fill_panel(s,x,y,w,h,fill):
    sp=s.shapes.add_shape(5,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); return sp
def _need(lines,w,scale,space,pad=0.34):
    h=0.26
    avail=(w-2*pad-0.04)*72.0
    for (t,sz,b,c) in lines:
        fsz=max(1.0,sz*scale); cpl=max(6,avail/(fsz*0.48)); nl=max(1,math.ceil(len(t)/cpl))
        h+=nl*(fsz*1.30/72.0)+space*scale/72.0
    return h
def _fit(lines,w,h,space):
    lo,hi=0.80,1.70
    for _ in range(28):
        m=(lo+hi)/2.0
        if _need(lines,w,m,space)<=h*0.93: lo=m
        else: hi=m
    return lo
def panel(s,x,y,w,h,fill,lines,space=10,pad=0.34,anchor=MSO_ANCHOR.MIDDLE):
    """Auto-fills: scales fonts so the text block fills the panel -> no internal empty space."""
    lines=[ln for ln in lines if ln[0]!=""]
    sp=_fill_panel(s,x,y,w,h,fill); tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(pad); tf.margin_right=Inches(pad); tf.margin_top=Inches(0.1); tf.margin_bottom=Inches(0.1)
    scale=_fit(lines,w,h,space)
    for i,(txt,sz,bold,col) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space*scale); p.alignment=PP_ALIGN.LEFT
        r=p.add_run(); r.text=txt; r.font.size=Pt(min(round(sz*scale,1),30)); r.font.bold=bold; r.font.color.rgb=col; sf(r,"Calibri")
    return sp

def L(t,sz,b=False,c=INK): return (t,sz,b,c)

def two_stack(s,A,B,split=0.5,fillA=TINT,fillB=TINT2,spaceA=10,spaceB=10):
    gap=0.16; Htot=CY1-CY0-gap; hA=Htot*split; hB=Htot-hA
    panel(s,LX,CY0,RW,hA,fillA,A,space=spaceA)
    panel(s,LX,CY0+hA+gap,RW,hB,fillB,B,space=spaceB)

def three_stack(s,A,B,C,fills=(REDT,TINT,TINT2),splits=(0.34,0.33)):
    gap=0.16; Htot=CY1-CY0-2*gap; hA=Htot*splits[0]; hB=Htot*splits[1]; hC=Htot-hA-hB
    panel(s,LX,CY0,RW,hA,fills[0],A); panel(s,LX,CY0+hA+gap,RW,hB,fills[1],B)
    panel(s,LX,CY0+hA+hB+2*gap,RW,hC,fills[2],C)

def chart_row(s,chart,Aleft,Bright,fillA=TINT,fillB=TINT2,chart_w=9.0):
    p=os.path.join(FIG,chart)
    if os.path.exists(p):
        x=(13.333-chart_w)/2; s.shapes.add_picture(p,Inches(x),Inches(1.12),width=Inches(chart_w))
    pt=4.78; ph=CY1-pt
    panel(s,LX,pt,6.13,ph,fillA,Aleft,space=8)
    panel(s,6.79,pt,6.13,ph,fillB,Bright,space=8)

def appendix(s,columns):
    """Balanced card grid: every snippet becomes a colored card; cards are distributed evenly
    across columns and tile each column top-to-bottom, so no column is empty or sparse."""
    cards=[]
    for col in columns:
        header,hcol,fill,items=col
        for (ref,snip,tag) in items:
            cards.append((snip, f"{ref}  ·  {header}", hcol, fill))
    N=len(cards)
    if N==0: return
    ncols = 1 if N<=2 else (2 if N<=4 else 3)
    base=N//ncols; extra=N%ncols
    counts=[base+(1 if i<extra else 0) for i in range(ncols)]
    cw=(RW-(ncols-1)*0.30)/ncols
    idx=0
    for ci in range(ncols):
        x=LX+ci*(cw+0.30); cnt=counts[ci]
        ch=CY1-CY0; slot=ch/cnt
        for r in range(cnt):
            snip,reftag,hcol,fill=cards[idx]; idx+=1
            y=CY0+r*slot+0.06; hh=slot-0.12
            _fill_panel(s,x,y,cw,hh,fill)
            tb=s.shapes.add_textbox(Inches(x+0.22),Inches(y),Inches(cw-0.44),Inches(hh)); tf=tb.text_frame
            tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
            tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
            # auto-scale snippet to fill the card slot
            ssz=max(15.0,min(24.0, hh*72.0*0.30))
            p1=tf.paragraphs[0]; p1.alignment=PP_ALIGN.RIGHT; p1.space_after=Pt(5)
            r1=p1.add_run(); r1.text=snip; r1.font.size=Pt(round(ssz,1)); r1.font.color.rgb=INK; sf(r1,"Calibri")
            p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.LEFT
            r2=p2.add_run(); r2.text=reftag; r2.font.size=Pt(round(max(11.0,ssz*0.58),1)); r2.font.bold=True; r2.font.color.rgb=hcol; sf(r2,"Calibri")

def _rect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(max(h,0.012)))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); return sp
def _tb(s,x,y,w,h,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.03); tf.margin_right=Inches(0.03); tf.margin_top=Inches(0.01); tf.margin_bottom=Inches(0.01)
    for i,(t,sz,b,c) in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; sf(r,"Calibri")
    return tb
def ebar(s,bx,by,bw,bh,title,cats,series,ymax=None,legend=True,fmt="{:.0f}"):
    """Editable grouped bar chart. series=[(name,color,[vals]),...]. cats=[str,...] (Arabic ok)."""
    _tb(s,bx,by,bw,0.36,[(title,14,True,NAVY)])
    multi=len(series)>1
    ytop=by+0.40+(0.32 if (legend and multi) else 0.0)
    cat_h=0.62; base=by+bh-cat_h; plot_h=base-ytop
    n=len(cats); ns=len(series)
    allv=[v for _,_,vs in series for v in vs]; mx=max(allv) if allv else 1
    ymax=ymax or (mx*1.20 if mx>0 else 1)
    gw=bw/n; inner=gw*0.72; barw=inner/ns
    _rect(s,bx,base,bw,0.016,GREY)
    for ci in range(n):
        gx=bx+ci*gw+(gw-inner)/2
        for si,(nm,col,vs) in enumerate(series):
            v=vs[ci]; h2=plot_h*(v/ymax); xx=gx+si*barw
            cc=col[ci] if isinstance(col,list) else col
            if h2>0.02: _rect(s,xx+barw*0.07,base-h2,barw*0.86,h2,cc)
            _tb(s,xx-0.12,base-h2-0.30,barw*0.86+0.24,0.27,[(fmt.format(v),11.5,True,NAVY)])
        _tb(s,bx+ci*gw,base+0.04,gw,cat_h,[(cats[ci],12.5,False,NAVY)],anchor=MSO_ANCHOR.TOP)
    if legend and multi:
        lx=bx+0.1; ly=by+0.42
        for nm,col,_ in series:
            _rect(s,lx,ly,0.17,0.17,col); _tb(s,lx+0.2,ly-0.05,2.7,0.28,[(nm,10.5,False,INK)],align=PP_ALIGN.LEFT)
            lx+=2.85
def finding2(s,cA,cB,panelA,panelB,fillA=TINT,fillB=TINT2):
    """Two editable charts (top) + two panels (bottom). cA/cB = dict(title,cats,series,ymax?,legend?,fmt?)."""
    def draw(box,c):
        ebar(s,box[0],box[1],box[2],box[3],c["title"],c["cats"],c["series"],
             c.get("ymax"),c.get("legend",True),c.get("fmt","{:.0f}"))
    draw((0.42,1.16,6.13,3.40),cA); draw((6.79,1.16,6.13,3.40),cB)
    pt=4.70; ph=CY1-pt
    panel(s,0.42,pt,6.13,ph,fillA,panelA,space=8); panel(s,6.79,pt,6.13,ph,fillB,panelB,space=8)
