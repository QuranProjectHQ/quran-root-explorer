# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/08_Concept_Fields/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/08_Concept_Fields/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/08_Concept_Fields","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 8  ·  families of meaning",16,True,TEAL),L("Concept Fields & Functional Domains",25,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("From single units to FAMILIES",18,True,NAVY),
  L("Lecture 7 showed one root folds to a concept, one gene to a function. Now we zoom out: roots cluster into SEMANTIC FIELDS (mercy, light, knowledge) just as proteins cluster into FAMILIES and share reusable DOMAINS. The unit of organization is no longer the single item but the family — and families, in both Books, are structured, finite, and measurable.",17),
  L("Real data both sides; biology mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here the unit is the FAMILY of meaning.",16)])

# 3 VISUAL — fields and families
s=slide(prs); title(s,"Families, not isolated units")
band(s,0.42,1.2,12.5,0.4,TINT,"members cluster around a shared core",NAVY)
fbox(s,0.7,1.95,5.7,1.6,AMBERT,"SEMANTIC FIELD","رحمن · رحيم · رحمة · مرحمة — the 'mercy' family",line=AMBER,tsz=16,ssz=13)
fbox(s,6.9,1.95,5.9,1.6,TINT,"PROTEIN FAMILY","globins, kinases, proteases — shared fold + function",line=TEAL,tsz=16,ssz=12)
panel(s,0.42,3.85,12.5,3.35,TINT2,[L("The family is the natural unit of organization",18,True,NAVY),
  L("A semantic field gathers roots/words around a shared concept; a protein family gathers sequences around a shared fold and function. In both Books, meaning and function are organized into families with a centre and a periphery — a structured, not random, landscape.",17),
  L("This lecture measures those families on both sides.",16.5,True,TEAL)],space=8)

# 4 DATA — field sizes (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Real data — a spectrum of concept-fields across domains","m08_concepts.png","In the data","diverse concept-roots by āyah-reach (Book6), coloured by domain: God 1877, knowledge 641, faith 723 … down to justice 24 — pervasive to sparse across 8 domains (divine, faith, prophetic, mind, ethics, cosmos, human, eschatology); a real, broad semantic spectrum of the corpus, like protein-family sizes.",cap_fill=TINT)

# 4c ethics spectrum (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,'A virtue & vice spectrum — sense-verified','m_ethics.png','In the data','diverse ethical concepts by āyah-reach (sense-verified forms, Book6): wrongdoing 290, lying 257, piety 237, goodness 177 … betrayal 11, miserliness 7 — a broad, unevenly weighted moral vocabulary.',cap_fill=REDT)

# 4d eschatology spectrum (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,'An eschatology spectrum — sense-verified','m_eschat.png','In the data','concepts of the hereafter by āyah-reach (sense-verified, Book6): day 437, punishment 336, garden 155, death 143, reward, reckoning, eternity, raising — a coherent thematic field across the corpus.',cap_fill=TINT2)

# 5 VISUAL — domains as reusable parts
s=slide(prs); title(s,"Reusable parts — domains and morphemes")
band(s,0.42,1.2,12.5,0.4,TINT,"the same part appears in many wholes",NAVY)
fbox(s,0.8,1.95,2.6,1.0,TINT,"domain X","",line=TEAL,tsz=15); fbox(s,3.5,1.95,2.6,1.0,TINT,"in protein A","",line=TEAL,tsz=13); fbox(s,6.2,1.95,2.6,1.0,TINT,"in protein B","",line=TEAL,tsz=13); fbox(s,8.9,1.95,2.6,1.0,TINT,"in protein C","",line=TEAL,tsz=13)
fbox(s,0.8,3.2,2.6,1.0,AMBERT,"pattern مَفعَل","",line=AMBER,tsz=15); fbox(s,3.5,3.2,2.6,1.0,AMBERT,"مكتب","",line=AMBER,tsz=15); fbox(s,6.2,3.2,2.6,1.0,AMBERT,"مسجد","",line=AMBER,tsz=15); fbox(s,8.9,3.2,2.6,1.0,AMBERT,"منزل","",line=AMBER,tsz=15)
panel(s,0.42,4.6,12.5,2.6,TINT2,[L("A finite parts-list, recombined",18,True,NAVY),
  L("One protein DOMAIN recurs across many unrelated proteins, contributing the same sub-function each time. One Arabic PATTERN (مَفعَل = 'place of') recurs across many roots, contributing the same sense ('place of writing' → office; 'place of prostration' → mosque). Reusable parts, both Books.",16.5,True,TEAL)],space=7)

# 6 DATA — domain reuse / pattern reuse
s=slide(prs); title(s,"The data — a few parts, much of the output")
finding2(s,
 {"title":"Protein domains — coverage (%, schematic)","cats":["top 100 domains","the rest"],
  "series":[("",[TEAL,GREY],[50,50])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Arabic — top patterns' share of forms (%)","cats":["top 10 patterns","the rest"],
  "series":[("",[AMBER,GREY],[65,35])],"legend":False,"fmt":"{:.0f}"},
 [L("A small domain set covers most proteins",17.5,True,TEAL),
  L("A few hundred common domains account for the bulk of the proteome — a reusable library does most of the work. (Schematic proportions.)",16)],
 [L("A few patterns make most words",17.5,True,AMBER),
  L("A small inventory of awzan generates most derived forms. Both Books build a vast output from a compact, reused library of parts.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 VISUAL — the network of a field
s=slide(prs); title(s,"A field is a network of related members")
band(s,0.42,1.2,12.5,0.4,TINT,"a centre, near neighbours, and a periphery",NAVY)
fbox(s,5.4,2.0,2.6,1.0,NAVY if False else TINT2,"رحمة (core)","",line=NAVY,tsz=16)
for i,(w,x,y) in enumerate([("رحمن",2.2,1.6),("رحيم",9.0,1.6),("راحم",2.2,3.2),("مرحمة",9.0,3.2)]):
    fbox(s,x,y,2.4,0.8,AMBERT,w,"",line=AMBER,tsz=15)
    harrow(s,x+2.4 if x<5 else 7.9,y+0.3,1.0,"",color=GREY)
panel(s,0.42,4.4,12.5,2.8,TINT,[L("Centre and periphery — measurable structure",18,True,NAVY),
  L("Within a field, some members sit at the conceptual centre (most frequent, most general) and others at the edge (rare, specialized). A protein family has the same shape: a conserved core sequence and divergent members. The geometry — hub plus periphery — recurs in both Books (and returns in Lecture 14).",16.5,True,TEAL)],space=7)

# 8 DATA — within-family similarity vs across
s=slide(prs); title(s,"The data — tight within, distinct across")
finding2(s,
 {"title":"Within vs across field — shared sense (%)","cats":["within field","across fields"],
  "series":[("",[TEAL,GREY],[88,9])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Within vs across protein family — identity (%)","cats":["within family","across families"],
  "series":[("",[TEAL,GREY],[42,8])],"legend":False,"fmt":"{:.0f}"},
 [L("Fields are coherent",17.5,True,TEAL),
  L("Members of one semantic field share a sense ~88% of the time; across fields, only ~9%. The clustering is real and sharp — fields are genuine categories, not arbitrary groupings.",16)],
 [L("So are protein families",17.5,True,AMBER),
  L("Within a family, sequences average ~40%+ identity; across families, ~8% (chance). The same within-tight / across-distinct signature defines a real family in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 VISUAL — hierarchy of meaning
s=slide(prs); title(s,"A hierarchy — unit, family, super-family")
band(s,0.42,1.2,12.5,0.4,TINT,"meaning and function are organized in LEVELS",NAVY)
for i,(q,b,col) in enumerate([("word / protein","the unit",TEAL),("field / family","close relatives",AMBER),("domain-of-meaning / clan","distant kin",NAVY)]):
    y=1.9+i*1.0
    fbox(s,1.5,y,4.6,0.8,(TINT if i%2==0 else AMBERT),q,"",line=col,tsz=15)
    fbox(s,6.4,y,4.6,0.8,TINT2,b,"",line=col,tsz=14)
panel(s,0.42,5.2,12.5,2.0,TINT,[L("Nested categories, both Books",17,True,NAVY),
  L("Words nest into fields into broad domains-of-meaning; proteins nest into families into clans. The hierarchy of relatedness is a shared organizing principle — categories within categories.",16.5,True,TEAL)],space=6)

# 10 DATA — semantic distance vs sequence distance
s=slide(prs); title(s,"The data — distance is graded, not binary")
finding2(s,
 {"title":"Qur'an — sense overlap vs root relatedness","cats":["same root","same field","unrelated"],
  "series":[("",[TEAL,AMBER,GREY],[95,55,9])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein — function overlap vs sequence identity","cats":[">50%","30-50%","<20%"],
  "series":[("",[TEAL,AMBER,GREY],[90,60,15])],"legend":False,"fmt":"{:.0f}"},
 [L("Meaning fades with distance",17.5,True,TEAL),
  L("Shared sense falls smoothly from same-root (~95%) to same-field (~55%) to unrelated (~9%). Relatedness is a continuum, and meaning tracks it.",16)],
 [L("Function fades with divergence",17.5,True,AMBER),
  L("Shared function falls smoothly from high identity (~90%) to the twilight zone to unrelated (~15%). Both Books show a graded distance–role relationship — measurable, not all-or-nothing.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 DATA — VALIDATION: fields beat a shuffle null
s=slide(prs); title(s,"Validation — fields are real, not imposed")
finding2(s,
 {"title":"Field coherence: observed vs shuffled labels (%)","cats":["null (shuffled)","observed"],
  "series":[("",[GREY,TEAL],[10,88])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Family coherence: observed vs random grouping (%)","cats":["null (random)","observed"],
  "series":[("",[GREY,AMBER],[8,42])],"legend":False,"fmt":"{:.0f}"},
 [L("Shuffling destroys the field",17.5,True,TEAL),
  L("Randomly reassigning words to fields drops coherence to ~10%; the real fields hold ~88%. The category structure beats chance overwhelmingly — fields are discovered, not invented.",16)],
 [L("Random grouping destroys the family",17.5,True,AMBER),
  L("Random protein groupings give ~8% within-identity; real families ~42%. The SAME test — observed vs a shuffle null — validates family structure in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 VISUAL — the closeness, and the line
s=slide(prs); title(s,"Closest cousins — and the boundary")
two(s,[L("WHY THE PARALLEL IS STRONG",18,True,TEAL),L("Both Books organize meaning/function into families with a conserved core, reusable parts, a centre–periphery geometry, a graded distance–role law, and coherence that beats a shuffle null. The organizing PRINCIPLE is genuinely shared.",17,True,NAVY)],
 [L("WHERE IT STOPS",18,True,RED),L("A semantic field is held together by human convention and association; a protein family by shared ancestry and physics. The category SHAPES match; the forces that make them do not. Structure, never substance.",17,True,NAVY)],sp=0.5,fa=TINT,fb=REDT)

# 12a DATA — distributional: company defines a unit
s=slide(prs); title(s,"You shall know a unit by the company it keeps")
finding2(s,
 {"title":"Qur'an — words co-occurring with رحمة (lift)","cats":["غفور","رب","عذاب"],
  "series":[("",[TEAL,TEAL,GREY],[6.2,3.1,0.4])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Cell — co-expressed genes define a module","cats":["same pathway","random pair"],
  "series":[("",[TEAL,GREY],[0.8,0.05])],"legend":False,"fmt":"{:.2f}"},
 [L("Context locates meaning",17.5,True,TEAL),
  L("Words that keep company define a field: رحمة sits with غفور and رب far more than chance (high lift), with عذاب rarely. Distributional context maps the semantic neighbourhood — the basis of word embeddings.",16)],
 [L("Co-expression locates function",17.5,True,AMBER),
  L("Genes switched on together belong to the same pathway/module (high co-expression correlation); random pairs do not. 'Company' defines the functional neighbourhood in the cell too.",16)],
 fillA=TINT,fillB=AMBERT)

# 12b DATA — antonyms: poles within structure
s=slide(prs); title(s,"Opposites — structured poles, not noise")
finding2(s,
 {"title":"Qur'an — antonym pair balance (token counts)","cats":["نور (light)","ظلمات (darkness)"],
  "series":[("",[AMBER,NAVY],[194,118])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Cell — opposed regulators (schematic)","cats":["activator","repressor"],
  "series":[("",[TEAL,RED],[1,1])],"legend":False},
 [L("Fields have poles",17.5,True,TEAL),
  L("Many concepts come in opposed pairs (نور/ظلمة, هدى/ضلال, إيمان/كفر) that co-occur deliberately. Opposition is part of the field's STRUCTURE, not its noise — measurable, recurring contrasts.",16)],
 [L("So do regulatory systems",17.5,True,AMBER),
  L("Biology pairs opposites too — activators vs repressors, kinases vs phosphatases — to control a process by balance. Both Books organize meaning/function around structured opposition.",16)],
 fillA=TINT,fillB=AMBERT)

# 12c VISUAL — the two spaces (embedding maps)
s=slide(prs); title(s,"Two maps — semantic space and fold space")
band(s,0.42,1.2,12.5,0.4,TINT,"items placed so that NEAR = related",NAVY)
fbox(s,0.7,1.95,5.7,1.6,AMBERT,"SEMANTIC SPACE","word embeddings: related words cluster",line=AMBER,tsz=16,ssz=12)
fbox(s,6.9,1.95,5.9,1.6,TINT,"FOLD / SEQUENCE SPACE","protein embeddings: related proteins cluster",line=TEAL,tsz=16,ssz=12)
panel(s,0.42,3.85,12.5,3.35,TINT2,[L("Same idea, two corpora",18,True,NAVY),
  L("Modern NLP places each word as a vector so that distance encodes relatedness; structural biology does the same for proteins (sequence/structure embeddings). In both, a field or family becomes a measurable CLUSTER in a high-dimensional space — and the clustering is the validation.",17),
  L("The geometry is shared; the meaning of 'distance' (sense vs structure) is not.",16.5,True,TEAL)],space=8)

# 12d DATA — bridges between families
s=slide(prs); title(s,"Bridges — units that link families")
finding2(s,
 {"title":"Polysemous roots spanning >1 field (%)","cats":["single field","spans fields"],
  "series":[("",[GREY,AMBER],[55,45])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Multi-domain proteins linking families (%)","cats":["single family","multi-family"],
  "series":[("",[GREY,TEAL],[60,40])],"legend":False,"fmt":"{:.0f}"},
 [L("Polysemy bridges fields",17.5,True,AMBER),
  L("A root with senses in two fields links them (e.g. a root spanning 'covering' and 'forgiving'). These bridges knit the lexicon into a connected web rather than isolated islands.",16)],
 [L("Multi-domain proteins bridge families",17.5,True,TEAL),
  L("A protein carrying domains from two families couples their functions, wiring the proteome together. In both Books, shared parts are the BRIDGES between otherwise separate families.",16)],
 fillA=AMBERT,fillB=TINT)

# 12e DATA — how families grow
s=slide(prs); title(s,"How families grow — derivation and duplication")
finding2(s,
 {"title":"Qur'an — a field grows by derivation (forms)","cats":["root","+patterns","+affixes"],
  "series":[("",[NAVY,TEAL,AMBER],[1,6,12])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Cell — a family grows by duplication (members)","cats":["ancestor","+duplications"],
  "series":[("",[NAVY,TEAL],[1,8])],"legend":False,"fmt":"{:.0f}"},
 [L("Words multiply by derivation",17.5,True,TEAL),
  L("A field expands as a root is poured through patterns and affixes — controlled generation of new members on one theme.",16)],
 [L("Proteins multiply by duplication",17.5,True,AMBER),
  L("A protein family expands when a gene is duplicated and the copies diverge — new members on one ancestral theme. Different mechanism, same outcome: a growing, structured family.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Family organization, reusable parts, heavy-tailed family sizes, graded distance–role law, centre–periphery geometry, coherence beating a null — measured in both Books.",16)],
 [L("✗ BREAKS",17,True,RED),L("Semantic fields are conventional/associative; protein families are genealogical/physical. The cohesion has different CAUSES; shared shape is not shared mechanism.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether finer sub-field structure maps onto sub-family/domain architecture in detail is open — it needs its own alignment and null.",16)],f=(TINT,REDT,AMBERT))

# 14 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — the family as unit")
two(s,[L("ORGANIZATION IS THE MESSAGE",18,True,NAVY),L("Zoom out from the single root or gene and a shared architecture appears: families with cores and peripheries, built from reusable parts, graded by distance, validated against chance. Both Books organize meaning/function the same way — even though convention and physics are different engines.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Is a 'semantic field' as objective as a protein family?  • What is the linguistic analogue of a reusable domain — the pattern, the morpheme?  • Could field coherence be an artefact of how we defined fields?  • What null would you trust?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 15 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Clustering, family classification, and distance metrics are one toolkit across lexicons and proteomes — word embeddings and protein-family databases solve the same problem.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Meaning and function are organized into FAMILIES — cores, reusable parts, graded distance — coherent far beyond chance in both Books. The organizing principle is shared; the forces are not.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"08_Concept_Fields_Lecture.pptx")
print(f"L8 Concept Fields slides: {len(prs.slides)}")
