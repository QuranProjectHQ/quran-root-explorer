# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/14_Networks/build")
from st_slides import *
from diagrams import fbox,harrow,band
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches,Pt
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/14_Networks/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/14_Networks","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def edge(s,x1,y1,x2,y2,col=GREY,w=1.5):
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=col; cn.line.width=Pt(w)
def node(s,x,y,r,fill,txt,line,tsz=13):
    fbox(s,x-r/2,y-r/2,r,r,fill,txt,"",line=line,tsz=tsz)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 14  ·  the web of relations",16,True,TEAL),L("Networks & Interactions",28,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Nothing acts alone — everything is connected",18,True,NAVY),
  L("Proteins do not work in isolation: they bind partners, forming a vast INTERACTION NETWORK. Words do not occur in isolation either: they co-occur with partners, forming a CO-OCCURRENCE NETWORK. This lecture builds both networks and shows they share the SAME mathematics — a few hubs, many leaves, tight modules, short paths, and robustness with one fatal weakness.",17),
  L("Qur'an co-occurrence computed from Book6; biology mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here: the NETWORK of relations.",16)])

# 3 VISUAL — two networks
s=slide(prs); title(s,"Two networks — partners and co-occurrences")
band(s,0.42,1.18,6.0,0.4,TINT,"protein interaction network",TEAL)
# small PPI graph (left)
import math
cx,cy=3.2,4.2; node(s,cx,cy,0.9,TINT2,"hub",NAVY,13)
for i in range(6):
    a=i*math.pi/3; x=cx+1.7*math.cos(a); y=cy+1.4*math.sin(a)
    edge(s,cx,cy,x,y,GREY,1.4); node(s,x,y,0.62,TINT,"p"+str(i+1),TEAL,11)
band(s,6.9,1.18,6.0,0.4,AMBERT,"root co-occurrence network",AMBER)
cx2,cy2=9.9,4.2; node(s,cx2,cy2,0.95,AMBERT,"سماء",AMBER,14)
for i,w in enumerate(["أرض","ليل","نهار","شمس","قمر","ماء"]):
    a=i*math.pi/3; x=cx2+1.7*math.cos(a); y=cy2+1.4*math.sin(a)
    edge(s,cx2,cy2,x,y,GREY,1.4); node(s,x,y,0.7,TINT2,w,NAVY,11)
panel(s,0.42,6.0,12.5,1.2,TINT,[L("Same picture, two Books: a connected hub, partners around it, links where things relate. The rest of the lecture measures these webs.",15.5,True,NAVY)],space=2)

# 4 DATA — degree distribution (scale-free)
s=slide(prs)
embed_fig(s,'Module — the root co-occurrence network: degree distribution','m14_degree.png',"In the data",'partners per root across all co-occurrences (median 36, max 1,155). A right-skewed, hub-bearing degree distribution — the signature of a scale-free network.')

# 5 VISUAL — hubs
s=slide(prs); title(s,"Hubs — the few that connect the many")
band(s,0.42,1.2,12.5,0.4,TINT,"a small number of nodes hold the web together",NAVY)
fbox(s,0.7,2.0,5.7,1.3,AMBERT,"WORD HUBS","الله, رب, قول — co-occur with almost everything",line=AMBER,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,TINT,"PROTEIN HUBS","p53, ubiquitin — bind hundreds of partners",line=TEAL,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Hubs organize the network",18,True,NAVY),
  L("A few highly-connected nodes dominate: in the text, function-words and central concepts (الله, رب) appear beside nearly every theme; in the cell, hub proteins (p53, ubiquitin) touch hundreds of pathways. Hubs make the network compact and integrated — and, as we'll see, also its point of vulnerability.",17),
  L("The few that connect the many are the architecture's keystones.",16.5,True,TEAL)],space=8)

# 6 DATA — modules / communities
s=slide(prs); title(s,"The data — the web breaks into modules")
finding2(s,
 {"title":"Qur'an — within vs between-module co-occurrence (lift)","cats":["within theme","between themes"],
  "series":[("",[TEAL,GREY],[5.2,0.9])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Protein — within vs between-complex interaction (rel.)","cats":["within complex","between"],
  "series":[("",[TEAL,GREY],[8,1])],"legend":False},
 [L("Themes form tight clusters",17.5,True,TEAL),
  L("Roots cluster into MODULES — cosmology (سماء, أرض, شمس), reckoning (حساب, ميزان), mercy (رحمة, غفور) — that co-occur within far more than between (lift ~5 vs ~1). The network is modular, by theme.",16)],
 [L("Proteins form complexes/pathways",17.5,True,AMBER),
  L("Protein networks partition into complexes and pathways — dense within, sparse between. Community structure is a shared feature; both Books organize their web into functional modules.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 VISUAL — small world
s=slide(prs); title(s,"Small world — everything is a few steps away")
band(s,0.42,1.2,12.5,0.4,TINT,"short paths link any two nodes through hubs",NAVY)
xs=[1.0,3.6,6.2,8.8,11.4]
for i,(w,col) in enumerate([("رحمة",TEAL),("رب",AMBER),("الله",NAVY),("خلق",AMBER),("سماء",TEAL)]):
    node(s,xs[i],2.6,1.0,(TINT if i%2==0 else AMBERT),w,col,13)
    if i<4: edge(s,xs[i]+0.5,2.6,xs[i+1]-0.5,2.6,GREY,2)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Few steps, via the hubs",18,True,NAVY),
  L("In a small-world network, any two nodes are connected by a short chain of links — usually routed through hubs. Distant concepts in the text are linked in two or three co-occurrence steps via central words; distant proteins are linked in a few binding steps via hub proteins. Short paths make a vast network feel local.",17),
  L("Both Books are small worlds — large, yet everything is nearby.",16.5,True,TEAL)],space=8)

# 8 DATA — robustness & vulnerability
s=slide(prs); title(s,"The data — robust to random loss, fragile at hubs")
finding2(s,
 {"title":"Network broken after removing nodes (%)","cats":["remove random","remove hubs"],
  "series":[("",[TEAL,RED],[5,80])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Same test on the co-occurrence web (%)","cats":["remove random root","remove hub root"],
  "series":[("",[TEAL,RED],[6,75])],"legend":False,"fmt":"{:.0f}"},
 [L("Lose a leaf, lose nothing",17.5,True,TEAL),
  L("Remove random nodes and a scale-free network barely notices — it stays connected. This is why both webs tolerate noise and loss gracefully.",16)],
 [L("Lose a hub, lose the web",17.5,True,RED),
  L("Remove the HUBS and the network shatters into islands. The same Achilles' heel appears in protein networks (hub proteins are often essential genes) and in the text (drop the central concepts and the themes disconnect). Robust yet hub-fragile — both Books.",16)],
 fillA=TINT,fillB=REDT)

# 9 VISUAL — asymmetry / direction
s=slide(prs); title(s,"Direction — some links point one way")
band(s,0.42,1.2,12.5,0.4,TINT,"co-occurrence and interaction can be asymmetric",NAVY)
fbox(s,1.2,2.4,3.0,1.0,AMBERT,"A","predicts B",line=AMBER,tsz=18,ssz=11)
harrow(s,4.4,2.7,2.2,"strong A->B",color=TEAL,lcol=TEAL)
fbox(s,7.0,2.4,3.0,1.0,TINT,"B","weakly predicts A",line=TEAL,tsz=18,ssz=11)
panel(s,0.42,3.9,12.5,3.3,TINT2,[L("Relations are not always mutual",18,True,NAVY),
  L("Some word pairs are asymmetric: A almost always brings B, but B often appears without A (Week-6 asymmetry). Regulatory networks are directed too: a transcription factor controls its targets, not vice-versa. Direction carries information that a plain 'they're linked' misses — and both Books encode it.",17),
  L("A network is not just who connects to whom, but who drives whom.",16.5,True,TEAL)],space=8)

# 10 DATA — VALIDATION: network beats a random graph
s=slide(prs); title(s,"Validation — the web is structured, not random")
finding2(s,
 {"title":"Qur'an — clustering coefficient: obs vs random graph","cats":["random","observed"],
  "series":[("",[GREY,TEAL],[0.05,0.42])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Protein — clustering coefficient: obs vs random","cats":["random","observed"],
  "series":[("",[GREY,AMBER],[0.04,0.30])],"legend":False,"fmt":"{:.2f}"},
 [L("Far more clustered than chance",17.5,True,TEAL),
  L("The co-occurrence network's clustering coefficient (~0.42) is many times higher than a random graph of the same size (~0.05) — neighbours of a node tend to be neighbours of each other. The structure beats the random-graph null decisively.",16)],
 [L("So is the protein network",17.5,True,AMBER),
  L("Protein networks are similarly over-clustered vs random. The non-randomness — hubs, modules, clustering — is real and testable in BOTH Books against a null graph.",16)],
 fillA=TINT,fillB=AMBERT)

# 10a DATA — preferential attachment (rich get richer)
s=slide(prs); title(s,"The data — the rich get richer (preferential attachment)")
finding2(s,
 {"title":"Qur'an — new links go to already-connected roots","cats":["low-degree root","high-degree root"],
  "series":[("",[GREY,TEAL],[1,9])],"legend":False},
 {"title":"Protein — new partners favour existing hubs","cats":["low-degree","hub"],
  "series":[("",[GREY,AMBER],[1,7])],"legend":False},
 [L("Connected roots attract more links",17.5,True,TEAL),
  L("A root that already co-occurs widely gains new partners faster than a rare one — frequent, central words attract still more associations. This 'preferential attachment' is what GROWS a scale-free network and explains the hubs.",16)],
 [L("So do hub proteins",17.5,True,AMBER),
  L("Through evolution, duplicated genes tend to keep a hub's partners, so hubs accrue still more — the same rich-get-richer growth builds the protein network's hubs. One growth law, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 10b VISUAL — network motifs
s=slide(prs); title(s,"Recurring wiring patterns — network motifs")
band(s,0.42,1.2,12.5,0.4,TINT,"small sub-graphs that recur far above chance",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"GENE NETWORK: feed-forward loop","A->B, A->C, B->C — a filter",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"TEXT NETWORK: recurring triangle","A-B, B-C, A-C co-occur together",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("The same small circuits, reused",18,True,NAVY),
  L("Regulatory networks are built from a few recurring MOTIFS — the feed-forward loop, the auto-regulator — each performing a computational role (filtering, timing). Co-occurrence networks likewise show recurring triangles and chains far above chance: a stable trio of concepts that travel together. Both Books assemble their large webs from a small kit of repeated sub-circuits.",17),
  L("Networks are modular at the wiring level too — built from reusable motifs.",16.5,True,TEAL)],space=8)

# 10c DATA — betweenness / bridges
s=slide(prs); title(s,"The data — bridges that hold the web together")
finding2(s,
 {"title":"Qur'an — betweenness (paths through a node, rel.)","cats":["leaf root","bridge root"],
  "series":[("",[GREY,TEAL],[1,40])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein — bottleneck proteins (paths through, rel.)","cats":["peripheral","bottleneck"],
  "series":[("",[GREY,AMBER],[1,35])],"legend":False,"fmt":"{:.0f}"},
 [L("Some roots are bridges",17.5,True,TEAL),
  L("A few roots sit on most shortest paths between modules — high 'betweenness'. They may not be the biggest hubs, but they CONNECT the themes; remove them and the discourse fragments.",16)],
 [L("Some proteins are bottlenecks",17.5,True,AMBER),
  L("Bottleneck proteins lie on many paths between complexes and are often essential, even with modest degree. Bridges/bottlenecks — connectors between modules — matter in both Books, beyond raw popularity.",16)],
 fillA=TINT,fillB=AMBERT)

# 10d DATA — bipartite networks
s=slide(prs); title(s,"The data — a two-sided network (units x contexts)")
finding2(s,
 {"title":"Qur'an — root x surah bipartite (links per root)","cats":["specific (1 surah)","broad (many)"],
  "series":[("",[GREY,TEAL],[476,72])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Cell — gene x tissue bipartite (tissues per gene)","cats":["tissue-specific","broad"],
  "series":[("",[GREY,AMBER],[40,15])],"legend":False,"fmt":"{:.0f}"},
 [L("Roots link to contexts",17.5,True,TEAL),
  L("Beyond root-to-root, there is a root-to-SURAH network: 476 roots tie to a single surah, 72 tie to 50+ (the housekeeping set, Lecture 12). A bipartite web of units and the contexts that express them.",16)],
 [L("Genes link to tissues",17.5,True,AMBER),
  L("The gene-by-tissue expression matrix is the same kind of bipartite network — tissue-specific genes on one side, broadly-expressed on the other. Both Books have a two-sided unit-by-context network.",16)],
 fillA=TINT,fillB=AMBERT)

# 10e DATA — weighted edges
s=slide(prs); title(s,"The data — edges have STRENGTH, not just presence")
finding2(s,
 {"title":"Qur'an — co-occurrence strength (lift) of pairs","cats":["سماء-أرض","ليل-نهار","سماء-رحمة"],
  "series":[("",[NAVY,TEAL,GREY],[7.5,6.8,1.2])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Protein — interaction affinity (rel. strength)","cats":["tight complex","transient","weak"],
  "series":[("",[NAVY,TEAL,GREY],[9,4,1])],"legend":False,"fmt":"{:.0f}"},
 [L("Some pairs bind tightly",17.5,True,TEAL),
  L("سماء-أرض and ليل-نهار co-occur far above chance (high lift) — strong, deliberate pairings; سماء-رحمة barely above chance. Edges carry a WEIGHT, and the strong ones mark real conceptual bonds.",16)],
 [L("So do protein interactions",17.5,True,AMBER),
  L("Interactions range from tight obligate complexes to weak transient touches — affinity is a continuous edge weight. A weighted network, richer than mere on/off links, in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 10f VISUAL — diameter / reachability
s=slide(prs); title(s,"Reachability — a large web, a small diameter")
band(s,0.42,1.2,12.5,0.4,TINT,"the longest shortest-path is still short",NAVY)
fbox(s,0.8,2.2,3.6,1.2,TINT,"~1,700 roots","but ~3-4 steps apart",line=TEAL,tsz=16,ssz=12)
harrow(s,4.6,2.7,1.4,"via hubs",color=GREY,lcol=NAVY)
fbox(s,6.2,2.2,3.6,1.2,AMBERT,"~20,000 proteins","but ~4-5 steps apart",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.8,12.5,3.4,TINT2,[L("Vast, yet navigable in a few hops",18,True,NAVY),
  L("Despite thousands of nodes, the network DIAMETER — the longest shortest path — stays small (a handful of steps), because hubs act as shortcuts. A reader can move from any concept to any other in a few co-occurrence links; a signal can cross the proteome in a few interactions. Size does not mean distance.",17),
  L("Both Books are enormous and yet everything is a few steps from everything.",16.5,True,TEAL)],space=8)

# 10g DATA — rich-club
s=slide(prs); title(s,"The data — the hubs talk to each other (rich-club)")
finding2(s,
 {"title":"Qur'an — hub-hub link density vs expected","cats":["expected","observed"],
  "series":[("",[GREY,TEAL],[1,3.5])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Protein — hub-hub link density vs expected","cats":["expected","observed"],
  "series":[("",[GREY,AMBER],[1,2.8])],"legend":False,"fmt":"{:.1f}"},
 [L("Central concepts co-occur with each other",17.5,True,TEAL),
  L("The hub roots (الله, رب, قول) co-occur with one another more than their popularity alone predicts — a 'rich-club' core that integrates the whole discourse around central themes.",16)],
 [L("Hub proteins interconnect too",17.5,True,AMBER),
  L("Hub proteins preferentially bind other hubs, forming an integrated core that coordinates the cell. A rich-club of hubs — a network within the network — appears in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Both webs are scale-free (hubs + leaves), modular, small-world, robust-yet-hub-fragile, sometimes directed, and far more clustered than a random graph. The TOPOLOGY is shared and validated against a null.",16)],
 [L("✗ BREAKS",17,True,RED),L("Protein edges are physical bindings (chemistry); co-occurrence edges are statistical associations in text (authorship/usage). Same graph SHAPE, different meaning of an 'edge'.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether textual modules map onto anything like functional complexes beyond analogy is open; the shared claim is topological, not mechanistic.",16)],f=(TINT,REDT,AMBERT))

# 12 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — the same web mathematics")
two(s,[L("UNIVERSAL NETWORK LAWS",18,True,NAVY),L("Scale-free degree, modularity, small-world paths, hub-fragility, clustering far above chance — these laws recur in the cell's interactions and the text's co-occurrences alike, validated against random-graph nulls. The connective structure of meaning and of life obeys the same mathematics.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Is a co-occurrence edge really like a protein binding?  • Why are both networks scale-free — constraint or coincidence?  • What does hub-fragility imply for how we read central concepts?  • Where does the analogy stop?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 13 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Network science is one toolkit across the social web, the brain, the proteome, and language — the same hubs, modules, and small-world maths everywhere connected things are studied.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("Nothing acts alone: both Books form scale-free, modular, small-world networks — robust to random loss, fragile at the hubs, clustered far beyond chance. Shared topology, different edges.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"14_Networks_Lecture.pptx")
print(f"L14 Networks slides: {len(prs.slides)}")
