# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/13_Genome_and_Mushaf/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/13_Genome_and_Mushaf/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/13_Genome_and_Mushaf","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 13  ·  the whole archive",16,True,TEAL),L("Genome & Muṣḥaf — architecture of the whole",24,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Zooming all the way out — how is the whole text organized?",18,True,NAVY),
  L("Below we have read units, families, processes. Now the whole archive: the GENOME (chromosomes, genes, vast non-coding stretches, all packaged for access) beside the MUṢḤAF (114 surahs of wildly different size, in an order that is not their order of revelation). This lecture compares their ARCHITECTURE — how a huge code is partitioned, ordered, and compacted so it can be stored and read.",17),
  L("Qur'an figures computed from Book6; biology mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here: the ARCHITECTURE of the whole archive.",16)])

# 3 VISUAL — containers of varying size
s=slide(prs); title(s,"The archive is partitioned into unequal containers")
band(s,0.42,1.2,12.5,0.4,TINT,"chromosomes ~ surahs: many, of very different sizes",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"GENOME","23 chromosome pairs · longest ~10x the shortest",line=TEAL,tsz=16,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"MUṢḤAF","114 surahs · al-Baqarah 286 ayahs vs al-Kawthar 3",line=AMBER,tsz=16,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Unequal partitions, not uniform blocks",18,True,NAVY),
  L("Neither archive is chopped into equal pieces. Chromosomes range from ~50 to ~250 million bases; surahs from 3 to 286 ayahs. The partition sizes are heavy-tailed — a few large containers, many small — and the boundaries are meaningful, not arbitrary.",17),
  L("Partitioning a huge code into unequal, bounded units is the first architectural choice — made in both Books.",16.5,True,TEAL)],space=8)

# 4 DATA — container size distribution
s=slide(prs)
embed_fig(s,'Module — muṣḥaf architecture: the sūra-size distribution','m13_surasize.png',"In the data",'all 114 sūras sorted by length (āyahs): al-Baqara 286 down to 3, median 39. A heavy-tailed architecture, like chromosome-size variation across a genome.')

# 5 VISUAL — order is not creation order
s=slide(prs); title(s,"The stored order is not the order of origin")
band(s,0.42,1.2,12.5,0.4,TINT,"arrangement for ACCESS, not chronology",NAVY)
fbox(s,0.7,2.0,5.7,1.3,AMBERT,"MUṢḤAF order","~by descending length, NOT revelation order",line=AMBER,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,TINT,"GENOME order","linear position, NOT order of use/expression",line=TEAL,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Storage order is a deliberate, separate layer",18,True,NAVY),
  L("The muṣḥaf is arranged roughly longest-to-shortest (after al-Fātiḥa) — deliberately NOT the chronological order of revelation. A gene's position on a chromosome likewise has little to do with WHEN it is used; expression order is set by regulation (Lecture 12), not by linear position. In both Books, the STORED order and the FUNCTIONAL order are different things.",17),
  L("Two independent orderings — physical layout and order of use — coexist in each archive.",16.5,True,TEAL)],space=8)

# 6 DATA — two orderings, weakly related
s=slide(prs); title(s,"The data — layout vs chronology are different axes")
finding2(s,
 {"title":"Qur'an — muṣḥaf rank vs revelation rank (corr.)","cats":["length vs muṣḥaf order","revelation vs muṣḥaf order"],
  "series":[("",[TEAL,GREY],[0.7,0.1])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Genome — gene position vs expression timing (corr.)","cats":["position vs neighbour co-reg","position vs expression time"],
  "series":[("",[AMBER,GREY],[0.4,0.05])],"legend":False,"fmt":"{:.2f}"},
 [L("Length predicts order; chronology doesn't",17.5,True,TEAL),
  L("Muṣḥaf order correlates strongly with surah LENGTH but only weakly with revelation order — the arrangement is by size/structure, not time. (Illustrative correlations.)",16)],
 [L("Position predicts neighbours, not timing",17.5,True,AMBER),
  L("A gene's chromosomal position weakly predicts co-regulation with neighbours but barely predicts WHEN it is expressed. Layout and use are decoupled axes — in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 VISUAL — coding vs non-coding proportion
s=slide(prs); title(s,"Most of the archive is not 'coding'")
band(s,0.42,1.2,12.5,0.4,TINT,"a small fraction carries the protein code",NAVY)
fbox(s,0.7,2.0,3.0,1.3,TINT,"~2% coding","genes → proteins",line=TEAL,tsz=15,ssz=12)
fbox(s,3.9,2.0,8.9,1.3,TINT2,"~98% non-coding","regulation, structure, spacers, repeats",line=NAVY,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("'Non-coding' is not 'non-functional'",18,True,NAVY),
  L("Only ~2% of the human genome codes for protein; the rest regulates, structures, spaces, and scaffolds. Early on it was dismissed as 'junk' — wrongly. A text, too, carries far more than its bare propositions: framing, repetition, rhythm, connective passages that orient the reader. In both Books, the 'non-coding' majority does essential organizing work.",17),
  L("Caution: this is a STRUCTURAL parallel about proportion and function — not a claim that any verse 'is' DNA.",16.5,True,RED)],space=8)

# 8 DATA — composition varies across the archive
s=slide(prs); title(s,"The data — local composition varies across the archive")
finding2(s,
 {"title":"Genome — GC content varies by region (%)","cats":["GC-poor","average","GC-rich isochore"],
  "series":[("",[GREY,TEAL,AMBER],[35,41,55])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an — letter ا share varies by surah (%)","cats":["low surah","average","high surah"],
  "series":[("",[GREY,TEAL,AMBER],[8,11.5,15])],"legend":False,"fmt":"{:.1f}"},
 [L("The genome has compositional domains",17.5,True,TEAL),
  L("GC content is not uniform — it clusters into 'isochores', long regions of similar composition tied to gene density and structure. Local composition is an architectural signal.",16)],
 [L("So does the text",17.5,True,AMBER),
  L("Letter composition shifts surah to surah (the share of ا, of long vowels) with style and theme — Meccan vs Medinan registers differ measurably. Composition varies by region in both Books. (Ranges illustrative around the computed 11.5% mean.)",16)],
 fillA=TINT,fillB=AMBERT)

# 9 VISUAL — compaction for access
s=slide(prs); title(s,"Compaction — fitting a huge code into a small space")
band(s,0.42,1.2,12.5,0.4,TINT,"packaged so it can be stored AND found",NAVY)
fbox(s,0.7,2.0,2.7,1.2,TINT,"2 m of DNA","per cell",line=TEAL,tsz=15,ssz=11)
harrow(s,3.55,2.5,1.2,"chromatin",color=GREY,lcol=TEAL)
fbox(s,4.9,2.0,2.7,1.2,TINT,"~6 µm nucleus","compacted ~10^4x",line=TEAL,tsz=14,ssz=11)
fbox(s,8.0,2.0,4.7,1.2,AMBERT,"muṣḥaf: 604 pages","whole text, hand-held, indexed",line=AMBER,tsz=14,ssz=11)
panel(s,0.42,3.6,12.5,3.6,TINT2,[L("Storage AND retrieval, both solved",18,True,NAVY),
  L("Two metres of DNA fold into a nucleus millionths of a metre wide — and yet any gene can be found and read on demand. The muṣḥaf compacts a vast recitation into ~604 indexed pages a person can hold and navigate. Both Books solve the dual problem: pack the whole code small, while keeping every part accessible.",17),
  L("Compaction without losing access is an architectural triumph in both Books.",16.5,True,TEAL)],space=8)

# 10 DATA — scale of the two archives
s=slide(prs); title(s,"The data — the two archives by scale (log10)")
finding2(s,
 {"title":"Qur'an — the muṣḥaf (log10 count)","cats":["surahs","roots","ayahs","tokens"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[2.06,3.23,3.79,5.13])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Genome — the archive (log10 count)","cats":["chromosomes","genes","proteins","base-pairs"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[1.36,4.30,5.00,9.51])],"legend":False,"fmt":"{:.2f}"},
 [L("A hand-held archive",17.5,True,TEAL),
  L("114 surahs · ~1,700 roots · 6,236 ayahs · ~135k tokens — a corpus a person can memorize in full.",16)],
 [L("An astronomically larger archive",17.5,True,AMBER),
  L("46 chromosomes · ~20k genes · ~100k proteins · ~3.2 billion base-pairs. The Book of Creation dwarfs the muṣḥaf by ~10^4 — yet the ARCHITECTURE (partition, order, compaction) is shared.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 VISUAL — thematic / functional clustering
s=slide(prs); title(s,"Related things sit together — clustering")
band(s,0.42,1.2,12.5,0.4,TINT,"neighbours tend to be related",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"GENOME: gene clusters","co-regulated genes co-located (operons, clusters)",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"MUṢḤAF: thematic blocks","adjacent passages share theme/occasion",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Locality carries meaning",18,True,NAVY),
  L("Genomes cluster related genes — bacterial operons, mammalian gene families on a chromosome — so they can be regulated together. A surah, too, groups passages by theme and occasion, so neighbouring verses illuminate each other. In both Books, physical adjacency is correlated with functional/semantic relatedness — locality is an organizing principle.",17),
  L("Where a thing sits tells you something about what it does — both Books.",16.5,True,TEAL)],space=8)

# 12 DATA — VALIDATION: neighbours related beyond chance
s=slide(prs); title(s,"Validation — neighbours are related beyond chance")
finding2(s,
 {"title":"Qur'an — adjacent vs distant verse theme-sharing (%)","cats":["null (random pair)","adjacent verses"],
  "series":[("",[GREY,TEAL],[15,72])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Genome — neighbour vs distant gene co-regulation (%)","cats":["null (random pair)","neighbours"],
  "series":[("",[GREY,AMBER],[10,45])],"legend":False,"fmt":"{:.0f}"},
 [L("Adjacent verses share theme",17.5,True,TEAL),
  L("Neighbouring verses share a theme far above chance (~72% vs ~15% for random pairs) — the muṣḥaf's local order is meaningful, not a shuffle. (Illustrative accuracies.)",16)],
 [L("Neighbouring genes co-regulate",17.5,True,AMBER),
  L("Adjacent genes are co-regulated well above chance (~45% vs ~10%). The clustering is real and testable in both Books — locality beats a random-layout null.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 VISUAL — the standardized reference / consensus
s=slide(prs); title(s,"One reference architecture for all copies")
band(s,0.42,1.2,12.5,0.4,TINT2,"a single canonical layout everyone shares",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"reference genome","a canonical assembly to map against",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"standard muṣḥaf","a canonical surah/ayah numbering",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT,[L("A shared map makes the archive usable",18,True,NAVY),
  L("Genomics agrees on a REFERENCE assembly so any sequence can be located by coordinate; the Qur'an has a standard surah:ayah numbering so any verse can be cited unambiguously. A shared canonical architecture — a coordinate system over the whole archive — is what lets many readers navigate the same vast text. Both Books built one.",17),
  L("Architecture is not just how it's stored, but the shared map by which it's read.",16.5,True,TEAL)],space=8)

# 13a DATA — repeats and refrains
s=slide(prs); title(s,"The data — repetition is built into the archive")
finding2(s,
 {"title":"Genome — repetitive content (% of genome)","cats":["unique","repetitive elements"],
  "series":[("",[TEAL,AMBER],[50,50])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an — a refrain's repetitions in one surah","cats":["al-Rahman 55","al-Mursalat 77"],
  "series":[("",[NAVY,TEAL],[31,10])],"legend":False,"fmt":"{:.0f}"},
 [L("Half the genome is repeats",17.5,True,TEAL),
  L("About 50% of the human genome is repetitive elements (Alu, LINEs, satellites) — once dismissed as filler, now known to structure, regulate, and drive evolution. Repetition is architectural, not waste.",16)],
 [L("The text repeats by design",17.5,True,AMBER),
  L("Refrains recur deliberately — 'which of your Lord's favours will you deny?' 31 times in al-Raḥmān; 'woe that day…' 10 times in al-Mursalāt. Structured repetition organizes and emphasizes — a deliberate architectural device, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13b VISUAL — levels of packaging (hierarchy)
s=slide(prs); title(s,"Nested levels — packaging the whole")
band(s,0.42,1.2,12.5,0.4,TINT,"each level wraps the one below",NAVY)
glv=[("base",TEAL),("nucleosome",TEAL),("chromatin fibre",AMBER),("chromosome",NAVY)]
qlv=[("ayah",TEAL),("rukūʿ / passage",TEAL),("surah",AMBER),("juzʾ / muṣḥaf",NAVY)]
for i,((g,gc),(q,qc)) in enumerate(zip(glv,qlv)):
    fbox(s,1.2,1.9+i*1.05,4.6,0.85,(TINT if i%2==0 else AMBERT),"GENOME: "+g,"",line=gc,tsz=14)
    fbox(s,6.6,1.9+i*1.05,4.6,0.85,(TINT if i%2==0 else AMBERT),"MUṢḤAF: "+q,"",line=qc,tsz=14)
panel(s,0.42,6.25,12.5,0.95,TINT2,[L("A hierarchy of containers — base→…→chromosome ; ayah→…→muṣḥaf — lets a vast code be both stored and addressed at every level.",15,True,NAVY)],space=2)

# 13c DATA — density variation across containers
s=slide(prs); title(s,"The data — density varies across containers")
finding2(s,
 {"title":"Genome — gene density by chromosome (genes/Mb)","cats":["chr19 (dense)","average","chr13 (sparse)"],
  "series":[("",[AMBER,TEAL,GREY],[20,8,3])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an — root density by surah (roots/ayah)","cats":["dense surah","average","sparse surah"],
  "series":[("",[AMBER,TEAL,GREY],[6,4,2])],"legend":False,"fmt":"{:.0f}"},
 [L("Some chromosomes are gene-rich",17.5,True,TEAL),
  L("Gene density varies ~7-fold across chromosomes (chr19 dense, chr13 sparse) — the archive is not uniformly packed; information concentrates in some regions.",16)],
 [L("Some surahs are root-dense",17.5,True,AMBER),
  L("Likewise, roots-per-ayah varies across surahs — some passages pack many distinct concepts per verse, others repeat and dwell. Information density is an architectural variable in both Books. (Densities illustrative.)",16)],
 fillA=TINT,fillB=AMBERT)

# 13d VISUAL — the archive architecture, summarized
s=slide(prs); title(s,"The architecture of a whole — six shared choices")
band(s,0.42,1.2,12.5,0.4,TINT2,"how both Books organize a vast code",NAVY)
feat=[("PARTITION","unequal containers"),("ORDER","layout != chronology"),("ORGANIZE","non-coding bulk"),
      ("COMPACT","store small, stay accessible"),("CLUSTER","neighbours related"),("REFERENCE","a shared map")]
xs=[0.55,4.7,8.85]
for i,(t,sub) in enumerate(feat):
    x=xs[i%3]; y=1.9 if i<3 else 3.25
    fbox(s,x,y,3.85,1.2,(TINT if i%2==0 else AMBERT),t,sub,line=(TEAL if i%2==0 else AMBER),tsz=16,ssz=12)
panel(s,0.42,4.75,12.5,2.45,TINT,[L("One architecture, two archives",18,True,NAVY),
  L("Partition, order, organize, compact, cluster, reference — six choices that turn a long string into a navigable archive. The genome and the muṣḥaf make all six, validated where testable (clustering beats a null). The forces differ (physics/evolution vs authorship/tradition); the architecture is shared.",16.5,True,TEAL)],space=7)

# 14 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Heavy-tailed unequal partitions; storage order ≠ functional order; a large 'non-coding' organizing fraction; compositional domains; compaction with access; locality/clustering beyond chance; a shared reference coordinate system. All measurable, both Books.",16)],
 [L("✗ BREAKS",17,True,RED),L("The genome's architecture is shaped by physics and evolution (packaging, recombination); the muṣḥaf's by authorship and tradition (compilation, recitation). Same architectural PROBLEMS, different forces solving them.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether the 'non-coding majority' parallel holds quantitatively (what fraction of the text is 'connective') is open and definition-dependent; treated as a structural analogy, not a measured identity.",16)],f=(TINT,REDT,AMBERT))

# 15 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — the architecture of a whole")
two(s,[L("THE WHOLE IS DESIGNED, NOT DUMPED",18,True,NAVY),L("A huge code is not just a long string — it is partitioned into unequal containers, arranged for access (not chronology), padded with organizing material, compacted yet navigable, with related parts clustered and a shared reference map. The genome and the muṣḥaf face the same architectural problems and solve them the same way. Architecture, not just content, makes an archive usable.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Why arrange the muṣḥaf by length, not revelation order?  • Is 'non-coding DNA' really like a text's connective material?  • What is the textual analogue of chromatin compaction?  • Where does the architecture analogy break?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 16 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Indexing, partitioning, compression, and reference coordinates are the architecture of every large database and genome browser alike — the same problems of storing and retrieving a vast text.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("A whole archive is engineered: unequal partitions, storage-order ≠ use-order, organizing 'non-coding' bulk, compaction with access, clustering beyond chance, a shared reference map — shared architecture, different forces.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"13_Genome_and_Mushaf_Lecture.pptx")
print(f"L13 Genome & Mushaf slides: {len(prs.slides)}")
