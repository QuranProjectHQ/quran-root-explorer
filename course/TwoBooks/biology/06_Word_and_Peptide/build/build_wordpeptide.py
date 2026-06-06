# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/06_Word_and_Peptide/build")
from st_slides import *
from diagrams import fbox,harrow,band,chain
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/06_Word_and_Peptide/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/06_Word_and_Peptide","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 6  ·  the polymer level",16,True,TEAL),L("Word & Peptide — sequence in the longer chain",26,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Pillar 2, scaled up — from the triplet to the whole chain",18,True,NAVY),
  L("A protein is a linear chain of amino acids from a ~20-unit alphabet; its 'primary structure' is just the SEQUENCE. A Qur'anic word is a linear chain of letters from a ~28-letter alphabet. Lecture 5 proved order matters in the triplet; here we read both as longer STRINGS — and the same string tools (alignment, edit distance, motifs) apply to each.",17),
  L("Still below full meaning; biology figures mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language. The Book of SCRIPTURE — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation. The Book of CREATION — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source (Allah); both are āyāt. Here we read both as linear CHAINS.",16)])

# 3 VISUAL — two chains
s=slide(prs); title(s,"Two linear polymers")
band(s,0.42,1.18,12.5,0.42,TINT,"PROTEIN — a chain of amino acids (~20-letter alphabet)",TEAL)
chain(s,1.0,1.95,["L","A","G","V","E","S","K","D","Y"],col=TEAL)
band(s,0.42,3.4,12.5,0.42,AMBERT,"WORD — a chain of letters (~28-letter alphabet)",AMBER)
chain(s,1.0,4.15,["ا","ل","ح","م","د","ل","ل","ه"],col=AMBER)
panel(s,0.42,5.4,12.5,1.8,TINT2,[L("Both are strings",17,True,NAVY),
  L("Read in ONE direction (N→C for proteins; right→left for Arabic), each is a sequence of units from a small alphabet. At this scale a word IS a string — and so is a peptide.",16.5,True,TEAL)],space=6)

# 4 DATA — composition (skewed), consistent %
s=slide(prs); title(s,"The data — composition is skewed in both")
finding2(s,
 {"title":"Letter frequency — Qur'an (% of letters)","cats":["ا","ل","ن","م","و"],
  "series":[("",[NAVY,TEAL,TEAL,AMBER,AMBER],[18.0,11.5,8.4,8.1,7.5])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Amino-acid frequency — proteins (%)","cats":["Leu","Ala","Gly","Val","Trp"],
  "series":[("",[TEAL,TEAL,AMBER,AMBER,RED],[9.7,8.3,7.1,6.9,1.1])],"legend":False,"fmt":"{:.1f}"},
 [L("A few units dominate",17.5,True,TEAL),
  L("Arabic letters are far from uniform: ا, ل, ن lead; others trail. A heavy-headed distribution over 28 letters.",16)],
 [L("So do amino acids",17.5,True,AMBER),
  L("Proteins are biased too — Leu/Ala/Gly common, Trp rare (~1%). Both alphabets have a skewed, non-random composition that shapes every chain.",16)],
 fillA=TINT,fillB=AMBERT)

# 5 DATA — length distributions
s=slide(prs)
embed_fig(s,'Real data — the full word-length distribution','m06_wordlen.png',"In the data",'letters per word-token across all 135,366 tokens (Book6 segmentation, mean 2.45): short function-words dominate, a tail of longer content-words — the same right-skewed length profile proteins show.')

# 6 DATA — Pillar 2 at word scale: anagram words
s=slide(prs); title(s,"Order at the word scale — anagram words")
finding2(s,
 {"title":"Word-forms sharing a letter-set (%)","cats":["unique letter-set","has anagram-sibling"],
  "series":[("",[GREY,TEAL],[60.9,39.1])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Protein: identical composition, different fold","cats":["composition","folds (order-dependent)"],
  "series":[("",[GREY,RED],[1,2])],"legend":False},
 [L("39% of words have an anagram-sibling",17.5,True,TEAL),
  L("Computed across 7,236 word-forms: 39.1% share their exact multiset of letters with another distinct word. Even at the whole-word scale, ORDER — not the bag of letters — fixes the word.",16)],
 [L("Same residues, different protein",17.5,True,AMBER),
  L("Two proteins can share amino-acid composition yet fold and function differently because their SEQUENCE differs. Composition is necessary, never sufficient — in both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 7 VISUAL — alignment
s=slide(prs); title(s,"Sequence alignment — the same operation, both Books")
band(s,0.42,1.18,12.5,0.42,TINT,"align two strings; mark matches and mismatches",TEAL)
chain(s,2.0,1.95,["ك","ت","ب"],col=TEAL); chain(s,2.0,2.95,["ك","ذ","ب"],col=AMBER)
fbox(s,6.6,1.95,5.6,1.6,TINT2,"2 of 3 positions match","a 'point' difference at position 2",line=NAVY,tsz=16,ssz=12)
panel(s,0.42,4.0,12.5,3.2,TINT,[L("Alignment is character-level, not semantic",18,True,NAVY),
  L("Comparing two protein sequences (alignment / BLAST) finds matched and mismatched positions — exactly what the course's root-normalization does when it lines up two strings. No meaning is consulted; only characters and positions.",17),
  L("This is the engine under the whole course: match strings first, interpret later.",16.5,True,TEAL)],space=9)

# 8 DATA — sequence identity & edit distance
s=slide(prs); title(s,"The data — identity and edit distance")
finding2(s,
 {"title":"Sequence identity needed to infer relatedness (%)","cats":["clearly related","twilight zone","unrelated"],
  "series":[("",[TEAL,AMBER,GREY],[35,20,10])],"legend":False},
 {"title":"Edit distance — كتب to related strings","cats":["كذب","كتاب","قلب"],
  "series":[("",[TEAL,AMBER,RED],[1,1,3])],"legend":False},
 [L("Bioinformatics has a threshold",17.5,True,TEAL),
  L("Two proteins above ~35% sequence identity are confidently homologous; ~20% is the 'twilight zone'. Relatedness is read off the STRING, quantitatively — the same logic the course uses on roots.",16)],
 [L("Edit distance counts the steps",17.5,True,AMBER),
  L("كتب→كذب is one edit; كتب→كتاب is one insertion; كتب→قلب is three. Levenshtein distance measures string similarity for words exactly as for sequences — one shared toolkit.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 DATA — local patterns: motifs
s=slide(prs); title(s,"Recurring sub-strings — motifs in both Books")
finding2(s,
 {"title":"Qur'an — a recurring formula (occurrences)","cats":["الذين آمنوا","بسم الله"],
  "series":[("",[TEAL,AMBER],[258,114])],"legend":False},
 {"title":"Protein — conserved sequence motifs (illustrative)","cats":["zinc finger","kinase","signal peptide"],
  "series":[("",[TEAL,AMBER,GREY],[3,2,1])],"legend":False},
 [L("Texts repeat sub-strings",17.5,True,TEAL),
  L("Fixed formulae recur across the Qur'an — e.g. 'those who believe' appears hundreds of times. Recurring sub-sequences are a measurable feature of the string, found by the same scan that finds protein motifs. (Counts approximate.)",16)],
 [L("Proteins repeat motifs too",17.5,True,AMBER),
  L("Short conserved motifs (a zinc finger, a kinase site, a signal peptide) recur across unrelated proteins and signal function. Motif discovery is one algorithm, two Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 10 VISUAL — secondary structure ~ template patterns (awzan)
s=slide(prs); title(s,"Local templates — awzan and secondary structure")
band(s,0.42,1.2,12.5,0.4,TINT,"a fixed TEMPLATE shapes the chain, both Books",NAVY)
fbox(s,0.7,1.95,5.7,1.5,AMBERT,"ARABIC wazn: فَاعِل","template C-a-C-i-C → فاعل, كاتب, ناصر",line=AMBER,tsz=16,ssz=12)
fbox(s,6.9,1.95,5.9,1.5,TINT,"PROTEIN: α-helix / β-sheet","local sequence → a repeating fold motif",line=TEAL,tsz=16,ssz=12)
panel(s,0.42,3.75,12.5,3.45,TINT2,[L("A pattern imposed on a sequence",18,True,NAVY),
  L("Arabic morphology pours a root into a fixed vowel-consonant TEMPLATE (wazn) — فاعل, مفعول, استفعال — generating a family of words with shared shape. A protein's local sequence likewise prefers a repeating secondary-structure template (helix, sheet).",17),
  L("In both, a small set of templates patterns the chain into recognizable, reusable shapes.",16.5,True,TEAL)],space=8)

# 11 DATA — the templates by the numbers
s=slide(prs); title(s,"The data — a few templates, much of the output")
finding2(s,
 {"title":"Arabic — common patterns (share of forms, approx.)","cats":["فاعِل","مفعول","فعّال","other"],
  "series":[("",[TEAL,AMBER,NAVY,GREY],[22,15,8,55])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein — secondary structure (% of residues)","cats":["helix","sheet","coil"],
  "series":[("",[TEAL,AMBER,GREY],[34,21,45])],"legend":False},
 [L("A handful of awzan recur",17.5,True,TEAL),
  L("A small inventory of patterns (فاعِل, مفعول, فعّال …) accounts for a large share of derived word-forms — templated reuse, like a design library. (Shares approximate.)",16)],
 [L("Two folds dominate the chain",17.5,True,AMBER),
  L("Across proteins, ~34% of residues sit in α-helices and ~21% in β-sheets — two templates pattern most of the structure. Few templates, much of the output, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 12 DATA — information per position (entropy)
s=slide(prs); title(s,"The data — not every position is equal")
finding2(s,
 {"title":"Arabic root — variability by position (distinct letters)","cats":["pos 1","pos 2","pos 3"],
  "series":[("",[TEAL,AMBER,TEAL],[26,27,25])],"legend":False},
 {"title":"Protein motif — conservation by position (bits)","cats":["core pos","variable pos"],
  "series":[("",[TEAL,GREY],[4.0,0.8])],"legend":False,"fmt":"{:.1f}"},
 [L("Positions carry different loads",17.5,True,TEAL),
  L("Across roots, each slot draws on a different sub-set of letters with different spread — the positions are not interchangeable, echoing Lecture 5's positional bias.",16)],
 [L("Conserved vs free positions",17.5,True,AMBER),
  L("In a protein motif, some positions are near-invariant (high information, ~4 bits) while others tolerate any residue (~0.8 bits). A sequence-logo measures exactly which positions matter — usable on Arabic patterns too.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 VISUAL — directionality & palindromes
s=slide(prs); title(s,"Direction matters — and palindromes in both")
band(s,0.42,1.2,12.5,0.4,TINT,"a string read backwards is a different string",NAVY)
fbox(s,0.7,1.95,5.7,1.5,AMBERT,"Arabic reads right → left","reverse a word → usually a non-word; some palindromes (e.g. رمر)",line=AMBER,tsz=15,ssz=12)
fbox(s,6.9,1.95,5.9,1.5,TINT,"DNA reads 5' → 3'","reverse-complement; palindromic sites (e.g. GAATTC) cut by enzymes",line=TEAL,tsz=15,ssz=12)
panel(s,0.42,3.75,12.5,3.45,TINT2,[L("Orientation is part of the information",18,True,NAVY),
  L("A protein chain has a direction (N→C); DNA has a strand (5'→3'); Arabic is read right-to-left. Reversing a sequence generally destroys it — but special PALINDROMES read the same both ways, and in DNA these mark real functional sites (restriction enzymes, regulatory elements).",17),
  L("Direction and symmetry are measurable string properties — shared between text and sequence.",16.5,True,TEAL)],space=8)

# 14 DATA — Zipf of words vs protein abundance
s=slide(prs); title(s,"The data — a few strings dominate usage")
finding2(s,
 {"title":"Qur'an — most frequent words (token count)","cats":["من","الله","في","لا","ما"],
  "series":[("",[NAVY,TEAL,TEAL,AMBER,AMBER],[2763,2699,1185,1124,1010])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Cell — protein abundance is heavy-tailed (log10 copies)","cats":["top","median","rare"],
  "series":[("",[NAVY,TEAL,GREY],[7,4,2])],"legend":False},
 [L("Word usage is heavy-tailed",17.5,True,TEAL),
  L("A few function words (من, الله, في …) dominate the token stream; thousands of words appear rarely — Zipf's law, the statistics of natural text.",16)],
 [L("So is protein abundance",17.5,True,AMBER),
  L("A cell holds millions of copies of a few proteins and a handful of others — abundance spans ~10^7. Both Books deploy their strings with the same steep, heavy-tailed economy.",16)],
 fillA=TINT,fillB=AMBERT)

# 15 VISUAL — the shared string toolkit
s=slide(prs); title(s,"One toolkit — strings, both Books")
band(s,0.42,1.2,12.5,0.4,TINT2,"the same algorithms read a word and a peptide",NAVY)
tools=[("ALIGN","match positions"),("EDIT DIST.","count changes"),("MOTIF SCAN","find sub-strings"),
       ("n-GRAM","local statistics"),("ENTROPY","info per position"),("SEARCH","find a query")]
xs=[0.55,4.7,8.85]
for i,(t,sub) in enumerate(tools):
    x=xs[i%3]; y=1.85 if i<3 else 3.25
    fbox(s,x,y,3.85,1.2,(TINT if i%2==0 else AMBERT),t,sub,line=(TEAL if i%2==0 else AMBER),tsz=16,ssz=12)
panel(s,0.42,4.75,12.5,2.45,TINT,[L("Bioinformatics = computational linguistics, applied to molecules",18,True,NAVY),
  L("Every method here — alignment, edit distance, motif discovery, n-gram models, entropy, search — was built for one kind of string and works on the other. The word and the peptide are the same KIND of object, so one toolbox serves both.",16.5,True,TEAL)],space=7)

# 16 DATA — composition is not sequence (longer chains)
s=slide(prs); title(s,"At every length — composition is not the message")
finding2(s,
 {"title":"Strings from one letter-multiset (4 letters)","cats":["composition","orderings (4!)"],
  "series":[("",[GREY,RED],[1,24])],"legend":False},
 {"title":"Orderings of a chain (log10)","cats":["3 units","6 units","20 units"],
  "series":[("",[TEAL,AMBER,RED],[0.78,2.86,18.4])],"legend":False,"fmt":"{:.1f}"},
 [L("More units, far more orderings",17.5,True,TEAL),
  L("Four letters admit 24 orderings; six admit 720; twenty admit ~10^18. As chains lengthen, the order-space explodes — and only a vanishing fraction are real words or working proteins.",16)],
 [L("So sequence carries almost all the information",17.5,True,AMBER),
  L("Composition fixes the parts; the astronomically larger choice of ORDER fixes the message. The longer the chain, the more decisively this is true — both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 16b DATA — searching a string against a database
s=slide(prs); title(s,"Finding a string — search in both Books")
finding2(s,
 {"title":"Qur'an concordance — a query's hits","cats":["الله","رحيم","صبر"],
  "series":[("",[NAVY,TEAL,AMBER],[2699,114,103])],"legend":False,"fmt":"{:.0f}"},
 {"title":"BLAST significance vs identity (-log10 E)","cats":["40% identity","25% identity"],
  "series":[("",[TEAL,GREY],[20,2])],"legend":False},
 [L("Concordance = ranked string search",17.5,True,TEAL),
  L("Querying a word returns every occurrence, ranked and located — a concordance. It is the textual twin of a database sequence search: same operation, different corpus.",16)],
 [L("BLAST scores a match against chance",17.5,True,AMBER),
  L("A sequence search reports an E-value — how many matches this good you'd expect by chance in a database of this size. Strong identity (~40%) gives a tiny E (significant); weak identity fades into noise. Significance, again, is judged against a null.",16)],
 fillA=TINT,fillB=AMBERT)

# 17 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Both are linear strings from a small alphabet, skewed composition, characteristic lengths, recurring motifs, templated patterns, and order-dominated information — all read by one string toolkit, no semantics.",16)],
 [L("✗ BREAKS",17,True,RED),L("Proteins FOLD into 3-D shapes that give function; a written word does not physically fold. Scale and chemistry differ entirely. The match is the STRING, not the substance.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether recitation's acoustic shape is the word's 'fold' is suggestive; the right 'alphabet' (28 letters? with/without vowels?) is a modeling choice.",16)],f=(TINT,REDT,AMBERT))

# 18 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — and the next tier")
two(s,[L("THE STRING, THEN THE MEANING",18,True,NAVY),L("At the character scale, a word and a peptide are the same kind of object: a string whose information lives in ORDER. The course's power comes from working here FIRST — matching, aligning, scanning strings — before any interpretation. Pillar 2 holds at every length.",17,True,TEAL)],
 [L("NEXT TIER (Lecture 7)",18,True,AMBER),L("Move up one level — root ↔ PROTEIN — and a semantic/functional dimension enters: a root carries a CONCEPT as a protein carries a FUNCTION via its fold. That lecture reads meaning; this one read only characters.",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 19 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("String and sequence methods — alignment, BLAST, motif discovery, NLP — are one toolkit across texts and proteins alike; what you learn on one transfers to the other.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("At every length a word and a peptide are the same kind of object: a linear string whose meaning lives in ORDER, not composition. One toolkit reads both.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"06_Word_and_Peptide_Lecture.pptx")
print(f"L6 Word & Peptide slides: {len(prs.slides)}")
