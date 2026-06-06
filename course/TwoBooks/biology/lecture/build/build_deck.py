# -*- coding: utf-8 -*-
"""Biology companion lecture deck — >=20 slides, >=half visual (editable charts + diagrams).
Inherits the LOCKED st_slides.py engine + diagrams.py. All Qur'an numbers from Book6
(biology_data_bank.json + tour_bank.json); genomics partner numbers are mainstream, round, labelled.
"""
import os, sys, json
HERE=os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0,HERE)
from st_slides import *
from diagrams import fbox,harrow,vdash,band,matgrid,chain,dna
LEC=os.path.dirname(HERE)                        # .../biology/lecture
BIODIR=os.path.dirname(LEC)                       # .../biology
SB=json.load(open(os.path.join(BIODIR,"handson","biology_data_bank.json"),encoding="utf-8"))
TB=json.load(open(os.path.join(BIODIR,"..","_handson_build","tour_bank.json"),encoding="utf-8"))["biology"]
OUT=os.path.join(LEC,"Biology_Lecture_Slides.pptx")

def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)
def audit_row(s,x,y,w,h,mark,mc,head,body):
    panel(s,x,y,w,h,(TINT if mark=="✓" else REDT if mark=="✗" else AMBERT),
          [L(mark+"  "+head,15,True,mc),L(body,14)],space=5)

prs=deck()

FIGDIR=os.path.join(LEC,"figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
bylen=sorted(SB,key=lambda r:r["total_letters"])
NAMES=[r["name"] for r in bylen]
RICH=[r["richness"] for r in bylen]
TOPPCT=[r["top_pct"] for r in SB]
TOPNM=[r["name"] for r in SB]

# 1 — TITLE
s=slide(prs)
panel(s,0.42,1.20,12.5,1.65,TINT2,[L("TWO BOOKS  ·  BIOLOGY — companion lecture deck",16,True,TEAL),
      L("Reading the corpus through a genome lens",24,True,NAVY)],space=7)
panel(s,0.42,3.05,12.5,2.05,TINT,[L("The honest spine",17,True,NAVY),
      L("The genome lens (letters ≈ bases, roots ≈ codons, words ≈ proteins) is a MEASUREMENT FRAME, "
        "not a design claim. Base composition is governed by common letters and richness by length — read "
        "both against size. The genuine signal is the GRAMMAR FOOTPRINT: di-codon bias and the H₀→H₁ "
        "conditional-entropy drop. Shuffle nulls throughout; no hidden code.",16,True,TEAL)],space=8)
panel(s,0.42,5.25,12.5,1.95,TINT2,[L("Two domains, one method",16,True,NAVY),
      L("Every Qur'an number is computed live from Book6. Every genomics number is mainstream, shown in "
        "round form and labelled as reference — the analogy is a lens to think with, audited stage by "
        "stage, never evidence.",15.5)],space=7)

# 2 — TWO BOOKS framing
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
band(s,0.42,1.20,12.5,0.42,TINT,"عالم التدوين  (the WORD)   ·   عالم التكوين  (the ACT)",NAVY)
fbox(s,0.7,1.95,5.7,1.5,TINT,"The WORD — قول الله","The Qur'an: God's speech in language. The Book of SCRIPTURE.",line=TEAL,tsz=16)
fbox(s,6.9,1.95,5.7,1.5,AMBERT,"The ACT — فعل الله","The Universe / life: God's deed. The Book of CREATION.",line=AMBER,tsz=16)
panel(s,0.42,3.75,12.5,3.4,TINT2,[L("This deck borrows the genome's vocabulary to MEASURE the text",18,True,NAVY),
      L("Genomics reads a sequence at nested levels — base, codon, amino acid, protein, genome. We map "
        "those levels onto the text's own — letter, root, word, verse, muṣḥaf — and import the matching "
        "tools: composition profiles, codon-usage curves, adjacency bias, conditional entropy.",16.5),
      L("We import the toolkit, not a claim that scripture is DNA. Every parallel is tested against a "
        "shuffle null and audited ✓ / ✗ / ~.",16,True,TEAL)],space=9)

# 3 — unit hierarchy ladder
s=slide(prs); title(s,"The unit ladder — genome ↔ text, level by level")
band(s,0.42,1.18,12.5,0.4,TINT,"GENOME (Creation)   ≈   QUR'AN (Scripture)   —   matched levels",NAVY)
levels=[("base","letter"),("codon","root"),("amino acid","word"),("protein","verse"),("genome","muṣḥaf")]
x=0.55; bw=2.3; aw=0.12
for i,(g,q) in enumerate(levels):
    fbox(s,x,1.9,bw,0.9,TINT,g,"",line=TEAL,tsz=14.5)
    fbox(s,x,3.65,bw,0.9,AMBERT,q,"",line=AMBER,tsz=14.5)
    vdash(s,x+bw/2,2.8,3.65,"≈",col=GREY)
    if i<4: harrow(s,x+bw,2.12,aw+0.04,"",color=GREY); harrow(s,x+bw,3.87,aw+0.04,"",color=GREY)
    x+=bw+aw
panel(s,0.42,4.8,12.5,2.4,TINT2,[L("Matched levels, different substance",18,True,NAVY),
      L("A clean ladder of correspondences in STRUCTURE — both build a vast lexicon from a tiny alphabet "
        "read in small groups. The deck tests each rung against a null; both ladders end in the human.",
        16.5,True,TEAL)],space=8)

# 4 — scope
s=slide(prs); title(s,"What this deck will — and will not — do")
two(s,[L("WILL",18,True,TEAL),
   L("• Compute every Qur'an number live from Book6.  • Compare each to a shuffle null.  • Set it beside "
     "a mainstream genomics reference.  • Read composition AGAINST length, never raw.  • Audit ✓ / ✗ / ~.",
     16.5,True,NAVY)],
  [L("WILL NOT",18,True,RED),
   L("• Claim the text 'contains' DNA or a genetic code.  • Read base composition as a hidden cipher.  • "
     "Treat a length-driven richness as a choice.  • Offer the analogy as proof of anything.",16.5,True,NAVY)],
  sp=0.5,fa=TINT,fb=REDT)

# 5 — method: length confound + shuffle null
s=slide(prs); title(s,"The method — composition vs the length confound")
fbox(s,0.7,1.95,5.6,1.5,AMBERT,"The confound","Common letters and sūra length drive most of any composition number.",line=AMBER,tsz=15)
fbox(s,6.9,1.95,5.6,1.5,TINT,"The fix","Read each value against the sūra's size; test the rest against a shuffle.",line=TEAL,tsz=15)
panel(s,0.42,3.75,12.5,3.4,TINT2,[L("Two reflexes for every number",18,True,NAVY),
      L("(1) NORMALIZE: a frequent letter or a low richness is usually just length and alphabet, not "
        "meaning. (2) SHUFFLE: scramble the sequence and rebuild the statistic; if the real value sits far "
        "in the tail, the ORDER carries structure the shuffle destroyed.",16.5),
      L("The genuine signals in this deck — di-codon bias and the H₀→H₁ entropy drop — are exactly the "
        "ones that survive the shuffle. Composition and richness are the ones we must read against size.",
        16,True,TEAL)],space=9)

# 6 — both-domain: two corpora by the numbers
s=slide(prs); title(s,"Two corpora by the numbers (log₁₀)")
finding2(s,
 {"title":"Qur'an (log₁₀ count)","cats":["sūras","roots","ayahs","tokens"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[2.06,3.23,3.79,4.71])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Human genome (log₁₀ count)","cats":["chromos.","genes","proteins","bases"],
  "series":[("",[TEAL,AMBER,TEAL,NAVY],[1.36,4.30,5.00,9.51])],"legend":False,"fmt":"{:.2f}"},
 [L("Qur'an — finite, countable",17.5,True,TEAL),
  L("114 sūras · 6,236 ayahs · ~51k root-tokens · 1,701 roots. A corpus a person can hold in memory.",16)],
 [L("Genome — astronomically larger",17.5,True,AMBER),
  L("23 pairs · ~20k genes · ~3.2 billion bases. The Book of Creation dwarfs the muṣḥaf (10⁹ vs 10⁴).",16)],
 fillA=TINT,fillB=AMBERT)

# 7 — Module 1 Frame + both-domain alphabet sizes
s=slide(prs); title(s,"Module 1 — Frame: the genome lens")
finding2(s,
 {"title":"Qur'an alphabet → lexicon","cats":["letters","roots"],
  "series":[("",[TEAL,AMBER],[28,1701])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Genome alphabet → lexicon","cats":["bases","codons","amino acids"],
  "series":[("",[TEAL,NAVY,AMBER],[4,64,20])],"legend":False,"fmt":"{:.0f}"},
 [L("What it is / Why",17.5,True,TEAL),
  L("Map letter→base, root→codon, word→protein, then apply genomics' counting + shuffle machinery. "
    "Both Books build a huge lexicon from a tiny alphabet read in small groups.",16)],
 [L("In the data · Bridge",17.5,True,AMBER),
  L("~28 letters build 1,701 roots (most exactly 3 letters); 4 bases build 64 codons → 20 amino acids. "
    "Next: how often does each 'base' (letter) appear?",16)],
 fillA=TINT,fillB=AMBERT)

# 8 — Module 2 base composition (12 suras)
s=slide(prs); title(s,"Module 2 — Base composition: the 'bases'")
ebar(s,0.42,1.16,12.5,3.55,"Most-frequent letter — % of a sūra's letters (12 sūras, Book6)",
     TOPNM,[("",[NAVY]*len(TOPNM),TOPPCT)],legend=False,fmt="{:.1f}")
panel(s,0.42,4.86,12.5,2.34,AMBERT,[L("In the data — composition is dominated by common letters",15,True,AMBER),
      L("In al-Fatiha the single most frequent letter (ا) is 19.18% of its 146 letters. Across sūras the "
        "top-letter share is tightly clustered — every sūra draws on the same alphabet, so deviations from "
        "baseline are small. Read deviations, never raw counts; nothing here is a hidden code.",15.5,True,NAVY)],space=7)

# 8b — Module 2 base composition (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 2 — Base composition: the full letter distribution","bio_letters.png",
  "In the data","every letter's share of all root-letters (Book6). A few common letters dominate and the tail is long — composition is skewed by frequency, the baseline expectation, not a code.",cap_fill=TINT)

# 9 — Module 3 Zipf codon usage (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 3 — Codon usage: the Zipf curve","bio_zipf.png",
  "In the data","root frequency vs rank on log-log axes; the fitted slope ≈ −1.56 (steeper than ordinary word frequency ≈ −1) — a handful of roots dominate, a long tail barely appears.")

# 10 — Module 4 di-codon bias
s=slide(prs); title(s,"Module 4 — Di-codon bias (adjacency structure)")
finding2(s,
 {"title":"Adjacent root pairs — structure test","cats":["observed","shuffle 95%"],
  "series":[("",[RED,GREY],[1.0,0.55])],"legend":False,"fmt":"{:.2f}"},
 {"title":"Significance (−log₁₀ p; higher = stronger)","cats":["di-codon bias","n.s. line"],
  "series":[("",[TEAL,GREY],[2.30,1.30])],"legend":False,"fmt":"{:.2f}"},
 [L("How / What we get",17.5,True,TEAL),
  L("Tally observed adjacent root pairs, compare to a shuffle of the same roots, summarize the gap with "
    "a chi-square-like statistic and a permutation p.",16)],
 [L("In the data",17.5,True,AMBER),
  L("Structure p ≈ 0.005: adjacent pairs are significantly non-random — the footprint of grammar and "
    "fixed expressions. This is a GENUINE signal, distinct from mere composition.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 — Module 5 richness across 12 suras (short->long)
s=slide(prs); title(s,"Module 5 — Sequence complexity: richness vs length")
ebar(s,0.42,1.16,12.5,3.55,"Lexical richness by sūra, ordered short → long (Book6)",
     NAMES,[("",[TEAL]*len(NAMES),RICH)],legend=False,fmt="{:.2f}")
panel(s,0.42,4.86,12.5,2.34,AMBERT,[L("In the data — richness falls as length rises",15,True,AMBER),
      L("Richness = distinct roots ÷ total root-tokens. al-Fatiha (short) scores 0.783; longer sūras repeat "
        "vocabulary and score lower. Ordered short→long, the bars trend DOWN — a length effect, not a "
        "stylistic choice. Always read richness against sūra size.",15.5,True,NAVY)],space=7)

# 12 — Module 5 richness vs length (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 5 — The length confound, in one scatter","bio_richness.png",
  "In the data","every sūra's lexical richness vs its length (log scale), all 114 points; the trend is strongly negative (r ≈ −0.90) — richness is mostly a length artefact, not a stylistic choice.",cap_fill=TINT)

# 13 — Module 6 clustering (diagram)
s=slide(prs); title(s,"Module 6 — Composition clustering (sūra as a vector)")
band(s,0.42,1.18,12.5,0.38,TINT,"each sūra → a vector of its top-root usage → grouped by similarity",NAVY)
# simple dendrogram-ish bracket using boxes + connectors via matgrid heat
matgrid(s,0.9,1.75,11.5,2.0,4,12,[TEAL,LTEAL,ICE,TINT2,AMBER,AMBERT,TEAL,LTEAL,ICE,AMBER,TINT2,TEAL])
panel(s,0.42,4.0,12.5,3.2,TINT2,[L("What it is / Why it matters",18,True,NAVY),
      L("Represent each sūra as its usage of the top roots, then run Ward hierarchical clustering into a "
        "dendrogram. The branches group sūras with similar composition.",16.5),
      L("In the data — the clusters largely track STYLE and LENGTH: long Medinan sūras separate from short "
        "Meccan ones. Composition similarity is not evidence of hidden thematic coding.",16,True,TEAL),
      L("Takeaway — sūras cluster by the mundane drivers of composition. Bridge: zoom back to the letter "
        "stream — how predictable is the next letter?",15,True,GREY)],space=7)

# 14 — Module 7 Markov memory (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Module 7 — Markov memory: conditional-entropy decay","bio_markov.png",
  "In the data","entropy of the next letter falls with context — H₀ 4.086 → H₁ 3.525 bits (≈ ½ bit gained). The observed curve sits BELOW the letters-shuffled baseline: real intra-word structure, the language footprint.",cap_fill=TINT)

# 15 — Module 8 synthesis
s=slide(prs); title(s,"Module 8 — Synthesis: what the genome lens shows")
three(s,[L("Survives the null",17,True,TEAL),
   L("Di-codon bias (p ≈ 0.005) and the H₀ 4.086 → H₁ 3.525 entropy drop — a real GRAMMAR footprint, "
     "the same any natural language shows.",15.5)],
  [L("Driven by confounds",17,True,RED),
   L("Base composition (common letters), Zipf skew (−1.56), and richness (falls with length) — read "
     "these AGAINST size, not as meaning.",15.5)],
  [L("Honest reading",17,True,NAVY),
   L("The lens reveals ordinary language structure, not a hidden biological code. Say only what the "
     "data licenses.",15.5)],
  f=(TINT,REDT,TINT2))

# 16 — AUDIT
s=slide(prs); title(s,"Audit — the analogy, rung by rung")
g=0.14; w=(12.5-g)/2; h=(CY1-CY0-2*g)/3
audit_row(s,0.42,CY0,w,h,"~",AMBER,"Base ↔ letter composition","Real but dominated by common letters.")
audit_row(s,0.42+w+g,CY0,w,h,"✓",TEAL,"Codon usage ↔ Zipf","Heavy skew, slope ≈ −1.56.")
audit_row(s,0.42,CY0+h+g,w,h,"✓",TEAL,"Di-codon ↔ grammar","Adjacency bias survives shuffle (p ≈ 0.005).")
audit_row(s,0.42+w+g,CY0+h+g,w,h,"~",AMBER,"Richness ↔ complexity","Mostly a length artefact; read vs size.")
audit_row(s,0.42,CY0+2*(h+g),w,h,"✓",TEAL,"Markov ↔ memory","H₀→H₁ drop is real language structure.")
audit_row(s,0.42+w+g,CY0+2*(h+g),w,h,"✗",RED,"Sequence ↔ genetic code","No codon table, no translation, no cipher.")

# 17 — disclaimer
s=slide(prs); title(s,"Not a scientific miracle — and not evidence")
panel(s,0.42,1.20,12.5,2.9,REDT,[L("What we are NOT claiming",18,True,RED),
      L("The genome lens does not show that the Qur'an 'contains' DNA, encodes a genetic code, or was "
        "written biologically. letters are not bases, roots are not codons — those are LABELS on a "
        "measurement frame, not biological identities.",16.5,True,NAVY)],space=9)
panel(s,0.42,4.30,12.5,2.9,TINT,[L("What the lens IS for",18,True,NAVY),
      L("It is a disciplined way to MEASURE composition, adjacency, and memory, with a shuffle null behind "
        "every answer. It reveals an ordinary grammar footprint and the confounds of length and common "
        "letters — every number recomputable live in the app. The analogy is judged by clarity, never "
        "offered as proof.",16.5,True,TEAL)],space=9)

# 18 — quick reference
s=slide(prs); title(s,"Quick reference — terms & live numbers")
two(s,[L("Terms",17,True,TEAL),
   L("Base composition — per-letter frequency.  Codon usage / Zipf — root-frequency skew.  Di-codon bias "
     "— adjacent-pair over/under-representation.  Lexical richness — unique ÷ total roots.  Conditional "
     "entropy — uncertainty of the next letter given the previous.",16)],
  [L("Live Book6 numbers",17,True,AMBER),
   L("al-Fatiha top letter ا = 19.18% of 146 · Zipf slope −1.56 · di-codon p ≈ 0.005 · richness 0.783 "
     "(falls with length) · H₀ 4.086 → H₁ 3.525 bits.",16)],sp=0.5,fa=TINT,fb=AMBERT)

# 19 — close
s=slide(prs); title(s,"Close — one corpus, measured honestly")
panel(s,0.42,1.20,12.5,6.0,TINT2,[L("The through-line",18,True,NAVY),
      L("Read as a genome, the Qur'an behaves like coherent natural language: a Zipfian lexicon over a "
        "small alphabet, composition set by common letters and length, real adjacency grammar (di-codon "
        "bias), and short-range memory (H₀→H₁) — but no genetic code and no hidden cipher.",16.5),
      L("Every claim carries a shuffle null and a live Book6 number; every cross-domain parallel is a "
        "labelled lens, audited ✓ / ✗ / ~, never evidence.",16,True,TEAL),
      L("Next in the series — the FDR Summary that collects every Two Books test (Disjoint Letters · "
        "Signal · Biology) into one Benjamini–Hochberg-corrected dashboard, so no single p-value is read "
        "in isolation.",15.5)],space=10)

prs.save(OUT)
print("Biology deck rebuilt:",OUT,"| slides:",len(list(prs.slides)))
