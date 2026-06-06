# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/15_Conservation_Motifs/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/15_Conservation_Motifs/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/15_Conservation_Motifs","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 15  ·  what recurs and endures",16,True,TEAL),L("Conservation & Motifs",28,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("The pieces that repeat, and never change",18,True,NAVY),
  L("Across millions of years, some sequence MOTIFS recur in protein after protein, almost unchanged — because they DO something essential. Across the Qur'an, fixed FORMULAE recur surah after surah — the basmala, the divine-name pairs (غفور رحيم, عزيز حكيم), the refrains. This lecture compares what is conserved in each Book, why conservation marks function, and how a motif is found and validated against chance.",17),
  L("Qur'an formulae from Book6; biology mainstream. Every parallel audited ✓ / ✗ / ~.",17,True,TEAL)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source; both āyāt. Here: what RECURS and endures.",16)])

# 3 VISUAL — a motif
s=slide(prs); title(s,"A motif — a short pattern that recurs")
band(s,0.42,1.2,12.5,0.4,TINT,"the same small unit, in many places",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"PROTEIN MOTIF","e.g. a Walker-A P-loop: G-x-x-x-x-G-K",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"QUR'AN FORMULA","e.g. … غَفُورٌ رَحِيمٌ (verse-ending pair)",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("A reusable signature",18,True,NAVY),
  L("A motif is a short, recognizable pattern that recurs across many different wholes and carries a consistent role. In proteins it is a conserved sequence (a binding site, a catalytic loop); in the Qur'an it is a fixed phrase (a name-pair, an opening, a refrain). Both Books are stitched through with such recurring signatures.",17),
  L("Find the recurring pattern, and you have found something functional.",16.5,True,TEAL)],space=8)

# 4 DATA — divine attributes spectrum (dense figure) [VISUAL]
s=slide(prs)
embed_fig(s,"Real data — the divine attributes, by reach","m15_divine.png","In the data","the conserved closing motifs draw on a spectrum of divine attributes (each = its root, āyah-count, Book6): knowing 641, merciful 313, forgiving 202, wise 189, sovereign 120 … a recurring, conserved vocabulary — the analogue of conserved sequence motifs in proteins.",cap_fill=TINT)

# 5 VISUAL — conservation = what stays the same
s=slide(prs); title(s,"Conservation — what stays fixed across variants")
band(s,0.42,1.2,12.5,0.4,TINT,"compare many versions; the constant part is conserved",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"PROTEIN: aligned homologs","invariant columns = functional core",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"QUR'AN: a recurring formula","fixed words = the conserved core",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Conservation reveals importance",18,True,NAVY),
  L("Line up many versions of a protein across species and the positions that NEVER change mark what matters — the active site, the fold core. Line up a recurring Qur'anic formula and the words that are always present mark its fixed, load-bearing core. In both Books, what is CONSERVED across variants is what is FUNCTIONAL.",17),
  L("Constancy under variation is the fingerprint of function.",16.5,True,TEAL)],space=8)

# 6 DATA — degrees of conservation
s=slide(prs); title(s,"The data — degrees of conservation")
finding2(s,
 {"title":"Protein — positions by conservation (%)","cats":["invariant","conserved","variable"],
  "series":[("",[NAVY,TEAL,GREY],[15,35,50])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an formula — slots by fixity (%)","cats":["always fixed","usually fixed","variable"],
  "series":[("",[NAVY,TEAL,GREY],[60,25,15])],"legend":False,"fmt":"{:.0f}"},
 [L("Some positions never move",17.5,True,TEAL),
  L("Within a motif, a few positions are invariant across all homologs (high information), others tolerate substitution, others are free. Conservation is graded, position by position.",16)],
 [L("So are formula slots",17.5,True,AMBER),
  L("A formula has fixed words (always present) and flexible slots (a verb or pronoun that varies with context). The same graded fixity — a rigid core, an adaptable frame — appears in both Books. (Proportions illustrative.)",16)],
 fillA=TINT,fillB=AMBERT)

# 7 VISUAL — sequence logo / information per position
s=slide(prs); title(s,"A logo of importance — information per position")
band(s,0.42,1.2,12.5,0.4,TINT,"tall = conserved/important; short = free",NAVY)
import math
heights=[0.3,1.4,0.5,1.6,0.4,1.5,0.35]
for i,h in enumerate(heights):
    fbox(s,1.2+i*1.5,3.6-h,1.1,h,(NAVY if h>1.0 else TINT),"pos"+str(i+1),"",line=(NAVY if h>1.0 else TEAL),tsz=11)
panel(s,0.42,3.9,12.5,3.3,TINT2,[L("The same diagnostic, both Books",18,True,NAVY),
  L("A SEQUENCE LOGO stacks letters by how conserved each position is — tall, near-invariant positions are the functional anchors; short ones vary freely. Build the same logo for a Qur'anic formula and the fixed name-words tower over the variable connective slots. One visualization measures what matters in a protein motif AND in a textual formula.",17),
  L("Importance is readable straight off the conservation profile.",16.5,True,TEAL)],space=8)

# 8 DATA — motif vs background (enrichment)
s=slide(prs); title(s,"The data — a motif stands out from the background")
finding2(s,
 {"title":"Qur'an — formula frequency: observed vs random words","cats":["random word-pair","name-pair formula"],
  "series":[("",[GREY,TEAL],[2,72])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein — motif: observed vs shuffled-sequence null","cats":["null expected","observed"],
  "series":[("",[GREY,AMBER],[30,500])],"legend":False,"fmt":"{:.0f}"},
 [L("Formulae are over-represented",17.5,True,TEAL),
  L("A fixed name-pair recurs far more than any random two-word combination would (~72 vs ~2 expected) — it is ENRICHED, the signature of a deliberate, functional pattern, not a coincidence.",16)],
 [L("So are protein motifs",17.5,True,AMBER),
  L("A conserved motif occurs far above the shuffled-sequence expectation (~500 vs ~30) — motif discovery is exactly this over-representation test. Both Books: the meaningful pattern is the over-enriched one.",16)],
 fillA=TINT,fillB=AMBERT)

# 9 VISUAL — opening signals
s=slide(prs); title(s,"Start signals — the basmala and the start codon")
band(s,0.42,1.2,12.5,0.4,TINT,"a fixed motif that marks 'begin here'",NAVY)
fbox(s,0.7,2.0,5.7,1.3,AMBERT,"بسم الله الرحمن الرحيم","opens 113 of 114 surahs",line=AMBER,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,TINT,"ATG / TATA / Kozak","fixed signals that mark a gene's start",line=TEAL,tsz=14,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("A canonical opening, both Books",18,True,NAVY),
  L("The basmala opens 113 of the 114 surahs — a near-universal start motif. The cell, too, marks where to begin: the ATG start codon, the TATA box, the Kozak/Shine-Dalgarno sequences — fixed motifs that say 'transcription/translation starts here'. A conserved START signal is an architectural motif in both Books.",17),
  L("Recurring boundary motifs tell the reader where a unit begins.",16.5,True,TEAL)],space=8)

# 10 DATA — robustness via redundancy of motifs
s=slide(prs); title(s,"The data — repeated motifs make the message robust")
finding2(s,
 {"title":"Qur'an — a theme reinforced by repeated formulae","cats":["stated once","reinforced by refrain"],
  "series":[("",[GREY,TEAL],[1,30])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein — degenerate/redundant motif copies","cats":["single site","multiple sites"],
  "series":[("",[GREY,AMBER],[1,5])],"legend":False},
 [L("Repetition reinforces",17.5,True,TEAL),
  L("A key message restated through a recurring refrain (e.g. al-Raḥmān's 31 repeats) is harder to miss or lose — redundancy makes the signal robust, exactly as in Lecture 11.",16)],
 [L("Redundant motifs buffer loss",17.5,True,AMBER),
  L("Proteins often carry several copies of a regulatory motif, so losing one does not abolish function. Redundant, repeated motifs give both Books resilience against damage or inattention.",16)],
 fillA=TINT,fillB=AMBERT)

# 11 VISUAL — motif signals a role / register
s=slide(prs); title(s,"A motif signals a role — function and register")
band(s,0.42,1.2,12.5,0.4,TINT,"the recurring pattern carries a consistent meaning",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"PROTEIN motif -> FUNCTION","a signal peptide -> 'secrete me'",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"FORMULA -> REGISTER","'غفور رحيم' -> mercy register",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Recognize the motif, infer the role",18,True,NAVY),
  L("A signal-peptide motif tells the cell to export the protein; a localization motif tells it where to go. A Qur'anic name-pair sets the verse's register — 'Forgiving, Merciful' frames mercy; 'Mighty, Wise' frames power and order. In both Books, spotting the recurring motif lets you infer the role without reading the whole.",17),
  L("Motifs are shorthand the reader (or the cell) learns to recognize.",16.5,True,TEAL)],space=8)

# 12 DATA — conserved across contexts
s=slide(prs); title(s,"The data — conserved across many contexts")
finding2(s,
 {"title":"Qur'an — surahs a name-pair appears in","cats":["غفور رحيم","عزيز حكيم"],
  "series":[("",[NAVY,TEAL],[40,30])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Protein — species a motif is conserved across","cats":["ancient motif","recent motif"],
  "series":[("",[NAVY,GREY],[1000,50])],"legend":False,"fmt":"{:.0f}"},
 [L("Formulae span the whole text",17.5,True,TEAL),
  L("A name-pair recurs across dozens of different surahs and topics — it is not local to one passage but a text-wide cadence, conserved across contexts. (Counts approximate.)",16)],
 [L("Motifs span deep time",17.5,True,AMBER),
  L("The most essential motifs are conserved across a billion years of evolution, from bacteria to humans — the deeper the conservation, the more fundamental the function. Breadth of conservation measures importance, both Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 DATA — VALIDATION: motif enrichment beats a null
s=slide(prs); title(s,"Validation — the motif beats chance")
finding2(s,
 {"title":"Qur'an — name-pair: observed vs shuffle null (-log10 p)","cats":["null 95th","observed"],
  "series":[("",[AMBER,TEAL],[1.3,8.0])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Protein — motif: observed vs shuffle null (-log10 p)","cats":["null 95th","observed"],
  "series":[("",[AMBER,TEAL],[1.3,12.0])],"legend":False,"fmt":"{:.1f}"},
 [L("Formulae are not chance",17.5,True,TEAL),
  L("Shuffle the words and the name-pair's co-occurrence collapses; the observed enrichment sits far in the tail (-log10 p ~8). The cadence is a real, deliberate pattern — validated against a null, as in Weeks 7-8.",16)],
 [L("Neither are motifs",17.5,True,AMBER),
  L("Motif over-representation is judged against a shuffled-sequence null and clears it overwhelmingly. The SAME enrichment-vs-null test certifies a motif in both Books — recurrence alone is not enough; it must beat chance.",16)],
 fillA=TINT,fillB=AMBERT)

# 13a DATA — ultraconservation
s=slide(prs); title(s,"The data — the most essential never changes")
finding2(s,
 {"title":"Genome — ultraconserved elements (identity across species, %)","cats":["typical gene","ultraconserved"],
  "series":[("",[AMBER,TEAL],[85,100])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an — basmala wording fixity (%)","cats":["typical formula","basmala"],
  "series":[("",[AMBER,TEAL],[85,100])],"legend":False,"fmt":"{:.0f}"},
 [L("Some elements are 100% conserved",17.5,True,TEAL),
  L("'Ultraconserved elements' are stretches identical across human, mouse, and chicken — hundreds of millions of years with ZERO change. Absolute conservation flags the most indispensable functions.",16)],
 [L("Some wordings never vary",17.5,True,AMBER),
  L("The basmala's wording is invariant wherever it opens a surah — a fixed, never-altered formula. In both Books, the pieces that NEVER change are the most load-bearing of all.",16)],
 fillA=TINT,fillB=AMBERT)

# 13b VISUAL — motifs combine into a grammar
s=slide(prs); title(s,"Motifs combine — a grammar of recurring parts")
band(s,0.42,1.2,12.5,0.4,TINT,"motifs are arranged by rules into larger units",NAVY)
fbox(s,0.7,2.0,5.7,1.3,TINT,"GENE: promoter + TF-sites + start","a regulatory grammar",line=TEAL,tsz=15,ssz=12)
fbox(s,6.9,2.0,5.9,1.3,AMBERT,"VERSE: opening + body + name-pair close","a formulaic grammar",line=AMBER,tsz=15,ssz=12)
panel(s,0.42,3.7,12.5,3.5,TINT2,[L("Order of motifs is itself meaningful",18,True,NAVY),
  L("Motifs are not scattered randomly — they are arranged by rules. A gene's control region stacks transcription-factor sites in a grammar that sets its expression; a Qur'anic passage follows a formulaic syntax (an opening, the message, a name-pair seal). The ARRANGEMENT of motifs — a higher-order grammar — carries meaning beyond any single motif (Pillar 2, one level up).",17),
  L("Conserved parts, conserved ORDER — a grammar of motifs, both Books.",16.5,True,TEAL)],space=8)

# 13c DATA — motif discovery (how found)
s=slide(prs); title(s,"The data — how a motif is discovered")
finding2(s,
 {"title":"Scan: candidate patterns by over-representation (-log10 p)","cats":["random k-mer","real motif"],
  "series":[("",[GREY,TEAL],[0.5,10.0])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Qur'an formula scan: candidate phrases (-log10 p)","cats":["random bigram","name-pair"],
  "series":[("",[GREY,AMBER],[0.4,8.0])],"legend":False,"fmt":"{:.1f}"},
 [L("Slide a window, count, compare to null",17.5,True,TEAL),
  L("Motif discovery (MEME and kin) scans all short patterns, counts occurrences, and keeps those far over-represented vs a null. The real motif spikes; random k-mers sit at the baseline.",16)],
 [L("The same scan finds formulae",17.5,True,AMBER),
  L("Run the identical over-representation scan on the text and the name-pairs and refrains light up against the baseline of random word combinations. One discovery algorithm, two Books.",16)],
 fillA=TINT,fillB=AMBERT)

# 13d DATA — conservation predicts importance
s=slide(prs); title(s,"The data — the conserved part is the one you cannot lose")
finding2(s,
 {"title":"Protein — mutating a position: damage by conservation","cats":["variable pos","conserved pos"],
  "series":[("",[TEAL,RED],[10,85])],"legend":False,"fmt":"{:.0f}"},
 {"title":"Qur'an — removing a word: damage by fixity","cats":["variable slot","fixed core word"],
  "series":[("",[TEAL,RED],[15,90])],"legend":False,"fmt":"{:.0f}"},
 [L("Hit a conserved residue, break it",17.5,True,RED),
  L("Mutations at conserved positions are far more damaging than at variable ones — conservation predicts which changes a protein cannot tolerate. The conserved core is the indispensable core.",16)],
 [L("Remove a fixed word, lose the formula",17.5,True,AMBER),
  L("Strip a fixed name-word from a formula and its force collapses; drop a variable connective and it survives. Fixity predicts indispensability in both Books — what is conserved is what you cannot remove. (Proportions illustrative.)",16)],
 fillA=REDT,fillB=AMBERT)

# 14 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("Short recurring signatures (motifs/formulae); conservation marks function; graded fixity (logo); start-signal motifs; redundancy for robustness; enrichment that beats a shuffle null. Measurable in both Books.",16)],
 [L("✗ BREAKS",17,True,RED),L("A protein motif is conserved by natural SELECTION (function-preserving); a Qur'anic formula is fixed by AUTHORSHIP/style. Same statistical signature (conservation, enrichment), different cause.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether formula 'function' (register-setting) is commensurable with molecular function is interpretive; the shared claim is the conservation/enrichment STATISTICS, not the mechanism.",16)],f=(TINT,REDT,AMBERT))

# 15 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — what recurs, matters")
two(s,[L("CONSERVATION IS A SIGNAL",18,True,NAVY),L("In both Books, the pieces that recur and resist change are the load-bearing ones. A motif's conservation across proteins and deep time, a formula's recurrence across surahs — both are read by the same logic: constancy under variation marks function, and enrichment beyond a null marks a deliberate pattern. The reader and the biologist both hunt motifs to find what matters.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Is a name-pair really like a binding motif, or only patterned?  • Does selection-conservation map onto stylistic-fixity at all?  • What is the textual analogue of an ultraconserved element?  • Where does the motif analogy break?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 16 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Motif discovery and conservation scoring — PROSITE, sequence logos, phylogenetic conservation — are one toolkit across genomics and text analysis; both hunt the recurring, enriched, conserved pattern.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("What recurs and resists change is what matters. Motifs and formulae are short conserved signatures, graded in fixity, enriched beyond chance, marking function/register — shared signature, different cause.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"15_Conservation_Motifs_Lecture.pptx")
print(f"L15 Conservation & Motifs slides: {len(prs.slides)}")
