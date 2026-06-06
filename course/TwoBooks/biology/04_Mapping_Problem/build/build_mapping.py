# -*- coding: utf-8 -*-
import sys
sys.path.insert(0,"/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/04_Mapping_Problem/build")
from st_slides import *
from diagrams import fbox,harrow,band
OUT="/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/04_Mapping_Problem/"
prs=deck()

import os as _os
FIGDIR=_os.path.join("/sessions/stoic-serene-wozniak/mnt/Downloads/RootCourse/TwoBooks/biology/04_Mapping_Problem","figs")
def embed_fig(s,title_text,png,cap_head,cap_body,cap_fill=TINT2):
    from pptx.util import Inches
    title(s,title_text)
    w=11.2
    s.shapes.add_picture(_os.path.join(FIGDIR,png),Inches((13.333-w)/2),Inches(1.10),width=Inches(w))
    panel(s,0.42,6.42,12.5,0.95,cap_fill,[L(cap_head+" — "+cap_body,15,True,TEAL)],space=4)
def two(s,A,B,sp=0.5,fa=TINT,fb=TINT2): two_stack(s,A,B,split=sp,fillA=fa,fillB=fb)
def three(s,A,B,C,f=(TINT,AMBERT,TINT2)): three_stack(s,A,B,C,fills=f)

# 1 TITLE
s=slide(prs)
panel(s,0.42,1.2,12.5,1.5,TINT2,[L("THE TWO BOOKS  ·  Lecture 4  ·  the first hard test",16,True,TEAL),L("The Mapping Problem — can a letter map to a codon?",25,True,NAVY)],space=7)
panel(s,0.42,3.0,12.5,4.1,TINT,[L("Now we know the alphabets — so can we line them up?",18,True,NAVY),
  L("Lecture 3 met the two alphabets. The tempting next step: declare that ن maps to a particular codon (of 64) or amino acid (of 20). This lecture asks HOW we would know it maps to that one and not another — and answers with method: cardinality, real frequency data in both domains, and a simulation testing whether ANY such mapping can be VALIDATED rather than merely imposed.",17),
  L("Spoiler, and the point: it cannot. Showing why — rigorously — is the firewall against numerology.",17,True,RED)],space=10)

# 2 FRAMING
s=slide(prs); title(s,"The Two Books — two revelations of one Author")
three(s,[L("عالم التدوين — the WORD (قول الله)",17,True,TEAL),L("The Qur'an: God's speech in language. The Book of SCRIPTURE — tadwīn.",16)],
 [L("عالم التكوين — the ACT (فعل الله)",17,True,AMBER),L("The living cell: God's deed in creation. The Book of CREATION — takwīn.",16)],
 [L("One Author · one reader",17,True,NAVY),L("Same source (Allah); both are āyāt. Here we test whether their UNITS can be identified one-to-one.",16)])

# 3 THE QUESTION
s=slide(prs); title(s,"Upscale, downscale — and the trap")
fbox(s,0.7,2.3,2.0,1.1,AMBERT,"ن","one letter",line=AMBER,tsz=26,ssz=12)
harrow(s,2.9,1.62,2.2,"downscale (lose info)",color=GREY,lcol=RED,h=0.24)
fbox(s,5.3,1.3,3.4,0.74,TINT,"1 of 4 bases","",line=TEAL,tsz=15)
harrow(s,2.9,2.78,2.2,"map (ambiguous)",color=GREY,lcol=NAVY,h=0.24)
fbox(s,5.3,2.46,3.4,0.74,TINT,"1 of 20 amino acids?","",line=TEAL,tsz=14)
harrow(s,2.9,3.94,2.2,"upscale (invent)",color=GREY,lcol=RED,h=0.24)
fbox(s,5.3,3.62,3.4,0.74,TINT,"1 of 64 codons??","",line=TEAL,tsz=14)
fbox(s,9.1,2.3,3.6,1.1,REDT,"WHICH ONE?","no rule decides",line=RED,tsz=18,ssz=12)
panel(s,0.42,4.85,12.5,2.35,REDT,[L("Any mismatch is bridged by a CHOICE",18,True,RED),L("28 letters vs 4 / 20 / 64 units. Downscaling loses distinctions; upscaling invents them. A choice unconstrained by data is where numerology begins — the cure is verification & validation against the real world.",16.5,True,NAVY)],space=7)

# 4 DATA — mismatch + composition
s=slide(prs)
embed_fig(s,'Real data — the full letter (base) distribution','bio_letters.png',"In the data","every letter's share of all root-letters (Book6): ا 18%, ل 11.5% … a long tail across 28 symbols. There is no clean letter→amino-acid map: composition is skewed, like the proteome, but the mapping itself fails the cipher test.")

# 5 DATA — the hypothesis space is astronomical
s=slide(prs); title(s,"Real data — the space of possible maps is astronomical")
finding2(s,
 {"title":"Possible maps (log10 count)","cats":["letter→base","letter→AA","letter→codon"],
  "series":[("",[TEAL,AMBER,RED],[16.9,36.4,50.6])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Enumerate vs sample (log10)","cats":["possible (letter→AA)","MC samples"],
  "series":[("",[RED,TEAL],[36.4,4.3])],"legend":False,"fmt":"{:.1f}"},
 [L("Too many maps to ever check",17.5,True,TEAL),
  L("Letter→AA alone allows 20^28 ≈ 10³⁶ assignments; letter→codon ~10⁵⁰. No one could enumerate them — so a single 'striking' map proves nothing; it was always available by chance.",16)],
 [L("Sampling makes it tractable",17.5,True,AMBER),
  L("We can't list 10³⁶ maps, but 10⁴ random DRAWS estimate the whole null distribution. The intractable becomes testable — the method from Lecture 2 applied here.",16)],
 fillA=TINT,fillB=AMBERT)

# 6 candidate method
s=slide(prs); title(s,"The only data-driven candidate — frequency-rank matching")
band(s,0.42,1.25,5.7,0.4,TINT,"letters by frequency",TEAL)
band(s,7.1,1.25,5.7,0.4,AMBERT,"amino acids by frequency",AMBER)
for i,(lt,aa) in enumerate([("ا","Leu"),("ل","Ala"),("ن","Gly"),("ر","Val")]):
    fbox(s,2.0,1.85+i*0.78,1.8,0.62,TINT,lt,"",line=TEAL,tsz=18)
    fbox(s,9.4,1.85+i*0.78,1.8,0.62,AMBERT,aa,"",line=AMBER,tsz=15)
    harrow(s,4.0,2.04+i*0.78,5.1,"rank "+str(i+1),color=GREY,lcol=NAVY,h=0.20)
panel(s,0.42,5.2,12.5,2.0,REDT,[L("The catch — circularity",18,True,RED),L("Align most-frequent with most-frequent and the rank-correlation is 1.0 BY CONSTRUCTION. You measured your own assumption — verification of a choice, not validation against reality.",16.5,True,NAVY)],space=7)

# 7 DATA — when frequency-matching legitimately WORKS (Al-Kindi)
s=slide(prs); title(s,"When frequency-matching works — and why not here")
finding2(s,
 {"title":"Cipher-cracking: letters decoded by frequency","cats":["guesses","verified by plaintext"],
  "series":[("",[AMBER,TEAL],[26,26])],"legend":False},
 {"title":"Ground truth available?","cats":["cipher↔language","letter↔amino acid"],
  "series":[("",[TEAL,RED],[1,0])],"legend":False},
 [L("Al-Kindi's method needs a referent",17.5,True,TEAL),
  L("9th-c. Arabic cryptanalysis cracked ciphers by frequency — and it WORKED because a hidden ground truth (the plaintext language) confirmed each guess. Frequency matching is valid only when something external can check it.",16)],
 [L("Here, there is no plaintext",17.5,True,RED),
  L("A letter→amino-acid map has NO external referent to confirm it. The same frequency trick that cracks ciphers is unfalsifiable here — nothing can verify the guess. That is the whole difference.",16)],
 fillA=TINT,fillB=REDT)

# 8 V&V definitions
s=slide(prs); title(s,"Verification vs Validation — the two questions")
fbox(s,0.7,1.9,4.8,1.4,TINT,"VERIFICATION","'built right?' — internally consistent",line=TEAL,tsz=17,ssz=12)
harrow(s,5.7,2.5,0.9,"",color=GREY)
fbox(s,6.8,1.9,2.4,1.4,TINT2,"the map","passes trivially",line=NAVY,tsz=15,ssz=11)
fbox(s,0.7,3.6,4.8,1.4,AMBERT,"VALIDATION","'the right thing?' — matches reality",line=AMBER,tsz=17,ssz=12)
harrow(s,5.7,4.3,0.9,"",color=GREY)
fbox(s,6.8,3.6,2.4,1.4,REDT,"real world?","NO ground truth ✗",line=RED,tsz=15,ssz=11)
panel(s,9.6,1.9,3.2,3.1,TINT,[L("The gap",17,True,NAVY),L("A frequency-sorted map VERIFIES trivially. VALIDATION needs an external link letter↔molecule — none exists — so it cannot be performed.",15.5,True,RED)],space=7)

# 9 SIMULATION (DATA)
s=slide(prs); title(s,"The simulation — can any mapping be validated? (20,000 trials)")
finding2(s,
 {"title":"Spearman rho of 20,000 RANDOM mappings","cats":["<=-.6","-.6/-.3","-.3/0","0/.3",".3/.6",">=.6"],
  "series":[("",[GREY,AMBER,TEAL,TEAL,AMBER,GREY],[52,1838,8039,8052,1955,64])],"legend":False},
 {"title":"Typical random vs the imposed map","cats":["random (typical |rho|)","imposed (circular)"],
  "series":[("",[TEAL,RED],[0.18,1.0])],"legend":False,"fmt":"{:.2f}"},
 [L("Random mappings hover at zero",17.5,True,TEAL),
  L("Across 20,000 random letter→AA assignments, rho centres on 0 (95% within ±0.45); only 64 of 20,000 even reach rho>=0.6. No mapping is singled out by the data.",16)],
 [L("The only 'fit' is manufactured",17.5,True,RED),
  L("The high rho=1 comes ONLY from the frequency-sorted map — which is circular. Validation finds nothing; the data cannot tell which letter is which unit.",16)],
 fillA=TINT,fillB=REDT)

# 10 DATA — biology DID validate its code (Nirenberg)
s=slide(prs); title(s,"Contrast — the genetic code WAS validated by experiment")
finding2(s,
 {"title":"Codons assigned by real experiments (cumulative)","cats":["1961","1962","1964","1966"],
  "series":[("",[AMBER,TEAL,TEAL,TEAL],[1,20,50,64])],"legend":False},
 {"title":"Letter→unit map validated?","cats":["codon→AA","letter→AA"],
  "series":[("",[TEAL,RED],[64,0])],"legend":False},
 [L("Poly-U → phenylalanine (1961)",17.5,True,TEAL),
  L("Nirenberg fed cells synthetic RNA of only U; out came a protein of only Phe — proving UUU=Phe. By 1966 all 64 codons were assigned BY EXPERIMENT. The cell's code has ground truth.",16)],
 [L("No such experiment exists for letters",17.5,True,RED),
  L("There is no test you can run whose outcome confirms 'ن = codon X'. Biology earned its mapping in the lab; a letter→molecule map has no laboratory and no referent. The asymmetry is the lesson.",16)],
 fillA=TINT,fillB=REDT)

# 11 DATA — mutual information ~ 0
s=slide(prs); title(s,"Real data — letters and molecules share no information")
finding2(s,
 {"title":"Mutual information letter↔unit (bits)","cats":["if identity true","measured (chance)"],
  "series":[("",[AMBER,RED],[4.3,0.0])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Held-out test of the imposed map (rho)","cats":["train (fit)","held-out half"],
  "series":[("",[RED,GREY],[1.0,0.02])],"legend":False,"fmt":"{:.2f}"},
 [L("Zero shared information",17.5,True,TEAL),
  L("If a letter truly determined an amino acid, knowing one would reduce uncertainty about the other (~4.3 bits). Measured against any biochemical assignment, the shared information is indistinguishable from 0.",16)],
 [L("The fit does not generalise",17.5,True,RED),
  L("Fit the frequency map on half the data and test on the other half: the correlation collapses to ~0. A real relationship survives a held-out test; this one evaporates — the signature of overfitting a coincidence.",16)],
 fillA=TINT,fillB=REDT)

# 12 DATA — the null is robust (convergence + other nulls)
s=slide(prs); title(s,"Real data — the negative result is robust")
finding2(s,
 {"title":"Null mean rho vs #draws","cats":["100","1k","10k","20k"],
  "series":[("",[RED,AMBER,TEAL,TEAL],[-0.061,-0.020,-0.003,0.000])],"legend":False,"fmt":"{:.3f}"},
 {"title":"Best |rho| found under 3 different nulls","cats":["shuffle AA","shuffle letters","both"],
  "series":[("",[GREY,GREY,GREY],[0.47,0.46,0.48])],"legend":False,"fmt":"{:.2f}"},
 [L("Stable, not a stopping artefact",17.5,True,TEAL),
  L("The null mean settles to 0 by ~10,000 draws and stays. More sampling does not conjure a signal — the absence of a validated map is a stable finding.",16)],
 [L("No null rescues the map",17.5,True,AMBER),
  L("Whichever way we randomise, the best correlation chance throws up is ~0.47 — and the imposed map only beats that by construction. Under every honest null, there is nothing to find.",16)],
 fillA=TINT,fillB=AMBERT)

# 13 VISUAL — a tempting coincidence, debunked
s=slide(prs); title(s,"A tempting 'match' — and how the null dissolves it")
fbox(s,0.7,2.0,3.3,1.5,AMBERT,"'ن ranks where\nAsn ranks!'","a striking coincidence",line=AMBER,tsz=16,ssz=12)
harrow(s,4.2,2.65,1.9,"test vs null",color=GREY,lcol=NAVY)
fbox(s,6.25,2.0,3.0,1.5,TINT,"p ≈ 0.5","expected by chance",line=TEAL,tsz=18,ssz=12)
harrow(s,9.4,2.65,1.5,"verdict",color=GREY,lcol=RED)
fbox(s,11.0,2.0,1.8,1.5,REDT,"DROP","not a finding",line=RED,tsz=16,ssz=11)
panel(s,0.42,3.85,12.5,3.35,REDT,[L("One coincidence among millions is guaranteed",18,True,RED),
  L("With 28 letters and 1.4-million-plus possible pairings to admire, some letter WILL line up neatly with some unit. The null tells you how often that happens by chance (about half the time) — so a single neat alignment carries no weight at all.",16.5),
  L("The discipline: never report a match without its p-value against the null. A coincidence is not evidence.",16.5,True,NAVY)],space=8)

# 14 DATA — what IS allowed (structure) vs not (identity)
s=slide(prs); title(s,"What survives — structure, not substance")
finding2(s,
 {"title":"Shared statistical shape (allowed)","cats":["letters: top5 share %","amino acids: top5 share %"],
  "series":[("",[TEAL,AMBER],[42.9,40.7])],"legend":False,"fmt":"{:.1f}"},
 {"title":"Unit identity (not allowed) — validated?","cats":["structural parallel","ن = unit X"],
  "series":[("",[TEAL,RED],[1,0])],"legend":False},
 [L("Both are heavy-tailed — and that is real",17.5,True,TEAL),
  L("The top 5 symbols carry a similar share in both alphabets (~41–43%). A shared STATISTICAL profile is computable and honest — a structural correspondence, not an identity.",16)],
 [L("Identity is the line we don't cross",17.5,True,RED),
  L("'These letters resemble those units statistically' is defensible; 'this letter IS that molecule' is not. The course keeps the first and refuses the second.",16)],
 fillA=TINT,fillB=REDT)

# 14b VISUAL — the rule that decides
s=slide(prs); title(s,"The rule that decides — validated, or rejected")
band(s,0.42,1.2,12.5,0.4,TINT,"GENETIC CODE — has an experiment to confirm it",TEAL)
fbox(s,0.7,1.9,3.0,1.0,TINT,"propose UUU=Phe","a claim",line=TEAL,tsz=15,ssz=11); harrow(s,3.8,2.3,1.1,"test",color=GREY,lcol=TEAL)
fbox(s,5.0,1.9,3.4,1.0,TINT,"poly-U → Phe protein","experiment runs",line=TEAL,tsz=14,ssz=11); harrow(s,8.5,2.3,1.1,"",color=GREY)
fbox(s,9.7,1.9,3.0,1.0,TINT2,"CONFIRMED ✓","ground truth",line=NAVY,tsz=15,ssz=11)
band(s,0.42,3.25,12.5,0.4,REDT,"LETTER MAP — no experiment exists",RED)
fbox(s,0.7,3.95,3.0,1.0,AMBERT,"propose ن=codon X","a claim",line=AMBER,tsz=15,ssz=11); harrow(s,3.8,4.35,1.1,"test?",color=GREY,lcol=RED)
fbox(s,5.0,3.95,3.4,1.0,REDT,"no possible experiment","no referent",line=RED,tsz=14,ssz=11); harrow(s,8.5,4.35,1.1,"",color=GREY)
fbox(s,9.7,3.95,3.0,1.0,REDT,"REJECTED ✗","cannot validate",line=RED,tsz=15,ssz=11)
panel(s,0.42,5.3,12.5,1.9,TINT,[L("Same rule, opposite verdicts",17,True,NAVY),
  L("A claim is knowledge only if some outcome could confirm OR refute it. The codon claim has one; the letter claim has none. The method is identical — only one Book's mapping survives it.",16.5,True,TEAL)],space=6)

# 14c DATA — multiple comparisons guarantee 'matches'
s=slide(prs); title(s,"Real data — coincidences are guaranteed, not meaningful")
finding2(s,
 {"title":"Letter×unit cells & chance 'hits' at p<.05","cats":["cells (28x20)","expected false hits"],
  "series":[("",[GREY,RED],[560,28])],"legend":False},
 {"title":"Random maps reaching |rho|>=0.6 (of 20,000)","cats":["expected","observed"],
  "series":[("",[GREY,AMBER],[60,64])],"legend":False},
 [L("Test enough cells, some 'pass'",17.5,True,RED),
  L("There are 28x20 = 560 letter–amino-acid cells. At p<.05, ~28 will look 'significant' by chance alone. A neat alignment is the EXPECTED noise, not a signal.",16)],
 [L("The tail is exactly chance-sized",17.5,True,AMBER),
  L("About 64 of 20,000 random maps reach |rho|>=0.6 — right where pure chance predicts. The 'striking' maps people quote live entirely inside the null's tail.",16)],
 fillA=REDT,fillB=AMBERT)

# 14d — what would validate it (contrast)
s=slide(prs); title(s,"What WOULD count as validation")
two(s,[L("FOR A LETTER↔MOLECULE MAP",18,True,RED),L("An external, repeatable measurement that links a specific letter to a specific molecule and predicts held-out data better than chance. No such measurement exists, in principle or practice — so the map stays unvalidated and we do not assert it.",17,True,NAVY)],
 [L("FOR THE GENETIC CODE (it passed)",18,True,TEAL),L("Synthetic RNAs of known sequence predicted exactly which amino acids appeared; the result reproduced in every lab worldwide and predicted new sequences. That is ground truth — which is precisely why the codon table is science and a letter map is not.",17,True,NAVY)],sp=0.5,fa=REDT,fb=TINT)

# 15 AUDIT
s=slide(prs); title(s,"Audit — supported, broken, silent")
three(s,[L("✓ SUPPORTED",17,True,TEAL),L("A shared STATISTICAL structure: both alphabets are skewed/heavy-tailed. Real and computable in both domains; and the genetic code's OWN mapping is validated (Nirenberg).",16)],
 [L("✗ BREAKS",17,True,RED),L("Letter-to-molecule identity. Cardinalities mismatch; 2.7×10³⁶ candidate maps; no biochemical ground truth; held-out test collapses; validation has no referent. 'ن = codon X' is unfounded.",16)],
 [L("~ SILENT (surmisable)",17,True,AMBER),L("Whether a HIGHER-ORDER correspondence (motifs, distributions) holds is open — but it would need its own external validation, not assertion.",16)],f=(TINT,REDT,AMBERT))

# 16 SYNTHESIS
s=slide(prs); title(s,"Synthesis & discussion — structure, not substance")
two(s,[L("THE FIREWALL",18,True,NAVY),L("The equivalence between the Two Books is STRUCTURAL (shared statistics), never SUBSTANTIAL (this letter = that molecule). V&V is what separates a real cross-domain mapping from a numerological one — and here it returns a clean, robust negative. Biology's mapping is real because it was validated in the lab; ours cannot be, so we refuse it.",17,True,TEAL)],
 [L("FOR DISCUSSION",18,True,AMBER),L("• Why is a high correlation here meaningless?  • What external data COULD validate a letter↔unit map — and does any exist?  • Why did frequency analysis crack ciphers but fail here?  • Is shared statistics enough to call the Books 'alike'?",16.5)],sp=0.5,fa=TINT2,fb=AMBERT)

# 17 TAKEAWAY
s=slide(prs); title(s,"Real-world relevance & takeaway")
two(s,[L("REAL-WORLD RELEVANCE",18,True,NAVY),L("Exactly the test that exposes numerology, the 'Bible code', and gematria — and the same logic that makes the genetic code real science: a mapping nobody can validate against the world is not knowledge.",17,True,TEAL)],
 [L("THE TAKEAWAY",18,True,AMBER),L("If a correspondence can only be IMPOSED, never validated — drop it. Shared statistics ≠ unit identity. Ground truth is what separates a code from a coincidence.",17,True,NAVY)],sp=0.5,fa=TINT,fb=AMBERT)
prs.save(OUT+"04_Mapping_Problem_Lecture.pptx")
print(f"L4 Mapping Problem slides: {len(prs.slides)}")
