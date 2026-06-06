# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
WK="/sessions/kind-compassionate-feynman/mnt/RootCourse/week04"; FD=os.path.join(WK,"figs")
NAVY=RGBColor(0x1E,0x27,0x61); TEAL=RGBColor(0x0E,0x9D,0x8C); INK=RGBColor(0x1E,0x29,0x3B); GREY=RGBColor(0x55,0x60,0x70); RED=RGBColor(0xA2,0x3B,0x3B)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); SH=prs.slide_height
blank=prs.slide_layouts[6]
def sf(run,name):
    run.font.name=name; rPr=run._r.get_or_add_rPr()
    for tag in ("a:cs","a:ea"):
        el=rPr.find(qn(tag))
        if el is None: el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set("typeface",name)
def slide():
    s=prs.slides.add_slide(blank); b=s.shapes.add_shape(1,0,0,Inches(0.16),SH); b.fill.solid(); b.fill.fore_color.rgb=TEAL; b.line.fill.background(); return s
def title(s,t,sz=30):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.28),Inches(12.5),Inches(0.95)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=NAVY; sf(r,"Georgia"); return s
def body(s,lines,top=1.3,size=19,width=12.3,left=0.65,space=7):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(7.5-top-0.25)); tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,bold,col,sz) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(space)
        r=p.add_run(); r.text=txt; r.font.size=Pt(sz or size); r.font.bold=bold; r.font.color.rgb=col or INK; sf(r,"Calibri")
def pic(s,name,left,top,width=None,height=None):
    p=os.path.join(FD,name)
    if os.path.exists(p):
        kw={}
        if width:kw['width']=Inches(width)
        if height:kw['height']=Inches(height)
        s.shapes.add_picture(p,Inches(left),Inches(top),**kw)
def B(t,bold=False,col=None,sz=None): return (t,bold,col,sz)

# 1 TITLE
s=slide()
tb=s.shapes.add_textbox(Inches(0.6),Inches(2.0),Inches(12),Inches(3.4)); tf=tb.text_frame; tf.word_wrap=True
for txt,sz,bold,col,fn in [("WEEK 4",22,True,TEAL,"Georgia"),("Co-occurrence",46,True,NAVY,"Georgia"),
  ("Which root shares the most ayahs with a target — and why the raw count fools us.",21,False,GREY,"Calibri"),
  ("Our first PAIR measure and first FIND task.  Worked target: صلو (prayer).",19,False,GREY,"Calibri"),
  ("Every value computed from Book6 (6,236 ayahs).",15,False,GREY,"Calibri")]:
    p=tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]; p.space_after=Pt(8); r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; sf(r,fn)

# 2 recap
s=slide(); title(s,"Recap, and today's new question")
body(s,[B("Week 1 — how MUCH a root is named  (frequency).",False,GREY),
 B("Week 2 — WHERE it sits and how it spreads  (distribution).",False,GREY),
 B("Week 3 — what SHAPE it wears and who it keeps company with  (forms & partners).",False,GREY),
 B("Today — our first PAIR measure: how often do two roots share an ayah?",True,NAVY),
 B("And our first FIND task: among several candidates, which shares the MOST with a target?",True,NAVY),
 B("The catch: ranking by the raw shared count is not just imperfect — it is systematically WRONG.",True,RED),
 B("It is like ranking your friendships by who appears in the most of your photos —",False),
 B("the winner is whoever is in EVERY photo, not your closest friend.",False),
 B("Worked target: صلو (prayer). By the end you will know its true closest companion.",False,TEAL)],top=1.35,size=19,space=10)

# 3 joint
s=slide(); title(s,"Counting shared ayahs — the joint count")
pic(s,"fig_cooccurrence_concept.png",7.4,1.7,width=5.6)
body(s,[B("Co-occurrence = ayahs that contain BOTH roots.",True,NAVY,20),
 B("The overlap of two circles in a Venn diagram.",False,GREY),
 B("صلو (prayer): 90 ayahs.",False),
 B("زكو (zakat): 56 ayahs.",False),
 B("Shared by both: 28 ayahs.",True,RED,20),
 B("Deliberately blunt: we ignore WHERE in the verse they sit, or their grammar —",False,None,16),
 B("only that both are present.",False,None,16),
 B("It is the rawest notion of “together.” We sharpen it next.",True,TEAL),
 B("This is the ingredient — not yet the answer.",False,GREY)],top=1.4,size=18,width=6.6,space=7)

# 4 raw slate
s=slide(); title(s,"Raw co-occurrence with prayer")
pic(s,"fig_raw_joint.png",7.3,1.7,width=5.8)
body(s,[B("Rank صلو's candidates by raw shared ayahs:",True,NAVY,20),
 B("•  قوم (establish) — 52   (appears in 597 ayahs)",False),
 B("•  ءله (God) — 44   (appears in 1,879 ayahs)",False),
 B("•  زكو (zakat) — 28   (appears in only 56)",False,TEAL),
 B("•  ءتي (give) — 27",False),
 B("Raw says prayer's companion is “establish” or “God.”",True),
 B("Unremarkable — and, it turns out, an artefact of counting.",False,GREY),
 B("Notice: the raw leaders are the roots that appear",False,None,16),
 B("almost everywhere. That is the clue to the trap.",False,None,16)],top=1.4,size=18,width=6.7,space=8)

# 5 frequency confound
s=slide(); title(s,"Why the raw count misleads — the frequency confound")
body(s,[B("Picture a celebrity photographed at EVERY event in town.",True,NAVY),
 B("Flip through the year's photos: hundreds show the celebrity beside your friend —",False),
 B("not because they're close, but because the celebrity is in every single photo.",False),
 B("Now an ordinary person in only 50 photos all year — and your friend is in 28 of them.",True,TEAL),
 B("THAT is closeness. The small number is the loud one.",True,TEAL),
 B("ءله (God) is this corpus's celebrity: named in 1,879 ayahs — a quarter of the book.",False,RED),
 B("Of course it shares 44 with prayer; it shares dozens with nearly everything.",False),
 B("زكو appears in only 56 ayahs in the whole Qur'an — and 28 (half) are with prayer.",True,NAVY),
 B("So a big shared count can be nothing but the candidate being everywhere.",False,GREY),
 B("Same trap as Week 2's length confound — now it's frequency.",False,GREY)],top=1.35,size=18,space=8)

# 6 fix
s=slide(); title(s,"The fix — observed vs expected-by-chance")
pic(s,"fig_observed_vs_expected.png",7.4,1.8,width=5.7)
body(s,[B("Compare what we OBSERVED to what we'd EXPECT by pure chance.",True,NAVY,20),
 B("Expected = freq(A) × freq(B) ÷ 6,236.",True),
 B("صلو–ءله:  90 × 1,879 ÷ 6,236 ≈ 27 expected,  observe 44.",False),
 B("        → barely above chance.",False,RED),
 B("صلو–زكو:  90 × 56 ÷ 6,236 ≈ 0.8 expected,  observe 28.",False),
 B("        → far above chance.",False,TEAL),
 B("Measure = how many TIMES observed beats expected.",True,NAVY),
 B("A crowded room: two guests who keep meeting in the same corner",False,None,16),
 B("far more than random milling predicts are seeking each other out.",False,None,16),
 B("Ask not “how many shared?” but “how many MORE than chance?”",True,TEAL)],top=1.4,size=18,width=6.7,space=6)

# 7 controlled
s=slide(); title(s,"Controlled for frequency — the true companion")
pic(s,"fig_controlled_ratio.png",7.4,1.9,width=5.7)
body(s,[B("Rank the SAME candidates by × over chance:",True,NAVY,20),
 B("•  زكو (zakat) — ×34.6",False,TEAL,20),
 B("•  نفق (spend) — ×8.1",False),
 B("•  قوم (establish) — ×6.0",False),
 B("•  ءتي (give) — ×3.8",False),
 B("•  ءله (God) — ×1.6   (basically chance)",False,RED),
 B("Zakat, third on the raw list, wins by a landslide.",True,NAVY),
 B("Of zakat's 56 ayahs, 28 — exactly half — are with prayer.",True),
 B("Its bond is narrow and DEEP; قوم's is wide and THIN.",False,GREY)],top=1.45,size=18,width=6.7,space=8)

# 8 flip
s=slide(); title(s,"The flip — the “unlearn”")
pic(s,"fig_flip.png",1.5,1.45,width=10.3)
body(s,[B("Raw says قوم / God.  Controlled says zakat.",True,NAVY,20),
 B("Prayer's closest companion is CHARITY — أقيموا الصلاة وآتوا الزكاة (“establish prayer and give zakat”).",False,None,17),
 B("You could only see it by refusing to reward roots for being common.",False,GREY,16)],top=6.05,size=18,left=1.3,width=10.8,space=4)

# 9 find-task
s=slide(); title(s,"The find-task — how to rank")
body(s,[B("Given a target and a slate of candidates:",True,NAVY,20),
 B("1.  For each candidate, compute × over chance  (observed ÷ expected).",False),
 B("2.  Throw out any pair with fewer than 5 shared ayahs — too little evidence to trust.",False),
 B("      (the same small-sample caution as Week 2's support floor.)",False,GREY,16),
 B("3.  Read the TOP of the controlled list — never the raw list.",True,TEAL),
 B("The whole habit in one line:  rank by SURPRISE, not by SIZE.",True,NAVY,21),
 B("Across the corpus the pattern is monotonous:",False),
 B("•  the RAW winner is almost always a generic frequent root (God, say, establish);",False,RED),
 B("•  the CONTROLLED winner is the specific bond that actually teaches you something.",False,TEAL),
 B("The raw list is the same boring names every time; the controlled list is where the Qur'an's associations live.",False,GREY,16)],top=1.35,size=18,space=8)

# 10 second case
s=slide(); title(s,"A second case — honest scales")
pic(s,"fig_kayl_wazn.png",7.3,1.8,width=5.8)
body(s,[B("كيل (measure) and وزن (weight) share only 6 ayahs —",True,NAVY,20),
 B("but at ×137 over chance.",True,TEAL,22),
 B("The Qur'an's warning to the المطففين —",False),
 B("those who give short measure and short weight in the market.",False),
 B("Rare, specific pairs the raw count would have buried",False),
 B("under God and “say.”",False),
 B("More controlled bonds:",True,NAVY),
 B("•  عهد ↔ نقض  (covenant & breaking) — ×96",False),
 B("•  سجد ↔ ركع  (prostrate & bow) — ×46",False),
 B("•  قرض ↔ حسن  (a goodly loan, قرضًا حسنًا) — ×30",False)],top=1.4,size=18,width=6.7,space=7)

# 11 across targets
s=slide(); title(s,"The pattern holds across targets")
pic(s,"fig_targets_overview.png",7.5,1.4,width=5.6)
body(s,[B("Twelve targets, the same lesson every time:",True,NAVY,20),
 B("the RAW winner is a generic frequent root —",False,RED),
 B("ءله, قول, قوم, جري —",False,RED),
 B("the CONTROLLED winner is the specific bond:",False,TEAL),
 B("•  prayer → zakat",False),
 B("•  measure → weight",False),
 B("•  covenant → breaking",False),
 B("•  orphan → the needy",False),
 B("•  prostrate → bow",False),
 B("The control is not a trick tuned to one pair —",False,GREY),
 B("it is how the corpus's real structure shows itself.",False,GREY)],top=1.4,size=18,width=6.8,space=6)

# 12 clusters
s=slide(); title(s,"Co-occurrence also reveals clusters")
pic(s,"fig_paradise_cluster.png",7.4,1.9,width=5.7)
body(s,[B("A bond is not always a single pair — sometimes a constellation.",True,NAVY,20),
 B("Paradise (جنن) lights up with:",True),
 B("•  عدن (Eden) — ×32",False,TEAL),
 B("•  تحت / جري (rivers flowing beneath) — ×24",False,TEAL),
 B("•  خلد (eternity)",False,TEAL),
 B("جنات تجري من تحتها الأنهار —",False),
 B("“gardens beneath which rivers flow.”",False,GREY),
 B("Co-occurrence lights up the whole image at once —",False),
 B("the seed of the networks and motifs of later weeks.",False,GREY)],top=1.45,size=18,width=6.7,space=8)

# 13 limits
s=slide(); title(s,"What co-occurrence can't say (yet)")
body(s,[B("GIVES:  who travels with whom, fairly across frequencies.",True,NAVY,20),
 B("Does NOT give:",True,RED,20),
 B("•  Trust — our control divides out frequency but still ignores ayah LENGTH;",False),
 B("      a long verse hands out co-occurrences freely. Week 5 adds a length-aware null,",False,None,16),
 B("      tests it by simulation, and sorts pairs into tiers (real / borderline / spurious).",False,None,16),
 B("•  Direction — it is symmetric; it can't say which root pulls in the other.",False),
 B("      That asymmetry, P(A|B) vs P(B|A), is Week 6.",False,None,16),
 B("•  Cause or agreement — two roots share verses to be praised together OR contrasted;",False),
 B("      remember last week's antonym partners (ءمن ↔ كفر).",False,None,16),
 B("So pair every co-occurrence claim with a labeled reading.",True,TEAL)],top=1.35,size=18,space=7)

# 14 two-sentence + exercise
s=slide(); title(s,"Two-sentence reading & this week's exercise")
body(s,[B("FACT",True,NAVY,20),
 B("“Zakat shares 28 of its 56 ayahs with prayer — ×34.6 over chance, the top bond on prayer's slate.”",False),
 B("INTERPRETATION",True,TEAL,20),
 B("“I read prayer and charity as one devotion with two faces — toward God, and toward people.”",False),
 B("A classmate should tell instantly which sentence the computer could check, and which is yours.",False,GREY),
 B("EXERCISE",True,NAVY,20),
 B("•  Your target + its candidate slate → compute each candidate's × over chance.",False),
 B("•  Apply the support floor (joint ≥ 5); name the controlled winner.",False),
 B("•  One fact + one labeled interpretation. Submit the night before class.",False)],top=1.35,size=18,space=9)

# 15 close
s=slide(); title(s,"Rank by surprise, not by size")
body(s,[B("Frequency told us how much; distribution where; forms & partners what-shape and who-with.",False,GREY),
 B("Co-occurrence now measures who-with — fairly — by asking not “how many shared?”",True,NAVY),
 B("but “how many MORE than chance?”",True,NAVY,21),
 B("The course's spine, a fourth time: control before you conclude.",True,TEAL),
 B("",False),
 B("Next week — making the trust rigorous:",True,NAVY),
 B("a pair's lift, a length-aware null tested by simulation,",False),
 B("and the tiers that tell us which bonds are real.",False)],top=1.7,size=20,space=10)
prs.save(os.path.join(WK,"Week4_Slides.pptx")); print("deck rebuilt:",len(prs.slides._sldIdLst),"slides")
